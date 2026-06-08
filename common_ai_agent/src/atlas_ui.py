"""
src/atlas_ui.py — Atlas frontend server for common_ai_agent

Serves the static frontend bundle at common_ai_agent/frontend/atlas/ and
bridges it to the existing main.py agent loop via:

  • GET  /                       → frontend/atlas/index.html
  • GET  /<asset>                → frontend/atlas/<asset>     (jsx, css, js)
  • WS   /ws/agent               → bidirectional event stream

Activation:
    UI_MODE=atlas  python src/textual_main.py         (Windows)
    UI_MODE=atlas  python3 src/textual_main.py        (macOS/Linux)
    or directly:   python -m src.atlas_ui --port 8765 (Windows)

This mirrors web_ui.py (SSE) but uses WebSockets so the frontend can both
push prompts AND receive token / stage / tool / cost / todo events.

Author: Atlas frontend
"""

from __future__ import annotations

import argparse
import asyncio
import collections
import contextlib
import contextvars
import errno
import faulthandler
import hashlib
import html as html_lib
import importlib.util
import io
import json
import os
import queue
import re
import signal
import shlex
import socket
import subprocess
import sys
import threading
import time
from urllib.parse import quote
from pathlib import Path
from typing import Any, Callable, Final, Optional, TYPE_CHECKING

_THIS_DIR = Path(__file__).resolve().parent
_REPO_ROOT = _THIS_DIR.parent
for _p in (str(_THIS_DIR), str(_REPO_ROOT)):
    if _p not in sys.path:
        sys.path.insert(0, _p)

from core.atlas_context import AtlasContext
from core.import_evidence_assets import extract_markdown_import_assets
from core.prompt_input import normalize_prompt_images
from core.visual_evidence import (
    VisualSurface,
    render_office_visual_surfaces,
    render_pdf_visual_surfaces,
    visual_evidence_prompt,
)


def _configure_utf8_process_io() -> None:
    """Keep Atlas server/worker console I/O from crashing on Windows code pages."""

    os.environ.setdefault("PYTHONUTF8", "1")
    os.environ.setdefault("PYTHONIOENCODING", "utf-8:replace")
    for stream in (sys.stdout, sys.stderr):
        try:
            stream.reconfigure(errors="replace")
        except Exception:
            pass


_configure_utf8_process_io()


def _install_stack_dump_signal() -> None:
    """Allow `kill -USR1 <pid>` to dump stuck Atlas server stacks."""

    sigusr1 = getattr(signal, "SIGUSR1", None)
    if sigusr1 is None:
        return
    try:
        faulthandler.register(sigusr1, all_threads=True)
    except Exception:
        pass


_install_stack_dump_signal()


_DISCONNECT_ERRNOS = {
    errno.ECONNRESET,
    errno.ECONNABORTED,
    errno.EPIPE,
}
_DISCONNECT_WINERRORS = {
    10053,  # WSAECONNABORTED
    10054,  # WSAECONNRESET: remote host forcibly closed the connection
    10058,  # WSAESHUTDOWN
}
_DISCONNECT_TEXT_FRAGMENTS = (
    "connection reset by peer",
    "forcibly closed",
    "broken pipe",
    "connection aborted",
    "transport is closing",
)


def _is_disconnect_os_error(exc: BaseException) -> bool:
    """Classify OS-level socket resets as normal client disconnects."""

    if isinstance(exc, (ConnectionResetError, BrokenPipeError, EOFError)):
        return True
    if isinstance(exc, OSError):
        if getattr(exc, "errno", None) in _DISCONNECT_ERRNOS:
            return True
        if getattr(exc, "winerror", None) in _DISCONNECT_WINERRORS:
            return True
        msg = str(exc).lower()
        if any(fragment in msg for fragment in _DISCONNECT_TEXT_FRAGMENTS):
            return True
    return False


# `from __future__ import annotations` turns every type annotation into
# a string. FastAPI's `get_type_hints()` then needs to resolve those
# strings in the *module globals*. The inner endpoint functions live
# inside create_app() (they import fastapi locally), so without a
# module-level alias of `Request`, the annotation `request: Request`
# becomes an unresolvable ForwardRef and pydantic v2 falls back to
# treating `request` as a query parameter (→ 422 on every POST).
# This conditional import keeps the script usable even when fastapi is
# missing — `Request` just becomes None and the endpoint won't be
# registered (the local import inside create_app sys.exits first).
if TYPE_CHECKING:
    from fastapi import Request
else:
    try:
        from fastapi import Request  # noqa: F401  (runtime forward-ref target)
    except ImportError:
        class Request:  # fallback name for annotations when FastAPI is absent
            pass

from core.atlas_exec_policy import (
    EXEC_MODE_ORCHESTRATOR,
    apply_exec_mode_env,
    current_exec_mode,
    exec_policy_payload,
    initial_workflow_for_exec_mode,
    normalize_exec_mode,
)
from core.scm import (
    configured_scm_provider,
    resolve_scm_adapter,
    scm_provider_allows_missing_git_dir,
)

# ── Paths ──────────────────────────────────────────────────────────
HERE         = Path(__file__).resolve().parent
SOURCE_ROOT  = HERE.parent                            # common_ai_agent/ (source)
FRONTEND     = SOURCE_ROOT / "frontend" / "atlas"


def _history_content_text(content: Any) -> str:
    """Return readable text from an OpenAI-style message content field."""

    if isinstance(content, str):
        return content
    if isinstance(content, list):
        parts: list[str] = []
        for item in content:
            if isinstance(item, dict):
                text = item.get("text") or item.get("content") or ""
                if text:
                    parts.append(str(text))
            elif item is not None:
                parts.append(str(item))
        return "\n".join(parts)
    if content is None:
        return ""
    return str(content)


def _history_message_preview(message: dict[str, Any], *, limit: int = 240) -> str:
    role = str(message.get("role") or "unknown")
    text = _history_content_text(message.get("content"))
    text = re.sub(r"\s+", " ", text).strip()
    if len(text) > limit:
        text = text[: max(0, limit - 1)].rstrip() + "..."
    return f"{role}: {text or '(empty)'}"


def _prompt_ack_frames(
    *,
    msg_id: str,
    text_preview: str,
    session_id: str,
    ok: bool,
    queued: bool = False,
    handled: str = "",
    duplicate: bool = False,
    error: str = "",
) -> list[dict[str, Any]]:
    """Build the ordered WS ack frames for a prompt acceptance outcome.

    AC-2 contract (see tests/test_atlas_worker_warmup_input.py):

    - ``agent_received`` is the TRANSPORT ack consumed by backend.js to disarm
      its retransmit timer. It must be emitted ONLY when the prompt is
      genuinely accepted (``ok`` is True) — never up-front before delivery is
      known. Emitting it unconditionally lied about delivery and defeated the
      only auto-retry, the original silent-loss bug.
    - ``agent_accepted`` is the DELIVERY ack and is ALWAYS emitted, carrying
      ``ok``/``queued``/``duplicate``/``handled``/``error`` so the client can
      decide whether to hold the input for a manual resend.

    When ``ok`` is False (the worker dropped the prompt) NO ``agent_received``
    is produced, so backend.js's pending-ack timer fires and retransmits the
    prompt — which auto-lands once the worker is warm. A duplicate (already-seen
    msg_id) is still a genuine accept (the original landed), so it is acked
    ``ok:true,duplicate`` WITH ``agent_received``; msg_id dedup keeps the worker
    from double-running.
    """
    frames: list[dict[str, Any]] = []
    if ok:
        frames.append({
            "type": "agent_received",
            "msg_id": msg_id,
            "text_preview": text_preview,
            "session_id": session_id,
        })
    accepted: dict[str, Any] = {
        "type": "agent_accepted",
        "msg_id": msg_id,
        "text_preview": text_preview,
        "session_id": session_id,
        "ok": bool(ok),
        "queued": bool(queued),
        "duplicate": bool(duplicate),
    }
    if handled:
        accepted["handled"] = handled
    if error:
        accepted["error"] = error
    frames.append(accepted)
    return frames


def _parse_compact_history_signal(signal: str) -> tuple[int, bool, str]:
    """Parse COMPACT_HISTORY signal options from the slash-command registry."""

    options = signal.split(":", 1)[1].strip() if ":" in signal else ""
    keep_recent = 4
    dry_run = False
    instruction = ""
    if not options:
        return keep_recent, dry_run, instruction
    keep_match = re.search(r"\bkeep=(\d+)\b", options)
    if keep_match:
        keep_recent = max(0, int(keep_match.group(1)))
    elif "dry_run=true" in options:
        dry_run = True
    else:
        instruction = options
    return keep_recent, dry_run, instruction


def _write_history_json(path: Path, messages: list[dict[str, Any]]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    tmp = path.with_suffix(path.suffix + ".tmp")
    tmp.write_text(json.dumps(messages, ensure_ascii=False, indent=2), encoding="utf-8")
    tmp.replace(path)


def _clear_history_file(path: Path, signal: str) -> tuple[str, list[dict[str, Any]]]:
    messages = _load_history_json(path)
    keep_pairs = 0
    if ":" in signal:
        try:
            keep_pairs = max(0, int(signal.split(":", 1)[1]))
        except ValueError:
            keep_pairs = 0
    system_msgs = [m for m in messages if m.get("role") == "system"]
    active_msgs = [m for m in messages if m.get("role") != "system"]
    kept = active_msgs[-(keep_pairs * 2):] if keep_pairs > 0 else []
    cleared = system_msgs + kept
    _write_history_json(path, cleared)
    if keep_pairs > 0:
        return f"Conversation history cleared; kept last {keep_pairs} message pair(s).", cleared
    return "Conversation history cleared.", cleared


def _estimate_history_tokens(messages: list[dict[str, Any]]) -> int:
    try:
        from src.llm_client import estimate_message_tokens as _estimate_msg
    except Exception:
        try:
            from llm_client import estimate_message_tokens as _estimate_msg  # type: ignore
        except Exception:
            _estimate_msg = None
    if callable(_estimate_msg):
        return sum(_estimate_msg(m) for m in messages)
    return sum(max(1, len(json.dumps(m, ensure_ascii=False)) // 4) for m in messages)


def _max_context_tokens() -> int:
    try:
        import src.config as _cfg_context  # noqa: WPS433
    except Exception:
        try:
            import config as _cfg_context  # type: ignore  # noqa: WPS433
        except Exception:
            _cfg_context = None
    return int(getattr(_cfg_context, "MAX_CONTEXT_TOKENS", 0) or 0) if _cfg_context else 0


def _history_context_usage(messages: list[dict[str, Any]]) -> tuple[int, int]:
    """Return the same context total shape used by the /context tracker."""

    try:
        from core.context_tracker import get_tracker

        tracker = get_tracker()
        tracker.messages = messages
        first_system = next((m for m in messages if m.get("role") == "system"), None)
        tracker.update_system_prompt(
            _history_content_text(first_system.get("content")) if first_system else ""
        )
        tracker.update_messages(messages, exclude_system=True)
        return int(tracker.get_total_tokens()), int(getattr(tracker, "max_tokens", 0) or 0)
    except Exception:
        return _estimate_history_tokens(messages), _max_context_tokens()


def _emit_history_context_update(client_session: Any, messages: list[dict[str, Any]]) -> None:
    used, max_ctx = _history_context_usage(messages)
    if max_ctx:
        client_session.emit("context", used=used, max=max_ctx)




from src.atlas_runtime import _resolve_workflow_root  # Phase 6: needed before module-level WORKFLOW_ROOT computation
WORKFLOW_ROOT = _resolve_workflow_root()
# PROJECT_ROOT is the user's cwd at launch, NOT the source repo. This
# lets the user run `python ../path/to/textual_main.py` from any
# project directory and have the file API + scope operate on THAT dir.
PROJECT_ROOT = Path(os.getcwd()).resolve()

# Phase 1 refactor: SSOT rendering layer extracted to src/atlas_ssot_export.py.
# All 56 moved names are re-exported here for backward compatibility with
# tests and any module that still resolves them via atlas_ui.
from src.atlas_ssot_export import (
    _SSOT_EXPORT_SECTION_ORDER,
    _ssot_md_escape_cell,
    _ssot_md_scalar,
    _ssot_md_bit_range,
    _ssot_bit_range_info,
    _ssot_md_is_short_scalar,
    _ssot_md_yaml_block,
    _ssot_md_dict_table,
    _ssot_md_auto_table,
    _ssot_md_definition_list,
    _ssot_md_section_top_module,
    _ssot_md_section_sub_modules,
    _ssot_md_section_parameters,
    _ssot_md_section_io_list,
    _ssot_md_section_features,
    _ssot_md_section_registers,
    _ssot_md_section_raw_yaml,
    _ssot_md_section_generic,
    _SSOT_MD_SECTION_RENDERERS,
    _ssot_section_is_empty,
    _ssot_html_escape,
    _ssot_html_item_name,
    _ssot_html_submodules,
    _ssot_html_interfaces,
    _ssot_html_block_diagram,
    _ssot_fsm_machines,
    _ssot_mermaid_state_id,
    _ssot_mermaid_label,
    _ssot_html_fsm_mermaid,
    _ssot_html_fsm_section,
    _ssot_html_signal_values,
    _ssot_html_timing_section,
    _ssot_pick,
    _ssot_html_is_scalar,
    _ssot_html_scalar,
    _ssot_html_yaml_pre,
    _ssot_html_value_block,
    _ssot_html_rule_cards,
    _ssot_html_function_model,
    _ssot_html_cycle_model,
    _ssot_field_bit_info,
    _ssot_html_field_bits,
    _ssot_html_register_field_table,
    _ssot_html_register_block,
    _ssot_html_hex,
    _ssot_html_registers,
    _ssot_html_insert_after_section,
    _ssot_html_design_views,
    _ssot_to_markdown,
    _ssot_html_insert_after_top_module,
    _ssot_html_normalize_mermaid_fences,
    _ssot_html_mermaid_runtime,
    _ssot_resolve_custom_file,
    _ssot_html_render_custom_block,
    _ssot_html_custom_blocks_for,
    _ssot_to_html,
    set_project_root as _set_ssot_export_project_root,
)
_set_ssot_export_project_root(PROJECT_ROOT)

# Phase 2 refactor: SSOT DOCX rendering extracted to src/atlas_ssot_docx.py.
# All 36 moved names re-exported for backward compat with tests/in-tree callers.
from src.atlas_ssot_docx import (
    _ssot_docx_set_mono,
    _ssot_docx_add_kv,
    _ssot_docx_yaml_block,
    _ssot_docx_list_of_dicts,
    _ssot_docx_dict_block,
    _ssot_docx_render_section,
    _ssot_docx_cover_page,
    _ssot_docx_render_revision_history,
    _ssot_docx_apply_heading_numbering,
    _ssot_docx_set_footer,
    _ssot_docx_page_break,
    _ssot_docx_add_word_index,
    _ssot_docx_add_toc,
    _ssot_docx_add_list_of_tables,
    _ssot_docx_add_list_of_figures,
    _ssot_docx_caption,
    _ssot_docx_table_from_rows,
    _ssot_docx_render_top_module,
    _ssot_docx_render_parameters,
    _ssot_docx_render_features,
    _ssot_docx_render_register_detail,
    _ssot_docx_port_direction,
    _ssot_docx_port_rows,
    _ssot_docx_render_port_buckets,
    _ssot_docx_render_io_list,
    _ssot_docx_render_interrupts,
    _ssot_docx_render_fsm,
    _ssot_docx_render_cycle_model,
    _ssot_docx_render_error_handling,
    _ssot_docx_block_diagram_png,
    _ssot_docx_render_block_diagram,
    _ssot_docx_render_function_description,
    _ssot_docx_render_programming_model,
    _ssot_docx_render_programming_sequence,
    _SSOT_DOCX_APPENDIX_KEYS,
    _ssot_to_docx,
)

# Phase 3 refactor: history compactor extracted to src/atlas_compactor.py
from src.atlas_compactor import (
    _load_history_json,
    _compact_history_file,
    _default_web_compress_fn,
    _compact_history_llm,
)

# Phase 10 refactor (PoC): Q&A section labels moved to src/atlas_qa.py
from src.atlas_qa import (
    _ssot_qa_section,
    _status_group, _qa_slug, _ssot_q_pairs_from_questions,
)

# Phase 4 refactor: server runtime + CLI extracted to src/atlas_runtime.py
from src.atlas_runtime import (
    run_atlas_ui,
    _launch_admin_server,
    main,
)

# Phase 5 refactor: model options + reasoning-effort -> src/atlas_model_options.py
from src.atlas_model_options import (
    _REASONING_EFFORT_OPTIONS,
    _REASONING_EFFORT_ALIASES,
    _MODEL_OPTION_KEYS,
    _BASE_MODEL_OPTION_KEYS,
    _LEGACY_MODEL_OPTION_KEYS,
    _RUNTIME_MODEL_OPTION_KEY,
    _MODEL_CATALOG_ENV_KEYS,
    _PROFILE_MODEL_OPTION_PREFIX,
    _RAW_MODEL_OPTION_PREFIX,
    _normalize_reasoning_effort,
    _canonical_model_option_key,
    _catalog_model_option_rows,
    _model_option_rows,
    _set_runtime_model,
)
# Backwards compat alias — older code references ROOT.
ROOT         = SOURCE_ROOT

_HEALTHZ_COST_CACHE: dict[tuple[str, str, str, str], tuple[float, dict[str, Any]]] = {}
_HEALTHZ_COST_CACHE_LOCK = threading.RLock()


def _healthz_cost_cache_ttl() -> float:
    try:
        return max(0.0, float(os.environ.get("ATLAS_HEALTHZ_COST_CACHE_TTL", "30") or 30))
    except Exception:
        return 30.0


def _healthz_cost_cache_get(key: tuple[str, ...]) -> dict[str, Any] | None:
    ttl = _healthz_cost_cache_ttl()
    if ttl <= 0:
        return None
    now = time.monotonic()
    with _HEALTHZ_COST_CACHE_LOCK:
        hit = _HEALTHZ_COST_CACHE.get(key)
        if hit and (now - hit[0]) < ttl:
            return dict(hit[1])
    return None


def _healthz_cost_cache_set(key: tuple[str, ...], value: dict[str, Any]) -> None:
    ttl = _healthz_cost_cache_ttl()
    if ttl <= 0:
        return
    with _HEALTHZ_COST_CACHE_LOCK:
        _HEALTHZ_COST_CACHE[key] = (time.monotonic(), dict(value))


_atlas_active_session_cv = contextvars.ContextVar("atlas_active_session", default="")
_atlas_active_ip_cv = contextvars.ContextVar("atlas_active_ip", default="")
_agent_mode_override_cv = contextvars.ContextVar("agent_mode_override", default="")
_plan_mode_cv = contextvars.ContextVar("plan_mode", default="false")


def _active_session_value() -> str:
    current = (_atlas_active_session_cv.get() or "").strip()
    env_value = (os.environ.get("ATLAS_ACTIVE_SESSION", "") or "").strip()
    try:
        from core.atlas_multiuser import get_atlas_bridge_session_id
        bridge_session = (get_atlas_bridge_session_id() or "").strip()
    except Exception:
        bridge_session = ""
    if bridge_session and bridge_session not in {"default", "default/default", "default/default/default"}:
        return bridge_session
    if env_value and env_value not in {"default", "default/default", "default/default/default"} and env_value != current:
        return env_value
    # FastAPI request tasks can inherit the startup contextvar value
    # ("default/default/default"). /api/session/activate mirrors the real
    # active namespace into os.environ for cross-task visibility, so do not
    # let that placeholder shadow a newer backend activation.
    if current and (current not in {"default", "default/default", "default/default/default"} or not env_value):
        return current
    return env_value or current


def _atlas_todo_item_from_raw(item: dict[str, Any]) -> dict[str, Any]:
    todo = dict(item)
    content = todo.get("content") or todo.get("title") or todo.get("id") or ""
    todo["content"] = str(content)
    todo["activeForm"] = str(todo.get("activeForm") or todo.get("active_form") or content)
    todo["status"] = str(todo.get("status") or "pending")
    todo["priority"] = str(todo.get("priority") or "medium")
    todo["detail"] = str(todo.get("detail") or "")
    todo["criteria"] = str(todo.get("criteria") or "")
    return todo


def _atlas_todo_payload_from_raw(raw: Any, fallback: dict[str, Any] | None = None) -> dict[str, Any]:
    payload = dict(fallback or {})
    raw_items: Any = None
    if isinstance(raw, dict):
        for key in (
            "name",
            "description",
            "lock_additions",
            "source_plan",
            "source_task_count",
            "status_counts",
            "ui_grouping",
            "current_index",
        ):
            if key in raw:
                payload[key] = raw[key]
        raw_items = raw.get("todos")
        if not isinstance(raw_items, list):
            raw_items = raw.get("tasks")
    elif isinstance(raw, list):
        raw_items = raw
    if isinstance(raw_items, list):
        payload["todos"] = [
            _atlas_todo_item_from_raw(item)
            for item in raw_items
            if isinstance(item, dict)
        ]
    else:
        payload.setdefault("todos", [])
    return payload




def _persist_config_values(updates: dict[str, str]) -> None:
    """Persist simple KEY=value settings to common_ai_agent/.config."""
    _persist_key_values(SOURCE_ROOT / ".config", updates)


def _persist_env_values(updates: dict[str, str]) -> None:
    """Persist user-editable runtime settings to common_ai_agent/.env."""
    _persist_key_values(SOURCE_ROOT / ".env", updates)


def _persist_key_values(path: Path, updates: dict[str, str]) -> None:
    try:
        lines = path.read_text(encoding="utf-8").splitlines()
    except OSError:
        lines = []

    seen = set()
    key_re = "|".join(re.escape(k) for k in updates)
    out = []
    for line in lines:
        match = re.match(rf"^(\s*)({key_re})(\s*)=.*$", line)
        if match:
            key = match.group(2)
            out.append(f"{match.group(1)}{key}{match.group(3)}={updates[key]}")
            seen.add(key)
        else:
            out.append(line)
    for key, value in updates.items():
        if key not in seen:
            out.append(f"{key}={value}")
    path.write_text("\n".join(out).rstrip() + "\n", encoding="utf-8")


def _read_env_file_values() -> dict[str, str]:
    values: dict[str, str] = {}
    try:
        lines = (SOURCE_ROOT / ".env").read_text(encoding="utf-8").splitlines()
    except OSError:
        return values
    for raw in lines:
        line = raw.strip()
        if not line or line.startswith("#") or "=" not in line:
            continue
        key, value = line.split("=", 1)
        key = key.strip()
        value = value.strip()
        if "#" in value:
            value = value.split("#", 1)[0].strip()
        if key:
            values[key] = value
    return values




def _env_value(env_file: dict[str, str], key: str) -> str:
    return (env_file.get(key, os.environ.get(key, "")) or "").strip()


def _profile_from_env_values(name: str, env_file: dict[str, str]) -> dict[str, str]:
    profile = str(name or "").strip()
    if not profile:
        return {}
    pfx = f"PROFILE_{profile}_"
    model = _env_value(env_file, pfx + "MODEL")
    if not model:
        return {}
    return {
        "name": profile,
        "model": model,
        "base_url": _env_value(env_file, pfx + "BASE_URL") or _env_value(env_file, "LLM_BASE_URL"),
        "api_key": _env_value(env_file, pfx + "API_KEY") or _env_value(env_file, "LLM_API_KEY"),
    }


def _profile_name_from_option_key(key: str) -> str:
    raw = _canonical_model_option_key(key)
    if raw.startswith(_PROFILE_MODEL_OPTION_PREFIX):
        return raw[len(_PROFILE_MODEL_OPTION_PREFIX):].strip()
    return ""


def _is_model_slot_key(key: str) -> bool:
    raw = _canonical_model_option_key(key)
    return any(raw in group for group in (_MODEL_OPTION_KEYS, _BASE_MODEL_OPTION_KEYS, _LEGACY_MODEL_OPTION_KEYS))


def _split_model_catalog(raw: str) -> list[str]:
    return [
        part.strip()
        for part in re.split(r"[,\n]+", str(raw or ""))
        if part.strip()
    ]




def _model_option_index(key: str) -> int | None:
    raw = _canonical_model_option_key(key)
    for group in (_MODEL_OPTION_KEYS, _BASE_MODEL_OPTION_KEYS, _LEGACY_MODEL_OPTION_KEYS):
        if raw in group:
            return group.index(raw)
    return None


def _display_model_option_keys(env_file: dict[str, str]) -> tuple[str, str, str]:
    selected_key = (
        env_file.get("LLM_SELECTED_MODEL_KEY", "")
        or os.environ.get("LLM_SELECTED_MODEL_KEY", "")
    ).strip()
    if selected_key in _BASE_MODEL_OPTION_KEYS:
        return _BASE_MODEL_OPTION_KEYS
    if selected_key in _LEGACY_MODEL_OPTION_KEYS:
        return _LEGACY_MODEL_OPTION_KEYS
    if any((env_file.get(key, os.environ.get(key, "")) or "").strip() for key in _BASE_MODEL_OPTION_KEYS):
        return _BASE_MODEL_OPTION_KEYS
    if any((env_file.get(key, os.environ.get(key, "")) or "").strip() for key in _LEGACY_MODEL_OPTION_KEYS):
        return _LEGACY_MODEL_OPTION_KEYS
    return _MODEL_OPTION_KEYS


def _model_option_value(env_file: dict[str, str], index: int) -> str:
    """Return dropdown slot value, preferring LLM_MODEL_NAME* over old names."""
    for group in (_MODEL_OPTION_KEYS, _BASE_MODEL_OPTION_KEYS, _LEGACY_MODEL_OPTION_KEYS):
        key = group[index]
        value = (env_file.get(key, os.environ.get(key, "")) or "").strip()
        if value:
            return value
    return ""


def _refresh_config_after_persist() -> None:
    """Refresh config mtime cache so the next /healthz does not undo runtime settings."""
    for mod_name in ("src.config", "config"):
        mod = sys.modules.get(mod_name)
        if mod is None:
            try:
                mod = __import__(mod_name, fromlist=["*"])
            except Exception:
                mod = None
        if mod is not None:
            try:
                mod.reload_env()
            except Exception:
                pass


def _persist_reasoning_effort(effort: str) -> None:
    glm_thinking = "disabled" if effort == "none" else "enabled"
    _persist_config_values({
        "REASONING_MODE": effort,
        "REASONING_EFFORT": effort,
        "GLM_THINKING_TYPE": glm_thinking,
    })


def _set_runtime_reasoning_effort(effort: str) -> None:
    glm_thinking = "disabled" if effort == "none" else "enabled"
    os.environ["REASONING_MODE"] = effort
    os.environ["REASONING_EFFORT"] = effort
    os.environ["GLM_THINKING_TYPE"] = glm_thinking
    config_modules = []
    seen_module_ids = set()
    for mod_name in ("src.config", "config"):
        mod = sys.modules.get(mod_name)
        if mod is not None and id(mod) not in seen_module_ids:
            config_modules.append(mod)
            seen_module_ids.add(id(mod))
    if not config_modules:
        try:
            mod = __import__("src.config", fromlist=["*"])
            config_modules.append(mod)
            sys.modules.setdefault("config", mod)
        except Exception:
            try:
                mod = __import__("config", fromlist=["*"])
                config_modules.append(mod)
                sys.modules.setdefault("src.config", mod)
            except Exception:
                pass
    for mod in config_modules:
        if mod is not None:
            setattr(mod, "REASONING_MODE", effort)
            setattr(mod, "REASONING_EFFORT", effort)
            setattr(mod, "GLM_THINKING_TYPE", glm_thinking)






def _apply_selected_model_from_env() -> str:
    selected_key = _canonical_model_option_key(os.environ.get("LLM_SELECTED_MODEL_KEY", ""))
    selected_row = next((row for row in _model_option_rows() if row.get("key") == selected_key), None)
    if selected_row:
        model = str(selected_row.get("model") or "")
        if model:
            _set_runtime_model(model, selected_key)
            return model
    selected_index = _model_option_index(selected_key)
    if selected_index is not None:
        model = _model_option_value({}, selected_index)
        if model:
            _set_runtime_model(model, selected_key)
            return model
    return ""


def _active_ip_value() -> str:
    current = (_atlas_active_ip_cv.get() or "").strip()
    env_value = (os.environ.get("ATLAS_ACTIVE_IP", "") or "").strip()
    try:
        from core.atlas_multiuser import get_atlas_bridge_session_id
        bridge_session = (get_atlas_bridge_session_id() or "").strip()
    except Exception:
        bridge_session = ""
    bridge_parts = [part for part in bridge_session.split("/") if part]
    if len(bridge_parts) >= 3 and bridge_parts[-2] != "default":
        return bridge_parts[-2]
    if env_value and env_value != "default" and env_value != current:
        return env_value
    if current and (current != "default" or not env_value):
        return current
    return env_value or current


def _plan_mode_value() -> str:
    return _plan_mode_cv.get() or os.environ.get("PLAN_MODE", "false")


def _sync_env_to_context() -> None:
    """Copy contextvars back to os.environ for legacy main.py reads."""
    session_value = (_atlas_active_session_cv.get() or "").strip()
    session_env = (os.environ.get("ATLAS_ACTIVE_SESSION", "") or "").strip()
    if session_value and (
        session_value not in {"default", "default/default", "default/default/default"}
        or not session_env
    ):
        os.environ["ATLAS_ACTIVE_SESSION"] = session_value
    elif not session_env:
        os.environ["ATLAS_ACTIVE_SESSION"] = session_value

    ip_value = (_atlas_active_ip_cv.get() or "").strip()
    ip_env = (os.environ.get("ATLAS_ACTIVE_IP", "") or "").strip()
    if ip_value and (ip_value != "default" or not ip_env):
        os.environ["ATLAS_ACTIVE_IP"] = ip_value
    elif not ip_env:
        os.environ["ATLAS_ACTIVE_IP"] = ip_value

    os.environ["AGENT_MODE_OVERRIDE"] = _agent_mode_override_cv.get()
    os.environ["PLAN_MODE"] = _plan_mode_cv.get()


def _python_cmd() -> str:
    """Return the Python launcher for generated user-facing commands."""
    return "python" if os.name == "nt" else "python3"


try:
    from .workflow_stage_engine import _rtl_manifest_progress as _shared_rtl_manifest_progress
except Exception:
    try:
        from workflow_stage_engine import _rtl_manifest_progress as _shared_rtl_manifest_progress  # type: ignore
    except Exception:
        _shared_rtl_manifest_progress = None  # type: ignore

try:
    from core.session_names import normalize_session_name
except Exception:
    from session_names import normalize_session_name  # type: ignore


_DEFAULT_SESSION_PLACEHOLDERS = {"default/default", "default/default/default"}
_LOCKED_TRUTH_DRAFT_SENTINEL: Final[str] = "[ATLAS LOCKED TRUTH DRAFT MODE]"
_LOCKED_TRUTH_DRAFT_SKIP_IPS: Final[frozenset[str]] = frozenset(
    {"", "default", "soc", "user"}
)
_LOCKED_TRUTH_DRAFT_OVERLAY: Final[str] = """[ATLAS LOCKED TRUTH DRAFT MODE]
Active IP: {ip}

This IP does not have an approved locked-truth manifest yet.

Purpose: collect and refine locked truth only.

Do not run workflow stages.
Do not generate RTL/TB/SIM.
Do not call dispatch_workflow.
Do not modify implementation artifacts.

Conversation rule:
- Lock truth per requirement. An IP is implementation-locked only when every required requirement is locked or approved.
- Prefer short tiktaka. Discuss only the next missing or ambiguous decision.
- Prefer normal chat for one or two missing decisions.
- Ask one concise natural-language question at a time in normal chat.
- Ask only for lock-blocking ambiguity.
- Do not ask optional refinement questions when a reasonable default is available.
- If enough information exists, stop interviewing and offer to draft or summarize locked truth.
- If 3 or more lock-blocking decisions remain, suggest grill-me/deep interview in normal chat before asking more one-by-one questions.
- Do not launch ask_user automatically for grill-me; wait until the user explicitly asks for grill-me/deep interview or final approval.
- Use ask_user only for grill-me/deep-interview, final approval, or batched structured decisions.
- If structured approval is needed, call ask_user with the smallest useful question set.
- Do not leave lock-blocking open questions as a prose list; ask the next one directly or use ask_user for structured approval.
- Existing RTL/doc/SSOT artifacts are read-only candidate evidence, not authority. If they exist, use them only to extract draft decisions, conflicts, and questions.
- Do not mark a legacy-derived requirement locked unless the user explicitly approves that requirement.
- If an existing artifact conflicts with user answers or lacks approval, ask_user whether to keep, update, or ignore that candidate.
- After ask_user answers, do not dump the full draft. Reply only with captured decisions, remaining blockers, and the next question.
- Produce the full requirement/obligation/contract_ref/structural_contract/behavioral_contract/evidence draft only when the user explicitly asks to show/review/export it, or when asking for final approval.
- Use /draft-req to write draft req JSON files: requirements_index.json, obligations.json, contract_refs.json, structural_contracts.json, behavioral_contracts.json, and evidence_plan.json.
- Use /finalize-req to quality-review/repair those draft files and run check_contract_bundle.py --review-candidate. Finalize means ready_for_human_review, not locked.
- Use /lock-req only after explicit human approval. It runs lock_requirement_set.py --from-candidate and then check_contract_bundle.py.
- Do not say approved or locked unless /lock-req created or you can cite req/locked_truth.md and req/approval_manifest.json with status requirements_locked.
- If no writer/tool is available, say the approval was captured in chat only and ask the user to run /draft-req, /finalize-req, then /lock-req before /to-ssot.
- Keep normal replies under 12 lines unless the user asks for detail.

Wait for explicit user approval before implementation.

[USER MESSAGE]
{user_message}
"""


def _ensure_default_workspace_commands_registered() -> None:
    from core.slash_commands import get_registry
    from workflow.loader import (
        get_todo_template_registry,
        load_workspace,
        register_workspace_commands,
    )

    ws = load_workspace("default", project_root=_REPO_ROOT)
    if ws.todo_templates_dir:
        get_todo_template_registry().load_from_dir(ws.todo_templates_dir)
    register_workspace_commands(ws, get_registry())


def _normalized_session_or_empty(value: Any) -> str:
    session = normalize_session_name(str(value or ""))
    if session in _DEFAULT_SESSION_PLACEHOLDERS:
        return ""
    return session


_FAST_IDENTITY_PROMPTS = {
    "who",
    "whoareyou",
    "whoyou",
    "whoami",
    "너누구",
    "누구",
}
_FAST_GREETING_PROMPTS = {
    "hi",
    "hello",
    "hey",
    "안녕",
    "안녕하세요",
}


def _atlas_fast_prompt_candidate(text: str) -> str:
    raw = str(text or "").strip()
    if not raw:
        return ""
    if "\n" not in raw:
        return raw
    lines = [line.strip() for line in raw.splitlines() if line.strip()]
    if not lines:
        return ""
    if lines[0].startswith("[Atlas UI language preference]") or lines[0].startswith("[scope]"):
        return lines[-1]
    return raw


def _atlas_fast_prompt_kind(text: str) -> str:
    raw = _atlas_fast_prompt_candidate(text)
    if not raw or "\n" in raw or len(raw) > 48:
        return ""
    compact = re.sub(r"[\s\?!.。！？,，:;`'\"()\[\]{}<>]+", "", raw.lower())
    if compact in _FAST_IDENTITY_PROMPTS:
        return "identity"
    if compact in _FAST_GREETING_PROMPTS:
        return "greeting"
    return ""


def _atlas_session_route_parts(session_id: str) -> tuple[str, str, str]:
    parts = [part for part in normalize_session_name(str(session_id or "")).split("/") if part]
    if len(parts) >= 4:
        return parts[0], parts[2], parts[3]
    while len(parts) < 3:
        parts.append("default")
    return parts[0], parts[1], parts[2]


def _locked_truth_draft_ip(session_id: str, msg: dict[str, Any]) -> str:
    for key in ("ip", "scope", "ip_id"):
        candidate = normalize_session_name(str(msg.get(key) or "")).strip("/")
        if "/" in candidate:
            candidate = candidate.split("/")[-1]
        if candidate not in _LOCKED_TRUTH_DRAFT_SKIP_IPS and re.fullmatch(
            r"[A-Za-z][A-Za-z0-9_]*",
            candidate,
        ):
            return candidate
    _, ip, _ = _atlas_session_route_parts(session_id)
    if ip in _LOCKED_TRUTH_DRAFT_SKIP_IPS:
        return ""
    if not re.fullmatch(r"[A-Za-z][A-Za-z0-9_]*", ip):
        return ""
    return ip


def _locked_truth_active(project_root: Path, ip: str) -> bool:
    try:
        from core.locked_truth_guard import is_locked_truth_active
    except ModuleNotFoundError:
        from locked_truth_guard import is_locked_truth_active  # type: ignore
    return is_locked_truth_active(project_root, ip)


def _locked_truth_candidate_roots(project_root: Path, session_id: str) -> tuple[Path, ...]:
    roots: list[Path] = []
    seen: set[Path] = set()

    def add(root: Path) -> None:
        resolved = root.expanduser().resolve()
        if resolved not in seen:
            seen.add(resolved)
            roots.append(resolved)

    add(project_root)
    session = normalize_session_name(str(session_id or ""))
    if session:
        atlas_roots = [
            os.environ.get("ATLAS_ROOT", "").strip(),
            str(project_root),
        ]
        for raw_root in atlas_roots:
            if not raw_root:
                continue
            try:
                context = AtlasContext.from_session_key(session, atlas_root=raw_root)
            except ValueError:
                continue
            add(context.workspace_root)
    return tuple(roots)


def _locked_truth_active_for_session(project_root: Path, session_id: str, ip: str) -> bool:
    return any(
        _locked_truth_active(root, ip)
        for root in _locked_truth_candidate_roots(project_root, session_id)
    )


def _apply_locked_truth_draft_overlay(
    project_root: Path,
    session_id: str,
    msg: dict[str, Any],
    text: str,
) -> tuple[str, bool]:
    raw_text = str(text or "").strip()
    if not raw_text or raw_text.startswith("/") or raw_text.startswith(_LOCKED_TRUTH_DRAFT_SENTINEL):
        return raw_text, False
    ip = _locked_truth_draft_ip(session_id, msg)
    if not ip or _locked_truth_active_for_session(project_root, session_id, ip):
        return raw_text, False
    return _LOCKED_TRUTH_DRAFT_OVERLAY.format(ip=ip, user_message=raw_text), True


def _atlas_fast_identity_response(session_id: str, text: str) -> str:
    _, ip, workflow = _atlas_session_route_parts(session_id)
    kind = _atlas_fast_prompt_kind(text)
    prefix = "Hi. " if kind == "greeting" else ""
    route = f"ask:{workflow or 'default'}"
    ip_label = ip or "default"
    workflow_label = workflow or "default"
    if workflow_label == "ssot-gen":
        body = (
            f"I am the `ssot-gen` workflow agent for `{ip_label}`. "
            f"This route is `{route}`: I own the YAML contract only, mainly "
            f"`{ip_label}/yaml/{ip_label}.ssot.yaml`, plus requirement capture, "
            "schema validation, and `[SSOT HANDOFF]`. RTL, testbench, firmware, "
            "and simulation stay with downstream workflows."
        )
    elif workflow_label == "default":
        body = (
            f"I am the default Atlas agent for `{ip_label}`. "
            f"This route is `{route}`: I can inspect files, update artifacts, "
            "run validations, answer workspace questions, and route specialized "
            "work to workflow agents when needed."
        )
    else:
        body = (
            f"I am the `{workflow_label}` workflow agent for `{ip_label}`. "
            f"This chat is routed through `{route}`."
        )
    return prefix + body


def _atlas_emit_session_id() -> str:
    """Resolve the safest session id for backend-to-browser agent events.

    Prefer the per-agent bridge context over process-global env. The env
    fallback remains only for older code paths that have no context.
    """
    try:
        from core.atlas_multiuser import get_atlas_bridge_session_id
        bridge_session = _normalized_session_or_empty(get_atlas_bridge_session_id())
    except Exception:
        bridge_session = ""
    if bridge_session and bridge_session != "default":
        return bridge_session

    context_session = _normalized_session_or_empty(_atlas_active_session_cv.get())
    if context_session and context_session != "default":
        return context_session

    env_session = _normalized_session_or_empty(os.environ.get("ATLAS_ACTIVE_SESSION", ""))
    return bridge_session or context_session or env_session


# ── ask_user answer formatter ──────────────────────────────────────
def _format_answer(ans: dict[str, Any], options: list[dict[str, Any]]) -> str:
    """Render a UI answer payload back into a tool observation string."""
    selected_ids = ans.get("selected") or []
    custom = (ans.get("custom") or "").strip()
    label_by_id = {o.get("id"): o.get("label", o.get("id")) for o in options or []}
    selected_labels = [label_by_id.get(sid, sid) for sid in selected_ids]
    parts = []
    if selected_labels:
        parts.append("selected: " + ", ".join(selected_labels))
    if custom:
        parts.append("note: " + custom)
    if not parts:
        return "(user submitted with no selection)"
    return " · ".join(parts)


# ── SSOT export helpers (module scope) ─────────────────────────────
# Reverse direction of the /import flow: take <ip>/yaml/<ip>.ssot.yaml
# and render md / docx / html for human review and sign-off. Walked
# deterministically (no LLM, no subprocess). Exposed at module scope so
# the DoD smoke test and tests/test_ssot_export.py can call them
# directly without going through create_app(). The HTTP endpoint
# `api_ssot_export` (inside create_app) is a thin wrapper around these.

_SSOT_EXPORT_IP_NAME_RE = re.compile(r"^[A-Za-z][A-Za-z0-9_]*$")


def _ssot_export_valid_ip(name: str) -> bool:
    return bool(_SSOT_EXPORT_IP_NAME_RE.match(name or ""))


def _configured_ip_root(ip: str = "") -> Path | None:
    raw = os.environ.get("ATLAS_IP_ROOT", "").strip()
    if not raw:
        return None
    try:
        root = Path(os.path.expandvars(raw)).expanduser()
        if not root.is_absolute():
            root = PROJECT_ROOT / root
        root = root.resolve()
    except Exception:
        return None
    if not root.is_dir():
        return None
    if ip and _ssot_export_valid_ip(ip):
        if root.name == ip:
            return root
        for name in (f"{ip}.ssot.yaml", f"{ip}_ssot.yaml", f"{ip}.ssot.yml"):
            if (root / "yaml" / name).is_file():
                return root
        return None
    return root


def _project_ip_root(ip: str) -> Path:
    configured = _configured_ip_root(ip)
    if configured is not None:
        return configured
    nested = PROJECT_ROOT / ip
    if nested.exists():
        return nested
    if PROJECT_ROOT.name == ip and (PROJECT_ROOT / "yaml").is_dir():
        return PROJECT_ROOT
    return nested


def _ssot_yaml_path(ip: str) -> Path:
    if not _ssot_export_valid_ip(ip):
        raise ValueError(f"invalid ip name {ip!r}")
    ip_dir = _project_ip_root(ip)
    for name in (f"{ip}.ssot.yaml", f"{ip}_ssot.yaml", f"{ip}.ssot.yml"):
        candidate = ip_dir / "yaml" / name
        if candidate.is_file():
            return candidate
    return ip_dir / "yaml" / f"{ip}.ssot.yaml"


def _load_ssot_yaml(ip: str) -> dict:
    import yaml as _yaml  # type: ignore

    path = _ssot_yaml_path(ip)
    if not path.is_file():
        raise FileNotFoundError(str(path))
    try:
        text = path.read_text(encoding="utf-8")
        data = _yaml.safe_load(text)
    except Exception as exc:
        raise ValueError(f"invalid yaml: {exc}") from exc
    if data is None:
        return {}
    if not isinstance(data, dict):
        raise ValueError("invalid yaml: top-level must be a mapping")
    return data


def create_app():
    try:
        from fastapi import FastAPI, Request, WebSocket, WebSocketDisconnect
        from fastapi.responses import FileResponse, HTMLResponse, JSONResponse, Response
        from fastapi.staticfiles import StaticFiles
        from starlette.routing import WebSocketRoute
    except ImportError:
        print("ERROR: fastapi not installed. Run: pip install fastapi uvicorn websockets")
        sys.exit(1)
    from core.atlas_db import AtlasDB
    from core.atlas_multiuser import _MultiUserBridge, set_atlas_bridge_session_id

    if not FRONTEND.exists():
        print(f"ERROR: frontend bundle not found at {FRONTEND}")
        sys.exit(1)

    # Load common_ai_agent/.config even when atlas_ui.py is launched directly
    # instead of through textual_main.py. Shell-set env still wins because
    # src.config only fills missing keys at first load.
    try:
        try:
            from . import config as _atlas_boot_config  # noqa: F401
        except Exception:
            import config as _atlas_boot_config  # type: ignore  # noqa: F401
    except Exception:
        pass

    app = FastAPI(title="ATLAS · common_ai_agent")
    _multi_raw = os.environ.get("ATLAS_MULTI_USER", "1").strip().lower()
    _multi_user_env = _multi_raw not in ("0", "false", "no", "off")
    # Multi-user and process-per-session isolation default on so agent output,
    # command results, and main.py global state do not bleed across users.
    # Operators can still opt out explicitly with ATLAS_MULTI_USER=0 or
    # ATLAS_MULTI_USER_PROC=0.
    _proc_raw = os.environ.get("ATLAS_MULTI_USER_PROC", "1").strip().lower()
    _use_proc = _multi_user_env and _proc_raw not in ("0", "false", "no", "off")
    _strict_raw = os.environ.get("ATLAS_STRICT_SESSION_ROUTING", "0").strip().lower()
    _strict_routing = _strict_raw in ("1", "true", "yes", "on")
    _exec_raw = (
        os.environ.get("ATLAS_EXEC_MODE")
        or os.environ.get("ATLAS_DEFAULT_EXEC_MODE")
        or ""
    ).strip().lower()
    _orch_raw = os.environ.get("ATLAS_ORCHESTRATOR_MODE")
    _single_loop = os.environ.get("ATLAS_SINGLE_MAIN_LOOP", "").strip().lower() in (
        "1", "true", "yes", "on",
    )
    _orch_disabled = (
        _orch_raw is not None
        and _orch_raw.strip().lower() in ("0", "false", "no", "off")
    )
    _single_worker_raw = (
        os.environ.get("ATLAS_SINGLE_WORKER_PER_OWNER")
        or os.environ.get("ATLAS_SINGLE_WORKER_PER_USER")
        or ""
    ).strip().lower()
    if _single_worker_raw:
        _single_worker_per_owner = _single_worker_raw in ("1", "true", "yes", "on")
    else:
        # Multi-user process mode means isolation is owned by the full
        # user/ip/workflow namespace.  Collapsing to one worker per owner
        # makes workflow tabs kill each other and can route browser prompts
        # into the wrong SQLite queue after reconnects.  Keep the legacy
        # owner-singleton behavior only behind the explicit env flag above.
        _single_worker_per_owner = False

    def _multi_user_enabled() -> bool:
        raw = os.environ.get("ATLAS_MULTI_USER", "1").strip().lower()
        return raw not in ("0", "false", "no", "off")

    def _normalize_atlas_exec_mode(value: Any) -> str:
        return normalize_exec_mode(value)

    def _current_atlas_exec_mode() -> str:
        return current_exec_mode(os.environ)

    def _new_ip_initial_workflow(explicit: Any = "", exec_mode: Any = "") -> str:
        return initial_workflow_for_exec_mode(
            _normalize_atlas_exec_mode(exec_mode) or _current_atlas_exec_mode(),
            normalize_session_name(str(explicit or "")),
        )

    if _multi_user_env:
        print(f"[atlas] Multi-user enabled (process_per_session={'on' if _use_proc else 'off'})")
    # single_user collapses every WS-bound session_id onto "default" so
    # the agent thread's inbox and the WS handler's inbox are the same.
    bridge = _MultiUserBridge(
        single_user=not _multi_user_env,
        use_processes=_use_proc,
        strict_session_routing=_strict_routing,
        single_worker_per_owner=_single_worker_per_owner,
    )
    add_event_handler = getattr(app, "add_event_handler", None)
    if not callable(add_event_handler):
        add_event_handler = getattr(getattr(app, "router", None), "add_event_handler", None)
    if callable(add_event_handler):
        add_event_handler("shutdown", bridge.stop_all_processes)
    else:
        app.on_event("shutdown")(bridge.stop_all_processes)
    # Register the bridge so the ReAct loop's orchestrator chat
    # injector (built lazily inside main.py / agent_server.py before
    # this point) can resolve sessions for the chat watermark.
    try:
        from core.orchestrator_inject import register_bridge as _orch_register_bridge
        _orch_register_bridge(bridge)
    except Exception:
        pass
    try:
        from src.orchestrator.react_bridge import register_live_event_emitter as _orch_live_events

        def _emit_orchestrator_live(session_id: str, event: dict[str, Any]) -> None:
            sid = normalize_session_name(str(session_id or ""))
            if not sid:
                return
            bridge.emit("orchestrator_chat", session_id=sid, **dict(event or {}))

        _orch_live_events(_emit_orchestrator_live)
    except Exception:
        pass
    # Opt-in auto-start of the chat-responder bots. With
    # CHAT_RESPONDER_AUTOSTART=1 in .env, atlas_ui spawns one daemon
    # thread per IP room plus one for _global, so launching the UI is
    # enough to make the bot respond to teammate feedback. Threads die
    # with the parent. Without the env var, users still launch
    # responders manually via `python3 -m core.chat_responder <room>`.
    if os.environ.get("CHAT_RESPONDER_AUTOSTART", "").strip() in ("1", "true", "yes", "on"):
        try:
            from core.chat_responder import autostart_all as _chat_autostart
            _chat_autostart()
        except Exception as _e:
            print(f"[chat-responder] autostart failed: {_e}")
    # Session Flow / runtime usage rollup scheduler (Task 7 / B2 / RS-1): the
    # REAL out-of-band production trigger that folds per-session runtime rows into
    # the control-side session_flow_rollups + runtime_usage_rollups. Without this
    # the dashboard reads empty/stale in session mode (the admin read is
    # no-fanout by design and never folds on-read). The scheduler is its own
    # guard: it starts ONLY in session mode (or when forced), and is test-disabled
    # unless ATLAS_FLOW_ROLLUP_ENABLE=1 so pytest never spins the daemon. Interval
    # is ATLAS_FLOW_ROLLUP_INTERVAL_S (default 30s). NOTE: the pre-existing
    # runtime_usage_rollups had the SAME missing-trigger gap; this scheduler wires
    # BOTH folds (run_rollup_pass) so both are populated.
    try:
        from core.runtime_rollup import start_rollup_scheduler as _start_flow_rollup
        _flow_rollup_thread = _start_flow_rollup()
        if _flow_rollup_thread is not None:
            print("[atlas] Session Flow rollup scheduler started")
    except Exception as _e:
        print(f"[atlas] flow rollup scheduler not started: {_e}")
    clients: set[Any] = set()
    broadcaster_task: asyncio.Task | None = None

    def _session_emit_target(client_session: Any | None) -> str | None:
        return getattr(client_session, "session_id", None) if client_session is not None else None

    def _queue_prompt_for_session(client_session: Any | None, text: str) -> None:
        if client_session is not None:
            bridge.queue_prompt_for_session(client_session.session_id, text)
            return
        bridge.queue_prompt(text)

    def _input_history_path() -> Path:
        try:
            try:
                from . import config as _config
            except Exception:
                import config as _config  # type: ignore
            base = str(getattr(_config, "SESSION_DIR", "") or "").strip()
        except Exception:
            base = ""
        return Path(base) / "input_history.txt" if base else PROJECT_ROOT / ".session" / "input_history.txt"

    def _read_input_history(limit: int = 200) -> list[str]:
        path = _input_history_path()
        if not path.is_file():
            return []
        entries: list[str] = []
        cur: list[str] = []
        try:
            for raw in path.read_text(encoding="utf-8", errors="replace").splitlines():
                if raw.startswith("+"):
                    cur.append(raw[1:])
                elif raw.startswith("#"):
                    if cur:
                        entries.append("\n".join(cur))
                        cur = []
            if cur:
                entries.append("\n".join(cur))
        except OSError:
            return []
        return [e for e in entries if e.strip()][-max(1, min(int(limit or 200), 1000)):]

    def _append_input_history(text: str) -> None:
        body = str(text or "").replace("\r\n", "\n").replace("\r", "\n").strip()
        if not body:
            return
        path = _input_history_path()
        path.parent.mkdir(parents=True, exist_ok=True)
        with path.open("a", encoding="utf-8") as fh:
            fh.write(f"\n# {time.ctime()}\n")
            for line in body.split("\n"):
                fh.write("+" + line + "\n")

    async def _send_one(client, raw: str, timeout: float):
        # raw is the already-serialized JSON text — broadcasters
        # serialize ONCE up front and pass the same string to every
        # client instead of re-serializing per WS. send_text avoids
        # FastAPI re-encoding the message a second time too. Scale
        # timeout with message size; tiny frames floor at 4 s.
        try:
            await asyncio.wait_for(client.send_text(raw), timeout=timeout)
            return None
        except Exception:
            return client

    def _is_websocket_disconnect(exc: BaseException) -> bool:
        """Return True only for normal client-side websocket disconnects."""
        seen: set[int] = set()
        cur: BaseException | None = exc
        while cur is not None and id(cur) not in seen:
            seen.add(id(cur))
            if _is_disconnect_os_error(cur):
                return True
            if isinstance(cur, WebSocketDisconnect):
                return True
            cls_name = cur.__class__.__name__
            if cls_name in {
                "ClientDisconnected",
                "ConnectionClosed",
                "ConnectionClosedError",
                "ConnectionClosedOK",
            }:
                return True
            if cls_name == "RuntimeError":
                msg = str(cur).lower()
                if "disconnect" in msg or "websocket is not connected" in msg:
                    return True
            if cls_name == "AttributeError" and "transfer_data_task" in str(cur):
                return True
            cur = getattr(cur, "__cause__", None) or getattr(cur, "__context__", None)
        return False

    async def _close_websocket_quietly(client, *, code: int, reason: str) -> None:
        try:
            await client.close(code=code, reason=reason)
        except Exception as exc:
            if not _is_websocket_disconnect(exc):
                raise

    async def _broadcast_outbox():
        """Single consumer for bridge events, fan out to the event session's WS clients.

        Each websocket used to start its own consumer on the same queue. With
        a browser tab plus an automation client, events were load-balanced
        between clients instead of broadcast, so ask_user cards and slash
        output could disappear from one surface.

        The whole loop body is wrapped in try/except so a transient lookup
        miss (e.g. ``get_session`` raising KeyError for a session that was
        deleted between event emit and broadcast) never kills the task —
        a dead broadcaster silently strands every connected WS client and
        looks to the frontend like "Backend disconnected" even though the
        agent process is still alive.
        """
        while True:
            try:
                msg, session_id = await bridge.next_event()
                if msg is None:
                    continue
                # _ensure_session is non-raising; get_session raises KeyError
                # for unknown ids which would propagate up and kill the task.
                try:
                    session = bridge.get_session(session_id)
                except Exception:
                    session = bridge._ensure_session(session_id)
                snapshot = list(session.clients)
                if not snapshot:
                    continue
                # Serialize once for the whole fan-out and skip the
                # second encode-to-bytes the timeout sizing used to do
                # for every frame — len(raw) is the same order of
                # magnitude as the UTF-8 byte length for ASCII-heavy
                # JSON, and the 4-second floor swallows the rounding
                # error.
                raw = json.dumps(msg, ensure_ascii=False)
                size_kb = max(len(raw) / 1024, 1)
                timeout = max(4.0, size_kb * 0.25)
                results = await asyncio.gather(
                    *(_send_one(c, raw, timeout) for c in snapshot),
                    return_exceptions=True,
                )
                for stale_client in results:
                    if stale_client is None or isinstance(stale_client, BaseException):
                        continue
                    session.clients.discard(stale_client)
            except asyncio.CancelledError:
                raise
            except Exception as _bcast_err:
                # Never let the broadcaster die. A killed broadcaster
                # produces the exact "WS open but no events flowing" symptom
                # users see as "Backend disconnected" with the chat empty.
                import sys as _sys
                print(f"[broadcaster] swallowed: {_bcast_err!r}", file=_sys.stderr)
                continue

    def _ensure_broadcaster() -> None:
        nonlocal broadcaster_task
        if broadcaster_task is None or broadcaster_task.done():
            broadcaster_task = asyncio.create_task(_broadcast_outbox())

    # ── inlined-HTML cache ────────────────────────────────────────
    # The browser fetches /  on every reload — the 14 .jsx files we
    # inline are large (~1.1 MB total) and each read_text() is sync,
    # so doing it from inside `async def index()` blocks the event
    # loop and starves every other request (api, static assets, ws).
    # Cache the assembled HTML and rebuild only when any inlined
    # source file's mtime changes.
    _INLINE_INDEX_RE = re.compile(
        r'<script\s+type="text/babel"\s+(?P<attrs>[^>]*?)src="(?P<src>[^"]+)"[^>]*>\s*</script>'
    )
    _inline_cache: dict[str, dict[str, Any]] = {}

    def _scm_ui_override_ref() -> str:
        provider = configured_scm_provider()
        suffix = provider.upper()
        return (
            os.environ.get(f"ATLAS_SCM_UI_OVERRIDE_{suffix}", "").strip()
            or os.environ.get(f"ATLAS_{suffix}_SCM_UI_OVERRIDE", "").strip()
            or os.environ.get("ATLAS_SCM_UI_OVERRIDE", "").strip()
        )

    def _scm_ui_override_is_url(ref: str) -> bool:
        return bool(re.match(r"^https?://", str(ref or ""), re.I))

    def _scm_ui_override_path(ref: str) -> Path:
        path = Path(ref).expanduser()
        if not path.is_absolute():
            path = SOURCE_ROOT / path
        return path.resolve()

    def _scm_ui_override_local_path() -> Path | None:
        ref = _scm_ui_override_ref()
        if not ref or _scm_ui_override_is_url(ref):
            return None
        return _scm_ui_override_path(ref)

    def _scm_ui_override_version() -> str:
        path = _scm_ui_override_local_path()
        if path is not None:
            try:
                st = path.stat()
                return str(st.st_mtime_ns)
            except OSError:
                return str(int(time.time()))
        return str(int(time.time()))

    def _scm_ui_override_script_tag() -> str:
        ref = _scm_ui_override_ref()
        if not ref:
            return ""
        if _scm_ui_override_is_url(ref):
            src = ref
        else:
            src = f"/api/scm/ui/override.js?v={_scm_ui_override_version()}"
        return (
            '<script type="text/babel" data-presets="react" '
            'data-atlas-scm-ui-override="1" '
            f'src="{html_lib.escape(src, quote=True)}"></script>'
        )

    def _inject_scm_ui_override(html: str) -> str:
        tag = _scm_ui_override_script_tag()
        if not tag:
            return html
        marker = '<script type="text/babel" data-filename="workspace.jsx"'
        if marker in html:
            return html.replace(marker, tag + "\n" + marker, 1)
        return html.replace("</body>", tag + "\n</body>", 1)

    # 2026-05-30 retirement: the legacy .jsx inliner helpers
    # (_inline_html_cached / _html_with_atlas_boot_config) were removed when
    # index() and lobby() went VITE-ONLY. _inject_scm_ui_override (above) is
    # retained — it is still used by _vite_index_html below.

    def _vite_index_html(filename: str = "index.vite.html") -> str | None:
        """Serve a built Vite HTML document when ATLAS_FRONTEND_MODE=vite.

        Reads FRONTEND/dist/<filename> fresh on each call (the file is tiny)
        and injects the ATLAS_BOOT_CONFIG <script>.  Used for both the main
        app (index.vite.html) and the lobby (lobby.vite.html).
        Returns None when the built dist file is absent (or the dist is
        partial) so the caller can serve a clean "build missing" 503.
        """
        dist_index = FRONTEND / "dist" / filename
        if not dist_index.is_file():
            return None
        # QA #5: a PARTIAL dist (the .html present but assets/ missing or
        # empty — e.g. an interrupted build) would serve an HTML shell that
        # references non-existent hashed bundles → a blank page. Require a
        # populated assets/ dir; otherwise treat the dist as missing.
        dist_assets = FRONTEND / "dist" / "assets"
        if not dist_assets.is_dir() or not any(dist_assets.glob("*.js")):
            return None
        html = dist_index.read_text(encoding="utf-8")
        exec_mode = _current_atlas_exec_mode()
        policy = exec_policy_payload(exec_mode, env=os.environ)
        payload = {
            "run_mode": os.environ.get("ATLAS_RUN_MODE", "engineering"),
            "exec_mode": exec_mode,
            "exec_policy": policy,
            "multi_user": os.environ.get("ATLAS_MULTI_USER", "1"),
            "multi_user_proc": os.environ.get("ATLAS_MULTI_USER_PROC", "1"),
            "scm_provider": configured_scm_provider(),
            "scm_ui_override": bool(_scm_ui_override_ref()),
        }
        script = (
            "<script>window.ATLAS_BOOT_CONFIG="
            + json.dumps(payload, separators=(",", ":"))
            + ";window.ATLAS_DEFAULT_RUN_MODE=window.ATLAS_BOOT_CONFIG.run_mode;"
            + "window.ATLAS_DEFAULT_EXEC_MODE=window.ATLAS_BOOT_CONFIG.exec_mode;</script>"
        )
        html = html.replace("</head>", script + "\n</head>", 1)
        return _inject_scm_ui_override(html)

    _asset_cache: dict[str, dict[str, Any]] = {}

    def _cached_frontend_asset_response(rel_path: str) -> Response:
        """Serve critical frontend assets from memory.

        Starlette StaticFiles streams large files in chunks.  On the shared
        ATLAS backend process that can make first paint hang on vendor Babel
        long enough that the browser stays blank and the chat textarea never
        mounts.  These vendor files are small enough to cache as bytes and
        return in one response.
        """
        clean = rel_path.replace("\\", "/").lstrip("/")
        if ".." in clean.split("/"):
            return JSONResponse({"ok": False, "error": "invalid asset path"}, status_code=404)
        path = (FRONTEND / clean).resolve()
        try:
            path.relative_to(FRONTEND.resolve())
        except Exception:
            return JSONResponse({"ok": False, "error": "invalid asset path"}, status_code=404)
        if not path.is_file():
            return JSONResponse({"ok": False, "error": "asset not found"}, status_code=404)

        stat = path.stat()
        stamp = (stat.st_mtime_ns, stat.st_size)
        cached = _asset_cache.get(clean)
        if not cached or cached.get("stamp") != stamp:
            cached = {"stamp": stamp, "body": path.read_bytes()}
            _asset_cache[clean] = cached

        ext = path.suffix.lower()
        media_type = _MIME_OVERRIDES.get(ext)
        if not media_type:
            media_type = _mimetypes.guess_type(str(path))[0] or "application/octet-stream"
        return Response(
            content=cached["body"],
            media_type=media_type,
            headers={"Cache-Control": "public, max-age=31536000, immutable"},
        )

    @app.get("/")
    async def index():
        """Serve the built Vite index (VITE-ONLY).

        2026-05-30: the legacy .jsx frontend has been retired. index() now
        serves only the built Vite dist (frontend/atlas/dist/index.vite.html
        + dist/assets/). There is no legacy inliner fallback: when the dist
        is missing or partial (QA #5 guard inside _vite_index_html), serve a
        clean 503 telling the operator to build the frontend, rather than
        inlining now-deleted .jsx.
        """
        vite_html = _vite_index_html()
        if vite_html is not None:
            return HTMLResponse(vite_html)
        return HTMLResponse(
            "<!doctype html><meta charset=utf-8>"
            "<title>ATLAS frontend build missing</title>"
            "<body style=\"font-family:system-ui,sans-serif;padding:2rem\">"
            "<h1>ATLAS frontend build missing</h1>"
            "<p>Run <code>npm run build</code> in <code>frontend/atlas</code>.</p>",
            status_code=503,
        )

    @app.get("/vendor/{asset_path:path}")
    async def vendor_asset(asset_path: str):
        return _cached_frontend_asset_response(f"vendor/{asset_path}")

    @app.get("/api/scm/ui/override.js")
    async def api_scm_ui_override_js():
        ref = _scm_ui_override_ref()
        if not ref:
            return JSONResponse({"ok": False, "error": "SCM UI override is not configured"}, status_code=404)
        if _scm_ui_override_is_url(ref):
            return JSONResponse({"ok": False, "error": "SCM UI override is configured as a remote URL"}, status_code=404)
        path = _scm_ui_override_path(ref)
        if path.suffix.lower() not in (".js", ".jsx"):
            return JSONResponse({"ok": False, "error": "SCM UI override must be .js or .jsx"}, status_code=400)
        if not path.is_file():
            return JSONResponse({"ok": False, "error": "SCM UI override file not found"}, status_code=404)
        try:
            body = path.read_text(encoding="utf-8")
        except OSError as exc:
            return JSONResponse({"ok": False, "error": str(exc)}, status_code=500)
        return Response(
            content=body,
            media_type="text/babel; charset=utf-8",
            headers={"Cache-Control": "no-store, max-age=0"},
        )

    @app.get("/lobby")
    async def lobby():
        # VITE-ONLY (2026-05-30 retirement): serve the built lobby.vite.html
        # via the shared _vite_index_html reader; same dist-missing 503 as index().
        vite_html = _vite_index_html("lobby.vite.html")
        if vite_html is not None:
            return HTMLResponse(vite_html)
        return HTMLResponse(
            "<!doctype html><meta charset=utf-8>"
            "<title>ATLAS frontend build missing</title>"
            "<body style=\"font-family:system-ui,sans-serif;padding:2rem\">"
            "<h1>ATLAS frontend build missing</h1>"
            "<p>Run <code>npm run build</code> in <code>frontend/atlas</code>.</p>",
            status_code=503,
        )

    # Per-process startup epoch — bumps every time the backend
    # restarts. Surfaced in /api/version so the frontend polling loop
    # triggers a reload not only on frontend file edits (mtime track)
    # but also on a pure backend restart, even when no .jsx changed.
    _BACKEND_STARTED_AT = time.time()

    @app.get("/api/version")
    async def api_version():
        """Returns two pieces of "should the browser reload?" data:

        - `mtime`   — latest mtime across the frontend bundle (catches
                      .jsx / .js / .html edits)
        - `started` — backend process start epoch (catches a pure
                      backend reboot when only .py changed)

        The browser polls every few seconds and reloads when EITHER
        value bumps, so backend-reboot and frontend-reload paths both
        propagate without manual cache clears.
        """
        latest = 0.0
        try:
            for f in FRONTEND.iterdir():
                if f.is_file():
                    latest = max(latest, f.stat().st_mtime)
        except OSError:
            pass
        override_path = _scm_ui_override_local_path()
        if override_path is not None:
            try:
                latest = max(latest, override_path.stat().st_mtime)
            except OSError:
                pass
        return JSONResponse({"mtime": latest, "started": _BACKEND_STARTED_AT})

    @app.get("/api/pdk/status")
    async def api_pdk_status():
        """Expose the resolved PDK/liberty paths used by Python-launched flows."""
        try:
            try:
                import src.config as _cfg
            except Exception:
                import config as _cfg  # type: ignore
            try:
                _cfg.reload_env()
            except Exception:
                pass
            status_fn = getattr(_cfg, "pdk_status", None)
            if callable(status_fn):
                return JSONResponse(status_fn())
        except Exception as e:
            return JSONResponse({"ok": False, "error": str(e)}, status_code=500)
        return JSONResponse({"ok": False, "error": "config.pdk_status unavailable"}, status_code=500)

    @app.get("/api/llm/ping")
    async def api_llm_ping():
        """Cheap provider reachability probe (no token spend).

        Hits `${BASE_URL}/models` with the configured bearer key.
        200 → provider URL + auth + network all OK.
        4xx → auth or path issue (key expired, wrong base_url).
        5xx → provider outage.
        timeout → network / DNS / firewall.

        Used by the frontend boot handshake to surface LLM-side problems
        before the user types their first prompt. Does NOT validate the
        model can actually generate — that needs a real completion call,
        which costs tokens. /models is a list endpoint, 0 tokens.
        """
        try:
            import src.config as _cfg
        except Exception:
            try:
                import config as _cfg  # type: ignore
            except Exception:
                return JSONResponse({"ok": False, "error": "config not loadable"},
                                    status_code=500)
        base = (getattr(_cfg, "BASE_URL", "") or "").rstrip("/")
        api_key = getattr(_cfg, "API_KEY", "") or ""
        if not base:
            return JSONResponse({"ok": False, "error": "BASE_URL not configured"},
                                status_code=500)
        is_codex_backend = "chatgpt.com/backend-api/codex" in base.lower()
        url = base + ("/models?client_version=0.0.0" if is_codex_backend else "/models")

        def _probe():
            import http.client as _hc
            from urllib.parse import urlparse
            u = urlparse(url)
            host, port = u.hostname, (u.port or (443 if u.scheme == "https" else 80))
            conn_cls = _hc.HTTPSConnection if u.scheme == "https" else _hc.HTTPConnection
            conn = conn_cls(host, port, timeout=4)
            try:
                headers = {"Authorization": f"Bearer {api_key}"} if api_key else {}
                if is_codex_backend:
                    try:
                        from src.llm_client import build_api_headers as _build_api_headers
                    except Exception:
                        from llm_client import build_api_headers as _build_api_headers  # type: ignore
                    headers = _build_api_headers(api_key)
                conn.request("GET", u.path + (("?" + u.query) if u.query else ""),
                             headers=headers)
                resp = conn.getresponse()
                body = resp.read(2048).decode("utf-8", errors="replace")
                return resp.status, body
            finally:
                try: conn.close()
                except Exception: pass

        try:
            status, body = await asyncio.to_thread(_probe)
            # 2xx → fully OK.
            # 401/403 → real auth failure (key bad / expired).
            # 400 → server reachable but rejects the bare /models GET.
            #       Codex OAuth backend (chatgpt.com/backend-api/codex)
            #       demands a `client_version` query param and replies
            #       400 without it; the chat completions path still
            #       works fine. TCP + DNS + TLS all proved alive, so
            #       count as reachable.
            # 404 → endpoint /models not supported, but the server
            #       responded → network + DNS OK. Several providers
            #       (some Azure deployments) don't expose /models;
            #       treat as "reachable" too.
            # 5xx / other → upstream problem.
            ok = (200 <= status < 300) or status in (400, 404)
            reason = (
                "ok" if 200 <= status < 300
                else "auth failed" if status in (401, 403)
                else "server reachable, /models rejected the bare GET" if status == 400
                else "endpoint not exposed but server reachable" if status == 404
                else f"http {status}"
            )
            return JSONResponse({
                "ok": ok,
                "status": status,
                "reason": reason,
                "base_url": base,
                "provider": getattr(_cfg, "LLM_PROVIDER", ""),
                "model": getattr(_cfg, "MODEL", ""),
                "preview": body[:240],
            }, status_code=200 if ok else 502)
        except Exception as e:
            return JSONResponse({"ok": False, "error": str(e), "base_url": base},
                                status_code=502)

    def _request_active_session_for_user(request: Request) -> str:
        user = request.scope.get("user") or {}
        username = normalize_session_name(str(user.get("username") or ""))
        if username:
            try:
                active = normalize_session_name(bridge.active_session_for_owner(username))
            except Exception:
                active = ""
            return active or f"{username}/default/default/default"
        return normalize_session_name(_active_session_value() or "")

    @app.post("/api/control/stop")
    async def api_control_stop(request: Request):
        """HTTP fallback for the UI Stop button and Escape key.

        The primary control plane is the WebSocket, but control buttons
        should still work when the WS is reconnecting or its outbound queue
        is wedged behind a larger message.
        """
        target_session = _request_active_session_for_user(request)
        bridge.request_stop_for_session(target_session)
        try:
            session = bridge._ensure_session(target_session)
            session.agent_running = False
        except Exception:
            pass
        bridge.emit("agent_state", running=False, session_id=target_session)
        return JSONResponse({"ok": True, "action": "stop", "session_id": target_session})

    @app.post("/api/control/shutdown")
    async def api_control_shutdown(request: Request):
        """HTTP fallback for the UI Exit button.

        Exit terminates the active session worker only. Atlas UI is the
        backend server for every browser/user, so it must stay alive.
        """
        target_session = _request_active_session_for_user(request)
        bridge.exit_session(target_session)
        return JSONResponse({"ok": True, "action": "exit_session", "session_id": target_session})

    @app.post("/api/settings/reasoning-effort")
    async def api_settings_reasoning_effort(request: Request):
        try:
            body = await request.json()
        except Exception:
            body = {}
        try:
            effort = _normalize_reasoning_effort(
                body.get("effort") or body.get("reasoning_effort")
            )
        except ValueError as exc:
            return JSONResponse({
                "ok": False,
                "error": str(exc),
                "allowed": list(_REASONING_EFFORT_OPTIONS),
            }, status_code=400)
        try:
            _persist_reasoning_effort(effort)
            _refresh_config_after_persist()
            _set_runtime_reasoning_effort(effort)
            bridge.emit(
                "context",
                reasoning_effort=effort,
                session_id=_request_active_session_for_user(request),
            )
            return JSONResponse({"ok": True, "reasoning_effort": effort})
        except Exception as exc:
            return JSONResponse({"ok": False, "error": str(exc)}, status_code=500)

    @app.post("/api/settings/model")
    async def api_settings_model(request: Request):
        try:
            body = await request.json()
        except Exception:
            body = {}

        try:
            import src.config as _cfg_model  # noqa: WPS433
        except Exception:
            try: import config as _cfg_model  # noqa: WPS433
            except Exception: _cfg_model = None
        if _cfg_model is not None:
            try:
                _cfg_model.reload_env()
            except Exception:
                pass

        active = getattr(_cfg_model, "MODEL_NAME", "") if _cfg_model is not None else os.environ.get("LLM_MODEL_NAME", "")
        options = _model_option_rows(active)
        model_key = str(body.get("key") or body.get("model_key") or "").strip()
        requested_model = str(body.get("model") or "").strip()

        selected = None
        if model_key:
            selected = next((row for row in options if row["key"] == model_key), None)
        if selected is None and requested_model:
            selected = next((row for row in options if row["model"] == requested_model), None)
        if selected is None:
            return JSONResponse({
                "ok": False,
                "error": "unknown or empty model option",
                "model_options": options,
            }, status_code=400)
        if selected.get("runtime") == "true":
            return JSONResponse({
                "ok": True,
                "model": selected["model"],
                "selected_model_key": selected["key"],
                "model_options": options,
            })

        model = selected["model"]
        try:
            persist_updates = {
                "LLM_SELECTED_MODEL_KEY": selected["key"],
                "LLM_ACTIVE_MODEL_NAME": model,
                "LLM_ACTIVE_BASE_NAME": model,
            }
            if _is_model_slot_key(selected["key"]):
                persist_updates[selected["key"]] = model
            profile_name = _profile_name_from_option_key(selected["key"])
            if profile_name:
                persist_updates["LLM_PROFILE"] = profile_name
            _persist_env_values(persist_updates)
            _refresh_config_after_persist()
            _set_runtime_model(model, selected["key"])
            active_model = model
            if _cfg_model is not None:
                active_model = str(getattr(_cfg_model, "MODEL_NAME", "") or active_model)
            active_model = os.environ.get("LLM_MODEL_NAME", active_model) or active_model
            _persist_env_values({
                "LLM_ACTIVE_MODEL_NAME": active_model,
                "LLM_ACTIVE_BASE_NAME": active_model,
            })
            updated_options = _model_option_rows(active_model)
            updated_selected_key = next(
                (row["key"] for row in updated_options if row.get("selected") == "true"),
                selected["key"],
            )
            bridge.emit(
                "context",
                model=active_model,
                model_options=updated_options,
                selected_model_key=updated_selected_key,
                session_id=_request_active_session_for_user(request),
            )
            return JSONResponse({
                "ok": True,
                "model": active_model,
                "selected_model_key": updated_selected_key,
                "model_options": updated_options,
            })
        except Exception as exc:
            return JSONResponse({"ok": False, "error": str(exc)}, status_code=500)

    def _atlas_root_for_context() -> Path:
        return Path(os.environ.get("ATLAS_ROOT") or str(PROJECT_ROOT)).expanduser().resolve()

    def _validated_context_workspace_root(context: AtlasContext) -> tuple[Path | None, str]:
        if context.legacy:
            try:
                return context.workspace_root.resolve(), ""
            except OSError as exc:
                return None, str(exc)
        root = Path(context.atlas_root).expanduser().resolve()
        owner_root = root / context.user_name
        workspace_root = owner_root / context.workspace_session
        try:
            if owner_root.is_symlink() or workspace_root.is_symlink():
                return None, "workspace session path is a symlink"
            resolved = workspace_root.resolve(strict=False)
            resolved.relative_to(root)
        except ValueError:
            return None, "workspace session path escapes atlas root"
        except OSError as exc:
            return None, str(exc)
        return resolved, ""

    def _validated_workspace_root(
        owner_name: str,
        workspace_session: str,
        atlas_root: Path | str | None = None,
    ) -> tuple[Path | None, str]:
        try:
            context = AtlasContext(
                user_name=owner_name,
                workspace_session=workspace_session or "default",
                ip_name="default",
                workflow="default",
                atlas_root=atlas_root or _atlas_root_for_context(),
            )
        except ValueError as exc:
            return None, str(exc)
        return _validated_context_workspace_root(context)

    @app.get("/healthz")
    async def healthz(request: Request):
        info = {
            "ok": True,
            "frontend": str(FRONTEND),
            "source_root":  str(SOURCE_ROOT),     # where atlas_ui.py lives
            "project_root": str(PROJECT_ROOT),    # = user's cwd at launch
            "cwd": os.getcwd(),
        }
        # Identity is now derived from the authenticated user (cookie).
        # Single-user vs multi-user only affects whether multiple distinct
        # users can run concurrently — login is required either way.
        _multi_user_on = _multi_user_enabled()
        info["multi_user"] = _multi_user_on
        user = request.scope.get("user")
        username = user.get("username") if user else None
        info["user_session"] = username
        username_norm = normalize_session_name(str(username or ""))
        include_cost = str(request.query_params.get("cost", "1") or "1").strip().lower() not in {
            "0", "false", "no", "lite",
        }
        query_active_session = normalize_session_name(
            str(
                request.query_params.get("session_id")
                or request.query_params.get("session")
                or ""
            )
        )
        request_active_session = ""
        if username_norm:
            try:
                request_active_session = (
                    bridge.active_session_for_owner(username_norm)
                    or bridge.active_session_for_owner(_session_owner_with_model(username_norm))
                )
            except Exception:
                request_active_session = ""
        client_host = (request.client.host if request.client else "") or "127.0.0.1"
        if client_host.startswith("::ffff:"):
            client_host = client_host[7:]
        info["client_ip"] = client_host
        info["project_root_name"] = PROJECT_ROOT.name or ""
        # Expose the real model + context window so the sidebar doesn't
        # have to invent values. Pull from src.config (the per-process
        # frozen settings); if config isn't importable yet, fall through
        # to env vars.
        try:
            import src.config as _cfg  # noqa: WPS433
        except Exception:
            try: import config as _cfg  # noqa: WPS433
            except Exception: _cfg = None
        # Pick up any .env edits made while the server has been running so
        # the sidebar (and dispatch on the next call) reflect the latest
        # active model / profile without a restart. mtime-cached → cheap.
        if _cfg is not None:
            try:
                _cfg.reload_env()
            except Exception:
                pass
        if _cfg is not None:
            model = ""
            try:
                from src.llm_client import get_active_model
                model = get_active_model() or ""
            except Exception:
                pass
            if not model:
                model = (
                    getattr(_cfg, "MODEL_NAME", None)
                    or getattr(_cfg, "PRIMARY_MODEL", None)
                    or getattr(_cfg, "LLM_MODEL_NAME", "")
                )
            info["model"] = model
            info["model_options"] = _model_option_rows(model)
            info["selected_model_key"] = next(
                (row["key"] for row in info["model_options"] if row.get("selected") == "true"),
                "",
            )
            # Use the active dispatch model as the primary display value.
            # PRIMARY_MODEL can remain stale after --model/profile overrides
            # (for example glm in .env while --model deepseek is active),
            # which made the ATLAS sidebar look like it was calling the wrong
            # backend even though dispatch was already using MODEL_NAME.
            info["base_model"] = model or getattr(_cfg, "MODEL_NAME", "")
            if getattr(_cfg, "CURSOR_AGENT_ENABLE", False):
                info["base_url"] = "cursor-agent"
            elif getattr(_cfg, "CLAUDE_CLI_ENABLE", False):
                info["base_url"] = "claude-cli"
            else:
                info["base_url"] = getattr(_cfg, "BASE_URL", "")
            info["provider"] = getattr(_cfg, "LLM_PROVIDER", "")
            info["max_context"] = getattr(_cfg, "MAX_CONTEXT_TOKENS", 0)
            info["max_iterations"] = getattr(_cfg, "MAX_ITERATIONS", 0)
            # Surface the active reasoning effort/mode so the ATLAS sidebar
            # mirrors what textual_main.py shows in its model line.
            info["reasoning_effort"] = (
                getattr(_cfg, "REASONING_MODE", "")
                or getattr(_cfg, "REASONING_EFFORT", "")
                or os.environ.get("REASONING_EFFORT", "")
                or os.environ.get("REASONING_MODE", "")
                or ""
            )
            info["chat_feed_summary"] = bool(getattr(_cfg, "ATLAS_CHAT_FEED_SUMMARY", True))
            # Resolve the "active session" the user is looking at. When
            # the agent boots WITHOUT -w, ACTIVE_WORKSPACE is unset but
            # the session still maps to .session/default/. Prefer the
            # explicit workspace; fall back to the actual project the
            # session loader is using (config.ACTIVE_PROJECT) so the
            # UI can show "default" instead of an ambiguous "—".
            info["workspace"] = (os.environ.get("ACTIVE_WORKSPACE")
                                  or os.environ.get("WORKSPACE")
                                  or getattr(_cfg, "ACTIVE_PROJECT", "")
                                  or "default")
            # Canonical (session_id, ip, workflow) triple — the single
            # source of truth for which IP the user is editing. Updated
            # any time _set_active_ssot_ip runs (e.g. /new-ip <name>),
            # so polling /healthz from the frontend is enough to keep
            # the preview / SSOT / QA panels in sync without a custom
            # WS event.
            healthz_owner_guard = username_norm if _multi_user_on else ""

            def _owned_healthz_session(raw: str) -> str:
                normalized = normalize_session_name(str(raw or ""))
                if not normalized:
                    return ""
                if _multi_user_on and not healthz_owner_guard:
                    return ""
                if not healthz_owner_guard:
                    return normalized
                owner = normalized.split("/", 1)[0]
                allowed = {
                    healthz_owner_guard,
                    _session_owner_with_model(healthz_owner_guard),
                }
                return normalized if owner in allowed else ""

            if username_norm or (query_active_session and not _multi_user_on):
                active_session = (
                    _owned_healthz_session(query_active_session)
                    or _owned_healthz_session(request_active_session)
                    or _owned_healthz_session(_active_session_value())
                )
                if not active_session:
                    active_session = f"{_session_owner_with_model(username_norm)}/default/default/default"
                info["active_session"] = active_session
            else:
                info["active_session"] = (
                    request_active_session
                    or _active_session_value()
                    or _canonical_session_string()
                )
            active_session_path = normalize_session_name(str(info.get("active_session") or ""))
            active_context: AtlasContext | None = None
            if active_session_path:
                try:
                    active_context = AtlasContext.from_session_key(
                        active_session_path,
                        atlas_root=_atlas_root_for_context(),
                    )
                    workspace_root, workspace_error = _validated_context_workspace_root(active_context)
                    if workspace_error or workspace_root is None:
                        return JSONResponse({
                            "ok": False,
                            "error": workspace_error or "invalid workspace root",
                            "active_session": active_session_path,
                        }, status_code=400)
                    session_dir = active_context.session_dir
                    try:
                        if session_dir.is_symlink():
                            return JSONResponse({
                                "ok": False,
                                "error": "session path is a symlink",
                                "active_session": active_session_path,
                            }, status_code=400)
                        session_dir.resolve(strict=False).relative_to(workspace_root)
                    except ValueError:
                        return JSONResponse({
                            "ok": False,
                            "error": "session path escapes workspace root",
                            "active_session": active_session_path,
                        }, status_code=400)
                    except OSError as exc:
                        return JSONResponse({
                            "ok": False,
                            "error": str(exc),
                            "active_session": active_session_path,
                        }, status_code=400)
                    info["active_session"] = active_context.active_session_key
                    info["context_key"] = active_context.context_key
                    info["workspace_session"] = active_context.workspace_session
                    info["active_ip"] = active_context.ip_name
                    info["active_workflow"] = active_context.workflow
                    info["atlas_root"] = str(active_context.atlas_root)
                    info["workspace_root"] = str(workspace_root)
                    info["ip_root"] = str(workspace_root / active_context.ip_name)
                    if not active_context.legacy:
                        info["backend_project_root"] = str(PROJECT_ROOT)
                        info["project_root"] = str(workspace_root)
                        info["project_root_name"] = workspace_root.name or ""
                except Exception:
                    active_context = None
            if active_context is None:
                _active_parts = [
                    part for part in normalize_session_name(
                        str(info.get("active_session") or "")
                    ).split("/") if part
                ]
                info["active_ip"] = (
                    _active_parts[1]
                    if len(_active_parts) >= 2 and _active_parts[1]
                    else (_active_ip_value() or "default")
                )
                info["active_workflow"] = (
                    _active_parts[2]
                    if len(_active_parts) >= 3 and _active_parts[2]
                    else (os.environ.get("ATLAS_DEFAULT_WORKFLOW") or "default")
                )
            info["db_session_id"] = ""
            info["session_uid"] = ""
            info["session_label"] = ""
            if include_cost and active_session_path and user:
                try:
                    with AtlasDB() as _db:
                        _row = _db.get_session_for_user(str(user.get("id") or ""), active_session_path)
                    if _row:
                        info["db_session_id"] = str(_row.get("id") or "")
                        info["session_uid"] = str(_row.get("session_uid") or "")
                        _uid = info["session_uid"]
                        info["session_label"] = f"S-{_uid[:8]}" if _uid else active_session_path
                except Exception:
                    pass
            if active_session_path:
                session_dir = (
                    active_context.session_dir
                    if active_context is not None
                    else PROJECT_ROOT / ".session" / active_session_path
                )
                info["session_dir"] = str(session_dir)
                info["todo_file"] = str(session_dir / "todo.json")
                info["history_file"] = str(session_dir / "conversation.json")
            else:
                info["session_dir"] = str(getattr(_cfg, "SESSION_DIR", "") or "")
                info["todo_file"] = str(getattr(_cfg, "TODO_FILE", "") or "")
                info["history_file"] = str(getattr(_cfg, "HISTORY_FILE", "") or "")
            try:
                if active_session_path:
                    _session_obj = bridge.get_session(active_session_path)
                    info["agent_running"] = bool(getattr(_session_obj, "agent_running", False))
                    info["agent_alive"] = bool(getattr(_session_obj, "agent_alive", False))
                else:
                    info["agent_running"] = bool(bridge.agent_running)
                    info["agent_alive"] = bool(bridge.agent_alive)
            except Exception:
                info["agent_running"] = bool(getattr(bridge, "agent_running", False))
                info["agent_alive"] = bool(getattr(bridge, "agent_alive", False))
            # Per-model pricing (USD / 1M tokens) — input / cache / output.
            # Try the displayed/runtime model first, then fall back to
            # LLM_BASE_NAME for opaque deployment aliases.
            info["pricing"] = None
            try:
                from lib.model_pricing import get_active_pricing
                p = get_active_pricing(str(info.get("model") or ""))
                if p is not None:
                    info["pricing"] = {
                        "input": p.input, "cache": p.cache, "output": p.output,
                    }
            except Exception:
                pass
            # Live worker context comes from the selected worker/session
            # cost.json. Cost rows below are overwritten from DB as a
            # user+IP aggregate so one IP's spend stays visible across
            # orchestrator/worker workflow switches.
            # Hot path: /healthz polls every few seconds and the disk
            # read used to block the asyncio loop. Push the read into
            # a thread; the surrounding logic stays cheap.
            worker_cost_seen = False
            try:
                import json as _json
                from pathlib import Path as _P
                _sess = str(PROJECT_ROOT)
                _sess_str = normalize_session_name(str(info.get("active_session") or "")).strip("/")
                _candidates = []
                _health_session_dir = str(info.get("session_dir") or "").strip()
                if _health_session_dir:
                    _candidates.append(_P(_health_session_dir) / "cost.json")
                if _sess_str:
                    # Canonical 3-part path:
                    # .session/<owner>/<ip>/<workflow>/cost.json. Do not
                    # fall back to workflow/default ledgers here; context and
                    # cost must follow the active namespace exactly.
                    _candidates.append(_P(_sess) / ".session" / _sess_str / "cost.json")

                def _read_cost_file(c):
                    try:
                        return _json.loads(c.read_text(encoding="utf-8", errors="replace"))
                    except Exception:
                        return None

                def _pick_cost():
                    for c in _candidates:
                        if c.exists():
                            d = _read_cost_file(c)
                            if d is not None:
                                return d
                    return None
                d = await asyncio.to_thread(_pick_cost)
                if d is not None:
                    worker_cost_seen = True
                    # cost.json schema written by both lib/textual_ui.py
                    # and atlas_ui._emit_token:
                    #   {in_tok, cache_tok, out_tok, sum_tok, cost_usd,
                    #    model, updated_at}
                    worker_tokens_in = d.get("in_tok",    d.get("input",  0))
                    worker_tokens_cache = d.get("cache_tok", d.get("cached", 0))
                    worker_tokens_out = d.get("out_tok",   d.get("output", 0))
                    info["worker_tokens_in"] = worker_tokens_in
                    info["worker_tokens_cache"] = worker_tokens_cache
                    info["worker_tokens_out"] = worker_tokens_out
                    # Current context window usage = the LAST turn's
                    # input tokens. Cumulative tokens make sense for cost
                    # ledger but not for "Context X / max" because a 200K
                    # window can host many turns whose sum exceeds max.
                    info["tokens"] = d.get("last_in_tok", 0)
                    # Prefer the cost_usd written by _emit_token (uses
                    # the per-call pricing at the moment of the LLM
                    # request, so model switches mid-session don't
                    # retroactively re-price old calls). Only recompute
                    # from the live pricing if cost_usd isn't on disk.
                    disk_cost = d.get("cost_usd")
                    if disk_cost is not None:
                        info["worker_cost_usd"] = float(disk_cost)
                    elif info["pricing"]:
                        ti = worker_tokens_in or 0
                        tc = worker_tokens_cache or 0
                        to = worker_tokens_out or 0
                        ti_billable = max(0, ti - tc)
                        info["worker_cost_usd"] = (
                            ti_billable * info["pricing"]["input"]  / 1_000_000
                            + tc        * info["pricing"]["cache"]  / 1_000_000
                            + to        * info["pricing"]["output"] / 1_000_000
                        )
                else:
                    info["worker_tokens_in"] = 0
                    info["worker_tokens_cache"] = 0
                    info["worker_tokens_out"] = 0
                    info["tokens"] = 0
                    info["worker_cost_usd"] = 0.0
            except Exception:
                pass
            info.setdefault("worker_tokens_in", 0)
            info.setdefault("worker_tokens_cache", 0)
            info.setdefault("worker_tokens_out", 0)
            info.setdefault("worker_cost_usd", 0.0)

            info["tokens_in"] = 0
            info["tokens_cache"] = 0
            info["tokens_out"] = 0
            info["cost_usd"] = 0.0
            info["cost_scope"] = "user_ip"
            info["cost_user"] = username_norm or ""
            info["cost_ip"] = str(info.get("active_ip") or "")
            info["cost_calls"] = 0
            try:
                active_ip_name = str(info.get("active_ip") or "").strip()
                if include_cost and active_ip_name and active_ip_name != "default":
                    db_path = os.environ.get("ATLAS_DB_PATH") or str(Path.home() / ".common_ai_agent" / "atlas.db")
                    user_id = str(user.get("id") or "") if user else ""
                    cache_key = (db_path, user_id, username_norm, active_ip_name, active_session_path)
                    usage = _healthz_cost_cache_get(cache_key)
                    if usage is None:
                        with AtlasDB() as _db:
                            usage = _db.summarize_llm_usage_for_user_ip(
                                user_id=user_id,
                                username=username_norm,
                                ip=active_ip_name,
                                active_session=active_session_path,
                            )
                        _healthz_cost_cache_set(cache_key, usage)
                    info["tokens_in"] = usage.get("tokens_input", 0)
                    info["tokens_cache"] = usage.get("cache_read_tokens", 0)
                    info["tokens_out"] = usage.get("tokens_output", 0)
                    info["cost_usd"] = usage.get("cost_usd", 0.0)
                    info["cost_calls"] = usage.get("calls", 0)
                elif include_cost and worker_cost_seen:
                    # No active IP means there is no user+IP aggregate key.
                    # Keep legacy local visibility for default sessions only.
                    info["cost_scope"] = "worker_session_fallback"
                    info["tokens_in"] = info.get("worker_tokens_in", 0)
                    info["tokens_cache"] = info.get("worker_tokens_cache", 0)
                    info["tokens_out"] = info.get("worker_tokens_out", 0)
                    info["cost_usd"] = info.get("worker_cost_usd", 0.0)
            except Exception:
                if worker_cost_seen:
                    info["cost_scope"] = "worker_session_fallback"
                    info["tokens_in"] = info.get("worker_tokens_in", 0)
                    info["tokens_cache"] = info.get("worker_tokens_cache", 0)
                    info["tokens_out"] = info.get("worker_tokens_out", 0)
                    info["cost_usd"] = info.get("worker_cost_usd", 0.0)
        else:
            import os as _os
            info["model"] = _os.environ.get("LLM_MODEL_NAME", "") or _os.environ.get("MODEL_NAME", "")
            info["model_options"] = _model_option_rows(info["model"])
            info["selected_model_key"] = next(
                (row["key"] for row in info["model_options"] if row.get("selected") == "true"),
                "",
            )
            info["max_context"] = int(_os.environ.get("MAX_CONTEXT_TOKENS", "0") or "0")
            info["max_iterations"] = int(_os.environ.get("MAX_ITERATIONS", "0") or "0")
            info["workspace"] = _os.environ.get("ACTIVE_WORKSPACE", "") or _os.environ.get("WORKSPACE", "")
        return JSONResponse(info)

    # ── REAL project data API ────────────────────────────────────
    # File-system backed endpoints. All paths are confined to the user's
    # PROJECT_ROOT (= cwd at launch, computed at module import) and
    # rejected if they try to escape via .. or absolute paths. This is
    # NOT the source repo — when the user runs:
    #   cd Custom_IP && python ../brian_hw/common_ai_agent/src/textual_main.py
    # the file API operates on Custom_IP, not on common_ai_agent/.
    # We intentionally re-bind here as a local var so the module-level
    # PROJECT_ROOT survives even if the import gets reloaded weirdly.
    import sys as _sys_local
    _PROJECT_ROOT = globals().get("PROJECT_ROOT") or Path(os.getcwd()).resolve()
    PROJECT_ROOT = _PROJECT_ROOT
    MAX_READ_BYTES = 256 * 1024
    SKIP_DIRS = {".git", "__pycache__", "node_modules", ".session",
                 "ATLAS", "vendor", ".venv", ".pytest_cache",
                 # Internal workflow scaffolding — generated for the LLM as
                 # authoring inputs, not user-facing artifacts. Keep on disk
                 # (workflow scripts still read them) but hide from the file
                 # tree so `rtl/` only surfaces real RTL (.sv/.v) plus the
                 # canonical compile/lint evidence JSONs.
                 "authoring_packets", "stage_engine"}

    # Files that are produced by workflow scripts for downstream-stage
    # consumption (TODO trackers, traceability maps, blocker reports,
    # authoring plans, gate-evidence packets, manifest bookkeeping). They
    # are *not* user-facing artifacts — the relevant data already shows
    # up in the Gates tab / TODO panel / compile/lint evidence cards.
    # Hide them in `/api/files` listings and exclude them from .ZIP
    # downloads. They remain on disk; workflow scripts and the
    # SSOT-gen blocker inline flow still read them directly.
    SKIP_FILES = {
        "manifest.json", "decomposition.json",
        "import_manifest.json", "ssot_downstream_blockers.json",
        "rtl_authoring_plan.json", "rtl_authoring_status.md",
        "rtl_blocked.json", "rtl_blocked_resolved.json",
        "rtl_todo_plan.json", "rtl_todo_tracker.json",
        "rtl_traceability.json",
    }

    def _is_internal_artifact(name: str) -> bool:
        if name in SKIP_FILES:
            return True
        if name.startswith("rtl_gate_") and (name.endswith(".json") or name.endswith(".md")):
            return True
        return False

    def _is_file_tree_hidden_artifact(child: Path, listing_root: Path) -> bool:
        if _is_internal_artifact(child.name):
            return True
        try:
            rel = child.relative_to(listing_root).as_posix().lower()
        except ValueError:
            rel = child.name.lower()
        # Imported PDF/Office image extraction creates low-value cache files
        # under req/imports/images. Markdown/DOC previews still reference them
        # through /api/file/raw, but listing them beside originals makes users
        # click placeholder masks/backgrounds instead of the source document.
        return bool(re.search(
            r"(^|/)req/imports/images/[^/]+\.(?:png|jpe?g|gif|webp|bmp|svg|tiff?|ico)$",
            rel,
        ))

    def _safe(rel_path):
        rel = (rel_path or "").lstrip("/")
        candidate = (PROJECT_ROOT / rel).resolve()
        try:
            candidate.relative_to(PROJECT_ROOT)
        except ValueError:
            return None
        return candidate



    def _safe_ip_file_delete_target(ip: str, path: str) -> tuple[Path | None, str | None]:
        clean_ip = str(ip or "").strip().strip("/")
        clean_path = str(path or "").strip().strip("/")
        if not clean_ip or not clean_path:
            return None, "ip and path are required"
        ip_parts = [part for part in clean_ip.split("/") if part]
        path_parts = [part for part in clean_path.split("/") if part]
        if not ip_parts or not path_parts:
            return None, "ip and path are required"
        if any(part in {".", ".."} for part in path_parts):
            return None, "invalid path"
        if any(part.startswith(".") for part in path_parts):
            return None, "hidden/internal files cannot be deleted from the UI"
        if any(part in {".", ".."} or not re.match(r"^[A-Za-z0-9_.-]+$", part) for part in ip_parts):
            return None, "invalid ip"
        if path_parts[:len(ip_parts)] != ip_parts:
            return None, "path is outside the selected IP"

        project_root = PROJECT_ROOT.resolve()
        ip_root = PROJECT_ROOT.joinpath(*ip_parts)
        candidate = PROJECT_ROOT.joinpath(*path_parts)
        try:
            resolved_ip = ip_root.resolve()
            resolved_target = candidate.resolve()
            resolved_ip.relative_to(project_root)
            resolved_target.relative_to(resolved_ip)
        except (OSError, ValueError):
            return None, "path outside project root"
        if resolved_target == resolved_ip:
            return None, "cannot delete the IP root"
        if not resolved_ip.is_dir():
            return None, "IP not found"
        if candidate.is_dir():
            return None, "directory delete is not supported from the UI"
        if not candidate.is_file():
            return None, "file not found"
        return candidate, None


    # ── Multi-tenant filesystem read/write authorization (audit finding B1) ──
    # Every content-serving endpoint (files, source, vcd, ssot, coverage, …)
    # confined its path/ip param to PROJECT_ROOT via _safe() but did NO owner
    # check, so any authenticated user could read any other user's IP source,
    # waveforms and .session conversation logs. This injected gate closes that
    # hole using the existing IP ACLs (can_user_access_ip) + the per-user
    # .session/<owner>/ namespace. Fail closed: anything we cannot positively
    # authorize is denied. A single shared AtlasDB handle (like _chat_db) backs
    # the lookups; the per-thread conn cache + RLock make it thread-safe.
    from types import SimpleNamespace as _AuthzNS  # noqa: WPS433
    from core.atlas_db import AtlasDB as _AuthzAtlasDB  # noqa: WPS433
    from core.atlas_auth import is_admin_user as _authz_is_admin  # noqa: WPS433
    from core import atlas_fs_authz as _fsz  # noqa: WPS433

    # A fresh AtlasDB() is opened per gate call (the per-thread conn cache makes
    # this cheap and matches /api/ip/list), so no long-lived shared handle.

    def _authz_identity(request):
        try:
            user = request.scope.get("user") or {}
        except Exception:
            user = {}
        uid = str((user or {}).get("id") or "default").strip() or "default"
        uname = normalize_session_name(str((user or {}).get("username") or ""))
        try:
            is_admin = bool(_authz_is_admin(user)) or (user or {}).get("role") == "admin"
        except Exception:
            is_admin = (user or {}).get("role") == "admin"
        return uid, uname, bool(is_admin)

    def _authz_owner_aliases(uname):
        """Owner-segment aliases: the bare username plus the per-model owner
        ("<user>__<model>") used when ATLAS_SESSION_PER_MODEL is on."""
        try:
            model_owner = _session_owner_with_model(uname)
        except Exception:
            model_owner = uname
        return {v for v in (uname, model_owner) if v}

    def _authz_owned_ips_db(db, uid, uname, is_admin):
        """IP names the user owns via their .session/db-sessions namespace.

        The authoritative multi-user IP-ownership source — the SAME logic
        /api/ip/list uses (db.list_sessions(user_id) -> namespaces whose owner
        segment is the user). None = no restriction (single-user/admin); fail
        closed (empty set) on any error.
        """
        if not _multi_user_enabled() or is_admin:
            return None
        if not uid or uid == "default" or not uname:
            return set()
        try:
            allowed_owners = _authz_owner_aliases(uname)
            owned = set()
            for row in db.list_sessions(uid) or []:
                ns = normalize_session_name(str(row.get("namespace") or row.get("id") or ""))
                parts = [p for p in ns.split("/") if p]
                if len(parts) >= 3 and parts[0] in allowed_owners:
                    owned.add(parts[-2])
            return owned
        except Exception:
            return set()

    def _authz_resolve_rel(rel_path):
        """Resolve symlinks so the gate authorizes the REAL target, not the
        requested name — a symlink inside an owned IP must not reach a victim
        IP. Returns (effective_rel, escaped)."""
        try:
            target = _safe(rel_path)
            if target is None:
                return rel_path, True
            return target.resolve().relative_to(PROJECT_ROOT.resolve()).as_posix(), False
        except Exception:
            return rel_path, False

    def _authz_path(request, rel_path, permission="view"):
        """Return a 401/403/400 JSONResponse if the read/write is not allowed, else None."""
        uid, uname, is_admin = _authz_identity(request)
        eff, escaped = _authz_resolve_rel(rel_path)
        if escaped:
            return JSONResponse({"error": "invalid path"}, status_code=400)
        with _AuthzAtlasDB() as db:
            decision = _fsz.authorize_path(
                eff, user_id=uid, username=uname,
                owner_aliases=_authz_owner_aliases(uname), is_admin=is_admin,
                multi_user=_multi_user_enabled(), db=db, permission=permission,
                owned_ips=_authz_owned_ips_db(db, uid, uname, is_admin),
            )
        if decision.allow:
            return None
        return JSONResponse({"error": decision.reason or "forbidden"}, status_code=decision.status)

    def _authz_ip(request, ip, permission="view"):
        uid, uname, is_admin = _authz_identity(request)
        with _AuthzAtlasDB() as db:
            decision = _fsz.authorize_ip(
                ip, user_id=uid, username=uname, is_admin=is_admin,
                multi_user=_multi_user_enabled(), db=db, permission=permission,
                owned_ips=_authz_owned_ips_db(db, uid, uname, is_admin),
            )
        if decision.allow:
            return None
        return JSONResponse({"error": decision.reason or "forbidden"}, status_code=decision.status)

    def _authz_accessible_ips(request):
        """Set of IP names the caller may see (for filtering listings), or None = no restriction."""
        uid, uname, is_admin = _authz_identity(request)
        with _AuthzAtlasDB() as db:
            return _fsz.accessible_ip_names(
                user_id=uid, is_admin=is_admin, multi_user=_multi_user_enabled(),
                db=db, permission="view",
                owned_ips=_authz_owned_ips_db(db, uid, uname, is_admin),
            )

    _fs_authz = _AuthzNS(
        path=_authz_path,
        ip=_authz_ip,
        accessible_ips=_authz_accessible_ips,
        shared_roots=_fsz.SHARED_ROOTS,
        shared_root_files=_fsz.SHARED_ROOT_FILES,
    )

    # Phase 9: file API cluster moved to src/atlas_api_files.py.
    # Factory pattern — closure captures become explicit kwargs.
    from src.atlas_api_files import register_file_routes as _register_file_routes
    _register_file_routes(
        app,
        safe_path_fn=_safe,
        project_root=PROJECT_ROOT,
        skip_dirs=SKIP_DIRS,
        is_hidden_artifact_fn=_is_file_tree_hidden_artifact,
        max_read_bytes=MAX_READ_BYTES,
        safe_ip_delete_fn=_safe_ip_file_delete_target,
        bridge=bridge,
        fs_authz=_fs_authz,
    )

    def _lint_ip_candidates(ip: str) -> list[Path]:
        clean = str(ip or "").strip().strip("/")
        if not clean:
            return []
        parts = [p for p in clean.split("/") if p]
        if any(p in {".", ".."} or not re.match(r"^[A-Za-z0-9_.-]+$", p) for p in parts):
            return []

        candidates: list[Path] = []
        if len(parts) > 1:
            target = _safe(clean)
            if target is not None:
                candidates.append(target)
        else:
            leaf = parts[0]
            direct = _safe(leaf)
            if direct is not None:
                candidates.append(direct)
            for pattern in (f"*/{leaf}", f"*/*/{leaf}"):
                for match in PROJECT_ROOT.glob(pattern):
                    try:
                        match.relative_to(PROJECT_ROOT)
                    except ValueError:
                        continue
                    if any(part in SKIP_DIRS for part in match.parts):
                        continue
                    candidates.append(match)

        out: list[Path] = []
        seen: set[str] = set()
        for candidate in candidates:
            try:
                resolved = candidate.resolve()
                rel = resolved.relative_to(PROJECT_ROOT).as_posix()
            except (OSError, ValueError):
                continue
            if rel in seen or not resolved.is_dir():
                continue
            seen.add(rel)
            out.append(resolved)
        return out

    def _choose_lint_ip_dir(ip: str) -> Optional[Path]:
        candidates = _lint_ip_candidates(ip)
        if not candidates:
            return None
        with_report = [p for p in candidates if (p / "lint" / "dut_lint.json").is_file()]
        if with_report:
            return max(with_report, key=lambda p: (p / "lint" / "dut_lint.json").stat().st_mtime)
        with_filelist = [p for p in candidates if (p / "list" / f"{p.name}.f").is_file()]
        if with_filelist:
            return with_filelist[0]
        return candidates[0]

    def _read_lint_report(ip_dir: Path) -> tuple[dict, str]:
        report_path = ip_dir / "lint" / "dut_lint.json"
        if not report_path.is_file():
            return {}, "missing lint/dut_lint.json"
        try:
            return json.loads(report_path.read_text(encoding="utf-8")), ""
        except Exception as exc:
            return {}, f"invalid lint report: {exc}"

    def _lint_diagnostic_path(rel_ip: str, diag_file: Any) -> str:
        file_text = str(diag_file or "").strip()
        if not file_text:
            return ""
        try:
            if os.path.isabs(file_text):
                return Path(file_text).resolve().relative_to(PROJECT_ROOT).as_posix()
        except (OSError, ValueError):
            return file_text
        file_text = file_text.replace("\\", "/").lstrip("/")
        if file_text == rel_ip or file_text.startswith(rel_ip + "/"):
            return file_text
        return f"{rel_ip}/{file_text}"

    def _normalize_lint_tool_results(report: dict, rel_ip: str) -> list[dict]:
        raw_results = report.get("tool_results") if isinstance(report, dict) else []
        if not isinstance(raw_results, list):
            return []
        out: list[dict] = []
        for result in raw_results:
            if not isinstance(result, dict):
                continue
            next_result = dict(result)
            diagnostics = []
            for diag in result.get("diagnostics") or []:
                if not isinstance(diag, dict):
                    continue
                next_diag = dict(diag)
                next_diag["path"] = _lint_diagnostic_path(rel_ip, next_diag.get("file"))
                diagnostics.append(next_diag)
            next_result["diagnostics"] = diagnostics
            out.append(next_result)
        return out

    @app.get("/api/lint/report")
    async def api_lint_report(ip: str, top: str = "", refresh: int = 0):
        """Return the canonical DUT lint report, split by pyslang/Verilator.

        `refresh=1` runs workflow/lint/scripts/dut_lint_report.py first so
        the UI can regenerate the report without asking the user to leave ATLAS.
        """
        ip_dir = _choose_lint_ip_dir(ip)
        if ip_dir is None:
            return JSONResponse({"error": "IP directory not found", "ip": ip}, status_code=404)

        rel_ip = ip_dir.relative_to(PROJECT_ROOT).as_posix()
        run_info: dict[str, Any] | None = None
        if refresh:
            script = WORKFLOW_ROOT / "lint" / "scripts" / "dut_lint_report.py"
            cmd = [_python_cmd(), str(script), rel_ip, "--top", top or ip_dir.name]

            def _run_lint_report():
                return subprocess.run(
                    cmd,
                    cwd=PROJECT_ROOT,
                    text=True,
                    encoding="utf-8",
                    errors="replace",
                    stdout=subprocess.PIPE,
                    stderr=subprocess.STDOUT,
                    timeout=180,
                )

            try:
                proc = await asyncio.to_thread(_run_lint_report)
                run_info = {
                    "command": " ".join(cmd),
                    "returncode": proc.returncode,
                    "output": proc.stdout[-12000:] if proc.stdout else "",
                }
            except subprocess.TimeoutExpired as exc:
                run_info = {
                    "command": " ".join(cmd),
                    "returncode": 124,
                    "output": str(exc),
                }
            except Exception as exc:
                run_info = {
                    "command": " ".join(cmd),
                    "returncode": 1,
                    "output": str(exc),
                }

        report, error = _read_lint_report(ip_dir)
        report_path = ip_dir / "lint" / "dut_lint.json"
        log_path = ip_dir / "lint" / "dut_lint.log"
        tool_results = _normalize_lint_tool_results(report, rel_ip) if report else []
        return JSONResponse({
            "ip": ip,
            "resolved_ip": rel_ip,
            "top": top or ip_dir.name,
            "exists": bool(report),
            "error": error,
            "report_path": report_path.relative_to(PROJECT_ROOT).as_posix(),
            "log_path": log_path.relative_to(PROJECT_ROOT).as_posix(),
            "log_exists": log_path.is_file(),
            "tool": report.get("tool", "") if report else "",
            "passed": report.get("passed") if report else None,
            "errors": int(report.get("errors") or 0) if report else 0,
            "warnings": int(report.get("warnings") or 0) if report else 0,
            "suppression_violations": int(report.get("suppression_violation_count") or 0) if report else 0,
            "style_violations": int(report.get("style_violation_count") or 0) if report else 0,
            "command": report.get("command", "") if report else "",
            "timestamp": report.get("timestamp", "") if report else "",
            "rtl_files": report.get("rtl_files", []) if report else [],
            "tool_results": tool_results,
            "run": run_info,
        })

    def _choose_coverage_ip_dir(ip: str) -> Optional[Path]:
        candidates = _lint_ip_candidates(ip)
        if not candidates:
            return None
        artifact_names = (
            "coverage.json",
            "coverage_ssot.json",
            "coverage.info",
            "toggle.json",
            "merged.dat",
        )
        with_cov = [
            p for p in candidates
            if any((p / "cov" / name).is_file() for name in artifact_names)
            or any((p / "sim").glob("*.vcd"))
        ]
        if with_cov:
            def _latest_artifact(path: Path) -> float:
                mtimes: list[float] = []
                for name in artifact_names:
                    target = path / "cov" / name
                    if target.is_file():
                        mtimes.append(target.stat().st_mtime)
                mtimes.extend(p.stat().st_mtime for p in (path / "sim").glob("*.vcd") if p.is_file())
                return max(mtimes or [0.0])
            return max(with_cov, key=_latest_artifact)
        with_filelist = [p for p in candidates if (p / "list" / f"{p.name}.f").is_file()]
        if with_filelist:
            return with_filelist[0]
        return candidates[0]

    def _read_json_artifact(path: Path) -> tuple[dict, str]:
        if not path.is_file():
            return {}, ""
        try:
            data = json.loads(path.read_text(encoding="utf-8"))
        except Exception as exc:
            return {}, f"invalid json: {exc}"
        return data if isinstance(data, dict) else {}, ""

    def _coverage_metric(hit: Any, total: Any, target: Any = None) -> dict[str, Any]:
        try:
            hit_i = int(hit or 0)
        except (TypeError, ValueError):
            hit_i = 0
        try:
            total_i = int(total or 0)
        except (TypeError, ValueError):
            total_i = 0
        pct_val = round(100.0 * hit_i / total_i, 2) if total_i else None
        out: dict[str, Any] = {"hit": hit_i, "total": total_i, "pct": pct_val}
        if target is not None:
            try:
                out["target_pct"] = float(target)
                out["meets_target"] = pct_val is not None and pct_val >= out["target_pct"]
            except (TypeError, ValueError):
                out["target_pct"] = target
        return out

    def _parse_lcov_summary(path: Path) -> dict[str, Any]:
        empty = {
            "available": False,
            "path": path.relative_to(PROJECT_ROOT).as_posix() if path.exists() else "",
            "lines": _coverage_metric(0, 0),
            "branches": _coverage_metric(0, 0),
            "functions": _coverage_metric(0, 0),
            "files": [],
            "error": "",
        }
        if not path.is_file():
            return empty

        def new_record() -> dict[str, Any]:
            return {
                "source": "",
                "line_total": 0,
                "line_hit": 0,
                "line_seen": set(),
                "branch_total": 0,
                "branch_hit": 0,
                "branch_seen": set(),
                "function_total": 0,
                "function_hit": 0,
                "function_seen": set(),
            }

        records: list[dict[str, Any]] = []
        current = new_record()

        def finish_record() -> None:
            nonlocal current
            if not current["source"] and not current["line_seen"] and not current["branch_seen"] and not current["function_seen"]:
                current = new_record()
                return
            if current["line_total"] == 0 and current["line_seen"]:
                current["line_total"] = len(current["line_seen"])
                current["line_hit"] = sum(1 for _, hits in current["line_seen"] if hits > 0)
            if current["branch_total"] == 0 and current["branch_seen"]:
                current["branch_total"] = len(current["branch_seen"])
                current["branch_hit"] = sum(1 for _, taken in current["branch_seen"] if taken > 0)
            if current["function_total"] == 0 and current["function_seen"]:
                current["function_total"] = len(current["function_seen"])
                current["function_hit"] = sum(1 for _, hits in current["function_seen"] if hits > 0)
            records.append(current)
            current = new_record()

        try:
            lines = path.read_text(encoding="utf-8", errors="replace").splitlines()
        except OSError as exc:
            empty["error"] = str(exc)
            return empty

        for raw in lines:
            line = raw.strip()
            if not line:
                continue
            if line == "end_of_record":
                finish_record()
                continue
            if line.startswith("SF:"):
                current["source"] = line[3:]
                continue
            if line.startswith("DA:"):
                parts = line[3:].split(",")
                if len(parts) >= 2:
                    try:
                        current["line_seen"].add((int(parts[0]), int(parts[1])))
                    except ValueError:
                        pass
                continue
            if line.startswith("LF:"):
                try:
                    current["line_total"] = int(line[3:])
                except ValueError:
                    pass
                continue
            if line.startswith("LH:"):
                try:
                    current["line_hit"] = int(line[3:])
                except ValueError:
                    pass
                continue
            if line.startswith("BRDA:"):
                parts = line[5:].split(",")
                if len(parts) >= 4:
                    try:
                        current["branch_seen"].add((":".join(parts[:3]), 0 if parts[3] == "-" else int(parts[3])))
                    except ValueError:
                        pass
                continue
            if line.startswith("BRF:"):
                try:
                    current["branch_total"] = int(line[4:])
                except ValueError:
                    pass
                continue
            if line.startswith("BRH:"):
                try:
                    current["branch_hit"] = int(line[4:])
                except ValueError:
                    pass
                continue
            if line.startswith("FNDA:"):
                parts = line[5:].split(",", 1)
                if len(parts) == 2:
                    try:
                        current["function_seen"].add((parts[1], int(parts[0])))
                    except ValueError:
                        pass
                continue
            if line.startswith("FNF:"):
                try:
                    current["function_total"] = int(line[4:])
                except ValueError:
                    pass
                continue
            if line.startswith("FNH:"):
                try:
                    current["function_hit"] = int(line[4:])
                except ValueError:
                    pass
        finish_record()

        total_lines = sum(int(r["line_total"] or 0) for r in records)
        hit_lines = sum(int(r["line_hit"] or 0) for r in records)
        total_branches = sum(int(r["branch_total"] or 0) for r in records)
        hit_branches = sum(int(r["branch_hit"] or 0) for r in records)
        total_functions = sum(int(r["function_total"] or 0) for r in records)
        hit_functions = sum(int(r["function_hit"] or 0) for r in records)

        files = []
        for record in records:
            source = str(record["source"] or "")
            try:
                if os.path.isabs(source):
                    source = Path(source).resolve().relative_to(PROJECT_ROOT).as_posix()
            except (OSError, ValueError):
                pass
            files.append({
                "path": source,
                "lines": _coverage_metric(record["line_hit"], record["line_total"]),
                "branches": _coverage_metric(record["branch_hit"], record["branch_total"]),
                "functions": _coverage_metric(record["function_hit"], record["function_total"]),
            })

        return {
            "available": True,
            "path": path.relative_to(PROJECT_ROOT).as_posix(),
            "lines": _coverage_metric(hit_lines, total_lines),
            "branches": _coverage_metric(hit_branches, total_branches),
            "functions": _coverage_metric(hit_functions, total_functions),
            "files": files,
            "error": "",
        }

    def _coverage_filelist_entries(ip_dir: Path) -> list[str]:
        filelist = ip_dir / "list" / f"{ip_dir.name}.f"
        entries: list[str] = []
        if filelist.is_file():
            for raw in filelist.read_text(encoding="utf-8", errors="replace").splitlines():
                line = raw.split("//", 1)[0].strip()
                if not line or line.startswith("+"):
                    continue
                if line.endswith((".v", ".sv", ".vh", ".svh")) and "/tb/" not in line and not line.startswith("tb/"):
                    entries.append(line)
        if entries:
            return entries
        return [
            path.relative_to(ip_dir).as_posix()
            for path in sorted((ip_dir / "rtl").glob("**/*"))
            if path.is_file() and path.suffix.lower() in {".v", ".sv", ".vh", ".svh"}
        ]

    def _diagnostic_file_name(diag: Any, sm: Any) -> str:
        loc = getattr(diag, "location", None)
        if loc is None or sm is None:
            return ""
        for arg in (loc, getattr(loc, "buffer", None)):
            if arg is None:
                continue
            try:
                file_name = sm.getFileName(arg)
                if file_name:
                    return str(file_name)
            except Exception:
                pass
        for attr in ("fileName", "filename", "file"):
            value = getattr(loc, attr, None)
            if value:
                return str(value)
        return ""

    def _static_rtl_coverage(ip_dir: Path, rel_ip: str) -> dict[str, Any]:
        entries = _coverage_filelist_entries(ip_dir)
        paths = [ip_dir / rel for rel in entries if rel.endswith((".v", ".sv", ".vh", ".svh"))]
        missing = [str(path.relative_to(ip_dir)) for path in paths if not path.is_file()]
        existing = [path for path in paths if path.is_file()]

        counts = {
            "files": len(existing),
            "listed_files": len(paths),
            "missing_files": len(missing),
            "lines": 0,
            "nonempty_lines": 0,
            "modules": 0,
            "always_blocks": 0,
            "assigns": 0,
            "case_blocks": 0,
            "assertions": 0,
            "cover_statements": 0,
        }
        file_rows: list[dict[str, Any]] = []
        for path in existing:
            text = path.read_text(encoding="utf-8", errors="replace")
            lines = text.splitlines()
            row = {
                "path": path.relative_to(PROJECT_ROOT).as_posix(),
                "lines": len(lines),
                "modules": len(re.findall(r"\bmodule\s+[A-Za-z_][A-Za-z0-9_$]*", text)),
                "always_blocks": len(re.findall(r"\balways(?:_ff|_comb|_latch)?\b", text)),
                "assigns": len(re.findall(r"\bassign\b", text)),
                "case_blocks": len(re.findall(r"\bcase[zx]?\s*\(", text)),
                "assertions": len(re.findall(r"\bassert(?:\s+property)?\b", text)),
                "cover_statements": len(re.findall(r"\bcover(?:\s+property)?\b", text)),
            }
            file_rows.append(row)
            counts["lines"] += row["lines"]
            counts["nonempty_lines"] += sum(1 for line in lines if line.strip())
            for key in ("modules", "always_blocks", "assigns", "case_blocks", "assertions", "cover_statements"):
                counts[key] += int(row[key])

        diagnostics: list[dict[str, Any]] = []
        compile_error = ""
        pyslang_available = False
        try:
            from core.pyslang_compat import compile_files as compile_pyslang_files
            from core.pyslang_compat import diagnostic_is_error, diagnostic_line, diagnostic_message
            compiled = compile_pyslang_files(existing)
            pyslang_available = not str(compiled.error or "").startswith("pyslang import failed")
            sm = compiled.source_manager
            for diag in compiled.diagnostics or []:
                is_error = diagnostic_is_error(diag)
                diag_file = _diagnostic_file_name(diag, sm)
                diagnostics.append({
                    "severity": "error" if is_error else "warning",
                    "file": diag_file,
                    "path": _lint_diagnostic_path(rel_ip, diag_file),
                    "line": diagnostic_line(diag, sm),
                    "message": diagnostic_message(compiled.pyslang, diag, sm),
                })
            if compiled.error:
                compile_error = compiled.error
                if not diagnostics:
                    diagnostics.append({
                        "severity": "error",
                        "file": "",
                        "path": "",
                        "line": 0,
                        "message": compiled.error,
                    })
        except Exception as exc:
            compile_error = f"pyslang static analysis failed: {exc}"
            diagnostics.append({"severity": "error", "file": "", "path": "", "line": 0, "message": compile_error})

        errors = sum(1 for d in diagnostics if str(d.get("severity", "")).lower() == "error")
        warnings = len(diagnostics) - errors
        return {
            "tool": "pyslang",
            "kind": "static_elab",
            "available": pyslang_available,
            "passed": bool(existing) and not missing and errors == 0,
            "source": "list/{name}.f + pyslang compile + static RTL scan".format(name=ip_dir.name),
            "filelist": (ip_dir / "list" / f"{ip_dir.name}.f").relative_to(PROJECT_ROOT).as_posix()
                if (ip_dir / "list" / f"{ip_dir.name}.f").is_file() else "",
            "rtl_files": [path.relative_to(PROJECT_ROOT).as_posix() for path in existing],
            "missing": missing,
            "metrics": counts,
            "files": file_rows,
            "diagnostics": diagnostics[:100],
            "errors": errors,
            "warnings": warnings,
            "error": compile_error,
        }

    def _domain_matches(bin_id: str, item: Any, domain: str) -> bool:
        text = str(bin_id).lower()
        if isinstance(item, dict):
            text += " " + " ".join(
                str(item.get(key, "")).lower()
                for key in ("coverage_domain", "domain", "class", "source", "description")
            )
            plan = item.get("plan")
            if isinstance(plan, dict):
                text += " " + " ".join(str(plan.get(key, "")).lower() for key in ("coverage_domain", "domain", "class", "source", "description"))
        if domain == "function":
            return any(token in text for token in ("function", "functional", "transaction", "scenario", "fl"))
        return any(token in text for token in ("cycle", "protocol", "handshake", "latency", "fsm", "transition", "cl"))

    def _domain_coverage_from_report(report: dict, domain: str) -> dict[str, Any]:
        key = "function_coverage" if domain == "function" else "cycle_coverage"
        raw = report.get(key) if isinstance(report.get(key), dict) else {}
        out = {
            "domain": domain,
            "hit": int(raw.get("hit") or 0),
            "total": int(raw.get("total") or 0),
            "pct": raw.get("pct"),
            "target_pct": raw.get("target_pct"),
            "meets_target": raw.get("meets_target"),
            "source": raw.get("source") or ("function_model" if domain == "function" else "cycle_model"),
            "missing_bins": [],
            "bins": [],
        }
        bins = report.get("functional_bins") if isinstance(report.get("functional_bins"), dict) else {}
        for bin_id, item in bins.items():
            if not _domain_matches(str(bin_id), item, domain):
                continue
            hit = bool(item.get("hit") or item.get("covered") or item.get("raw_hit")) if isinstance(item, dict) else bool(item)
            row = {
                "id": str(bin_id),
                "hit": hit,
                "source": item.get("source", "") if isinstance(item, dict) else "",
                "description": item.get("description", "") if isinstance(item, dict) else "",
            }
            out["bins"].append(row)
            if not hit:
                out["missing_bins"].append(row)
        return out

    def _normalize_toggle_report(toggle: dict, path: Path) -> dict[str, Any]:
        if not toggle:
            return {
                "available": False,
                "path": path.relative_to(PROJECT_ROOT).as_posix() if path.exists() else "",
                "vcd": "",
                "metrics": _coverage_metric(0, 0),
                "nets": 0,
                "scopes": [],
            }
        total = int(toggle.get("total_bits") or 0)
        hit = int(toggle.get("toggled_bits") or 0)
        scopes = toggle.get("scopes") if isinstance(toggle.get("scopes"), list) else []
        scopes = sorted(
            [s for s in scopes if isinstance(s, dict)],
            key=lambda s: (float(s.get("pct") or 0.0), str(s.get("scope") or "")),
        )
        return {
            "available": True,
            "path": path.relative_to(PROJECT_ROOT).as_posix(),
            "vcd": str(toggle.get("vcd") or ""),
            "metrics": {**_coverage_metric(hit, total), "pct": round(float(toggle.get("pct") or 0.0), 2) if total else None},
            "nets": int(toggle.get("nets") or 0),
            "scopes": scopes[:20],
        }

    def _coverage_card(id_: str, label: str, available: bool, status: str, metrics: list[dict[str, Any]], **extra: Any) -> dict[str, Any]:
        payload = {
            "id": id_,
            "label": label,
            "available": available,
            "status": status,
            "metrics": metrics,
        }
        payload.update(extra)
        return payload

    # Phase 17: /api/coverage/report extracted to atlas_api_coverage_report.py.
    # All 9 closure captures are defined earlier in create_app() than this
    # call site (verified line numbers 293..2881), so no lambda forward-refs
    # are needed.
    from src.atlas_api_coverage_report import register_coverage_report_routes as _register_coverage_report_routes
    _register_coverage_report_routes(
        app,
        PROJECT_ROOT=PROJECT_ROOT,
        WORKFLOW_ROOT=WORKFLOW_ROOT,
        _python_cmd=_python_cmd,
        _safe=_safe,
        _parse_lcov_summary=_parse_lcov_summary,
        _choose_coverage_ip_dir=_choose_coverage_ip_dir,
        _coverage_card=_coverage_card,
        _coverage_metric=_coverage_metric,
        _domain_coverage_from_report=_domain_coverage_from_report,
        _normalize_toggle_report=_normalize_toggle_report,
        _read_json_artifact=_read_json_artifact,
        _static_rtl_coverage=_static_rtl_coverage,
    )

    # ── Foldable structure for PreviewPane (sv / yaml) ────────────
    # Returns line ranges the frontend wraps in <details>. Empty list
    # for unknown extensions — caller falls back to plain Prism.
    _FOLD_CACHE: "collections.OrderedDict[str, tuple[float, list]]" = collections.OrderedDict()
    _FOLD_CACHE_CAP = 32
    _FOLD_MAX_BYTES = 5 * 1024 * 1024   # 5 MB
    _FOLD_MAX_LINES = 10_000


    # ── VCD (waveform) endpoints — sim_debug workspace ────────────
    # VCD files can be MB+ so we bypass MAX_READ_BYTES with a separate
    # ceiling. Path resolution still goes through _safe() so the user
    # can't escape PROJECT_ROOT.
    MAX_VCD_BYTES = 32 * 1024 * 1024  # 32 MB
    # Routes registered via register_vcd_routes() below (see atlas_api_vcd.py).

    @app.post("/api/ip/create")
    async def api_ip_create(request: Request):
        """Create the on-disk IP scaffold used by the workspace file tree.

        `/new-ip <name>` still owns the chat-facing SSOT plan, but the UI
        needs a synchronous HTTP path so `+ IP` cannot leave only a
        `.session/<owner>/<ip>/...` namespace when the websocket prompt is
        delayed or disconnected.
        """
        try:
            body = await request.json()
        except Exception:
            body = {}
        name = str((body or {}).get("name") or "").strip()
        kind = str((body or {}).get("kind") or "TBD").strip() or "TBD"
        requested_exec_mode = _normalize_atlas_exec_mode((body or {}).get("exec_mode"))
        if (body or {}).get("exec_mode") is not None and not requested_exec_mode:
            return JSONResponse({"error": "exec_mode must be single-worker or orchestrator"}, status_code=400)
        workflow = _new_ip_initial_workflow(
            (body or {}).get("workflow") or (body or {}).get("initial_workflow"),
            requested_exec_mode,
        )
        if not name:
            return JSONResponse({"error": "name required"}, status_code=400)
        if not _valid_ip_name(name) or "/" in name or "\\" in name or ".." in name:
            return JSONResponse({"error": "invalid name"}, status_code=400)
        user = request.scope.get("user") or {}
        username = normalize_session_name(str(user.get("username") or ""))
        user_id = str(user.get("id") or "").strip()
        multi_user_on = _multi_user_enabled()
        if multi_user_on and (not username or not user_id):
            return JSONResponse({"error": "login required"}, status_code=401)
        def _owner_hint_from_create_payload() -> str:
            for key in ("session_id", "session", "namespace", "active_session"):
                raw = normalize_session_name(str((body or {}).get(key) or ""))
                parts = [part for part in raw.split("/") if part]
                if len(parts) >= 2 and parts[0] and parts[0] != "default":
                    return parts[0]
            for key in ("user_name", "username", "owner"):
                raw = normalize_session_name(str((body or {}).get(key) or ""))
                if raw and raw != "default":
                    return raw
            raw = normalize_session_name(str(os.environ.get("ATLAS_USER_SESSION_ID") or ""))
            return raw if raw and raw != "default" else ""

        workspace_session = normalize_session_name(
            str(
                (body or {}).get("workspace_session")
                or (body or {}).get("workspaceSession")
                or os.environ.get("ATLAS_WORKSPACE_SESSION")
                or "default"
            )
        )
        if not workspace_session:
            workspace_session = "default"
        atlas_root = _atlas_root_for_context()
        context: AtlasContext | None = None
        workspace_root = PROJECT_ROOT
        create_owner_hint = _owner_hint_from_create_payload()
        context_owner = username if multi_user_on else (create_owner_hint or ("" if username == "local-admin" else username))
        if context_owner:
            try:
                context = AtlasContext(
                    user_name=context_owner,
                    workspace_session=workspace_session or "default",
                    ip_name=name,
                    workflow=workflow,
                    atlas_root=atlas_root,
                )
            except ValueError as exc:
                return JSONResponse({"error": str(exc)}, status_code=400)
            validated_workspace_root, workspace_error = _validated_context_workspace_root(context)
            if workspace_error or validated_workspace_root is None:
                return JSONResponse({"error": workspace_error or "invalid workspace root"}, status_code=400)
            workspace_root = validated_workspace_root
            raw_target = workspace_root / name
            if raw_target.is_symlink():
                return JSONResponse({"error": "ip path is a symlink"}, status_code=400)
            target = raw_target.resolve(strict=False)
        else:
            target = (PROJECT_ROOT / name).resolve()
        try:
            target.relative_to(workspace_root.resolve())
        except ValueError:
            return JSONResponse({"error": "outside project root"}, status_code=400)
        if target.exists():
            return JSONResponse({
                "error": f'IP "{name}" already exists. Select it from IP_ID or choose another name.',
                "ip": name,
            }, status_code=409)
        session_namespace = context.active_session_key if context is not None else (f"{username}/{name}/{workflow}" if username else "")
        session_dir = context.session_dir if context is not None else ((PROJECT_ROOT / ".session" / username / name / workflow) if username else None)
        db_session: dict[str, Any] = {}
        workspace_row: dict[str, Any] = {}
        ip_row: dict[str, Any] = {}
        try:
            target.mkdir(parents=True, exist_ok=True)
            paths = _ensure_new_ip_structure(name, base_root=workspace_root)
            _ensure_ssot_draft(name, kind, base_root=workspace_root)
            if session_dir is not None:
                session_dir.mkdir(parents=True, exist_ok=True)
                conv = session_dir / "conversation.json"
                if not conv.exists():
                    conv.write_text("[]", encoding="utf-8")
            if multi_user_on and session_namespace:
                summary = {
                    "kind": "atlas_ip_scaffold",
                    "namespace": session_namespace,
                    "owner": context_owner,
                    "workspace_session": workspace_session or "default",
                    "context_key": context.context_key if context is not None else session_namespace,
                    "ip": name,
                    "workflow": workflow,
                    "workspace_root": str(workspace_root),
                    "ip_root": str(target),
                }
                with AtlasDB() as db:
                    workspace_row = db.upsert_workspace(
                        f"{PROJECT_ROOT.name or 'default'}/{workspace_session or 'default'}",
                        owner_user_id=user_id,
                        local_path=str(workspace_root.resolve()),
                    ) or {}
                    ip_row = db.upsert_ip_block(
                        str(workspace_row.get("id") or ""),
                        name,
                        ip_type=kind,
                        ssot_path=f"{name}/yaml/{name}.ssot.yaml",
                    ) or {}
                    db_session = db.upsert_runtime_session(
                        session_namespace,
                        user_id,
                        owner=context_owner,
                        ip=name,
                        workflow=workflow,
                        workspace_id=str(workspace_row.get("id") or ""),
                        ip_id=str(ip_row.get("id") or name),
                        project_id=name,
                        directory=str(session_dir) if session_dir is not None else "",
                        title=f"{name} / {workflow}",
                        status="active",
                        summary=summary,
                    )
        except Exception as exc:
            return JSONResponse({"error": f"failed to scaffold IP: {exc}"}, status_code=500)
        if session_namespace:
            try:
                bridge.activate_session(session_namespace)
                _atlas_active_session_cv.set(session_namespace)
                _atlas_active_ip_cv.set(name)
            except Exception:
                pass
        worker_warmup: dict[str, Any] = {}
        try:
            try:
                from atlas_api_jobs import schedule_worker_warmup  # noqa: WPS433
            except ImportError:
                from src.atlas_api_jobs import schedule_worker_warmup  # type: ignore  # noqa: WPS433

            worker_warmup = schedule_worker_warmup(
                ip=name,
                owner=context_owner,
                db_user_id=user_id,
                session_name=session_namespace,
                active_workflow=workflow,
                project_root_value=str(workspace_root),
                exec_mode=requested_exec_mode or _current_atlas_exec_mode(),
                reason="ip_create",
                background=True,
            )
        except Exception as exc:
            worker_warmup = {"enabled": False, "error": str(exc)}
        return JSONResponse({"ok": True,
                             "ip": name,
                             "created": True,
                             "path": str(target.relative_to(workspace_root.resolve())),
                             "ssot_path": f"{name}/yaml/{name}.ssot.yaml",
                             "paths": paths,
                             "session": session_namespace,
                             "active_session": session_namespace,
                             "context_key": context.context_key if context is not None else session_namespace,
                             "workspace_session": workspace_session or "default",
                             "workspace_root": str(workspace_root),
                             "ip_root": str(target),
                             "session_dir": str(session_dir) if session_dir is not None else "",
                             "workflow": workflow,
                             "exec_mode": requested_exec_mode or _current_atlas_exec_mode(),
                             "policy": exec_policy_payload(requested_exec_mode or _current_atlas_exec_mode(), env=os.environ),
                             "session_uid": str(db_session.get("session_uid") or ""),
                             "workspace_id": str(workspace_row.get("id") or ""),
                             "ip_block_id": str(ip_row.get("id") or ""),
                             "worker_warmup": worker_warmup})

    def _git_route_provider(value: Any = "") -> str:
        provider = str(value or "").strip().lower()
        return "" if provider in {"", "auto", "default"} else provider

    def _resolve_ip_path(
        name: str,
        provider: str = "",
        session_id: str = "",
    ) -> Path | tuple[None, JSONResponse]:
        """Validate a path-segment-style IP name and return its on-disk dir.

        Returns a Path on success, or a (None, JSONResponse) tuple on
        failure for the caller to forward."""
        clean = str(name or "").strip()
        if not clean or "/" in clean or "\\" in clean or ".." in clean:
            return None, JSONResponse({"error": "invalid ip name"}, status_code=400)
        base = PROJECT_ROOT.resolve()
        session = normalize_session_name(str(session_id or ""))
        if session:
            try:
                context = AtlasContext.from_session_key(
                    session,
                    atlas_root=os.environ.get("ATLAS_ROOT") or str(PROJECT_ROOT),
                )
                base = context.workspace_root.resolve()
            except ValueError as exc:
                return None, JSONResponse({"error": str(exc)}, status_code=400)
        target = (base / clean).resolve()
        try:
            target.relative_to(base)
        except ValueError:
            return None, JSONResponse({"error": "outside project root"}, status_code=400)
        if not target.is_dir():
            return None, JSONResponse({"error": "ip not found"}, status_code=404)
        if (
            not scm_provider_allows_missing_git_dir(_git_route_provider(provider) or configured_scm_provider())
            and not (target / ".git").is_dir()
        ):
            return None, JSONResponse({"error": "ip has no .git — create via /new-ip first"}, status_code=409)
        return target

    @app.post("/api/ip/{name}/git/commit")
    async def api_ip_git_commit(name: str, request: Request):
        """Stage and commit the IP's working tree with a user/agent message."""
        try:
            body = await request.json()
        except Exception:
            body = {}
        message = str((body or {}).get("message") or "").strip() or "commit"
        provider = _git_route_provider((body or {}).get("provider"))
        session_id = str((body or {}).get("session_id") or (body or {}).get("sessionId") or "")
        resolved = _resolve_ip_path(name, provider=provider, session_id=session_id)
        if isinstance(resolved, tuple):
            _, err = resolved
            return err
        target = resolved
        try:
            adapter = resolve_scm_adapter(target, provider=provider or None)
            out = await asyncio.to_thread(
                adapter.submit,
                message,
                add_all=True,
                allow_empty=True,
            )
            status = await asyncio.to_thread(adapter.status)
            return JSONResponse({
                "ok": out.ok,
                "ip": name,
                "hash": str(status.get("head_full") or status.get("head") or "")[:12],
                "stdout": out.stdout,
                "stderr": out.stderr,
                "error": out.error,
                "provider": out.provider,
                "returncode": out.returncode,
            })
        except Exception as exc:
            return JSONResponse({"error": str(exc)}, status_code=500)

    @app.get("/api/ip/{name}/git/log")
    async def api_ip_git_log(name: str, limit: int = 50, provider: str = "", session_id: str = ""):
        """Return the last N commits of the per-IP repo as JSON."""
        provider = _git_route_provider(provider)
        resolved = _resolve_ip_path(name, provider=provider, session_id=session_id)
        if isinstance(resolved, tuple):
            _, err = resolved
            return err
        target = resolved
        try:
            limit = max(1, min(int(limit or 50), 500))
            adapter = resolve_scm_adapter(target, provider=provider or None)
            log = await asyncio.to_thread(adapter.log, limit)
            commits = [{
                "hash": str(item.get("sha") or item.get("hash") or ""),
                "short": str(item.get("short") or ""),
                "author": str(item.get("author") or ""),
                "time": float(item.get("time") or 0),
                "subject": str(item.get("subject") or ""),
            } for item in log.get("commits", [])]
            payload = {
                "ip": name,
                "commits": commits,
                "provider": log.get("provider", adapter.provider),
            }
            if not log.get("ok", True):
                payload["error"] = log.get("error") or "scm log failed"
            return JSONResponse(payload)
        except Exception as exc:
            return JSONResponse({"error": str(exc)}, status_code=500)

    @app.get("/api/ip/{name}/git/url")
    async def api_ip_git_url(name: str, request: Request):
        """Return clone URLs for the per-IP bare repo at <root>/<name>.git.
        Honors BARE_GIT_OPTION — returns 404 when the bare wasn't built."""
        clean = str(name or "").strip()
        if not clean or "/" in clean or "\\" in clean or ".." in clean:
            return JSONResponse({"error": "invalid ip name"}, status_code=400)
        bare = (PROJECT_ROOT / f"{clean}.git").resolve()
        try:
            bare.relative_to(PROJECT_ROOT.resolve())
        except ValueError:
            return JSONResponse({"error": "outside project root"}, status_code=400)
        if not (bare / "HEAD").is_file():
            return JSONResponse({"error": "no bare repo — BARE_GIT_OPTION may be off, or scaffold hasn't run"},
                                status_code=404)
        host = request.headers.get("host") or "127.0.0.1:8765"
        scheme = "https" if request.url.scheme == "https" else "http"
        return JSONResponse({
            "ip": clean,
            "bare_path": str(bare),
            "clone": {
                "http": f"{scheme}://{host}/git/{clean}.git",
                "file": f"file://{bare}",
            },
        })

    @app.api_route("/git/{path:path}", methods=["GET", "POST"])
    async def git_http_backend_proxy(path: str, request: Request):
        """Smart-HTTP gateway over git-http-backend so the per-IP bare
        repos under PROJECT_ROOT are clone+push targets on the LAN.
        Honors BARE_GIT_OPTION — returns 404 when disabled."""
        from starlette.responses import Response as _StarResponse
        try:
            import config as _cfg_git
            if not getattr(_cfg_git, "BARE_GIT_OPTION", True):
                return JSONResponse({"error": "BARE_GIT_OPTION disabled"}, status_code=404)
        except Exception:
            pass
        # Git smart-HTTP authz. ATLAS_GIT_ANON_READ is ON BY DEFAULT, so anonymous
        # fetch/clone of any IP is allowed (AuthMiddleware lets /git/ through).
        # Set ATLAS_GIT_ANON_READ=0 to require auth + per-IP view for fetch
        # (fetch -> view, with cross-tenant reads blocked). PUSH (receive-pack)
        # ALWAYS requires an authenticated user with write access, either way.
        from core.atlas_auth import git_anon_read_enabled as _git_anon_read_enabled
        _repo_seg = (path or "").split("/", 1)[0]
        _git_ip = _repo_seg[:-4] if _repo_seg.endswith(".git") else _repo_seg
        _is_push = (
            "service=git-receive-pack" in (request.url.query or "")
            or str(path or "").endswith("git-receive-pack")
        )
        if _is_push:
            _git_denied = _fs_authz.ip(request, _git_ip, "write")
            if _git_denied is not None:
                return _git_denied
        elif not _git_anon_read_enabled():
            _git_denied = _fs_authz.ip(request, _git_ip, "view")
            if _git_denied is not None:
                return _git_denied
        # Find git-http-backend. Brew puts it under libexec/git-core.
        import subprocess as _sp_git
        backend = ""
        try:
            ep = _sp_git.run(["git", "--exec-path"], capture_output=True,
                             timeout=3, check=False)
            cand = (ep.stdout.decode("utf-8", "replace").strip() + "/git-http-backend")
            if Path(cand).is_file():
                backend = cand
        except Exception:
            pass
        if not backend:
            return JSONResponse({"error": "git-http-backend not found on host"},
                                status_code=501)
        auth_header = request.headers.get("authorization", "")
        query_string = request.url.query or ""
        if "service=git-receive-pack" in query_string and not auth_header.lower().startswith("basic "):
            return _StarResponse(
                content=b"Authentication required for git push\n",
                status_code=401,
                headers={"WWW-Authenticate": 'Basic realm="Atlas Git"'},
            )
        env = dict(os.environ)
        env.update({
            "GIT_PROJECT_ROOT": str(PROJECT_ROOT.resolve()),
            "GIT_HTTP_EXPORT_ALL": "1",
            "PATH_INFO": "/" + path,
            "REQUEST_METHOD": request.method,
            "QUERY_STRING": query_string,
            "CONTENT_TYPE": request.headers.get("content-type", ""),
            "CONTENT_LENGTH": request.headers.get("content-length", "0"),
            "REMOTE_ADDR": (request.client.host if request.client else "127.0.0.1"),
        })
        if auth_header.lower().startswith("basic "):
            try:
                import base64 as _base64_git_auth
                raw = _base64_git_auth.b64decode(
                    auth_header.split(None, 1)[1],
                    validate=True,
                ).decode("utf-8", "replace")
                username = raw.split(":", 1)[0].strip()
                if username:
                    # git-http-backend treats REMOTE_USER as the signal
                    # that HTTP auth succeeded. Without this, clients that
                    # retry with http://user:pass@host/... still look
                    # anonymous to receive-pack and can get Unauthorized.
                    env["REMOTE_USER"] = username
                    env["AUTH_TYPE"] = "Basic"
            except Exception:
                pass
        body = await request.body()
        proc = await asyncio.create_subprocess_exec(
            backend,
            env=env,
            stdin=asyncio.subprocess.PIPE,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
        )
        try:
            stdout, _stderr = await asyncio.wait_for(
                proc.communicate(body), timeout=120,
            )
        except asyncio.TimeoutError:
            proc.kill()
            return JSONResponse({"error": "git-http-backend timed out"},
                                status_code=504)
        sep = stdout.find(b"\r\n\r\n")
        if sep == -1:
            return _StarResponse(content=stdout, status_code=200)
        header_text = stdout[:sep].decode("latin-1", "replace")
        resp_body = stdout[sep + 4:]
        status = 200
        headers: dict[str, str] = {}
        for line in header_text.split("\r\n"):
            if line.lower().startswith("status:"):
                try: status = int(line.split(":", 1)[1].strip().split()[0])
                except Exception: status = 200
            elif ":" in line:
                k, v = line.split(":", 1)
                headers[k.strip()] = v.strip()
        return _StarResponse(content=resp_body, status_code=status, headers=headers)

    @app.get("/api/ip/{name}/git/graph")
    async def api_ip_git_graph(name: str, limit: int = 80, provider: str = "", session_id: str = ""):
        """ASCII graph of the per-IP commit history. Returns the raw
        `git log --graph --oneline --decorate --all` text plus a parsed
        commit list so the frontend can render either a monospaced graph
        or a structured list."""
        provider = _git_route_provider(provider)
        resolved = _resolve_ip_path(name, provider=provider, session_id=session_id)
        if isinstance(resolved, tuple):
            _, err = resolved
            return err
        target = resolved
        try:
            limit = max(1, min(int(limit or 80), 1000))
            adapter = resolve_scm_adapter(target, provider=provider or None)
            graph = await asyncio.to_thread(adapter.graph, limit)
            payload = {
                "ip": name,
                "graph": graph.get("graph", ""),
                "commits": graph.get("commits", []),
                "provider": graph.get("provider", adapter.provider),
            }
            if not graph.get("ok", True):
                payload["error"] = graph.get("error") or "scm graph failed"
            return JSONResponse(payload)
        except Exception as exc:
            return JSONResponse({"error": str(exc)}, status_code=500)

    @app.post("/api/ip/{name}/git/revert")
    async def api_ip_git_revert(name: str, request: Request):
        """Restore the working tree to a previous commit (hard reset)."""
        try:
            body = await request.json()
        except Exception:
            body = {}
        target_hash = str((body or {}).get("hash") or "").strip()
        provider = _git_route_provider((body or {}).get("provider"))
        session_id = str((body or {}).get("session_id") or (body or {}).get("sessionId") or "")
        resolved = _resolve_ip_path(name, provider=provider, session_id=session_id)
        if isinstance(resolved, tuple):
            _, err = resolved
            return err
        target = resolved
        try:
            adapter = resolve_scm_adapter(target, provider=provider or None)
            if adapter.provider == "git" and not re.match(r"^[0-9a-f]{7,40}$", target_hash, re.I):
                return JSONResponse({"error": "invalid hash"}, status_code=400)
            if adapter.provider != "git" and not re.match(r"^[0-9A-Za-z._/@#:+-]{1,160}$", target_hash):
                return JSONResponse({"error": "invalid revision"}, status_code=400)
            out = await asyncio.to_thread(adapter.hard_reset, target_hash)
            if out.returncode == 404:
                return JSONResponse({"error": out.error or "revision not in this ip's history"}, status_code=404)
            return JSONResponse({
                "ok": out.ok,
                "ip": name,
                "hash": target_hash,
                "stdout": out.stdout,
                "stderr": out.stderr,
                "error": out.error,
                "provider": out.provider,
                "returncode": out.returncode,
            })
        except Exception as exc:
            return JSONResponse({"error": str(exc)}, status_code=500)

    @app.get("/api/ip/list")
    async def api_ip_list(request: Request, session_id: str = ""):
        """List IPs that belong to a session namespace.

        In single-user desktop mode the dropdown should include real IP
        directories under PROJECT_ROOT, so a newly generated IP is
        selectable before any per-user session history exists. In
        multi-user mode the authoritative source remains
        `.session/<session_id>/<ip>/<workflow>/` to avoid cross-user IP
        leakage.
        """
        skip = set(SKIP_DIRS) | {
            ".session", ".git", ".venv", ".sisyphus", ".omc", "logs",
            "node_modules", "workflow", "src", "core", "lib", "tests",
            "frontend", "docs", "scripts",
        }
        ip_markers = {
            "yaml", "rtl", "tb", "sim", "lint", "syn", "sta", "sta-post",
            "pnr", "dft", "doc", "req", "list", "cov", "verify", "model",
        }
        by_name: dict[str, dict[str, Any]] = {}

        workspace_root_for_list: Path | None = None
        workspace_root_for_list_error = ""

        def _ssot_exists(name: str) -> bool:
            roots = []
            if workspace_root_for_list is not None:
                roots.append(workspace_root_for_list)
            roots.append(PROJECT_ROOT)
            for root in roots:
                yaml_dir = root / name / "yaml"
                if not yaml_dir.is_dir():
                    continue
                if (yaml_dir / f"{name}.ssot.yaml").is_file():
                    return True
                try:
                    if any(yaml_dir.glob("*.ssot.yaml")):
                        return True
                except OSError:
                    continue
            return False

        def _add_item(name: str, *, workflows=None, mtime: float = 0.0) -> None:
            if not name or name.startswith(".") or name in skip:
                return
            row = by_name.setdefault(name, {
                "name": name,
                "has_ssot": _ssot_exists(name),
                "workflows": [],
                "mtime": mtime,
            })
            row["has_ssot"] = bool(row.get("has_ssot")) or _ssot_exists(name)
            row["mtime"] = max(float(row.get("mtime") or 0.0), float(mtime or 0.0))
            if workflows:
                merged = set(row.get("workflows") or [])
                merged.update(str(w) for w in workflows if str(w or "").strip())
                row["workflows"] = sorted(merged)

        def _looks_like_project_ip(entry: Path) -> bool:
            if entry.is_symlink() or not entry.is_dir() or entry.name.startswith(".") or entry.name in skip:
                return False
            try:
                return any((entry / marker).is_dir() for marker in ip_markers)
            except OSError:
                return False

        def _workspace_catalog_ip_names(db: Any, db_user_id: str = "") -> set[str] | None:
            if workspace_root_for_list is None:
                return None
            try:
                workspace_path = str(workspace_root_for_list.resolve())
            except OSError:
                workspace_path = str(workspace_root_for_list)
            rows = db._fetchall(
                """
                SELECT i.ip_name, i.status
                  FROM workspaces w
                  JOIN ip_blocks i ON i.workspace_id = w.id
                 WHERE w.local_path = ?
                   AND (? = '' OR w.owner_user_id = ?)
                 ORDER BY i.ip_name
                """,
                (workspace_path, db_user_id, db_user_id),
            )
            names = {
                str(row["ip_name"] or "").strip()
                for row in rows
                if str(row["ip_name"] or "").strip()
                and str(row["status"] or "active").strip().lower() == "active"
            }
            return names or None

        user = request.scope.get("user") or {}
        username = normalize_session_name(str(user.get("username") or ""))
        requested = normalize_session_name(str(session_id or ""))
        requested_parts = [part for part in requested.split("/") if part]
        owner = (requested_parts[0] if requested_parts else "") or username
        requested_workspace_session = (
            requested_parts[1]
            if len(requested_parts) >= 4 or len(requested_parts) == 2
            else ""
        )
        multi_user_on = _multi_user_enabled()
        if multi_user_on and not username:
            return JSONResponse({"error": "login required", "items": [], "count": 0}, status_code=401)
        if multi_user_on and username and owner and owner != username:
            return JSONResponse({"error": "session owner mismatch", "items": []}, status_code=403)
        if owner and requested_workspace_session:
            workspace_root_for_list, workspace_root_for_list_error = _validated_workspace_root(
                owner,
                requested_workspace_session,
                _atlas_root_for_context(),
            )
        session_root = (PROJECT_ROOT / ".session" / owner).resolve() if owner else None
        try:
            if workspace_root_for_list_error:
                return JSONResponse({
                    "error": workspace_root_for_list_error,
                    "project_root": str(PROJECT_ROOT),
                    "workspace_root": str(workspace_root_for_list) if workspace_root_for_list is not None else "",
                    "session_id": owner or "",
                    "workspace_session": requested_workspace_session or "",
                    "items": [],
                    "count": 0,
                }, status_code=400)
            if multi_user_on:
                user_id = str(user.get("id") or "").strip()
                if not user_id:
                    return JSONResponse({"error": "login required", "items": [], "count": 0}, status_code=401)
                try:
                    model_owner = _session_owner_with_model(username)
                except Exception:
                    model_owner = username
                allowed_owners = {value for value in (username, model_owner) if value}
                with AtlasDB() as db:
                    catalog_ip_names = _workspace_catalog_ip_names(db, user_id)
                    session_rows = db.list_sessions(user_id)
                for row in session_rows:
                    namespace = normalize_session_name(
                        str(row.get("namespace") or row.get("id") or "")
                    )
                    parts = [part for part in namespace.split("/") if part]
                    if len(parts) < 3 or parts[0] not in allowed_owners:
                        continue
                    row_workspace_session = parts[1] if len(parts) >= 4 else ""
                    if requested_workspace_session:
                        if len(parts) >= 4:
                            if row_workspace_session != requested_workspace_session:
                                continue
                        elif requested_workspace_session != "default":
                            continue
                    name = parts[-2]
                    if catalog_ip_names is not None and name not in catalog_ip_names:
                        continue
                    workflow = parts[-1]
                    mtime = float(row.get("updated_at") or row.get("created_at") or 0.0)
                    ip_dir = None
                    if len(parts) >= 4:
                        ip_dir = PROJECT_ROOT / parts[0] / row_workspace_session / name
                    elif session_root is not None:
                        ip_dir = session_root / name
                    try:
                        if ip_dir is not None and ip_dir.is_dir():
                            mtime = max(mtime, ip_dir.stat().st_mtime)
                    except OSError:
                        pass
                    _add_item(name, workflows=[workflow], mtime=mtime)
                if catalog_ip_names is not None:
                    for name in catalog_ip_names:
                        _add_item(name)
                items = sorted(by_name.values(), key=lambda x: (-x["mtime"], x["name"]))
                return JSONResponse({
                    "project_root": str(PROJECT_ROOT),
                    "workspace_root": str(workspace_root_for_list) if workspace_root_for_list is not None else "",
                    "session_id": owner or "",
                    "workspace_session": requested_workspace_session or "",
                    "items": items,
                    "count": len(items),
                    "source": "db_ip_blocks" if catalog_ip_names is not None else "db_sessions",
                })
            if not multi_user_on:
                scan_root = workspace_root_for_list or PROJECT_ROOT
                catalog_ip_names = None
                if workspace_root_for_list is not None:
                    with AtlasDB() as db:
                        catalog_ip_names = _workspace_catalog_ip_names(db)
                if catalog_ip_names is not None:
                    for name in catalog_ip_names:
                        _add_item(name)
                elif scan_root.is_dir():
                    for entry in scan_root.iterdir():
                        if _looks_like_project_ip(entry):
                            _add_item(entry.name, mtime=entry.stat().st_mtime)
                items = sorted(by_name.values(), key=lambda x: (-x["mtime"], x["name"]))
                if workspace_root_for_list is not None:
                    return JSONResponse({
                        "project_root": str(PROJECT_ROOT),
                        "workspace_root": str(workspace_root_for_list),
                        "session_id": owner or "",
                        "workspace_session": requested_workspace_session or "",
                        "items": items,
                        "count": len(items),
                        "source": "db_ip_blocks" if catalog_ip_names is not None else "filesystem",
                    })
            if session_root is None or not session_root.is_dir():
                items = sorted(by_name.values(), key=lambda x: (-x["mtime"], x["name"]))
                return JSONResponse({
                    "project_root": str(PROJECT_ROOT),
                    "workspace_root": str(workspace_root_for_list) if workspace_root_for_list is not None else "",
                    "session_id": owner or "",
                    "workspace_session": requested_workspace_session or "",
                    "items": items,
                    "count": len(items),
                })
            try:
                session_root.relative_to((PROJECT_ROOT / ".session").resolve())
            except ValueError:
                return JSONResponse({"error": "session path escapes .session", "items": []}, status_code=400)
            for entry in session_root.iterdir():
                if not entry.is_dir():
                    continue
                name = entry.name
                if name.startswith(".") or name in skip:
                    continue
                workflows = sorted(
                    child.name for child in entry.iterdir()
                    if child.is_dir() and not child.name.startswith(".")
                )
                _add_item(name, workflows=workflows, mtime=entry.stat().st_mtime)
        except OSError as exc:
            return JSONResponse({"error": str(exc), "items": []}, status_code=500)
        items = sorted(by_name.values(), key=lambda x: (-x["mtime"], x["name"]))
        return JSONResponse({
            "project_root": str(PROJECT_ROOT),
            "workspace_root": str(workspace_root_for_list) if workspace_root_for_list is not None else "",
            "session_id": owner or "",
            "workspace_session": requested_workspace_session or "",
            "items": items,
            "count": len(items),
        })

    # Phase 15: /api/ssot-gates/{ip} (~497 lines) extracted to
    # atlas_api_ssot_gates.py. Same factory pattern as Phase 14.
    from src.atlas_api_ssot_gates import register_ssot_gates_routes as _register_ssot_gates_routes
    _register_ssot_gates_routes(app, PROJECT_ROOT=PROJECT_ROOT, _safe=_safe)

    # Phase 11b: sim_debug API cluster (5 endpoints + 2 helpers + _ELAB_CACHE)
    # moved to src/atlas_api_sim_debug.py via register_sim_debug_routes.
    # 3 closure captures injected as kwargs; everything else is local to
    # the moved bodies.
    from src.atlas_api_sim_debug import register_sim_debug_routes as _register_sim_debug_routes
    _register_sim_debug_routes(
        app,
        _safe=_safe,
        PROJECT_ROOT=PROJECT_ROOT,
        WORKFLOW_ROOT=WORKFLOW_ROOT,
        fs_authz=_fs_authz,
    )

    # ── Source endpoint — sim_debug signal→driver + cocotb test view ─
    # Accepts SV/V plus the text extensions that show up in the
    # cocotb tab (Python tests, sequences, agents, env, Makefile,
    # YAML/JSON, etc.). Rejects binaries to avoid shipping .vvp / .out
    # contents over WS.
    _SOURCE_EXTS = {
        ".sv", ".v", ".svh", ".vh",          # SystemVerilog
        ".py",                                # cocotb / Python testbench
        ".sdc", ".tcl", ".f",                 # constraints / filelists
        ".yaml", ".yml", ".json", ".md",      # config / docs
        ".txt", ".log", ".rpt",               # reports
        ".sh", ".bash",                       # scripts
        ".c", ".h", ".cpp", ".hpp",           # firmware
        ".xml",                               # results.xml
    }
    _SOURCE_NO_EXT_NAMES = {"Makefile", "makefile", "Dockerfile"}

    # Phase 11a: /api/source moved to src/atlas_api_sim_debug.py via factory.
    from src.atlas_api_sim_debug import register_source_route as _register_source_route
    _register_source_route(app, safe_path_fn=_safe, fs_authz=_fs_authz)

    # ── sim_debug elab module loader ─────────────────────────────
    # Lives at workflow/sim_debug/elab.py — co-located with the rest
    # of the sim_debug workspace (system_prompt.md, commands/, rules/,
    # scripts/). Loaded via importlib so we don't have to add
    # workflow/sim_debug/ to sys.path globally.

    # ── Elab endpoints (pyslang / Verilator / slang) — sim_debug hierarchy + trace ─




    # ── cocotb / TB env browsing — sim_debug "TB" tab ─────────────

    def _same_todo_path(left: Path | None, right: Path | None) -> bool:
        if left is None or right is None:
            return False
        try:
            return left.expanduser().resolve() == right.expanduser().resolve()
        except Exception:
            return str(left) == str(right)

    def _todo_session_for_request(
        request: Request,
        requested: str | None = None,
    ) -> tuple[str, JSONResponse | None]:
        session = normalize_session_name(requested or "")
        if not session:
            session = _request_active_session_for_user(request)
        user = request.scope.get("user") or {}
        username = normalize_session_name(str(user.get("username") or ""))
        owner = session.split("/", 1)[0] if session else ""
        if username and owner and owner != username:
            return "", JSONResponse({"error": "session owner mismatch"}, status_code=403)
        return session, None

    def _todo_file_for_session(session_name: str) -> Path | None:
        session = normalize_session_name(session_name or "")
        if not session:
            return None
        parts = [part for part in session.split("/") if part]
        if len(parts) > 4:
            return None
        if len(parts) == 4:
            try:
                context = AtlasContext.from_session_key(
                    session,
                    atlas_root=os.environ.get("ATLAS_ROOT") or str(PROJECT_ROOT),
                )
                return context.session_dir / "todo.json"
            except ValueError:
                return None
        return PROJECT_ROOT / ".session" / session / "todo.json"

    @app.post("/api/todos/clear")
    async def api_todos_clear(request: Request):
        """Clear the todo file for the requested active session."""
        try:
            body = await request.json()
        except Exception:
            body = {}
        session_name, denied = _todo_session_for_request(
            request,
            str((body or {}).get("session") or request.query_params.get("session") or ""),
        )
        if denied is not None:
            return denied
        session_todo = _todo_file_for_session(session_name)
        import os as _os
        _os.environ.pop("TODO_TEMPLATE_LOCK_ADDITIONS", None)
        _os.environ.pop("TODO_TEMPLATE_LOCK_NAME", None)
        try:
            import main as _main  # noqa: WPS433
            tt = getattr(_main, "todo_tracker", None)
            persist = getattr(tt, "_persist_path", None) if tt is not None else None
            live_path = Path(persist) if persist else None
            if (
                tt is not None
                and hasattr(tt, "todos")
                and (session_todo is None or _same_todo_path(live_path, session_todo))
            ):
                tt.todos = []
                if hasattr(tt, "current_index"):
                    tt.current_index = -1
                if hasattr(tt, "save"):
                    try: tt.save()
                    except Exception: pass
        except Exception:
            pass
        if session_todo is not None:
            try:
                session_todo.parent.mkdir(parents=True, exist_ok=True)
                session_todo.write_text('{"todos": []}', encoding="utf-8")
            except Exception:
                pass
        else:
            try:
                from pathlib import Path as _P
                for cand in ("current_todos.json",
                             str(_P.home() / ".common_ai_agent" / "current_todos.json")):
                    p = _P(cand)
                    if p.exists():
                        try: p.unlink()
                        except Exception: pass
            except Exception:
                pass
        return JSONResponse({"ok": True, "session": session_name})

    def _sync_live_tracker_from_session(session_todo: "Path | None") -> None:
        """Reload main.todo_tracker from the session file when it points at
        the same path the editing endpoints just wrote. Mirrors the live-sync
        the clear endpoint performs so an in-flight agent loop sees edits."""
        if session_todo is None:
            return
        try:
            import main as _main  # noqa: WPS433
            from lib.todo_tracker import TodoTracker as _TT
            tt = getattr(_main, "todo_tracker", None)
            persist = getattr(tt, "_persist_path", None) if tt is not None else None
            live_path = Path(persist) if persist else None
            if (
                tt is not None
                and hasattr(tt, "todos")
                and _same_todo_path(live_path, session_todo)
            ):
                fresh = _TT.load(session_todo)
                tt.todos = fresh.todos
                if hasattr(tt, "current_index"):
                    tt.current_index = fresh.current_index
                if hasattr(tt, "save"):
                    try: tt.save()
                    except Exception: pass
        except Exception:
            pass

    def _bind_live_tracker_to_session(
        session_name: str,
        session_todo: Path,
        tracker: Any,
    ) -> list[str]:
        errors: list[str] = []
        session_name = normalize_session_name(session_name)
        session_todo = session_todo.expanduser().resolve(strict=False)
        session_dir = session_todo.parent
        session_cfg = {
            "TODO_FILE": str(session_todo),
            "TODO_ERROR_FILE": str(session_dir / "todo_error.json"),
            "SESSION_DIR": str(session_dir),
            "ACTIVE_PROJECT": session_name,
            "ATLAS_ACTIVE_SESSION": session_name,
        }
        for cfg_name in ("config", "src.config"):
            try:
                cfg_mod = importlib.import_module(cfg_name)
            except ImportError as exc:
                errors.append(f"{cfg_name}: {exc}")
                continue
            for key, value in session_cfg.items():
                setattr(cfg_mod, key, value)
        try:
            import lib.todo_tracker as todo_tracker_mod  # noqa: WPS433
        except ImportError as exc:
            errors.append(f"lib.todo_tracker: {exc}")
        else:
            todo_tracker_mod.TODO_FILE = session_todo
        if hasattr(tracker, "_persist_path"):
            tracker._persist_path = session_todo
        for module_name in ("main", "src.main", "__main__"):
            module = sys.modules.get(module_name)
            if module is not None:
                setattr(module, "todo_tracker", tracker)
        return errors

    @app.post("/api/todos/add")
    async def api_todos_add(request: Request):
        """Append one todo to the requested session's todo file."""
        try:
            body = await request.json()
        except Exception:
            body = {}
        body = body or {}
        session_name, denied = _todo_session_for_request(
            request,
            str(body.get("session") or request.query_params.get("session") or ""),
        )
        if denied is not None:
            return denied
        if not session_name:
            return JSONResponse({"error": "no active session"}, status_code=400)
        content = str(body.get("content") or "").strip()
        if not content:
            return JSONResponse({"error": "content is required"}, status_code=400)
        detail = str(body.get("detail") or "").strip()
        if not detail:
            return JSONResponse({"error": "detail is required"}, status_code=400)
        criteria = str(body.get("criteria") or "").strip()
        if not criteria:
            return JSONResponse({"error": "criteria is required"}, status_code=400)
        session_todo = _todo_file_for_session(session_name)
        if session_todo is None:
            return JSONResponse({"error": "invalid session"}, status_code=400)
        try:
            from lib.todo_tracker import TodoTracker
            tracker = TodoTracker.load(session_todo)
            existing = tracker.to_dict().get("todos", [])
            new_todo = {
                "content": content,
                "activeForm": str(body.get("activeForm") or "").strip() or content,
                "status": "pending",
                "priority": str(body.get("priority") or "medium"),
                "detail": detail,
                "criteria": criteria,
            }
            # Optional `index` inserts the todo in the MIDDLE at that 0-based
            # position; out-of-range or omitted index appends to the end.
            raw_index = body.get("index")
            if raw_index is None or str(raw_index).strip() == "":
                existing.append(new_todo)
            else:
                try:
                    pos = int(raw_index)
                except (TypeError, ValueError):
                    pos = len(existing)
                if pos < 0:
                    pos = 0
                if pos > len(existing):
                    pos = len(existing)
                existing.insert(pos, new_todo)
            tracker.add_todos(existing)
            tracker.save()
        except Exception as e:
            return JSONResponse({"error": str(e)}, status_code=500)
        # Adding a todo is an editor action. Do not sync it into the live
        # in-flight worker tracker here, or an active loop may immediately
        # consume the new pending task without an explicit "continue" turn.
        return JSONResponse(_gate_for_workflow(tracker.to_dict(), session_name))

    @app.post("/api/todos/update")
    async def api_todos_update(request: Request):
        """Modify the todo at 0-based `index` in the requested session."""
        try:
            body = await request.json()
        except Exception:
            body = {}
        body = body or {}
        session_name, denied = _todo_session_for_request(
            request,
            str(body.get("session") or request.query_params.get("session") or ""),
        )
        if denied is not None:
            return denied
        if not session_name:
            return JSONResponse({"error": "no active session"}, status_code=400)
        try:
            index = int(body.get("index"))
        except (TypeError, ValueError):
            return JSONResponse({"error": "index must be an integer"}, status_code=400)
        session_todo = _todo_file_for_session(session_name)
        if session_todo is None:
            return JSONResponse({"error": "invalid session"}, status_code=400)
        try:
            from lib.todo_tracker import TodoTracker
            tracker = TodoTracker.load(session_todo)
            if not (0 <= index < len(tracker.todos)):
                return JSONResponse({"error": "index out of range"}, status_code=400)
            todo = tracker.todos[index]
            if "content" in body and body.get("content") is not None:
                new_content = str(body.get("content")).strip()
                if not new_content:
                    return JSONResponse({"error": "content cannot be empty"}, status_code=400)
                todo.content = new_content
            if "detail" in body and body.get("detail") is not None:
                todo.detail = str(body.get("detail"))
            if "criteria" in body and body.get("criteria") is not None:
                todo.criteria = str(body.get("criteria"))
            if "priority" in body and body.get("priority") is not None:
                todo.priority = str(body.get("priority"))
            if "activeForm" in body and body.get("activeForm") is not None:
                todo.active_form = str(body.get("activeForm"))
            if "approved_reason" in body and body.get("approved_reason") is not None:
                todo.approved_reason = str(body.get("approved_reason"))
            if "approvedReason" in body and body.get("approvedReason") is not None:
                todo.approved_reason = str(body.get("approvedReason"))
            if "rejection_reason" in body and body.get("rejection_reason") is not None:
                todo.rejection_reason = str(body.get("rejection_reason"))
            if "rejectionReason" in body and body.get("rejectionReason") is not None:
                todo.rejection_reason = str(body.get("rejectionReason"))
            if "state" in body and body.get("state") is not None:
                from lib.todo_transition import apply_todo_status_transition, normalize_todo_status
                raw_state = str(body.get("state")).strip()
                new_status = normalize_todo_status(raw_state)
                if new_status != getattr(todo, "status", None):
                    transition_reason = str(
                        body.get("reason")
                        or body.get("approved_reason")
                        or body.get("approvedReason")
                        or body.get("rejection_reason")
                        or body.get("rejectionReason")
                        or ""
                    )
                    transition = apply_todo_status_transition(
                        tracker,
                        index,
                        new_status,
                        reason=transition_reason,
                        plan_mode=os.environ.get("PLAN_MODE") == "true",
                    )
                    if not transition.ok:
                        return JSONResponse(
                            {"error": transition.message, "code": transition.code},
                            status_code=transition.http_status or 409,
                        )
                else:
                    todo.status = new_status
            if todo.status == "approved" and not str(getattr(todo, "approved_reason", "") or "").strip():
                return JSONResponse({"error": "approved_reason is required"}, status_code=400)
            if todo.status == "rejected" and not str(getattr(todo, "rejection_reason", "") or "").strip():
                return JSONResponse({"error": "rejection_reason is required"}, status_code=400)
            tracker.save()
        except Exception as e:
            return JSONResponse({"error": str(e)}, status_code=500)
        _bind_live_tracker_to_session(session_name, session_todo, tracker)
        return JSONResponse(_gate_for_workflow(tracker.to_dict(), session_name))

    @app.post("/api/todos/remove")
    async def api_todos_remove(request: Request):
        """Remove the todo at 0-based `index` in the requested session."""
        try:
            body = await request.json()
        except Exception:
            body = {}
        body = body or {}
        session_name, denied = _todo_session_for_request(
            request,
            str(body.get("session") or request.query_params.get("session") or ""),
        )
        if denied is not None:
            return denied
        if not session_name:
            return JSONResponse({"error": "no active session"}, status_code=400)
        try:
            index = int(body.get("index"))
        except (TypeError, ValueError):
            return JSONResponse({"error": "index must be an integer"}, status_code=400)
        session_todo = _todo_file_for_session(session_name)
        if session_todo is None:
            return JSONResponse({"error": "invalid session"}, status_code=400)
        try:
            from lib.todo_tracker import TodoTracker
            tracker = TodoTracker.load(session_todo)
            if not (0 <= index < len(tracker.todos)):
                return JSONResponse({"error": "index out of range"}, status_code=400)
            del tracker.todos[index]
            if tracker.current_index == index:
                tracker.current_index = -1
            elif tracker.current_index > index:
                tracker.current_index -= 1
            tracker.save()
        except Exception as e:
            return JSONResponse({"error": str(e)}, status_code=500)
        _bind_live_tracker_to_session(session_name, session_todo, tracker)
        return JSONResponse(_gate_for_workflow(tracker.to_dict(), session_name))

    def _tracker_stage(name: str) -> str:
        # `derive_rtl_todos.py` writes tracker name as "<ip>-rtl"; other
        # generators follow the same `<ip>-<stage>` convention so we can
        # cheaply read the trailing stage off the name.
        s = (name or "").rsplit("-", 1)
        return s[-1] if len(s) == 2 else ""

    def _active_workflow_stage(session_name: str = "") -> str:
        sess = session_name or _active_session_value() or ""
        parts = [p for p in sess.split("/") if p]
        if len(parts) < 3:
            return ""
        wf = parts[-1]
        return wf.split("-")[0] if "-" in wf else wf

    def _gate_for_workflow(d: dict | None, session_name: str = "") -> dict | None:
        # Tracker carries an SSOT-derived workflow stage in its name.
        # When the user is viewing a different workflow we hide the
        # tracker so a previous /gen-rtl run doesn't keep showing 20
        # auto-generated rtl-gen TODOs in sim_debug or other workflows.
        if not isinstance(d, dict):
            return d
        ts = _tracker_stage(d.get("name", ""))
        ws = _active_workflow_stage(session_name)
        if ts and ws and ts != ws:
            return {"todos": [], "auto_hidden": True,
                    "reason": f"tracker '{d.get('name')}' is for {ts}, active workflow is {ws}"}
        return d

    @app.get("/api/todos")
    async def api_todos(request: Request, session: str = ""):
        # Prefer the live tracker the agent is mutating in main.py — that's
        # the only way to see in-progress changes before they hit disk. Fall
        # back to the on-disk file if main hasn't initialized one yet. When
        # ATLAS has an active namespaced session, that session's todo.json is
        # the source of truth; a process-global live tracker may still point at
        # an older HOME-level current_todos.json from another IP.
        candidates: list[Path] = []
        active_session, denied = _todo_session_for_request(request, session)
        if denied is not None:
            return denied
        active_todo_path = _todo_file_for_session(active_session)

        try:
            import main as _main  # noqa: WPS433
            live = getattr(_main, "todo_tracker", None)
            live_persist_path = (
                getattr(live, "_persist_path", None)
                if live is not None
                else None
            )
            live_path = Path(live_persist_path) if live_persist_path else None
            if live is not None and getattr(live, "todos", None):
                if (
                    not active_todo_path
                    or _same_todo_path(live_path, active_todo_path)
                ):
                    return JSONResponse(_gate_for_workflow(live.to_dict(), active_session))
            if active_todo_path:
                candidates.append(active_todo_path)
            elif live_path:
                candidates.append(live_path)
        except Exception:
            pass
        # On-disk fallback. Two persistence paths exist in this repo:
        #   1. <PROJECT_ROOT>/current_todos.json    (relative TODO_FILE,
        #                                            agent's actual writes)
        #   2. ~/.common_ai_agent/current_todos.json (HOME default for
        #                                            stand-alone scripts)
        # When `import config` succeeds at module load, TodoTracker module
        # caches TODO_FILE as Path("current_todos.json") — relative — and
        # whether `.exists()` resolves depends on the server's cwd at
        # request time. Resolve them *both* explicitly so the panel never
        # silently falls through to "no todos" when the file is right
        # there in PROJECT_ROOT.
        try:
            import json as _json
            from lib.todo_tracker import TodoTracker
            if not active_todo_path:
                try:
                    import config as _cfg
                    cfg_todo = Path(str(getattr(_cfg, "TODO_FILE", "current_todos.json")))
                    candidates.append(cfg_todo if cfg_todo.is_absolute() else PROJECT_ROOT / cfg_todo)
                except Exception:
                    pass
                candidates.extend([
                    PROJECT_ROOT / "current_todos.json",
                    Path.cwd() / "current_todos.json",
                    Path.home() / ".common_ai_agent" / "current_todos.json",
                ])
            deduped: list[Path] = []
            seen_paths: set[str] = set()
            for cand in candidates:
                try:
                    key = str(cand.expanduser().resolve())
                except Exception:
                    key = str(cand)
                if key not in seen_paths:
                    seen_paths.add(key)
                    deduped.append(cand)
            picked = next((p for p in deduped if p.exists()), None)
            if picked is None:
                return JSONResponse({"todos": []})
            raw = None
            try:
                raw = _json.loads(picked.read_text(encoding="utf-8"))
            except Exception:
                raw = None
            tt = TodoTracker.load(picked)
            d = tt.to_dict()
            if d.get("todos"):
                return JSONResponse(_gate_for_workflow(d, active_session))
            d = _atlas_todo_payload_from_raw(raw, d)
            return JSONResponse(_gate_for_workflow(d, active_session))
        except Exception as e:
            return JSONResponse({"error": str(e)}, status_code=500)

    @app.get("/api/commands")
    async def api_commands():
        """List every slash command currently registered, including the
        workspace-specific ones (e.g. /grill-me, /to-ssot for ssot-gen).
        """
        try:
            _ensure_default_workspace_commands_registered()
            from core.slash_commands import get_registry as _gr
            reg = _gr()
        except Exception as e:
            return JSONResponse({"error": str(e)}, status_code=500)
        # The registry stores commands in an internal dict; read it
        # defensively through whatever public surface is available.
        cmds = []
        seen = set()
        for attr in ("commands", "_commands"):
            entries = getattr(reg, attr, None)
            if isinstance(entries, dict):
                for name, spec in entries.items():
                    canonical = spec.get("name", name) if isinstance(spec, dict) else name
                    if canonical in seen:
                        continue
                    seen.add(canonical)
                    if isinstance(spec, dict):
                        cmds.append({
                            "cmd":     "/" + canonical,
                            "name":    canonical,
                            "aliases": spec.get("aliases", []) or [],
                            "hint":    spec.get("description", "") or "",
                            "usage":   spec.get("usage", f"/{canonical}"),
                        })
                break
        # Inline workflow slashes (handled in this file, not in the
        # core registry) need to surface in the UI picker too — without
        # this merge, /to-ssot, /ssot-rtl, /grill-me, etc. silently
        # vanish from autocomplete even though typing them still works.
        _INLINE_WORKFLOW_HINTS = {
            "wf":              ("workflow",   "switch active workflow: /wf ssot-gen"),
            "workflow":        ("",           "switch active workflow: /workflow ssot-gen"),
            "ip":              ("use",        "switch active IP: /ip uart"),
            "use":             ("",           "switch active IP: /use uart"),
            "session":         ("",           "show or switch session: /session default"),
            "new-ip":          ("ni",         "create a new IP block: /new-ip <name>"),
            "import":          ("imp",        "import an external doc / spec into the active IP"),
            "grill-me":        ("grill,g",    "Q&A the LLM with the active SSOT context"),
            "to-ssot":         ("ssot,ts",    "approve grill answers → write SSOT yaml"),
            "verify-ssot":     ("verify_ssot,vs", "verify SSOT YAML shape, Preview fields, and gates"),
            "validate-yaml":   ("",           "validate the active SSOT yaml against schema"),
            "ssot-fl-model":   ("sfm",        "generate functional / cycle model from SSOT"),
            "ssot-equiv-goals":("equiv-goals,seg", "generate FL↔RTL equivalence goals"),
            "repair-equiv":    ("repair-equivalence,reqv", "repair equivalence-goal failures"),
            "ssot-rtl":        ("sr",         "generate RTL from SSOT"),
            "gen-rtl":         ("",           "generate RTL from SSOT"),
            "repair-rtl":      ("rrtl",       "repair RTL gate failures"),
            "lint":            ("l",          "run lint pass on the active IP"),
            "tb":              ("",           "generate testbench (auto-select format)"),
            "gen-tb":          ("gt",         "generate TB from SSOT contract ledger"),
            "ssot-tb":         ("stb",        "generate cocotb testbench from SSOT"),
            "ssot-tb-cocotb":  ("stb-cocotb", "generate cocotb testbench"),
            "ssot-tb-uvm":     ("stb-uvm",    "generate UVM testbench"),
            "ssot-tb-verilog": ("stb-verilog,ssot-tb-sv,stb-sv", "generate Verilog/SV testbench"),
            "sim":             ("s",          "run simulation"),
            "sim-debug":       ("sd",         "FL↔RTL mismatch debug pass"),
            "coverage":        ("cov",        "run / iterate coverage"),
            "contract-check":   ("contract",   "run requirement contract evidence gate"),
            "goal-audit":      ("audit,ga",   "audit goal coverage and evidence"),
            "signoff":         ("",           "final signoff review"),
            "feedback":        ("fb",         "send admin-visible feedback: /feedback <message>"),
        }
        for name, (alias_str, hint) in _INLINE_WORKFLOW_HINTS.items():
            if name in seen:
                continue
            seen.add(name)
            aliases = [a for a in alias_str.split(",") if a] if alias_str else []
            cmds.append({
                "cmd":     "/" + name,
                "name":    name,
                "aliases": aliases,
                "hint":    hint,
                "usage":   f"/{name}",
            })

        cmds.sort(key=lambda c: c["name"])
        return JSONResponse({"commands": cmds})

    @app.get("/api/input-history")
    async def api_input_history(limit: int = 200):
        try:
            history = _read_input_history(limit)
            return JSONResponse({
                "history": history,
                "path": _relative_project_path(_input_history_path()),
            })
        except Exception as e:
            return JSONResponse({"error": str(e), "history": []}, status_code=500)

    @app.post("/api/input-history")
    async def api_append_input_history(request: Request):
        try:
            payload = await request.json()
        except Exception:
            payload = {}
        text = str(payload.get("text") or "").strip()
        if not text:
            return JSONResponse({"ok": True, "stored": False})
        try:
            _append_input_history(text)
            return JSONResponse({"ok": True, "stored": True})
        except Exception as e:
            return JSONResponse({"ok": False, "error": str(e)}, status_code=500)

    # SSOT routes registered via register_ssot_routes() below (see atlas_api_ssot.py).
    # /api/session/activate registered via register_sessions_routes() (see atlas_api_sessions.py).

    # POST /api/ssot/qa/answer registered via register_ssot_routes() (see atlas_api_ssot.py).

    _soc_cache_lock = threading.Lock()
    _soc_build_lock = threading.Lock()
    _soc_cache: dict[tuple[str, str, str], tuple[float, dict[str, Any]]] = {}
    try:
        _SOC_CACHE_TTL_SEC = max(
            0.0,
            float(os.environ.get("ATLAS_SOC_CACHE_TTL_SEC", "15.0") or "15.0"),
        )
    except ValueError:
        _SOC_CACHE_TTL_SEC = 15.0

    def _soc_cache_get(key: tuple[str, str, str]) -> dict[str, Any] | None:
        if _SOC_CACHE_TTL_SEC <= 0:
            return None
        now = time.monotonic()
        with _soc_cache_lock:
            hit = _soc_cache.get(key)
            if not hit:
                return None
            ts, payload = hit
            if now - ts > _SOC_CACHE_TTL_SEC:
                _soc_cache.pop(key, None)
                return None
            return payload

    def _soc_cache_set(key: tuple[str, str, str], payload: dict[str, Any]) -> None:
        if _SOC_CACHE_TTL_SEC <= 0:
            return
        with _soc_cache_lock:
            _soc_cache[key] = (time.monotonic(), payload)

    # Phase 14: /api/soc + nested _build_module (~3000 lines) extracted to
    # atlas_api_soc.py. Factory takes the 6 closure captures as kwargs
    # whose names mirror the original locals — zero body edits needed.
    # Inline import (matches sim_debug factory wiring at line 4163) keeps
    # this surface self-contained inside create_app() so test-time monkey
    # patches of atlas_ui module globals still observe a single import path.
    #
    # Lambda forward-refs: _load_ssot_state and _valid_ip_name are defined
    # further down in create_app (lines ~4745/~4749 → which now hold the
    # factory call). Wrapping each callable in a lambda defers the name
    # resolution to call time (when api_soc actually handles a request),
    # by which point every closure local is bound. Same pattern as the
    # Phase 13d preview-pane forward-refs and Phase 10b atlas_qa wiring.
    from src.atlas_api_soc import register_soc_routes as _register_soc_routes
    # Capture the returned api_soc handler so cross-route Python-level
    # callers like api_progress (line ~4093) can keep invoking
    # api_soc(scope=…, ip=…) directly without going through HTTP.
    api_soc = _register_soc_routes(
        app,
        SKIP_DIRS=SKIP_DIRS,
        _load_ssot_state=lambda *a, **k: _load_ssot_state(*a, **k),
        _soc_build_lock=_soc_build_lock,
        _soc_cache_get=lambda *a, **k: _soc_cache_get(*a, **k),
        _soc_cache_set=lambda *a, **k: _soc_cache_set(*a, **k),
        _valid_ip_name=lambda *a, **k: _valid_ip_name(*a, **k),
        # PROJECT_ROOT + SOURCE_ROOT are atlas_ui module globals — both
        # finalized by main()/import time before create_app() runs, so
        # the captured value-at-init stays current for the route lifetime.
        PROJECT_ROOT=PROJECT_ROOT,
        SOURCE_ROOT=SOURCE_ROOT,
    )

    @app.get("/api/progress")
    def api_progress(scope: str = "", ip: str = ""):
        """Return SSOT-derived implementation progress for the Atlas sidebar.

        The heavy lifting already lives in /api/soc because the architect
        canvas needs the same SSOT/RTL/TB/sim evidence. This endpoint flattens
        that structure into a compact shape for the normal chat workspace:
        one selected module plus the full module list. All metrics are derived
        from the canonical leaf SSOT YAML and disk artifacts, not from fixed IP
        templates or assistant prose.
        """
        resp = api_soc(scope=scope, ip=ip)
        try:
            data = json.loads(resp.body.decode("utf-8"))
        except Exception as e:
            return JSONResponse({"error": f"soc progress parse: {e}", "modules": []}, status_code=500)

        modules: list[dict[str, Any]] = []
        for cluster in data.get("clusters", []) if isinstance(data, dict) else []:
            if not isinstance(cluster, dict):
                continue
            for mod in cluster.get("modules", []) or []:
                if not isinstance(mod, dict):
                    continue
                entry = {
                    "id": mod.get("id") or mod.get("name") or "",
                    "name": mod.get("name") or mod.get("id") or "",
                    "label": mod.get("label") or mod.get("name") or mod.get("id") or "",
                    "kind": mod.get("kind") or "",
                    "ip_dir": mod.get("ip_dir") or "",
                    "ssot_path": mod.get("ssot_path") or "",
                    "status": mod.get("status") or {},
                    "status_detail": mod.get("status_detail") or {},
                    "status_source": mod.get("status_source") or {},
                    "artifact_status": mod.get("artifact_status") or {},
                    "artifact_detail": mod.get("artifact_detail") or {},
                    "artifact_source": mod.get("artifact_source") or {},
                    "progress": mod.get("progress") or {},
                    "signoff": mod.get("signoff") or {},
                    "simple_summary": (
                        mod.get("simple_summary")
                        or (mod.get("signoff") or {}).get("simple_summary")
                        or {}
                    ),
                }
                modules.append(entry)

        want = (ip or scope or "").strip().strip("/")
        selected = None
        if want:
            selected = next((
                m for m in modules
                if want in {str(m.get("id") or ""), str(m.get("name") or ""), str(m.get("ip_dir") or "")}
            ), None)
            if selected is None:
                selected = next((
                    m for m in modules
                    if str(m.get("ip_dir") or "").startswith(want + "/")
                    or str(m.get("ssot_path") or "").startswith(want + "/")
                ), None)
        if selected is None and modules:
            selected = modules[0]

        return JSONResponse({
            "project": data.get("name") if isinstance(data, dict) else PROJECT_ROOT.name,
            "source": data.get("source") if isinstance(data, dict) else "",
            "scope": want,
            "selected": selected,
            "modules": modules,
            "module_count": len(modules),
        })

    # ── Jobs (HTTP-worker dispatch tracker) ────────────────────────
    # Routes + state live in src/atlas_api_jobs.py (phase 6 of split).
    # register_jobs_routes() is called below near the other registrars.

    _WORKFLOW_SLASHES = {
        "/wf", "/workflow",
        "/ip", "/use",
        "/session",
        "/new-ip", "/ni",
        "/import", "/imp",
        "/grill-me", "/grill", "/g",
        "/to-ssot", "/ssot", "/ts",
        "/validate-yaml",
        "/ssot-fl-model", "/sfm",
        "/ssot-equiv-goals", "/equiv-goals", "/seg",
        "/repair-equiv", "/repair-equivalence", "/reqv",
        "/ssot-rtl", "/sr", "/gen-rtl",
        "/repair-rtl", "/rrtl",
        "/lint", "/l",
        "/tb", "/gen-tb", "/gt",
        "/ssot-tb", "/stb",
        "/ssot-tb-cocotb", "/stb-cocotb",
        "/ssot-tb-uvm", "/stb-uvm",
        "/ssot-tb-verilog", "/stb-verilog", "/ssot-tb-sv", "/stb-sv",
        "/sim", "/s",
        "/sim-debug", "/sd",
        "/coverage", "/cov",
        "/contract-check", "/contract",
        "/goal-audit", "/audit", "/ga",
        "/signoff",
    }

    _STAGE_RUNNERS = {
        "ssot-rtl": {
            "workflow": "rtl-gen",
            "template": "ssot-rtl",
            "artifact_hint": "rtl/",
        },
        "ssot-fl-model": {
            "workflow": "fl-model-gen",
            "template": "ssot-fl-model",
            "artifact_hint": "model/ and cov/fcov_plan.json",
        },
        "ssot-equiv-goals": {
            "workflow": "fl-model-gen",
            "template": "ssot-equiv-goals",
            "artifact_hint": "verify/equivalence_goals.json",
        },
        "lint": {
            "workflow": "lint",
            "template": "lint-fix",
            "artifact_hint": "lint/dut_lint.json",
        },
        "ssot-tb": {
            "workflow": "tb-gen",
            "template": "ssot-tb-cocotb",
            "artifact_hint": "tb/cocotb/ and sim/",
        },
        "gen-tb": {
            "workflow": "tb-gen",
            "template": "ssot-tb-cocotb",
            "artifact_hint": "tb/tb_todo_plan.json and tb/cocotb/",
        },
        "ssot-tb-cocotb": {
            "workflow": "tb-gen",
            "template": "ssot-tb-cocotb",
            "artifact_hint": "tb/cocotb/ and sim/",
        },
        "ssot-tb-uvm": {
            "workflow": "tb-gen",
            "template": "ssot-tb-uvm",
            "artifact_hint": "tb/uvm/ and sim/",
        },
        "ssot-tb-verilog": {
            "workflow": "tb-gen",
            "template": "ssot-tb-verilog",
            "artifact_hint": "tb/tb_*.sv and sim/",
        },
        "sim": {
            "workflow": "sim",
            "template": "sim-debug",
            "artifact_hint": "sim/results.xml, sim/scoreboard_events.jsonl, and waveform/coverage artifacts",
        },
        "sim-debug": {
            "workflow": "sim_debug",
            "template": "sim-debug",
            "artifact_hint": "sim/fl_rtl_compare.json and sim/mismatch_classification.json",
        },
        "coverage": {
            "workflow": "coverage",
            "template": "coverage_iter",
            "artifact_hint": "cov/coverage.json and sim/coverage_report.md",
        },
        "contract-check": {
            "workflow": "contract-reflection",
            "template": "contract-check",
            "artifact_hint": "signoff/contract_check.json",
        },
        "goal-audit": {
            "workflow": "sim_debug",
            "template": "sim-debug",
            "artifact_hint": "sim/fl_rtl_goal_audit.json",
        },
        "signoff": {
            "workflow": "sim_debug",
            "template": "sim-debug",
            "artifact_hint": "ATLAS /api/progress signoff gate",
        },
    }

    _SSOT_REQUIRED_DECISIONS = [
        ("purpose", "IP purpose / one sentence behavior"),
        ("bus_interface", "bus interface and role, e.g. APB4 slave"),
        ("register_map", "register map, address offsets, access policies"),
        ("clock_reset", "clock/reset names, frequency, reset polarity"),
        ("interrupt", "interrupt behavior, or explicit none"),
        ("memory_map", "memory map/base address requirement, or explicit none"),
        ("parameters", "parameters and defaults, or explicit none"),
        ("submodule_structure", "leaf submodule hierarchy and ownership"),
        ("test_expectation", "minimum cocotb/pyuvm TB/SIM acceptance expectations"),
    ]

    _SSOT_IMPORT_SECTION_TODO_SPECS = [
        ("top_module", "00 Top Module Identity", ("purpose", "submodule_structure")),
        ("sub_modules", "01 Sub-Module List", ("submodule_structure", "purpose")),
        ("parameters", "02 Parameters", ("parameters", "clock_reset", "memory_map")),
        ("io_list", "03 IO List", ("bus_interface", "clock_reset", "interrupt")),
        ("features", "04 Main Features", ("purpose", "register_map", "memory_map", "interrupt")),
        ("dataflow", "05 Data Flow", ("purpose", "bus_interface", "memory_map", "submodule_structure")),
        ("function_model", "06 Function Model", ("purpose", "register_map", "memory_map", "test_expectation")),
        ("cycle_model", "07 Cycle Model", ("clock_reset", "bus_interface", "test_expectation")),
        ("clock_reset_domains", "08 Clock & Reset Domain", ("clock_reset",)),
        ("cdc_requirements", "09 CDC Requirements", ("clock_reset", "bus_interface")),
        ("rdc_requirements", "10 RDC Requirements", ("clock_reset",)),
        ("registers", "11 Registers", ("register_map", "bus_interface")),
        ("memory", "12 Memory Requirements", ("memory_map", "parameters")),
        ("interrupts", "13 Interrupt", ("interrupt", "register_map")),
        ("fsm", "14 FSM", ("submodule_structure", "purpose", "test_expectation")),
        ("timing", "15 Timing & Performance", ("clock_reset", "parameters", "test_expectation")),
        ("power", "16 Power Intent", ("clock_reset", "parameters")),
        ("security", "17 Security & Safety", ("purpose", "bus_interface", "register_map")),
        ("error_handling", "18 Error Handling", ("interrupt", "register_map", "test_expectation")),
        ("debug_observability", "19 Debug & Observability", ("test_expectation", "interrupt", "clock_reset")),
        ("integration", "20 Integration Contract", ("bus_interface", "memory_map", "submodule_structure")),
        ("dft", "21 DFT / DFD", ("test_expectation", "clock_reset")),
        ("synthesis", "22 Synthesis / Implementation Constraints", ("parameters", "clock_reset", "submodule_structure")),
        ("coding_rules", "23 Coding Rules", ("parameters", "submodule_structure")),
        ("reuse_modules", "24 Reuse Modules", ("submodule_structure", "purpose")),
        ("custom", "25 Custom Extensions", ("purpose", "parameters", "test_expectation")),
        ("dir_structure", "26 Dir Structure", ("submodule_structure", "purpose")),
        ("filelist", "27 Filelist", ("submodule_structure", "purpose")),
        ("test_requirements", "28 Test Requirements / DV Plan", ("test_expectation", "bus_interface", "register_map", "interrupt")),
        ("quality_gates", "29 Quality Gates / Pass Criteria", ("test_expectation", "clock_reset", "parameters")),
        ("traceability", "30 Traceability", ("purpose", "test_expectation", "submodule_structure")),
        ("workflow_todos", "31 Workflow TODOs / Downstream Task Contract", ("test_expectation", "submodule_structure", "purpose")),
        ("generation_flow", "32 Generation Flow", ("purpose", "test_expectation")),
    ]

    _SSOT_IMPORT_EXTENSIONS = {
        ".md", ".txt", ".rst", ".yaml", ".yml", ".json", ".sv", ".svh",
        ".v", ".vh", ".py", ".csv", ".tsv", ".xml", ".f", ".sdc",
        ".tcl", ".rpt", ".log", ".h", ".c", ".cpp",
        ".pdf", ".pptx", ".docx", ".html", ".htm",
        ".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".svg", ".tif", ".tiff",
    }
    # Suffixes treated as plain text (no markitdown invocation).
    _SSOT_IMPORT_PASSTHROUGH = {".md", ".txt", ".rst"}
    _SSOT_IMPORT_IMAGE_EXTENSIONS = {
        ".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".svg", ".tif", ".tiff",
    }
    # External Python interpreter used to run markitdown (atlas_ui itself
    # may be on 3.9 where markitdown >=0.1 won't install). The probe order
    # is OS-aware:
    #   - Windows: plain `python` / `python.exe` first, then `python3.12`,
    #     then `py -3.12` as a last resort.
    #   - macOS / Linux: python3.12, python3.11, python3.10 via PATH, then
    #     the Homebrew default location as a last resort.
    # ATLAS_MARKITDOWN_PYTHON env var overrides everything.
    _SSOT_MARKITDOWN_PY = os.environ.get("ATLAS_MARKITDOWN_PYTHON", "").strip()
    _SSOT_IMPORT_SKIP_DIRS = {
        ".git", ".session", ".omx", "__pycache__", "node_modules",
        ".pytest_cache", ".mypy_cache", ".ruff_cache",
    }

    def _valid_ip_name(name: str) -> bool:
        # "default" is the reserved default workspace/workflow sentinel, never a
        # real IP. Rejecting it here makes every IP-requiring command
        # (verify-ssot, to-ssot, grill-me, repair-*, refresh-wiki,
        # /api/ip/create, …) fall into its existing "no active IP" guard instead
        # of scaffolding or writing a bogus <root>/default/ tree when no IP is
        # selected. The codebase already pairs _valid_ip_name with an explicit
        # `!= "default"` in the SSOT-IP resolvers; this folds that in centrally.
        if (name or "") == "default":
            return False
        return bool(re.match(r"^[A-Za-z][A-Za-z0-9_]*$", name or ""))

    def _slash_head(text: str) -> str:
        return (text.strip().split(None, 1)[0] if text and text.strip() else "").lower()

    def _is_workflow_slash(text: str) -> bool:
        head = _slash_head(text)
        return head in _WORKFLOW_SLASHES

    def _split_slash(text: str) -> tuple[str, str]:
        raw = (text or "").strip()
        if not raw:
            return "", ""
        parts = raw.split(None, 1)
        return parts[0].lstrip("/").lower(), (parts[1] if len(parts) > 1 else "").strip()

    _SLASH_ANSI_RE = re.compile(
        r"\x1b\[[0-9;?]*[a-zA-Z]"
        r"|\x1b\][^\x07\x1b]*(?:\x07|\x1b\\)"
        r"|\[(?:\d{1,3};)*\d{0,3}m"
    )

    def _clean_slash_output(text: str) -> str:
        return _SLASH_ANSI_RE.sub("", str(text or "")).strip("\n")

    def _emit_slash_output(client_session: Any, text: str = "", *, finish: bool = True) -> None:
        cleaned = _clean_slash_output(text)
        if cleaned:
            client_session.emit("slash_output", text=cleaned, finish=bool(finish))
        # Slash commands are command-plane operations. They should not
        # leave the chat input in "agent is streaming" state unless the
        # command explicitly queues an agent task below.
        if finish and not getattr(client_session, "agent_running", False):
            client_session.emit("agent_state", running=False)
        client_session.emit("flush")

    def _req_command_workflow_root() -> Path | None:
        candidates = [
            Path(WORKFLOW_ROOT),
            Path(SOURCE_ROOT) / "workflow",
            _REPO_ROOT / "workflow",
        ]
        seen: set[Path] = set()
        for candidate in candidates:
            root = candidate.expanduser().resolve()
            if root in seen:
                continue
            seen.add(root)
            if (
                (root / "req-gen" / "scripts" / "lock_requirement_set.py").is_file()
                and (root / "req-gen" / "scripts" / "check_contract_bundle.py").is_file()
            ):
                return root
        return None

    # Phase 12b: 14 slash command handlers (1851 lines)
    # moved to src/atlas_slash_handlers.py via make_slash_handlers factory.
    # All callable deps wrapped in lambdas to defer name lookup (forward-refs).
    from src.atlas_slash_handlers import make_slash_handlers as _make_slash_handlers
    _slash = _make_slash_handlers(
        PROJECT_ROOT=PROJECT_ROOT,
        SOURCE_ROOT=SOURCE_ROOT,
        WORKFLOW_ROOT=WORKFLOW_ROOT,
        _SSOT_IMPORT_EXTENSIONS=_SSOT_IMPORT_EXTENSIONS,
        _STAGE_RUNNERS=_STAGE_RUNNERS,
        _active_ssot_ip=lambda *a, **k: _active_ssot_ip(*a, **k),
        _append_active_history=lambda *a, **k: _append_active_history(*a, **k),
        _append_session_message=lambda *a, **k: _append_session_message(*a, **k),
        _append_workflow_history=lambda *a, **k: _append_workflow_history(*a, **k),
        _atlas_active_session_cv=lambda *a, **k: _atlas_active_session_cv(*a, **k),
        _canonical_session_string=lambda *a, **k: _canonical_session_string(*a, **k),
        _cmd_refresh_wiki=lambda *a, **k: _cmd_refresh_wiki(*a, **k),
        _collect_import_files=lambda *a, **k: _collect_import_files(*a, **k),
        _command_ip=lambda *a, **k: _command_ip(*a, **k),
        _emit_ssot_approval_ready=lambda *a, **k: _emit_ssot_approval_ready(*a, **k),
        _emit_workflow_result=lambda *a, **k: _emit_workflow_result(*a, **k),
        _ensure_new_ip_structure=lambda *a, **k: _ensure_new_ip_structure(*a, **k),
        _ensure_ssot_draft=lambda *a, **k: _ensure_ssot_draft(*a, **k),
        _extract_import_candidates=lambda *a, **k: _extract_import_candidates(*a, **k),
        _generated=lambda *a, **k: _generated(*a, **k),
        _graph=lambda *a, **k: _graph(*a, **k),
        _ip_root=lambda *a, **k: _ip_root(*a, **k),
        _ip_root_for_session=lambda *a, **k: _ip_root_for_session(*a, **k),
        _load_ssot_state=lambda *a, **k: _load_ssot_state(*a, **k),
        _merge_import_candidates=lambda *a, **k: _merge_import_candidates(*a, **k),
        _missing_ssot_decisions=lambda *a, **k: _missing_ssot_decisions(*a, **k),
        _new_ip_initial_workflow=lambda *a, **k: _new_ip_initial_workflow(*a, **k),
        _new_ssot_state=lambda *a, **k: _new_ssot_state(*a, **k),
        _parse_import_args=lambda *a, **k: _parse_import_args(*a, **k),
        _parse_new_ip_args=lambda *a, **k: _parse_new_ip_args(*a, **k),
        _python_cmd=lambda *a, **k: _python_cmd(*a, **k),
        _queue_prompt_for_session=lambda *a, **k: _queue_prompt_for_session(*a, **k),
        _refresh_ip_wiki_pages=lambda *a, **k: _refresh_ip_wiki_pages(*a, **k),
        _relative_project_path=lambda *a, **k: _relative_project_path(*a, **k),
        _render_approved_ssot_spec=lambda *a, **k: _render_approved_ssot_spec(*a, **k),
        _render_new_ip_plan=lambda *a, **k: _render_new_ip_plan(*a, **k),
        _render_ssot_llm_qna_prompt=lambda *a, **k: _render_ssot_llm_qna_prompt(*a, **k),
        _run_command=lambda *a, **k: _run_command(*a, **k),
        _save_ssot_state=lambda *a, **k: _save_ssot_state(*a, **k),
        _script_project_root=lambda *a, **k: _script_project_root(*a, **k),
        _script_project_root_for_session=lambda *a, **k: _script_project_root_for_session(*a, **k),
        _set_active_ssot_ip=lambda *a, **k: _set_active_ssot_ip(*a, **k),
        _split_slash=lambda *a, **k: _split_slash(*a, **k),
        _ssot_session_for_ip=lambda *a, **k: _ssot_session_for_ip(*a, **k),
        _ssot_yaml_path=lambda *a, **k: _ssot_yaml_path(*a, **k),
        _ssot_yaml_path_for_session=lambda *a, **k: _ssot_yaml_path_for_session(*a, **k),
        _start_sim_human_gate_qna=lambda *a, **k: _start_sim_human_gate_qna(*a, **k),
        _valid_ip_name=lambda *a, **k: _valid_ip_name(*a, **k),
    )
    _handle_bang_shell_command = _slash['_handle_bang_shell_command']
    _handle_import_command = _slash['_handle_import_command']
    _handle_grill_me_command = _slash['_handle_grill_me_command']
    _handle_new_ip_command = _slash['_handle_new_ip_command']
    _handle_ip_command = _slash['_handle_ip_command']
    _handle_session_command = _slash['_handle_session_command']
    _handle_approval_command = _slash['_handle_approval_command']
    _handle_refresh_wiki_command = _slash['_handle_refresh_wiki_command']
    _handle_to_ssot_gate = _slash['_handle_to_ssot_gate']
    _handle_repair_ssot_command = _slash['_handle_repair_ssot_command']
    _handle_verify_ssot_command = _slash['_handle_verify_ssot_command']
    _handle_repair_rtl_command = _slash['_handle_repair_rtl_command']
    _handle_repair_equiv_command = _slash['_handle_repair_equiv_command']
    _run_stage_command = _slash['_run_stage_command']

    def _apply_slash_model_switch(target: str, client_session: Any) -> str:
        try:
            import src.config as _cfg_model  # noqa: WPS433
        except Exception:
            import config as _cfg_model  # type: ignore  # noqa: WPS433

        if target == "1":
            model = str(getattr(_cfg_model, "PRIMARY_MODEL", "") or "").strip()
            if not model:
                return "No PRIMARY_MODEL is configured."
            _set_runtime_model(model)
            _cfg_model.MODEL_NAME = model
            msg = f"Model switched to: {model}"
        elif target == "2":
            model = str(getattr(_cfg_model, "SECONDARY_MODEL", "") or "").strip()
            if not model:
                return "No SECONDARY_MODEL is configured."
            _set_runtime_model(model)
            _cfg_model.MODEL_NAME = model
            msg = f"Model switched to: {model}"
        elif target.startswith("profile:"):
            profile = target.split(":", 1)[1]
            if not _cfg_model.set_active_profile(profile):
                return f"Profile '{profile}' is not defined."
            _set_runtime_model(str(getattr(_cfg_model, "MODEL_NAME", "") or ""))
            msg = (
                f"Profile '{profile}' active -> "
                f"{getattr(_cfg_model, 'MODEL_NAME', '')} @ {getattr(_cfg_model, 'BASE_URL', '')}"
            )
        elif target.startswith("cli:"):
            backend = target.split(":", 1)[1]
            if not _cfg_model.activate_cli_backend(backend):
                return f"CLI backend '{backend}' is not available."
            _set_runtime_model(str(getattr(_cfg_model, "MODEL_NAME", "") or backend))
            try:
                from src.llm_client import get_active_model as _get_active_model
                label = _get_active_model()
            except Exception:
                label = backend
            msg = f"CLI backend active -> {label}"
        elif target.startswith("opencode:"):
            model = target.split(":", 1)[1]
            if not _cfg_model.activate_opencode_oauth(model):
                return (
                    f"'{model}' needs ChatGPT OAuth but no opencode credential is available.\n"
                    "Run: python -m src.opencode_backend login"
                )
            _set_runtime_model(str(getattr(_cfg_model, "MODEL_NAME", "") or ""))
            msg = (
                f"Opencode-OAuth active -> "
                f"{getattr(_cfg_model, 'MODEL_NAME', '')} @ {getattr(_cfg_model, 'BASE_URL', '')}"
            )
        else:
            model = target.strip()
            _set_runtime_model(model)
            _cfg_model.MODEL_NAME = model
            msg = f"Model switched to: {model}"

        model_now = str(getattr(_cfg_model, "MODEL_NAME", "") or os.environ.get("LLM_MODEL_NAME", ""))
        options = _model_option_rows(model_now)
        selected_key = next((row["key"] for row in options if row.get("selected") == "true"), "")
        client_session.emit(
            "context",
            model=model_now,
            model_options=options,
            selected_model_key=selected_key,
        )
        return msg

    @contextlib.contextmanager
    def _scoped_slash_session_config(active_slash_session: str):
        """Bind generic slash commands to the active Atlas session files."""
        if not active_slash_session:
            yield
            return
        import importlib as _importlib
        import sys as _sys

        cfg_modules = []
        old_cfg_values: list[tuple[Any, str, bool, Any]] = []
        old_tt_todo_file: tuple[bool, Any] | None = None
        try:
            for cfg_name in ("config", "src.config"):
                try:
                    cfg_mod = _importlib.import_module(cfg_name)
                except Exception:
                    cfg_mod = _sys.modules.get(cfg_name)
                if cfg_mod is not None and all(cfg_mod is not m for m in cfg_modules):
                    cfg_modules.append(cfg_mod)
            session_dir = _session_json_path(active_slash_session).parent
            session_dir.mkdir(parents=True, exist_ok=True)
            req_workflow_root = _req_command_workflow_root()
            session_cfg = {
                "TODO_FILE": str(session_dir / "todo.json"),
                "TODO_ERROR_FILE": str(session_dir / "todo_error.json"),
                "SESSION_DIR": str(session_dir),
                "ACTIVE_PROJECT": active_slash_session,
                "ATLAS_ACTIVE_SESSION": active_slash_session,
            }
            if req_workflow_root is not None:
                session_cfg["ATLAS_WORKFLOW_ROOT"] = str(req_workflow_root)
                # Atlas Runtime source-root note for split workspaces:
                # allowed tooling lives under $ATLAS_SOURCE_ROOT/workflow/...
                # do not ask the user to mount or copy workflow/;
                # keep generated IP artifacts in the project workspace instead.
            for cfg_mod in cfg_modules:
                for key, value in session_cfg.items():
                    had = hasattr(cfg_mod, key)
                    old = getattr(cfg_mod, key, None)
                    old_cfg_values.append((cfg_mod, key, had, old))
                    setattr(cfg_mod, key, value)
            try:
                import lib.todo_tracker as slash_tt
                old_tt_todo_file = (
                    hasattr(slash_tt, "TODO_FILE"),
                    getattr(slash_tt, "TODO_FILE", None),
                )
                slash_tt.TODO_FILE = session_dir / "todo.json"
            except Exception:
                old_tt_todo_file = None
            yield
        finally:
            if old_tt_todo_file is not None:
                try:
                    import lib.todo_tracker as slash_tt
                    had_tt, old_tt = old_tt_todo_file
                    if had_tt:
                        slash_tt.TODO_FILE = old_tt
                    else:
                        delattr(slash_tt, "TODO_FILE")
                except Exception:
                    pass
            for cfg_mod, key, had, old in reversed(old_cfg_values):
                try:
                    if had:
                        setattr(cfg_mod, key, old)
                    else:
                        delattr(cfg_mod, key)
                except Exception:
                    pass

    def _req_lifecycle_template_guard(context: Any, tmpl_name: str) -> str:
        req_templates = {"draft-req", "finalize-req", "lock-req"}
        if tmpl_name not in req_templates:
            return ""
        workflow_root = _req_command_workflow_root()
        if workflow_root is None:
            return f"Cannot load /{tmpl_name}: no verified workflow root contains req-gen scripts."
        ip_root = Path(getattr(context, "ip_root", ""))
        ip_name = str(getattr(context, "ip_name", "") or "active IP")
        req_dir = ip_root / "req"
        manifest = req_dir / "approval_manifest.json"
        if _locked_truth_active(context.workspace_root, ip_name):
            manifest_state = ""
            if manifest.exists():
                try:
                    payload = json.loads(manifest.read_text(encoding="utf-8"))
                    if isinstance(payload, dict):
                        manifest_state = str(payload.get("status") or payload.get("locked_truth_status") or "")
                except (OSError, json.JSONDecodeError) as exc:
                    return f"Cannot load /{tmpl_name}: req/approval_manifest.json is invalid: {exc}"
            state_text = f" ({manifest_state})" if manifest_state else ""
            return (
                f"Requirements already locked for {ip_name}{state_text}. "
                f"Refusing /{tmpl_name} because it would replace the current TODO list. "
                "Next action: /to-ssot."
            )
        if tmpl_name in {"finalize-req", "lock-req"}:
            candidate_files = [
                req_dir / "requirements_index.json",
                req_dir / "obligations.json",
                req_dir / "contract_refs.json",
                req_dir / "structural_contracts.json",
                req_dir / "behavioral_contracts.json",
                req_dir / "evidence_plan.json",
            ]
            missing = [path.name for path in candidate_files if not path.exists()]
            if missing:
                return f"/{tmpl_name} needs draft req files first. Missing: {', '.join(missing)}. Run /draft-req."
        return ""

    def _req_lifecycle_template_tasks(
        tmpl_name: str,
        tasks: list[dict[str, Any]],
        approver: str = "",
    ) -> list[dict[str, Any]]:
        workflow_root = _req_command_workflow_root()
        copied = json.loads(json.dumps(tasks))
        if tmpl_name not in {"draft-req", "finalize-req", "lock-req"} or workflow_root is None:
            return copied
        # The lock gate runs `--approved-by "${ATLAS_APPROVED_BY:?...}"` as a
        # deterministic subprocess. That env var is never exported into the worker
        # process that runs the gate, so the ${VAR:?} guard aborts the command
        # before the stamper runs and the lock can never complete. The session
        # owner (the human who invoked /lock-req) is the approver of record, so
        # bake it into the command string here — the same load-time substitution
        # already used for $ATLAS_WORKFLOW_ROOT — instead of relying on env that
        # does not cross the web→worker process boundary.
        approver_token = shlex.quote(approver.strip()) if approver and approver.strip() else ""
        for task in copied:
            command = task.get("command")
            if isinstance(command, str):
                command = command.replace("$ATLAS_WORKFLOW_ROOT", str(workflow_root))
                if approver_token:
                    command = re.sub(
                        r'"?\$\{ATLAS_APPROVED_BY[^}]*\}"?',
                        approver_token,
                        command,
                    )
                    command = command.replace('"$ATLAS_APPROVED_BY"', approver_token)
                    command = command.replace("$ATLAS_APPROVED_BY", approver_token)
                task["command"] = command
        return copied

    def _execute_generic_slash_command(text: str, client_session: Any) -> bool:
        """Run non-ATLAS slash commands immediately on the command plane."""
        raw = (text or "").strip()
        if not raw.startswith("/"):
            return False
        active_slash_session = normalize_session_name(
            str(getattr(client_session, "session_id", "") or "")
        )
        if raw.lower() == "/normal":
            result = "AGENT_MODE:normal"
        else:
            try:
                from core.slash_commands import get_registry as _get_slash_registry
                _ensure_default_workspace_commands_registered()
                _old_memory_user = os.environ.get("ATLAS_MEMORY_USER")
                _owner_for_memory = active_slash_session.split("/", 1)[0]
                if _owner_for_memory:
                    os.environ["ATLAS_MEMORY_USER"] = _owner_for_memory
                try:
                    with _scoped_slash_session_config(active_slash_session):
                        result = _get_slash_registry().execute(raw)
                finally:
                    if _old_memory_user is None:
                        os.environ.pop("ATLAS_MEMORY_USER", None)
                    else:
                        os.environ["ATLAS_MEMORY_USER"] = _old_memory_user
            except Exception as exc:
                _emit_slash_output(client_session, f"Error executing {raw.split(None, 1)[0]}: {exc}")
                return True
        if result is None:
            return False

        if result.startswith("INJECT_TODO_TEMPLATE:"):
            tmpl_name = result[len("INJECT_TODO_TEMPLATE:"):].strip()
            session_key = normalize_session_name(str(getattr(client_session, "session_id", "") or ""))
            try:
                from core.atlas_context_paths import AtlasContext
                from lib.todo_tracker import TodoTracker
                from workflow.loader import get_todo_template_registry

                context = AtlasContext.from_session_key(session_key, atlas_root=PROJECT_ROOT)
                req_guard = _req_lifecycle_template_guard(context, tmpl_name)
                if req_guard:
                    _emit_slash_output(client_session, req_guard)
                    return True
                registry = get_todo_template_registry()
                template = registry.get(tmpl_name) or {}
                tasks = _req_lifecycle_template_tasks(
                    tmpl_name,
                    registry.get_tasks(tmpl_name) or [],
                    approver=str(getattr(context, "user_name", "") or ""),
                )
                if not tasks:
                    available = ", ".join(registry.list()) or "(none)"
                    _emit_slash_output(
                        client_session,
                        f"Todo template '{tmpl_name}' not found. Available: {available}",
                    )
                    return True
                todo_path = context.session_dir / "todo.json"
                tracker = TodoTracker(persist_path=todo_path)
                tracker.add_todos(tasks)
                tracker.template_lock_additions = bool(template.get("lock_additions", True))
                tracker.template_name = tmpl_name
                # Auto-start: the pending→in_progress transition is normally an
                # LLM-driven todo_update tool call. If the model never makes that
                # first call (writes prose, works off-ledger, or forgets) the task
                # stays pending and the no-action watchdog stops the loop. Seed the
                # first task to in_progress here so the runtime — not the model —
                # owns the first transition. Per-template opt-in; global kill switch.
                _auto_start = (
                    bool(template.get("auto_start", False))
                    and os.getenv("ATLAS_DISABLE_TODO_AUTO_START", "").strip().lower()
                    not in ("1", "true", "yes")
                )
                if _auto_start and tracker.todos:
                    tracker.mark_in_progress(0)
                tracker.save()
                bind_errors = _bind_live_tracker_to_session(session_key, todo_path, tracker)
                if bool(template.get("lock_additions", True)):
                    os.environ["TODO_TEMPLATE_LOCK_ADDITIONS"] = "1"
                    os.environ["TODO_TEMPLATE_LOCK_NAME"] = tmpl_name
                else:
                    os.environ.pop("TODO_TEMPLATE_LOCK_ADDITIONS", None)
                    os.environ.pop("TODO_TEMPLATE_LOCK_NAME", None)
                desc = str(template.get("description", "") or "").strip()
                rel_path = todo_path
                try:
                    rel_path = todo_path.relative_to(PROJECT_ROOT)
                except ValueError:
                    pass
                _emit_slash_output(
                    client_session,
                    f"Loaded todo template '{tmpl_name}': {len(tasks)} tasks added.\n"
                    f"Todo: {rel_path}"
                    + (f"\nDescription: {desc}" if desc else "")
                    + (
                        "\nTodo runtime bind warning: " + "; ".join(bind_errors)
                        if bind_errors else ""
                    )
                    + ("\nAuto-start: driving task 1 now." if _auto_start else ""),
                )
                # Auto-start kick: enqueue a worker turn so the loop drives the
                # plan without waiting for a user message. Reuses the proven
                # INJECT_PROMPT path (bridge.submit_prompt_for_session). The
                # worker reads the session todo.json (task 1 already in_progress)
                # and advances through the tasks.
                if _auto_start:
                    _auto_prompt = str(template.get("auto_start_prompt", "") or "").strip() or (
                        "Begin executing the loaded TODO plan now. Task 1 is "
                        "already in_progress — drive it to completion with real "
                        "tool actions, then review and continue through the "
                        "remaining tasks. Do not wait for further confirmation."
                    )
                    try:
                        bridge.submit_prompt_for_session(
                            client_session.session_id, _auto_prompt
                        )
                        client_session.emit("agent_state", running=True)
                    except Exception as _auto_exc:
                        _emit_slash_output(
                            client_session,
                            f"Auto-start kick failed (load succeeded): {_auto_exc}",
                        )
            except ValueError:
                _emit_slash_output(
                    client_session,
                    "Todo template commands require canonical session "
                    "user/workspace/ip/workflow.",
                )
            return True

        if result.startswith("MODEL_SWITCH:"):
            msg = _apply_slash_model_switch(result.split(":", 1)[1], client_session)
            _emit_slash_output(client_session, msg)
            return True

        if result.startswith("AGENT_MODE:"):
            target = result.split(":", 1)[1]
            is_plan = target == "plan"
            _agent_mode_override_cv.set("plan_q" if is_plan else "normal")
            _plan_mode_cv.set("true" if is_plan else "false")
            os.environ["AGENT_MODE_OVERRIDE"] = "plan_q" if is_plan else "normal"
            os.environ["PLAN_MODE"] = "true" if is_plan else "false"
            if not is_plan:
                os.environ.pop("_PLAN_TODO_WRITE_COUNT", None)
            client_session.emit("mode_change", mode="plan" if is_plan else "normal")
            _emit_slash_output(
                client_session,
                "Plan mode: read-only until confirmation." if is_plan else "Normal mode: tools enabled.",
            )
            return True

        if result.startswith("EXECUTION_MODE:"):
            parts = result.split(":")
            mode = parts[1] if len(parts) > 1 else "agent"
            count = int(parts[2]) if len(parts) > 2 and parts[2].isdigit() else 1
            try:
                import src.config as _cfg_exec  # noqa: WPS433
            except Exception:
                import config as _cfg_exec  # type: ignore  # noqa: WPS433
            _cfg_exec.EXECUTION_MODE = mode
            _cfg_exec.STEP_BY_STEP_MODE = mode == "step"
            if mode == "chat":
                _cfg_exec.CHAT_MAX_ITERATIONS = count
            _emit_slash_output(client_session, f"Execution mode set to: {mode}")
            return True

        if result.startswith("WINDOW_MODE:") or result.startswith("COMPRESSION_MODE:"):
            _emit_slash_output(client_session, result)
            return True

        if result.startswith("COMPACT_HISTORY"):
            _compact_path = _session_json_path(active_slash_session)
            try:
                # Real LLM compaction — the same compress_history path the CLI /
                # Textual UI use (AI summary preserving paths / code / todos).
                message, updated = _compact_history_llm(_compact_path, result)
            except Exception as exc:
                # LLM unavailable / failed → deterministic local compactor so
                # /compact still shrinks history instead of breaking.
                try:
                    message, updated = _compact_history_file(_compact_path, result)
                    message = f"{message} (AI summary unavailable: {type(exc).__name__})"
                except Exception as exc2:
                    _emit_slash_output(client_session, f"Compact failed: {exc2}")
                    return True
            _emit_history_context_update(client_session, updated)
            _emit_slash_output(client_session, message)
            return True

        if result == "CLEAR_HISTORY" or result.startswith("CLEAR_HISTORY:"):
            try:
                message, updated = _clear_history_file(
                    _session_json_path(active_slash_session),
                    result,
                )
                _emit_history_context_update(client_session, updated)
                _emit_slash_output(client_session, message)
            except Exception as exc:
                _emit_slash_output(client_session, f"Clear failed: {exc}")
            return True

        if result.startswith("PLAN_AND_RUN:"):
            task = result[len("PLAN_AND_RUN:"):].strip()
            _agent_mode_override_cv.set("plan_q")
            _plan_mode_cv.set("true")
            os.environ["AGENT_MODE_OVERRIDE"] = "plan_q"
            os.environ["PLAN_MODE"] = "true"
            bridge.submit_prompt_for_session(client_session.session_id, task)
            client_session.emit("agent_state", running=True)
            _emit_slash_output(client_session, "Plan command accepted; agent task queued.", finish=False)
            return True

        if result.startswith("INJECT_PROMPT:"):
            prompt = result[len("INJECT_PROMPT:"):].strip()
            bridge.submit_prompt_for_session(client_session.session_id, prompt)
            client_session.emit("agent_state", running=True)
            _emit_slash_output(client_session, "Command accepted; agent task queued.", finish=False)
            return True

        if result.startswith("WORKSPACE_SWITCH:"):
            # Workspace switching rewrites main.py's in-memory prompt and
            # todo bindings. Queue that specialized state transition through
            # the agent control loop, but keep the chat command itself out
            # of the LLM prompt stream.
            was_running = bool(getattr(client_session, "agent_running", False))
            bridge.submit_prompt_for_session(client_session.session_id, raw)
            if not was_running:
                client_session.agent_running = False
            _emit_slash_output(
                client_session,
                "Workflow switch command accepted.",
                finish=not was_running,
            )
            return True

        if result in {"CLEAR_ALL", "GIT_CLEAR"} or result.startswith("TODO_REVERT:"):
            _emit_slash_output(
                client_session,
                "This slash command is destructive and is not executed directly from the web command plane.",
            )
            return True

        _emit_slash_output(client_session, result)
        return True

    def _session_json_path(session: str) -> Path:
        """Map any session string (1/2/3-part) to the canonical
        .session/<session_id>/<ip>/<workflow>/conversation.json path.

        Legacy (shorter) on-disk locations are auto-migrated on read.
        """
        clean = normalize_session_name(session or "")
        parts = [p for p in clean.split("/") if p]
        owner_default = os.environ.get("ATLAS_DEFAULT_SESSION_ID") or "default"
        ip_default = _active_ip_value() or "default"
        wf_default = os.environ.get("ATLAS_DEFAULT_WORKFLOW") or "default"
        owner = _resolve_session_owner() or owner_default
        legacy_candidates: list[Path] = []
        if len(parts) >= 4:
            owner, workspace_session, ip, wf = parts[0], parts[1], parts[2], parts[3]
            canon = PROJECT_ROOT / owner / workspace_session / ".session" / ip / wf
            legacy_candidates.append(PROJECT_ROOT / ".session" / owner / ip / wf)
        elif len(parts) >= 3:
            owner, ip, wf = parts[0], parts[1], parts[2]
            canon = PROJECT_ROOT / ".session" / owner / ip / wf
        elif len(parts) == 2:
            ip, wf = parts[0], parts[1]
            canon = PROJECT_ROOT / ".session" / owner / ip / wf
        elif len(parts) == 1:
            ip, wf = ip_default, parts[0]
            canon = PROJECT_ROOT / ".session" / owner / ip / wf
        else:
            ip, wf = ip_default, wf_default
            canon = PROJECT_ROOT / ".session" / owner / ip / wf
        # Migrate legacy 3-part / 2-part / 1-part dirs to canonical on first access.
        # The slash-command setup may create the canonical directory before the
        # command handler reads conversation.json, so use the history file
        # itself as the migration guard instead of the directory.
        history_path = canon / "conversation.json"
        if not history_path.exists():
            for legacy in (
                *legacy_candidates,
                PROJECT_ROOT / ".session" / ip / wf,    # 2-part: ip/wf
                PROJECT_ROOT / ".session" / wf,          # 1-part: wf only
            ):
                try:
                    if legacy == canon or not legacy.exists() or not legacy.is_dir():
                        continue
                    if not (legacy / "conversation.json").exists():
                        continue
                    canon.parent.mkdir(parents=True, exist_ok=True)
                    if not canon.exists():
                        legacy.rename(canon)
                    else:
                        canon.mkdir(parents=True, exist_ok=True)
                        for child in legacy.iterdir():
                            target = canon / child.name
                            if not target.exists():
                                child.rename(target)
                        try:
                            legacy.rmdir()
                        except OSError:
                            pass
                        break
                except OSError:
                    continue
        return history_path

    def _append_session_message(session: str, role: str, content: str) -> None:
        session = normalize_session_name(session)
        if not session:
            return
        path = _session_json_path(session)
        try:
            path.parent.mkdir(parents=True, exist_ok=True)
            try:
                msgs = json.loads(path.read_text(encoding="utf-8")) if path.exists() else []
                if not isinstance(msgs, list):
                    msgs = []
            except Exception:
                msgs = []
            msgs.append({"role": role, "content": content})
            path.write_text(json.dumps(msgs, ensure_ascii=False, indent=2), encoding="utf-8")
        except Exception:
            pass

    def _handle_fast_identity_prompt(text: str, client_session: Any) -> bool:
        """Answer tiny identity/greeting prompts without waking a workflow LLM."""
        if not _atlas_fast_prompt_kind(text):
            return False
        session_id = normalize_session_name(str(getattr(client_session, "session_id", "") or ""))
        if not session_id:
            return False
        user_text = _atlas_fast_prompt_candidate(text)
        reply = _atlas_fast_identity_response(session_id, user_text)
        _append_session_message(session_id, "user", user_text)
        _append_session_message(session_id, "assistant", reply)
        try:
            client_session.agent_running = False
        except Exception:
            pass
        client_session.emit("agent_state", running=True)
        client_session.emit("token", text=reply)
        client_session.emit("flush")
        client_session.emit("agent_state", running=False)
        return True

    def _append_active_history(role: str, content: str) -> None:
        """Mirror direct Web workflow command output into the currently
        hydrated chat history. Without this, commands like `/new-ip` can
        emit a visible WS event and then disappear when data.jsx reloads
        `/api/conversation` for the active workspace.
        """
        try:
            try:
                import src.config as _cfg_hist  # type: ignore
            except Exception:
                try:
                    import config as _cfg_hist  # type: ignore
                except Exception:
                    _cfg_hist = None
            if _cfg_hist is None:
                return
            hpath = Path(getattr(_cfg_hist, "HISTORY_FILE", "") or "")
            if not hpath:
                return
            hpath.parent.mkdir(parents=True, exist_ok=True)
            try:
                msgs = json.loads(hpath.read_text(encoding="utf-8")) if hpath.exists() else []
                if not isinstance(msgs, list):
                    msgs = []
            except Exception:
                msgs = []
            msgs.append({"role": role, "content": content})
            hpath.write_text(json.dumps(msgs, ensure_ascii=False, indent=2), encoding="utf-8")
        except Exception:
            pass

    def _append_workflow_history(workflow: str, role: str, content: str) -> None:
        """Persist a message in the workflow-level session as well as the
        per-IP session. The actual slash dispatcher reloads `.session/<wf>`
        after `/wf <name>`, so approved Web Q&A must be visible there before
        `/to-ssot` runs.
        """
        _append_session_message(workflow, role, content)

    def _ssot_state_path(ip: str) -> Path:
        return _ssot_session_dir(ip) / "state.json"

    def _load_ssot_state(ip: str) -> dict[str, Any]:
        paths = [_ssot_state_path(ip)]
        # Single-user/dev mode historically stored SSOT state under the
        # default owner or the legacy two-part path. Keep those reads in
        # single-user mode only so multi-user sessions do not see another
        # user's SSOT decisions.
        if not _multi_user_enabled():
            paths.extend([
                PROJECT_ROOT / ".session" / "default" / ip / "ssot-gen" / "state.json",
                _legacy_ssot_session_dir(ip) / "state.json",
            ])
        seen: set[Path] = set()
        for path in paths:
            if path in seen:
                continue
            seen.add(path)
            if not path.is_file():
                continue
            try:
                doc = json.loads(path.read_text(encoding="utf-8"))
                if isinstance(doc, dict):
                    return doc
            except Exception:
                continue
        return {}

    def _save_ssot_state(ip: str, state: dict[str, Any]) -> None:
        path = _ssot_state_path(ip)
        path.parent.mkdir(parents=True, exist_ok=True)
        state["updated_at"] = time.time()
        path.write_text(json.dumps(state, ensure_ascii=False, indent=2), encoding="utf-8")


    def _ssot_session_dir(ip: str, session: str | None = None) -> Path:
        # Canonical layout under --ui atlas: .session/<session_id>/<ip>/ssot-gen
        # Honor an explicitly-passed canonical session string (3-part). For all
        # other cases delegate to the central canonical resolver, so the
        # session_id is never silently dropped from the path.
        if session:
            clean = normalize_session_name(str(session))
            parts = [p for p in clean.split("/") if p]
            if len(parts) >= 3 and parts[-1] == "ssot-gen" and parts[-2] == ip:
                return PROJECT_ROOT / ".session" / clean
        return PROJECT_ROOT / ".session" / _canonical_session_string(ip, "ssot-gen")

    def _legacy_ssot_session_dir(ip: str) -> Path:
        return PROJECT_ROOT / ".session" / ip / "ssot-gen"

    # Phase 10b: simple Q&A I/O helpers moved to src/atlas_qa.py via factory.
    # Lambda wrappers defer name lookup so forward-refs (e.g. _canonical_session_string)
    # work — Python resolves the names from create_app's enclosing scope at CALL time.
    from src.atlas_qa import make_qa_helpers as _make_qa_helpers
    _qa = _make_qa_helpers(
        ssot_session_dir_fn=_ssot_session_dir,
        legacy_ssot_session_dir_fn=lambda ip: _legacy_ssot_session_dir(ip),
        normalize_session_name_fn=lambda s: normalize_session_name(s),
        project_root_fn=lambda: PROJECT_ROOT,
        active_session_value_fn=lambda: _active_session_value(),
        active_ip_value_fn=lambda: _active_ip_value(),
        valid_ip_name_fn=lambda n: _valid_ip_name(n),
        canonical_session_fn=lambda *a, **k: _canonical_session_string(*a, **k),
        load_ssot_state_fn=lambda ip: _load_ssot_state(ip),
        ssot_decisions_fn=lambda ip, state: _ssot_decisions(ip, state),
        required_decisions_fn=lambda: _SSOT_REQUIRED_DECISIONS,
    )
    _ssot_qa_path = _qa["path"]
    _load_ssot_qa_items = _qa["load"]
    _save_ssot_qa_items = _qa["save"]
    _active_ssot_qa_context = _qa["active_context"]
    _upsert_ssot_qa_items = _qa["upsert"]
    _ssot_qa_view = _qa["view"]
    _ssot_qa_sessions_view = _qa["sessions_view"]




    _upsert_ssot_qa_items = _qa["upsert"]
    _ssot_qa_view = _qa["view"]
    _ssot_qa_sessions_view = _qa["sessions_view"]







    def _ssot_yaml_path(ip: str, base_root: Path | None = None) -> Path:
        ip_dir = (base_root / ip) if base_root is not None else _ip_root(ip)
        for name in (f"{ip}.ssot.yaml", f"{ip}_ssot.yaml", f"{ip}.ssot.yml"):
            candidate = ip_dir / "yaml" / name
            if candidate.is_file():
                return candidate
        return ip_dir / "yaml" / f"{ip}.ssot.yaml"

    def _load_ssot_draft(ip: str, base_root: Path | None = None) -> dict[str, Any]:
        path = _ssot_yaml_path(ip, base_root=base_root)
        if not path.is_file():
            return {}
        try:
            import yaml as _yaml  # type: ignore

            doc = _yaml.safe_load(path.read_text(encoding="utf-8", errors="replace")) or {}
            return doc if isinstance(doc, dict) else {}
        except Exception:
            return {}

    def _load_ssot_yaml_from_path(path: Path) -> dict[str, Any]:
        import yaml as _yaml  # type: ignore

        if not path.is_file():
            raise FileNotFoundError(str(path))
        try:
            data = _yaml.safe_load(path.read_text(encoding="utf-8", errors="replace"))
        except Exception as exc:
            raise ValueError(f"invalid yaml: {exc}") from exc
        if data is None:
            return {}
        if not isinstance(data, dict):
            raise ValueError("ssot yaml must be a mapping/object")
        return data

    def _save_ssot_draft(ip: str, doc: dict[str, Any], base_root: Path | None = None) -> None:
        path = _ssot_yaml_path(ip, base_root=base_root)
        path.parent.mkdir(parents=True, exist_ok=True)
        try:
            import yaml as _yaml  # type: ignore
        except Exception as exc:
            raise RuntimeError(f"PyYAML is required to update SSOT draft: {exc}") from exc
        path.write_text(_yaml.safe_dump(doc, sort_keys=False, allow_unicode=True, width=120), encoding="utf-8")
        # Fire file_changed so the open SSOT preview / full view / file
        # tree auto-reload. _emit_tool_result only catches tool-driven
        # writes (write_file / replace_in_file / …); _save_ssot_draft
        # is a direct Python disk write (called by /api/ssot/qa/answer
        # and the ssot scaffold helpers), so it would otherwise look
        # like a silent update from the frontend's perspective.
        try:
            target_session = ""
            try:
                target_session = normalize_session_name(str(_ssot_session_for_ip(ip) or ""))
            except Exception:
                try:
                    target_session = normalize_session_name(str(_canonical_session_string(ip, "ssot-gen") or ""))
                except Exception:
                    target_session = ""
            payload = {
                "path": str(path),
                "tool": "ssot_save",
                "ip": ip,
                "workflow": "ssot-gen",
            }
            if target_session:
                payload["session"] = target_session
                payload["session_id"] = target_session
            bridge.emit("file_changed", **payload)
        except Exception:
            pass

    def _ssot_base_root_for_request(
        request: Request,
        ip: str,
        session_id: str,
    ) -> tuple[Path | None, JSONResponse | None]:
        session_name = normalize_session_name(str(session_id or ""))
        if not session_name:
            return None, None
        try:
            context = AtlasContext.from_session_key(
                session_name,
                atlas_root=_atlas_root_for_context(),
            )
        except Exception as exc:
            return None, JSONResponse({"ok": False, "error": f"invalid session_id: {exc}"}, status_code=400)
        if context.legacy:
            return None, None
        if context.ip_name and context.ip_name not in {"default", ip}:
            return None, JSONResponse({"ok": False, "error": "session ip mismatch"}, status_code=400)
        try:
            user = request.scope.get("user") or {}
        except Exception:
            user = {}
        user_id = str((user or {}).get("id") or "").strip()
        if not user_id or user_id == "default":
            return None, JSONResponse({"ok": False, "error": "login required"}, status_code=401)
        if str((user or {}).get("role") or "").strip().lower() != "admin":
            username = str((user or {}).get("username") or "").strip().strip("/")
            if username != context.user_name:
                return None, JSONResponse({"ok": False, "error": "session owner mismatch"}, status_code=403)
        root, error = _validated_context_workspace_root(context)
        if root is None:
            return None, JSONResponse({"ok": False, "error": error or "invalid workspace root"}, status_code=400)
        return root, None

    def _parse_ssot_feedback_path(raw: str) -> list[Any]:
        text = str(raw or "").strip().strip(".")
        if not text:
            return []
        if not re.match(r"^[A-Za-z_][A-Za-z0-9_-]*(?:\.(?:[A-Za-z_][A-Za-z0-9_-]*|[0-9]+))*$", text):
            raise ValueError(
                "path must use dot notation, e.g. top_module.review_note or registers.register_list.0.description"
            )
        return [int(part) if part.isdigit() else part for part in text.split(".")]

    def _set_ssot_feedback_path(root: dict[str, Any], tokens: list[Any], value: Any) -> None:
        if not tokens:
            return
        cur: Any = root
        for idx, token in enumerate(tokens):
            last = idx == len(tokens) - 1
            nxt = None if last else tokens[idx + 1]
            if isinstance(token, int):
                if not isinstance(cur, list):
                    raise ValueError("numeric path segments require a list parent")
                while len(cur) <= token:
                    cur.append([] if isinstance(nxt, int) else {})
                if last:
                    cur[token] = value
                    return
                if not isinstance(cur[token], (dict, list)):
                    cur[token] = [] if isinstance(nxt, int) else {}
                cur = cur[token]
                continue

            if not isinstance(cur, dict):
                raise ValueError("mapping path segments require an object parent")
            if last:
                cur[token] = value
                return
            if token not in cur or not isinstance(cur[token], (dict, list)):
                cur[token] = [] if isinstance(nxt, int) else {}
            cur = cur[token]

    def _get_ssot_feedback_path(root: dict[str, Any], tokens: list[Any]) -> Any:
        cur: Any = root
        for token in tokens:
            if isinstance(token, int):
                if not isinstance(cur, list) or token >= len(cur):
                    raise KeyError(token)
                cur = cur[token]
                continue
            if not isinstance(cur, dict) or token not in cur:
                raise KeyError(token)
            cur = cur[token]
        return cur

    def _ssot_feedback_tokens_to_path(tokens: list[Any]) -> str:
        return ".".join(str(token) for token in tokens)

    def _ssot_feedback_section(raw: str, fallback: str = "custom") -> str:
        section = str(raw or "").strip()
        known = {key for key, _label in _SSOT_EXPORT_SECTION_ORDER}
        return section if section in known else fallback

    def _ssot_feedback_slug(text: str) -> str:
        slug = re.sub(r"[^A-Za-z0-9_.-]+", "_", str(text or "").strip()).strip("_.-").lower()
        return slug[:80] or "feedback"

    def _full_ssot_tbd_template(ip: str, top_file: str, kind: str) -> dict[str, Any]:
        """Comprehensive TBD skeleton covering every REQUIRED_SECTIONS key.

        Each section ships with the right shape (dict vs list) and at least
        one TBD-marked example so /new-ip immediately shows what needs to
        be filled in. The Validation pane stays at 0/N until real content
        replaces the TBDs — that's intentional, the template is meant as
        a checklist, not a free pass through the gates.
        """
        return {
            "top_module": {
                "name": ip,
                "file": top_file,
                "type": "draft",
                "description": kind or "TBD",
                "version": "draft",
            },
            "sub_modules": [
                {"name": "TBD", "role": "TBD", "file": "TBD"},
            ],
            "decomposition": {
                "strategy": "TBD",
                "notes": "TBD",
            },
            "parameters": [
                {"name": "TBD", "type": "TBD", "default": "TBD", "description": "TBD"},
            ],
            "io_list": [
                {"name": "TBD", "direction": "TBD", "width": "TBD", "clock_domain": "TBD", "description": "TBD"},
            ],
            "features": [
                {"id": "F_TBD", "description": "TBD"},
            ],
            "dataflow": {
                "stages": ["TBD"],
                "notes": "TBD",
            },
            "function_model": {
                "state_variables": [{"name": "TBD", "width": "TBD", "reset": "TBD"}],
                "transactions": [{"id": "TX_TBD", "name": "TBD", "preconditions": ["TBD"], "outputs": ["TBD"], "side_effects": ["TBD"], "error_cases": ["TBD"]}],
                "invariants": [{"name": "TBD", "expr": "TBD", "description": "TBD"}],
            },
            "cycle_model": {
                "clock": "TBD",
                "reset": "TBD",
                "pipeline": ["TBD"],
                "handshake_rules": ["TBD"],
                "ordering": ["TBD"],
                "observability": ["TBD"],
                "executable": "TBD",
            },
            "rtl_contract": {
                "sample_condition": "TBD",
                "input_map": {"TBD": "TBD"},
                "output_map": {"TBD": "TBD"},
            },
            "clock_reset_domains": [
                {"name": "TBD", "clock": "TBD", "reset": "TBD", "polarity": "TBD"},
            ],
            "registers": [
                {"name": "TBD", "address": "TBD", "access": "TBD", "fields": [{"name": "TBD", "bits": "TBD", "access": "TBD", "description": "TBD"}]},
            ],
            "memory": {
                "blocks": [{"name": "TBD", "size": "TBD", "width": "TBD", "kind": "TBD"}],
            },
            "interrupts": [
                {"name": "TBD", "source": "TBD", "polarity": "TBD", "mask_register": "TBD"},
            ],
            "fsm": [
                {"name": "TBD", "states": ["TBD"], "transitions": [{"from": "TBD", "to": "TBD", "guard": "TBD"}]},
            ],
            "timing": {
                "target_freq_mhz": "TBD",
                "critical_paths": ["TBD"],
            },
            "power": {
                "domains": ["TBD"],
                "gating": "TBD",
            },
            "security": {
                "threats": ["TBD"],
                "mitigations": ["TBD"],
            },
            "error_handling": [
                {"id": "ERR_TBD", "trigger": "TBD", "action": "TBD"},
            ],
            "debug_observability": [
                {"signal": "TBD", "purpose": "TBD"},
            ],
            "integration": {
                "parent_bus": "TBD",
                "addressing": "TBD",
                "external_signals": ["TBD"],
            },
            "dft": {
                "scan_chain": "TBD",
                "test_modes": ["TBD"],
            },
            "synthesis": {
                "library": "TBD",
                "constraints": ["TBD"],
            },
            "pnr": {
                "floorplan_notes": "TBD",
                "placement_hints": ["TBD"],
            },
            "test_requirements": [
                {"id": "TR_TBD", "description": "TBD", "coverage_kind": "TBD"},
            ],
            "quality_gates": [
                {"id": "QG_TBD", "criterion": "TBD", "status": "TBD"},
            ],
            "traceability": {
                "requirements_to_features": {"REQ_TBD": ["F_TBD"]},
                "features_to_tests": {"F_TBD": ["TR_TBD"]},
            },
            "workflow_todos": [
                {"id": "TODO_TBD", "title": "TBD", "owner": "TBD", "status": "TBD"},
            ],
            "filelist": [
                {"path": top_file, "role": "top"},
            ],
            "coding_rules": [
                "TBD",
            ],
            "reuse_modules": [
                {"name": "TBD", "source": "TBD"},
            ],
            "dir_structure": {
                "rtl": "rtl/",
                "tb": "tb/",
                "sim": "sim/",
                "yaml": "yaml/",
                "wiki": "wiki/",
            },
            "generation_flow": {
                "steps": ["ssot-gen", "fl-model-gen", "rtl-gen", "tb-gen", "sim", "coverage"],
                "notes": "TBD",
            },
            "cdc_requirements": [
                {"from_clock": "TBD", "to_clock": "TBD", "scheme": "TBD"},
            ],
            "rdc_requirements": [
                {"from_reset": "TBD", "to_reset": "TBD", "scheme": "TBD"},
            ],
            "custom": {},
        }

    def _ensure_ssot_draft(ip: str, kind: str = "TBD", base_root: Path | None = None) -> dict[str, Any]:
        # Default top file path follows the canonical convention
        # `rtl/<ip>.sv` so /new-ip scaffolds an SSOT that already has
        # the synthesizable top wired to a name matching the IP. Without
        # this, drafts shipped without a `file` field and downstream
        # rtl-gen runs occasionally settled on `<ip>_wrapper.sv` as the
        # de facto top, which surprised reviewers expecting `<ip>.sv`.
        _default_top_file = f"rtl/{ip}.sv" if ip else "rtl/top.sv"
        doc = _load_ssot_draft(ip, base_root=base_root)
        if not doc:
            doc = _full_ssot_tbd_template(ip, _default_top_file, kind)
        # Pull every REQUIRED_SECTIONS key from the template so older drafts
        # that pre-date the full template auto-upgrade on the next /new-ip
        # or /ssot edit. setdefault guarantees existing real content is
        # never clobbered by TBD stubs.
        else:
            _tpl = _full_ssot_tbd_template(ip, _default_top_file, kind)
            for k, v in _tpl.items():
                doc.setdefault(k, v)
        top = doc.setdefault("top_module", {})
        if isinstance(top, dict):
            top.setdefault("name", ip)
            top.setdefault("file", _default_top_file)
            top.setdefault("type", "draft")
            top.setdefault("description", kind or "TBD")
            top.setdefault("version", "draft")
        custom = doc.setdefault("custom", {})
        if not isinstance(custom, dict):
            custom = {}
            doc["custom"] = custom
        workflow = custom.setdefault("atlas_workflow", {})
        if isinstance(workflow, dict):
            workflow.setdefault("status", "draft")
            workflow.setdefault("source", "atlas-ui")
            workflow["updated_at"] = time.time()
        custom.setdefault("atlas_decisions", {})
        custom.setdefault("atlas_decision_sources", {})
        custom.setdefault("atlas_imports", [])
        custom.setdefault("atlas_import_conflicts", [])
        _save_ssot_draft(ip, doc, base_root=base_root)
        return doc

    def _ssot_custom(ip: str, kind: str = "TBD") -> tuple[dict[str, Any], dict[str, Any]]:
        doc = _ensure_ssot_draft(ip, kind)
        custom = doc.setdefault("custom", {})
        if not isinstance(custom, dict):
            custom = {}
            doc["custom"] = custom
        return doc, custom

    def _infer_decisions_from_yaml(doc: dict[str, Any]) -> dict[str, str]:
        # The Validation pane needs to reflect what is actually written to
        # <ip>/yaml/<ip>.ssot.yaml on disk, not just the parallel
        # custom.atlas_decisions promise tracker. /to-ssot and direct
        # generators write top-level sections (top_module, registers,
        # parameters, ...) without ever touching atlas_decisions, so the
        # tracker shows 0/9 while the yaml has 17+ populated sections.
        # Infer a "filled" signal for each required decision from real
        # section content; explicit atlas_decisions still wins downstream.
        if not isinstance(doc, dict):
            return {}
        out: dict[str, str] = {}

        def _meaningful(value: Any) -> bool:
            if value in (None, "", [], {}):
                return False
            if isinstance(value, str):
                token = value.strip().lower()
                if not token:
                    return False
                if token in {"tbd", "todo", "draft", "?", "n/a"}:
                    return False
                return re.search(r"(^|[^a-z0-9])(tbd|todo|fixme|placeholder)([^a-z0-9]|$)", token) is None
            if isinstance(value, (list, tuple, set)):
                return any(_meaningful(v) for v in value)
            if isinstance(value, dict):
                return any(_meaningful(v) for v in value.values())
            return True

        def _section(*keys: str) -> Any:
            for k in keys:
                v = doc.get(k)
                if _meaningful(v):
                    return v
            return None

        def _count(v: Any) -> int:
            if isinstance(v, (list, tuple, set, dict)):
                return len(v)
            return 1 if _meaningful(v) else 0

        top = doc.get("top_module") if isinstance(doc.get("top_module"), dict) else {}
        purpose = ""
        if isinstance(top, dict):
            for k in ("description", "purpose", "summary"):
                cand = top.get(k)
                if isinstance(cand, str) and _meaningful(cand):
                    purpose = cand.strip()
                    break
        if purpose:
            out["purpose"] = purpose[:200]

        io_list = doc.get("io_list") if isinstance(doc.get("io_list"), dict) else {}
        interfaces = None
        if isinstance(io_list, dict):
            interfaces = io_list.get("interfaces")
        if not _meaningful(interfaces):
            interfaces = _section("interfaces")
        if _meaningful(interfaces):
            labels: list[str] = []
            iter_src = interfaces if isinstance(interfaces, list) else (
                list(interfaces.values()) if isinstance(interfaces, dict) else []
            )
            for ent in iter_src:
                if isinstance(ent, dict):
                    label = ent.get("name") or ent.get("type") or ent.get("role")
                    if label:
                        labels.append(str(label))
                elif isinstance(ent, str):
                    labels.append(ent)
            out["bus_interface"] = ", ".join(labels[:4]) if labels else "interfaces present"

        regs = _section("registers", "memoryMap", "memory_map")
        if _meaningful(regs):
            out["register_map"] = f"{_count(regs)} register entries"

        clocks_present = (
            _meaningful(_section("clock_reset_domains", "clock_reset", "clocks"))
            or (isinstance(io_list, dict) and (
                _meaningful(io_list.get("clocks")) or _meaningful(io_list.get("resets"))
            ))
        )
        if clocks_present:
            out["clock_reset"] = "clock_reset section present"

        if _meaningful(_section("interrupts")):
            out["interrupt"] = "interrupts section present"

        if _meaningful(_section("memory", "memory_map")):
            out["memory_map"] = "memory section present"

        params = _section("parameters")
        if not _meaningful(params) and isinstance(top, dict):
            params = top.get("parameters")
        if _meaningful(params):
            out["parameters"] = f"{_count(params)} parameters"

        if _meaningful(_section("sub_modules", "submodules", "decomposition")):
            out["submodule_structure"] = "submodule section present"

        if _meaningful(_section("test_requirements", "tb_plan", "test_plan", "tests")):
            out["test_expectation"] = "test section present"

        return out

    def _ssot_raw_top_keys(ip: str) -> set[str]:
        # Regex fallback when PyYAML can't parse the file. Real-world SSOT
        # files occasionally contain block-style entries that confuse the
        # parser; we still want the Validation pane to reflect the actual
        # set of top-level sections written to disk.
        path = _ssot_yaml_path(ip)
        if not path.is_file():
            return set()
        try:
            text = path.read_text(encoding="utf-8", errors="replace")
        except Exception:
            return set()
        keys: set[str] = set()
        for line in text.splitlines():
            if not line or line.startswith((" ", "\t", "#", "-")):
                continue
            stripped = line.rstrip()
            if not stripped.endswith(":") and ":" not in stripped:
                continue
            head = stripped.split(":", 1)[0].strip()
            if not head or not re.match(r"^[A-Za-z_][A-Za-z0-9_]*$", head):
                continue
            keys.add(head)
        return keys

    _SSOT_SECTION_ALIASES: dict[str, set[str]] = {
        "purpose":             {"top_module", "purpose", "overview"},
        "bus_interface":       {"io_list", "interfaces", "ports", "bus", "bus_interface"},
        "register_map":        {"registers", "register_map", "memoryMap", "memory_map"},
        "clock_reset":         {"clock_reset_domains", "clock_reset", "clocking", "clocks"},
        "interrupt":           {"interrupts", "interrupt"},
        "memory_map":          {"memory", "memory_map"},
        "parameters":          {"parameters"},
        "submodule_structure": {"sub_modules", "submodules", "decomposition", "architecture"},
        "test_expectation":    {"test_requirements", "tb_plan", "test_plan", "tests", "verification"},
    }

    def _decisions_from_top_keys(top_keys: set[str]) -> dict[str, str]:
        # When the yaml fails to parse, surface presence-only inference
        # using the section names we recovered via regex. This keeps the
        # Validation pane honest even on a structurally broken draft.
        if not top_keys:
            return {}
        out: dict[str, str] = {}
        for decision, aliases in _SSOT_SECTION_ALIASES.items():
            hit = aliases & top_keys
            if hit:
                out[decision] = f"{sorted(hit)[0]} section present"
        return out

    def _ssot_decisions(ip: str, state: dict[str, Any] | None = None) -> dict[str, str]:
        doc = _load_ssot_draft(ip)
        custom = doc.get("custom") if isinstance(doc.get("custom"), dict) else {}
        raw = custom.get("atlas_decisions") if isinstance(custom, dict) else {}
        if not isinstance(raw, dict) or not raw:
            legacy = state if isinstance(state, dict) else _load_ssot_state(ip)
            raw = legacy.get("decisions") if isinstance(legacy.get("decisions"), dict) else {}
        merged = {str(k): str(v).strip() for k, v in (raw or {}).items() if str(v or "").strip()}
        # Fill any decisions the promise tracker has not recorded yet from
        # whatever is actually present in the SSOT yaml on disk. Explicit
        # tracker entries win — inference only covers gaps.
        for k, v in _infer_decisions_from_yaml(doc).items():
            if k not in merged and v:
                merged[k] = v
        # If PyYAML failed (doc empty) but the file exists, fall back to
        # regex-detected top-level section names so the Validation pane
        # still shows the real on-disk state. Do not use this fallback on a
        # parseable draft: the /new-ip TBD scaffold has all top-level keys,
        # but those keys are not approved design decisions.
        if not doc:
            for k, v in _decisions_from_top_keys(_ssot_raw_top_keys(ip)).items():
                if k not in merged and v:
                    merged[k] = v
        return merged

    def _missing_ssot_decisions(ip: str, state: dict[str, Any] | None = None) -> list[str]:
        decisions = _ssot_decisions(ip, state)
        return [key for key, _ in _SSOT_REQUIRED_DECISIONS if not str(decisions.get(key) or "").strip()]

    def _record_ssot_decisions(
        ip: str,
        kind: str,
        updates: dict[str, str],
        sources: dict[str, list[dict[str, str]]] | None = None,
    ) -> tuple[list[str], list[dict[str, Any]]]:
        doc, custom = _ssot_custom(ip, kind)
        decisions = custom.get("atlas_decisions")
        if not isinstance(decisions, dict):
            decisions = {}
            custom["atlas_decisions"] = decisions
        decision_sources = custom.get("atlas_decision_sources")
        if not isinstance(decision_sources, dict):
            decision_sources = {}
            custom["atlas_decision_sources"] = decision_sources
        filled: list[str] = []
        conflicts: list[dict[str, Any]] = []
        source_map = sources or {}
        for key, value in updates.items():
            candidate = str(value or "").strip()
            if not candidate:
                continue
            existing = str(decisions.get(key) or "").strip()
            if existing:
                if re.sub(r"\s+", " ", existing).lower() != re.sub(r"\s+", " ", candidate).lower():
                    conflicts.append({
                        "key": key,
                        "existing": existing[:500],
                        "candidate": candidate[:500],
                        "sources": source_map.get(key, [])[:5],
                    })
                continue
            decisions[key] = candidate
            decision_sources[key] = source_map.get(key, [])[:8]
            filled.append(key)
        if conflicts:
            prior = custom.get("atlas_import_conflicts")
            if not isinstance(prior, list):
                prior = []
            custom["atlas_import_conflicts"] = prior + conflicts
        _save_ssot_draft(ip, doc)
        return filled, conflicts

    def _latest_pending_ssot_ip() -> str:
        root = PROJECT_ROOT / ".session"
        candidates: list[tuple[float, str]] = []
        if root.is_dir():
            for p in root.rglob("ssot-gen/state.json"):
                try:
                    doc = json.loads(p.read_text(encoding="utf-8"))
                    if isinstance(doc, dict) and not doc.get("approved"):
                        candidates.append((p.stat().st_mtime, p.parent.parent.name))
                except Exception:
                    continue
        candidates.sort(reverse=True)
        return candidates[0][1] if candidates else ""

    def _resolve_session_owner() -> str:
        """Extract the session_id (owner) from ATLAS_ACTIVE_SESSION env.

        Canonical session string layout for --ui atlas is always 3-part:
            <session_id>/<ip>/<workflow>
        This helper returns just the session_id (or empty string if no owner
        can be determined).

        Accepted current-env shapes:
          - 3+ parts:                       owner/ip/wf[/...] -> owner = parts[0]
          - 2 parts ending in 'default':    owner/default     -> owner = parts[0]
          - 1 part that is NOT an IP-like:  bare session_id   -> owner = parts[0]
        """
        current = normalize_session_name(str(_active_session_value() or ""))
        parts = [p for p in current.split("/") if p]
        if len(parts) >= 3:
            return parts[0]
        if len(parts) == 2 and parts[-1] == "default":
            return parts[0]
        if len(parts) == 1 and parts[0] and not _valid_ip_name(parts[0]):
            return parts[0]
        return ""

    def _session_owner_with_model(owner: str) -> str:
        """Optionally isolate session namespace per model.

        When enabled, the owner segment becomes:
            <owner>__<model_slug>
        so model runs cannot overwrite each other's .session trees.
        """
        base = str(owner or "default").strip() or "default"
        enabled = os.environ.get("ATLAS_SESSION_PER_MODEL", "0").strip().lower() in ("1", "true", "yes", "on")
        if not enabled:
            return base
        raw_model = (
            os.environ.get("LLM_ACTIVE_MODEL_NAME")
            or os.environ.get("MODEL_NAME")
            or os.environ.get("LLM_MODEL_NAME")
            or ""
        ).strip()
        if not raw_model:
            return base
        model_slug = re.sub(r"[^A-Za-z0-9_-]+", "_", raw_model).strip("_")
        if not model_slug:
            return base
        if base.endswith(f"__{model_slug}"):
            return base
        return f"{base}__{model_slug}"

    def _canonical_session_string(ip: str | None = None,
                                   workflow: str | None = None) -> str:
        """Return canonical 3-part session path string for --ui atlas:
            <session_id>/<ip>/<workflow>

        Any segment that is empty/None falls back to "default", so the
        on-disk layout is always 3 levels deep. Defaults can be overridden
        from CLI via `-s/-ip/-w` (env: ATLAS_DEFAULT_SESSION_ID,
        ATLAS_ACTIVE_IP, ATLAS_DEFAULT_WORKFLOW).
        """
        owner = _resolve_session_owner() or os.environ.get("ATLAS_DEFAULT_SESSION_ID") or "default"
        owner = _session_owner_with_model(owner)
        ip = ip or _active_ip_value() or "default"
        workflow = workflow or os.environ.get("ATLAS_DEFAULT_WORKFLOW") or "default"
        return f"{owner}/{ip}/{workflow}"

    def _canonical_session_dir(ip: str | None = None,
                                workflow: str | None = None) -> Path:
        """Filesystem dir for the canonical session path (3 segments)."""
        return PROJECT_ROOT / ".session" / _canonical_session_string(ip, workflow)

    def _legacy_session_candidates(ip: str | None,
                                    workflow: str | None) -> list[Path]:
        """Legacy on-disk locations used before path canonicalization.
        Used by readers to recover history from older runs."""
        ip_seg = ip or "default"
        wf_seg = workflow or "default"
        cands: list[Path] = []
        # 2-part: <ip>/<workflow>
        cands.append(PROJECT_ROOT / ".session" / ip_seg / wf_seg)
        # 1-part: <workflow> alone (very old)
        cands.append(PROJECT_ROOT / ".session" / wf_seg)
        # 1-part: <ip> alone
        cands.append(PROJECT_ROOT / ".session" / ip_seg)
        return cands

    def _migrate_legacy_session(ip: str | None,
                                 workflow: str | None) -> Path:
        """Locate canonical path; if absent and a legacy path exists, move it
        in-place. Returns the canonical path (whether or not migration ran)."""
        canon = _canonical_session_dir(ip, workflow)
        if canon.exists():
            return canon
        for legacy in _legacy_session_candidates(ip, workflow):
            try:
                if legacy.exists() and legacy != canon and legacy.is_dir():
                    canon.parent.mkdir(parents=True, exist_ok=True)
                    legacy.rename(canon)
                    break
            except OSError:
                continue
        return canon

    def _set_active_ssot_ip(ip: str, workflow: str = "ssot-gen") -> None:
        if not _valid_ip_name(ip):
            return
        workflow = normalize_session_name(str(workflow or "ssot-gen")) or "ssot-gen"
        _atlas_active_ip_cv.set(ip)
        _atlas_active_session_cv.set(_canonical_session_string(ip, workflow))

    def _infer_ip_from_project_root() -> str:
        """Infer an IP when ATLAS is pointed directly at an IP workspace."""
        yaml_dir = PROJECT_ROOT / "yaml"
        if yaml_dir.is_dir():
            for pattern in ("*.ssot.yaml", "*_ssot.yaml", "*.ssot.yml"):
                for path in sorted(yaml_dir.glob(pattern)):
                    name = path.name
                    if name.endswith(".ssot.yaml"):
                        candidate = name[: -len(".ssot.yaml")]
                    elif name.endswith("_ssot.yaml"):
                        candidate = name[: -len("_ssot.yaml")]
                    elif name.endswith(".ssot.yml"):
                        candidate = name[: -len(".ssot.yml")]
                    else:
                        candidate = path.stem
                    if _valid_ip_name(candidate):
                        return candidate
            if _valid_ip_name(PROJECT_ROOT.name):
                return PROJECT_ROOT.name
        nested = sorted({
            p.parents[1].name
            for pattern in ("*/yaml/*.ssot.yaml", "*/yaml/*_ssot.yaml", "*/yaml/*.ssot.yml")
            for p in PROJECT_ROOT.glob(pattern)
            if _valid_ip_name(p.parents[1].name)
        })
        if len(nested) == 1:
            return nested[0]
        return ""

    def _session_ip_from_namespace(session: str) -> str:
        session = normalize_session_name(str(session or ""))
        parts = [p for p in session.split("/") if p]
        if len(parts) >= 3 and _valid_ip_name(parts[-2]) and parts[-2] != "default":
            return parts[-2]
        if len(parts) == 1 and _valid_ip_name(parts[0]) and parts[0] != "default" and _ssot_state_path(parts[0]).is_file():
            return parts[0]
        return ""

    def _active_ssot_ip() -> str:
        env_ip = str(_active_ip_value() or "").strip()
        if _valid_ip_name(env_ip) and env_ip != "default":
            return env_ip
        session_ip = _session_ip_from_namespace(str(_active_session_value() or ""))
        if session_ip:
            return session_ip
        root_ip = _infer_ip_from_project_root()
        if _valid_ip_name(root_ip):
            return root_ip
        return _latest_pending_ssot_ip()

    def _first_ip_token(args: str) -> str:
        try:
            tokens = shlex.split(args or "")
        except ValueError:
            tokens = str(args or "").split()
        skip_next = False
        for tok in tokens:
            if skip_next:
                skip_next = False
                continue
            if tok in {"--mode", "--run-mode", "--preview", "--preview-contract", "--top"}:
                skip_next = True
                continue
            if tok.startswith("-"):
                continue
            return tok
        return ""

    def _command_ip(args: str = "", client_session: Any | None = None) -> str:
        explicit = _first_ip_token(args)
        if _valid_ip_name(explicit):
            return explicit
        session_ip = _session_ip_from_namespace(getattr(client_session, "session_id", ""))
        if _valid_ip_name(session_ip):
            return session_ip
        return _active_ssot_ip()

    def _ip_root(ip: str) -> Path:
        configured = _configured_ip_root(ip)
        if configured is not None:
            return configured
        nested = PROJECT_ROOT / ip
        if nested.exists():
            return nested
        if PROJECT_ROOT.name == ip and (PROJECT_ROOT / "yaml").is_dir():
            return PROJECT_ROOT
        return nested

    def _script_project_root(ip: str) -> Path:
        ip_dir = _ip_root(ip)
        if ip_dir == PROJECT_ROOT and PROJECT_ROOT.name == ip:
            return PROJECT_ROOT.parent
        try:
            resolved_ip = ip_dir.resolve()
            resolved_project = PROJECT_ROOT.resolve()
            if resolved_ip != resolved_project and resolved_ip.parent != resolved_project and (resolved_ip / "yaml").is_dir():
                return resolved_ip.parent
        except Exception:
            pass
        return PROJECT_ROOT

    def _context_for_client_session(client_session: Any | None) -> AtlasContext | None:
        session_key = normalize_session_name(str(getattr(client_session, "session_id", "") or ""))
        if not session_key:
            return None
        try:
            return AtlasContext.from_session_key(
                session_key,
                atlas_root=os.environ.get("ATLAS_ROOT") or str(PROJECT_ROOT),
            )
        except ValueError:
            return None

    def _ip_root_for_session(ip: str, client_session: Any | None = None) -> Path:
        context = _context_for_client_session(client_session)
        if context is not None and not context.legacy and context.ip_name == ip:
            root, error = _validated_context_workspace_root(context)
            if not error and root is not None:
                return root / ip
        return _ip_root(ip)

    def _script_project_root_for_session(ip: str, client_session: Any | None = None) -> Path:
        context = _context_for_client_session(client_session)
        if context is not None and not context.legacy and context.ip_name == ip:
            root, error = _validated_context_workspace_root(context)
            if not error and root is not None:
                return root
        return _script_project_root(ip)

    def _ssot_yaml_path_for_session(ip: str, client_session: Any | None = None) -> Path:
        ip_dir = _ip_root_for_session(ip, client_session)
        for name in (f"{ip}.ssot.yaml", f"{ip}_ssot.yaml", f"{ip}.ssot.yml"):
            candidate = ip_dir / "yaml" / name
            if candidate.is_file():
                return candidate
        return ip_dir / "yaml" / f"{ip}.ssot.yaml"

    def _render_new_ip_plan(ip: str, kind: str, state: dict[str, Any]) -> str:
        missing = _missing_ssot_decisions(ip, state)
        lines = [
            f"[SSOT PLAN] {ip}",
            f"kind: {kind or 'simple APB peripheral'}",
            "mode: structure only; no document import is run by /new-ip",
            "",
            "Created structure:",
            f"- {ip}/doc, {ip}/req, {ip}/yaml",
            f"- {ip}/rtl, {ip}/list, {ip}/tb/cocotb",
            f"- {ip}/tc, {ip}/sim, {ip}/cov, {ip}/lint",
            "",
            "SSOT decisions still needed before production YAML write:",
        ]
        for key, label in _SSOT_REQUIRED_DECISIONS:
            mark = "✓" if key not in missing else "·"
            lines.append(f"- {mark} `{key}`: {label}")
        lines += [
            "",
            "Short flow:",
            f"1. Put source docs, RTL, specs, logs, or notes anywhere under `{ip}/`",
            f"2. Run `/import {ip}` or `/import @path` to scan the workspace into SSOT section TODOs",
            f"3. Run `/to-ssot {ip}`; use `/grill-me` only for gaps or conflicts",
        ]
        if missing:
            lines.append("")
            lines.append("missing decisions: " + ", ".join(missing))
        return "\n".join(lines)

    def _parse_new_ip_args(args: str) -> tuple[str, str, list[str], str]:
        import shlex

        try:
            tokens = [t.strip().strip("\"'") for t in shlex.split(args or "", posix=False) if t.strip()]
        except ValueError as exc:
            return "", "", [], f"cannot parse /new-ip arguments: {exc}"
        if not tokens:
            return "", "", [], ""
        ip = tokens[0]
        import_paths: list[str] = []
        kind_tokens: list[str] = []
        idx = 1
        while idx < len(tokens):
            tok = tokens[idx]
            if tok in ("--import", "--doc", "--docs"):
                if idx + 1 >= len(tokens):
                    return "", "", [], f"missing value after {tok}"
                import_paths.append(tokens[idx + 1])
                idx += 2
                continue
            if tok.startswith("@"):
                import_paths.append(tok)
                idx += 1
                continue
            kind_tokens.append(tok)
            idx += 1
        kind = " ".join(kind_tokens).strip() or "TBD"
        return ip, kind, import_paths, ""

    def _ssot_session_for_ip(ip: str) -> str:
        current = normalize_session_name(str(_active_session_value() or ""))
        parts = [p for p in current.split("/") if p]
        if len(parts) >= 3 and parts[-2] == ip and parts[-1] == "ssot-gen":
            return current
        return _canonical_session_string(ip, "ssot-gen")

    def _render_ssot_llm_qna_prompt(ip: str, kind: str, state: dict[str, Any]) -> str:
        session = _ssot_session_for_ip(ip)
        imported = state.get("imported_artifacts") if isinstance(state.get("imported_artifacts"), list) else []
        imported_paths = [
            str(item.get("path") or "").strip()
            for item in imported
            if isinstance(item, dict) and str(item.get("path") or "").strip()
        ]
        missing = _missing_ssot_decisions(ip, state)
        import_context_paths = [
            f"{ip}/wiki/index.md",
            f"{ip}/wiki/_graph.json",
            f"{ip}/wiki/import-evidence.md",
            f"{ip}/req/import_manifest.json",
            f"{ip}/req/extracted_decisions.json",
            f"{ip}/req/imports/",
            *imported_paths[:24],
        ]
        seen_import_context: set[str] = set()
        import_context_paths = [
            path
            for path in import_context_paths
            if path and not (path in seen_import_context or seen_import_context.add(path))
        ]
        path_lines = "\n".join(f"- {p}" for p in import_context_paths) or "- (none recorded; inspect the IP directory and draft SSOT)"
        missing_line = ", ".join(missing) if missing else "(backend baseline decisions already filled; still inspect for SSOT TBD/conflicts)"
        return "\n".join([
            f"You are ssot-gen for IP `{ip}` in ATLAS UI.",
            f"Session: `{session}`",
            "Goal: create IP-specific SSOT Q&A from the current evidence, not from a fixed template.",
            "This is a general-IP flow. Do not assume APB/register-only/simple peripheral structure unless evidence says so.",
            "Use the per-IP wiki as the navigation index. Follow relevant wiki links and graph nodes for import history, notes, requirements, SSOT, RTL, model, verification, and logs.",
            "Do not stop at a file list. Inspect imported markdown, wiki import evidence, existing SSOT draft, and any other relevant IP artifacts before deciding what needs human clarification.",
            "",
            "Truth ownership model:",
            "- Human owns requirement/spec/interface/FL golden model/coverage goals/performance targets/sign-off.",
            "- LLM owns drafting, import analysis, QA generation, SSOT patch proposals, and downstream workflow_todos.",
            "- Do not change locked truth to make downstream RTL pass; make a change-request question instead.",
            "- TODOs are execution work, not substitutes for unresolved human decisions.",
            "",
            "Current backend baseline missing keys, for orientation only:",
            f"- {missing_line}",
            "",
            "Evidence paths imported or known:",
            path_lines,
            "",
            "Required action:",
            f"1. Read `{ip}/wiki/index.md`, `{ip}/wiki/_graph.json`, `{ip}/wiki/import-evidence.md`, `{ip}/req/import_manifest.json`, `{ip}/req/extracted_decisions.json`, `{ip}/req/imports/`, and `{ip}/yaml/{ip}.ssot.yaml` if they exist.",
            f"2. Follow relevant wiki links and inspect relevant files under `{ip}/req`, `{ip}/doc`, `{ip}/rtl`, `{ip}/model`, `{ip}/tb`, `{ip}/sim`, `{ip}/lint`, and `{ip}/logs` when the index or evidence points there.",
            "3. Detect unresolved SSOT decisions, contradictions, assumptions, TBD/null/placeholders, vague imported facts, and any truth that needs human approval or concretization.",
            "4. Generate ONLY the questions needed for this IP. The question set may be 0, 1, 4, 20, or more depending on complexity.",
            "5. If the answer is not an immediate blocker, use `record_ssot_qa(questions=[...])` to save deferred QA cards.",
            "6. Use `ask_user(questions=[...])` only when the answer blocks the next SSOT write or import pass.",
            "   Do not ask plain prose questions in chat. Both tools preserve SSOT QA metadata.",
            "7. Each question object must carry metadata so ATLAS can save it in SSOT QA preview:",
            "   - id: stable snake_case id",
            "   - section_id: canonical section bucket such as 00_overview, 03_interface, 06_registers, 18_verification, 19_workflow_todos, or a specific section number",
            "   - section_title: human-readable SSOT section title",
            "   - decision_key: stable key for the decision",
            "   - decision_label: short label",
            "   - qa_type: human_decision | clarification | change_request | execution_blocker",
            "   - question, subtitle, kind, options when useful",
            "   - criteria: pass/fail criteria for using the answer downstream",
            "   - source_refs: SSOT paths, doc paths, or RTL paths that caused the question",
            "8. Prefer section-specific QA cards. Group by SSOT section and ask concrete decisions, not generic template prompts.",
            "8. If downstream RTL needs explicit decomposition, write `workflow_todos.rtl-gen[]` with content/detail/criteria/source_refs.",
            "9. If no immediate answer is needed after recording deferred QA, say `[SSOT Q&A] deferred questions recorded` with a short evidence summary.",
            "10. If no human decision is needed at all, say `[SSOT Q&A] no generated questions required` and explain the evidence briefly.",
            "",
            "Important: fixed question templates are forbidden here. Derive the QA from this IP's evidence and current SSOT only.",
        ])

    def _render_approved_ssot_spec(ip: str, state: dict[str, Any]) -> str:
        decisions = _ssot_decisions(ip, state)
        req_dir = _ip_root(ip) / "req"

        def _read_json_file(path: Path) -> dict[str, Any]:
            try:
                doc = json.loads(path.read_text(encoding="utf-8"))
                return doc if isinstance(doc, dict) else {}
            except Exception:
                return {}

        approval = _read_json_file(req_dir / "approval_manifest.json")
        requirements_doc = _read_json_file(req_dir / "requirements_index.json")
        obligations_doc = _read_json_file(req_dir / "obligations.json")
        contracts_doc = _read_json_file(req_dir / "contract_refs.json")
        structural_doc = _read_json_file(req_dir / "structural_contracts.json")
        behavioral_doc = _read_json_file(req_dir / "behavioral_contracts.json")
        evidence_doc = _read_json_file(req_dir / "evidence_plan.json")
        locked_truth_active = str(approval.get("status") or approval.get("locked_truth_status") or "").strip() == "requirements_locked"

        lines = [
            f"[WEB TO SSOT SPEC] {ip}",
            f"kind: {state.get('kind') or 'generic IP block'}",
            (
                "source: Approved Locked Truth req/ bundle + Web UI Plan Mode / import evidence"
                if locked_truth_active else
                "source: Web UI Plan Mode + SSOT draft decisions + per-IP wiki/import evidence"
            ),
            "",
            (
                "Use req/ as the authority for /to-ssot. Existing YAML is only a Design Spec projection and may be a stale TBD skeleton."
                if locked_truth_active else
                "Use this with the per-IP wiki as the source of truth for /to-ssot. Do not invent over missing fields."
            ),
            "Required evidence navigation before writing SSOT:",
        ]
        if locked_truth_active:
            req_count = len(requirements_doc.get("requirements") or [])
            obl_count = len(obligations_doc.get("obligations") or [])
            contract_count = len(contracts_doc.get("contract_refs") or [])
            structural_count = len(structural_doc.get("contracts") or [])
            behavioral_count = len(behavioral_doc.get("contracts") or [])
            evidence_count = len(evidence_doc.get("evidence_plan") or [])
            lines.extend([
                f"- {ip}/req/approval_manifest.json (status=requirements_locked, bundle_sha256={approval.get('bundle_sha256') or '(missing)'})",
                f"- {ip}/req/requirements_index.json ({req_count} requirements)",
                f"- {ip}/req/obligations.json ({obl_count} obligations)",
                f"- {ip}/req/contract_refs.json ({contract_count} contract refs)",
                f"- {ip}/req/structural_contracts.json ({structural_count} structural contracts)",
                f"- {ip}/req/behavioral_contracts.json ({behavioral_count} behavioral contracts)",
                f"- {ip}/req/evidence_plan.json ({evidence_count} evidence plan entries)",
                f"- {ip}/req/locked_truth.md",
                f"- {ip}/yaml/{ip}.ssot.yaml only as prior projection; replace/repair it if it is TBD-filled.",
                "",
                "Locked Truth projection rule:",
                "- Derive SSOT Preview blockers from req/ before asking the user: top_module.description, io_list.interfaces[].ports[], registers.register_list[], clock_reset_domains, interrupts, scenarios, quality_gates, and traceability.",
                "- Project structural_contracts into io_list and clock/reset domains exactly: signal name, direction, width, sync/async timing, and clock/reset association.",
                "- Project behavioral_contracts into function_model, cycle_model, fsm/register behavior, test_requirements, and quality_gates with decision tables or equivalent machine rules.",
                "- Map requirements -> obligations -> contract_refs/structural_contracts/behavioral_contracts -> evidence_plan into source_refs/contract_refs/evidence_refs on Design Spec items.",
                "- If req/ does not specify memory, FSM, child submodules, DFT, power, or security behavior, write an explicit no-feature/external-owner/non-goal policy instead of a TBD placeholder.",
                "- Do not modify canonical req/*.json from ssot-gen; write only the YAML projection and SSOT-side validation artifacts.",
                "",
                "Locked requirements:",
            ])
            for row in requirements_doc.get("requirements") or []:
                if not isinstance(row, dict):
                    continue
                rid = str(row.get("requirement_id") or row.get("id") or "").strip()
                title = str(row.get("title") or "").strip()
                statement = str(row.get("statement") or row.get("requirement") or "").strip()
                obligations = row.get("obligation_refs") or row.get("obligations") or []
                lines.append(f"- {rid or '(unnamed)'}: {title} — {statement}")
                if obligations:
                    if all(isinstance(item, str) for item in obligations):
                        lines.append(f"  obligation_refs: {', '.join(str(item) for item in obligations[:8])}")
                    else:
                        lines.append(f"  obligations: {len(obligations)} inline item(s)")
        else:
            lines.extend([
                f"- {ip}/wiki/index.md",
                f"- {ip}/wiki/_graph.json",
                f"- {ip}/wiki/import-evidence.md",
                f"- {ip}/wiki/log.md",
                f"- {ip}/wiki/notes.md",
                f"- {ip}/req/imports/",
                f"- relevant files linked by the wiki/import evidence under {ip}/req, {ip}/doc, {ip}/rtl, {ip}/model, {ip}/tb, {ip}/sim, {ip}/lint, and {ip}/logs",
                "Generation rule: reconcile current decisions with import history, wiki summaries, source excerpts, conflicts, and downstream TODO evidence.",
            ])
        for key, _label in _SSOT_REQUIRED_DECISIONS:
            lines.append(f"- {key}: {decisions.get(key) or '(missing)'}")
        return "\n".join(lines)

    def _emit_ssot_approval_ready(ip: str, state: dict[str, Any], missing: list[str] | None = None) -> None:
        decisions = _ssot_decisions(ip, state)
        miss = missing if missing is not None else _missing_ssot_decisions(ip, state)
        bridge.emit(
            "ssot_approval_ready",
            ip=ip,
            kind=state.get("kind") or "TBD",
            status=state.get("status") or ("approved" if state.get("approved") else "answered"),
            approved=bool(state.get("approved")),
            missing=miss,
            decisions=decisions,
            approve_cmd=f"approve {ip}",
            generate_cmd=f"/to-ssot {ip}",
        )

    def _answer_text(answer: dict[str, Any], question: dict[str, Any]) -> str:
        custom = str(answer.get("custom") or "").strip()
        if custom:
            return custom
        selected = answer.get("selected") or []
        by_id = {str(o.get("id")): str(o.get("label") or o.get("id"))
                 for o in (question.get("options") or []) if isinstance(o, dict)}
        labels = [by_id.get(str(s), str(s)) for s in selected]
        return ", ".join([x for x in labels if x]).strip()

    def _new_ssot_state(ip: str, kind: str = "TBD") -> dict[str, Any]:
        return {
            "ip": ip,
            "kind": kind,
            "approved": False,
            "approved_at": 0,
            "status": "planned",
            "active_session": _ssot_session_for_ip(ip),
            "last_step": "new-ip",
            "created_at": time.time(),
        }

    def _scaffold_ip_wiki(ip: str, base_root: Path | None = None) -> None:
        """Seed <ip>/wiki/{index.md, log.md, notes.md} idempotently.

        Karpathy-style: `[[link]]` ToC into the IP tree + append-only log + free-form notes.
        Re-runs of /new-ip leave existing pages untouched.
        """
        root = base_root or PROJECT_ROOT
        wiki_dir = root / ip / "wiki"
        wiki_dir.mkdir(parents=True, exist_ok=True)
        today = time.strftime("%Y-%m-%d", time.gmtime())
        seeds = {
            "index.md": (
                f"# {ip} IP Wiki\n\n"
                f"Per-IP knowledge base. Read with `wiki_query(ip=\"{ip}\")` or run\n"
                f"`python3 \"$ATLAS_WORKFLOW_ROOT/wiki/build_graph.py\" --ip {ip} --root \"$ATLAS_PROJECT_ROOT\"` to refresh the index.\n\n"
                "## Status snapshot\n\n"
                "Populated by `$ATLAS_WORKFLOW_ROOT/wiki/build_graph.py --ip <ip> --root $ATLAS_PROJECT_ROOT`; the synthetic\n"
                "`[[ssot]]`, `[[fl_model]]`, `[[cl_model]]`, `[[rtl]]`, `[[filelist]]`,\n"
                "`[[lint]]`, `[[tb]]`, `[[sim]]`, `[[coverage]]`, `[[audit]]`, and\n"
                "`[[last_run]]` nodes carry status/digest fields.\n\n"
                "## Tree\n\n"
                f"- [[notes]] — free-form owner/manager notes\n"
                f"- [[log]] — append-only event log\n"
                f"- requirements at `../req/`\n"
                f"- SSOT YAML at `../yaml/{ip}.ssot.yaml`\n"
                f"- function/cycle model at `../model/`\n"
                f"- RTL at `../rtl/` (filelist `../list/{ip}.f`)\n"
                f"- testbench at `../tb/`\n"
                f"- sim evidence at `../sim/`\n"
                f"- lint/coverage at `../lint/` and `../cov/`\n"
                f"- run logs at `../logs/`\n"
            ),
            "log.md": (
                "# Wiki Log\n\n"
                f"## [{today}] new-ip | scaffolded {ip} wiki\n"
            ),
            "notes.md": (
                f"# {ip} Notes\n\n"
                "Free-form notes for the IP owner and manager. Tooling does not\n"
                "write to this file; only humans (and chat agents acting on a\n"
                "user request) edit it. Cross-link with `[[ssot]]`, `[[rtl]]`,\n"
                "`[[sim]]`, etc. to compound context across sessions.\n"
            ),
        }
        for name, content in seeds.items():
            target = wiki_dir / name
            if not target.exists():
                try:
                    target.write_text(content, encoding="utf-8")
                except Exception:
                    pass

    def _compact_import_wiki_text(value: Any, limit: int = 700) -> str:
        text = " ".join(str(value or "").strip().split())
        if len(text) <= limit:
            return text
        return text[: max(0, limit - 3)].rstrip() + "..."

    def _refresh_ip_wiki_graph(ip: str, base_root: Path | None = None) -> None:
        try:
            import subprocess

            script = WORKFLOW_ROOT / "wiki" / "build_graph.py"
            root = base_root or PROJECT_ROOT
            subprocess.run(
                [
                    "python3",
                    str(script),
                    "--ip",
                    ip,
                    "--project-root",
                    str(root),
                    "--quiet",
                ],
                check=False,
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
                timeout=10,
            )
        except Exception:
            pass

    def _write_import_wiki_index(
        ip: str,
        kind: str,
        artifacts: list[dict[str, Any]],
        candidates: dict[str, str],
        sources: dict[str, list[dict[str, str]]],
        filled: list[str],
        conflicts: list[dict[str, Any]],
        todo_summary: dict[str, Any],
        next_action: str,
    ) -> None:
        try:
            _scaffold_ip_wiki(ip)
            wiki_dir = PROJECT_ROOT / ip / "wiki"
            wiki_dir.mkdir(parents=True, exist_ok=True)
            updated = time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())
            page = wiki_dir / "import-evidence.md"
            lines: list[str] = [
                "---",
                "type: reference",
                "tags: [import, ssot, requirements]",
                f"updated: {updated[:10]}",
                "---",
                "",
                f"# {ip} Import Evidence",
                "",
                "Index for files imported into the SSOT workflow. This page is generated by `/import` so agents can find imported requirements without re-scanning the full IP tree.",
                "",
                "## Summary",
                "",
                f"- IP: `{ip}`",
                f"- Kind: `{kind}`",
                f"- Last import: `{updated}`",
                f"- Imported source directory: `{ip}/req/imports/`",
                f"- Imported artifacts: {len(artifacts)}",
                f"- Candidate facts: {len(candidates)}",
                f"- Filled decisions: {len(filled)}",
                f"- Conflicts: {len(conflicts)}",
                f"- Next action: `{next_action}`",
                "",
                "## Imported files",
                "",
            ]
            if artifacts:
                lines.extend(["| path | bytes | derived md |", "| --- | ---: | --- |"])
                for artifact in artifacts[:80]:
                    path = _compact_import_wiki_text(artifact.get("path"), 240).replace("|", "\\|")
                    size = artifact.get("size_bytes")
                    if not isinstance(size, int):
                        size = artifact.get("bytes")
                    md_path = _compact_import_wiki_text(artifact.get("md_path"), 240).replace("|", "\\|")
                    lines.append(f"| `{path}` | {size if isinstance(size, int) else ''} | `{md_path}` |")
            else:
                lines.append("- No imported artifacts recorded.")
            lines.extend(["", "## Candidate SSOT facts", ""])
            if candidates:
                for key in sorted(candidates):
                    lines.extend([f"### {key}", "", _compact_import_wiki_text(candidates.get(key), 1200), ""])
            else:
                lines.append("- No candidate facts extracted.")
            lines.extend(["", "## Source excerpts", ""])
            if sources:
                for key in sorted(sources):
                    rows = sources.get(key) or []
                    lines.extend([f"### {key}", ""])
                    for row in rows[:8]:
                        path = _compact_import_wiki_text(row.get("path"), 240)
                        excerpt = _compact_import_wiki_text(row.get("excerpt"), 700)
                        lines.append(f"- `{path}`: {excerpt}")
                    lines.append("")
            else:
                lines.append("- No source excerpts recorded.")
            lines.extend(["", "## Filled decisions", ""])
            if filled:
                for key in filled:
                    lines.append(f"- `{key}`")
            else:
                lines.append("- None.")
            lines.extend(["", "## Conflicts", ""])
            if conflicts:
                for conflict in conflicts[:40]:
                    field = _compact_import_wiki_text(conflict.get("field"), 120)
                    existing = _compact_import_wiki_text(conflict.get("existing"), 360)
                    imported = _compact_import_wiki_text(conflict.get("imported"), 360)
                    lines.extend([f"### {field}", "", f"- Existing: {existing}", f"- Imported: {imported}", ""])
            else:
                lines.append("- None.")
            lines.extend(["", "## Workflow todo summary", ""])
            if isinstance(todo_summary, dict) and todo_summary:
                for key in sorted(todo_summary):
                    lines.append(f"- `{key}`: `{todo_summary.get(key)}`")
            else:
                lines.append("- None.")
            lines.extend([
                "",
                "## Agent usage",
                "",
                f"- Deep Interview should start from `{ip}/wiki/index.md` and `{ip}/wiki/_graph.json`, then read this `[[import-evidence]]` page, `[[log]]`, `[[notes]]`, `{ip}/req/imports/`, and any linked artifacts that look relevant.",
                f"- Deep Interview should ask only the concretization questions needed after comparing import history, candidate facts, SSOT draft, requirements, RTL/model, and verification evidence.",
                f"- To SSOT should use the wiki index, import evidence, candidate facts, conflicts, and source excerpts as evidence, then write `{ip}/yaml/{ip}.ssot.yaml`.",
            ])
            page.write_text("\n".join(lines).rstrip() + "\n", encoding="utf-8")

            index_path = wiki_dir / "index.md"
            if index_path.exists():
                index_text = index_path.read_text(encoding="utf-8")
                if "[[import-evidence]]" not in index_text:
                    index_text = index_text.rstrip() + (
                        "\n\n## Import Evidence\n\n"
                        "- [[import-evidence]] - imported files, extracted SSOT facts, conflicts, and next actions.\n"
                    )
                    index_path.write_text(index_text, encoding="utf-8")

            log_path = wiki_dir / "log.md"
            log_line = (
                f"\n## [{updated}] import | {len(artifacts)} artifacts\n\n"
                f"- Page: [[import-evidence]]\n"
                f"- Filled decisions: {len(filled)}\n"
                f"- Conflicts: {len(conflicts)}\n"
                f"- Next action: `{next_action}`\n"
            )
            with log_path.open("a", encoding="utf-8") as handle:
                handle.write(log_line)
            _refresh_ip_wiki_graph(ip)
        except Exception:
            pass

    def _ensure_new_ip_structure(ip: str, base_root: Path | None = None) -> list[str]:
        root = base_root or PROJECT_ROOT
        dirs = [
            "doc",
            "req",
            "yaml",
            "rtl",
            "list",
            "tb/cocotb",
            "tc",
            "sim",
            "cov",
            "lint",
            "wiki",
        ]
        created: list[str] = []
        for rel in dirs:
            path = root / ip / rel
            path.mkdir(parents=True, exist_ok=True)
            created.append(f"{ip}/{rel}")
        _scaffold_ip_wiki(ip, base_root=root)
        # Copy the FULL workflow engine into <ip>/workflow/ (every stage's
        # scripts, prompts, system_prompt.md, rules, todo templates, shared
        # scripts/ + prompts/, flow guides) AND generate the wiki/_generated/
        # runbook + ip_knowledge pages — the same scaffold the tool path runs.
        # Without this a UI-created IP had no <ip>/workflow/ scripts on disk, so
        # the runbook's `workflow/<stage>/scripts/...` commands had nothing to
        # run. Then rebuild the wiki graph. Best-effort: never block creation.
        try:
            from core.tools import _scaffold_ip_workflow as _scaffold_wf
            _scaffold_wf(str(root / ip), ip, [], [], [])
        except Exception:
            pass
        _refresh_ip_wiki_graph(ip, base_root=root)
        # Per-IP git repo. Each IP gets its OWN .git so the agent's
        # write_file / replace_in_file calls can auto-commit and the
        # user has a per-IP history independent of the outer project
        # repo. Idempotent — `git init` on an existing repo is a no-op.
        _ip_root = root / ip
        _git_dir = _ip_root / ".git"
        _gitignore = _ip_root / ".gitignore"
        # Ignore the heavy/derived artifacts so per-IP git history stays
        # small and reviewable. write it before `git init` so the first
        # `git add .` doesn't sweep up sim/*.vcd etc.
        if not _gitignore.exists():
            try:
                _gitignore.write_text(
                    "# Atlas-managed — generated/binary artifacts excluded\n"
                    "__pycache__/\n*.pyc\n.DS_Store\n*.log\n"
                    "sim/build/\nsim/results/\nsim/work/\n"
                    "sim/*.vcd\nsim/*.fst\nsim/*.shm/\nsim/*.vpd\nsim/*.wlf\n"
                    "tb/cocotb/__pycache__/\ntb/cocotb/results.xml\n"
                    "cov/build/\ncov/*.dat\ncov/*.ucdb\n"
                    "lint/build/\nlint/*.rpt\nlint/*.tmp\n"
                    ".vscode/\n.idea/\n",
                    encoding="utf-8",
                )
            except Exception:
                pass
        try:
            import subprocess as _sp_init
            if not _git_dir.is_dir():
                _sp_init.run(
                    ["git", "init", "-q", "-b", "main"],
                    cwd=str(_ip_root),
                    capture_output=True,
                    timeout=10,
                )
                # Pin the committer to a benign default so `git commit`
                # doesn't fail with "Please tell me who you are" on
                # fresh boxes that have no global git config.
                for _k, _v in (("user.email", "atlas@local"),
                                ("user.name",  "Atlas Agent")):
                    _sp_init.run(
                        ["git", "config", _k, _v],
                        cwd=str(_ip_root),
                        capture_output=True,
                        timeout=5,
                    )
                # Initial commit so subsequent auto-commits have a
                # parent — `git commit --allow-empty` keeps it clean
                # even when the dirs are empty at scaffold time.
                _sp_init.run(
                    ["git", "add", "--", ".gitignore"],
                    cwd=str(_ip_root),
                    capture_output=True,
                    timeout=5,
                )
                _sp_init.run(
                    ["git", "commit", "--allow-empty",
                     "-m", f"atlas: scaffold {ip}"],
                    cwd=str(_ip_root),
                    capture_output=True,
                    timeout=10,
                )
            # BARE_GIT_OPTION: also stand up a central bare repo as a
            # sibling of the IP (<root>/<ip>.git). Wire the working
            # repo's origin to it + install hooks so that:
            #   • atlas's per-edit auto-commits propagate to the bare
            #     via post-commit `git push origin main`
            #   • external pushes (`git push http://host:port/git/<ip>.git`)
            #     land in the bare, and a post-receive hook fast-forwards
            #     the atlas working tree so the agent sees the change
            try:
                import config as _cfg_bare
                _bare_on = getattr(_cfg_bare, "BARE_GIT_OPTION", True)
            except Exception:
                _bare_on = True
            if _bare_on:
                import shlex as _shlex_bare
                _bare_dir = root / f"{ip}.git"
                _post_commit  = _ip_root / ".git" / "hooks" / "post-commit"
                _post_receive = _bare_dir / "hooks" / "post-receive"
                try:
                    if not (_bare_dir / "HEAD").is_file():
                        _sp_init.run(
                            ["git", "init", "--bare", "-q", "-b", "main",
                             str(_bare_dir)],
                            capture_output=True, timeout=10,
                        )
                    # git-http-backend refuses receive-pack by default;
                    # enable it so /git/<ip>.git accepts `git push` over
                    # smart HTTP. Idempotent.
                    _sp_init.run(
                        ["git", "config", "http.receivepack", "true"],
                        cwd=str(_bare_dir), capture_output=True, timeout=5,
                    )
                    # Wire / re-wire working repo's origin to the bare.
                    _sp_init.run(
                        ["git", "remote", "remove", "origin"],
                        cwd=str(_ip_root), capture_output=True, timeout=5,
                    )
                    _sp_init.run(
                        ["git", "remote", "add", "origin",
                         str(_bare_dir.resolve())],
                        cwd=str(_ip_root), capture_output=True, timeout=5,
                    )
                    # Initial push so the bare has the scaffold commit.
                    _sp_init.run(
                        ["git", "push", "-u", "-q", "origin", "main"],
                        cwd=str(_ip_root), capture_output=True, timeout=15,
                    )
                    # post-commit hook in working repo → push to bare.
                    _post_commit.parent.mkdir(parents=True, exist_ok=True)
                    _post_commit.write_text(
                        "#!/bin/sh\n"
                        "# Atlas — mirror each working-tree commit to the central bare.\n"
                        "git push -q origin HEAD >/dev/null 2>&1 || true\n",
                        encoding="utf-8",
                    )
                    _post_commit.chmod(0o755)
                    # post-receive hook in bare → fast-forward the
                    # working tree so external pushes land in the live
                    # IP dir the agent reads from.
                    _post_receive.parent.mkdir(parents=True, exist_ok=True)
                    _post_receive.write_text(
                        "#!/bin/sh\n"
                        "# Atlas — fast-forward the working tree when an external push arrives.\n"
                        f"WORK={_shlex_bare.quote(str(_ip_root.resolve()))}\n"
                        "( cd \"$WORK\" && unset GIT_DIR && "
                        "  git fetch -q origin && "
                        "  git reset -q --hard origin/main ) >/dev/null 2>&1 || true\n",
                        encoding="utf-8",
                    )
                    _post_receive.chmod(0o755)
                except Exception:
                    pass
        except Exception:
            # Best-effort — never block /new-ip on a git failure.
            pass
        return created

    def _auto_commit_for_path(path: Path | str, tool: str = "edit") -> None:
        """After a write/replace tool call, auto-commit the change in
        the nearest enclosing per-IP git repo. Silent-best-effort: any
        git error is swallowed so a write to a non-IP location never
        blocks the agent. Walks up from `path` looking for `.git`,
        stops at PROJECT_ROOT — never escapes the project boundary."""
        try:
            import subprocess as _sp_ac
            p = Path(path).resolve()
            project_root = PROJECT_ROOT.resolve()
            try:
                p.relative_to(project_root)
            except ValueError:
                return
            # Walk up looking for .git, but stop before / and don't
            # land on the outer project's git repo (we only want per-IP).
            cur = p if p.is_dir() else p.parent
            git_root: Path | None = None
            while cur != project_root and cur != cur.parent:
                if (cur / ".git").is_dir() and cur != project_root:
                    git_root = cur
                    break
                cur = cur.parent
            if git_root is None:
                return
            rel = ""
            try:
                rel = str(p.relative_to(git_root))
            except ValueError:
                rel = str(p.name)
            _sp_ac.run(
                ["git", "add", "--", "."],
                cwd=str(git_root),
                capture_output=True,
                timeout=15,
            )
            # `git commit` returns non-zero when there's nothing to
            # commit — that's fine, just means the file content didn't
            # change. We don't surface it.
            _sp_ac.run(
                ["git", "commit",
                 "-m", f"{tool}: {rel}",
                 "--allow-empty-message"],
                cwd=str(git_root),
                capture_output=True,
                timeout=15,
            )
        except Exception:
            pass

    def _relative_project_path(path: Path) -> str:
        try:
            return str(path.resolve().relative_to(PROJECT_ROOT.resolve()))
        except Exception:
            return str(path)

    def _strip_import_marker(raw: str) -> str:
        token = str(raw or "").strip().strip("\"'")
        return token[1:] if token.startswith("@") else token

    def _safe_import_path(raw: str) -> Path | None:
        token = _strip_import_marker(raw)
        if not token:
            return None
        if os.name != "nt":
            token = token.replace("\\", "/")
        try:
            p = Path(token).expanduser()
            if not p.is_absolute():
                p = PROJECT_ROOT / token
            resolved = p.resolve()
            resolved.relative_to(PROJECT_ROOT.resolve())
            return resolved
        except Exception:
            return None

    def _resolve_import_path(ip: str, raw: str) -> tuple[Path | None, str]:
        token = _strip_import_marker(raw)
        if not token:
            return None, "empty import path"
        candidates = [token]
        maybe_path = Path(token)
        if not maybe_path.is_absolute():
            norm = token.replace("\\", "/")
            if not norm.startswith(f"{ip}/"):
                candidates.append(f"{ip}/{token}")
        first_safe: Path | None = None
        for candidate in candidates:
            p = _safe_import_path(candidate)
            if p is None:
                continue
            if first_safe is None:
                first_safe = p
            if p.exists():
                return p, ""
        if first_safe is None:
            return None, f"unsafe import path: {token}"
        return first_safe, f"import path not found: {_relative_project_path(first_safe)}"

    def _parse_import_args(args: str) -> tuple[str, list[str], str]:
        import shlex

        try:
            raw_tokens = shlex.split(args or "", posix=False)
        except ValueError as exc:
            return "", [], f"cannot parse /import arguments: {exc}"
        tokens = []
        for raw in raw_tokens:
            tok = _strip_import_marker(raw)
            if tok:
                tokens.append(tok)
        ip = ""
        paths: list[str] = []
        idx = 0
        while idx < len(tokens):
            tok = tokens[idx]
            if tok in ("--ip", "-i"):
                if idx + 1 >= len(tokens):
                    return "", [], "missing value after --ip"
                ip = tokens[idx + 1]
                idx += 2
                continue
            paths.append(tok)
            idx += 1
        if not ip and len(paths) > 1 and _valid_ip_name(paths[0]):
            maybe_ip = paths[0]
            if _ssot_state_path(maybe_ip).is_file() or (PROJECT_ROOT / maybe_ip).exists():
                ip = maybe_ip
                paths = paths[1:]
        if not ip:
            ip = _active_ssot_ip()
        if not _valid_ip_name(ip):
            return "", [], (
                "[SSOT IMPORT] no active IP found\n"
                "usage: /new-ip <ip_name> first, then /import [path ...]\n"
                "or: /import --ip <ip_name> [path ...]"
            )
        return ip, paths, ""

    def _default_import_roots(ip: str) -> list[Path]:
        ip_dir = _ip_root(ip)
        return [ip_dir] if ip_dir.exists() else []

    def _collect_import_files(ip: str, raw_paths: list[str]) -> tuple[list[Path], list[str]]:
        roots: list[Path] = []
        errors: list[str] = []
        if raw_paths:
            for raw in raw_paths:
                p, err = _resolve_import_path(ip, raw)
                if err:
                    errors.append(err)
                if p is not None and p.exists():
                    roots.append(p)
        else:
            roots = _default_import_roots(ip)

        files: list[Path] = []
        seen: set[Path] = set()
        for root in roots:
            candidates = [root]
            if root.is_dir():
                candidates = sorted(root.rglob("*"), key=lambda p: p.as_posix())
            for p in candidates:
                try:
                    rp = p.resolve()
                    rel_parts = rp.relative_to(PROJECT_ROOT.resolve()).parts
                except Exception:
                    continue
                if any(part in _SSOT_IMPORT_SKIP_DIRS or part.startswith(".") for part in rel_parts[:-1]):
                    continue
                if not rp.is_file() or rp.suffix.lower() not in _SSOT_IMPORT_EXTENSIONS:
                    continue
                if rp in seen:
                    continue
                seen.add(rp)
                files.append(rp)
                if len(files) >= 256:
                    errors.append("import file limit reached at 256 files")
                    return files, errors
        return files, errors

    def _clean_import_line(line: str) -> str:
        line = re.sub(r"\s+", " ", str(line or "").strip())
        line = line.lstrip("#/*- ").rstrip("*/ ")
        return line[:260]

    def _snippet_lines(text: str, pattern: str, *, limit: int = 5) -> list[str]:
        out: list[str] = []
        seen: set[str] = set()
        for raw in text.splitlines():
            line = _clean_import_line(raw)
            if len(line) < 4 or line in seen:
                continue
            if re.search(pattern, line, re.IGNORECASE):
                seen.add(line)
                out.append(line)
                if len(out) >= limit:
                    break
        return out

    def _purpose_lines(ip: str, path: Path, text: str) -> list[str]:
        out: list[str] = []
        module_names = re.findall(r"\bmodule\s+([A-Za-z_][A-Za-z0-9_]*)\b", text)
        if module_names:
            out.append("RTL modules: " + ", ".join(module_names[:8]))
        for raw in text.splitlines():
            line = _clean_import_line(raw)
            if len(line) < 12:
                continue
            if re.search(r"\b(purpose|overview|summary|objective|function|module|ip)\b", line, re.IGNORECASE):
                out.append(line)
            elif ip.lower() in line.lower():
                out.append(line)
            elif path.suffix.lower() in {".md", ".txt", ".rst"} and not out:
                out.append(line)
            if len(out) >= 5:
                break
        return out

    def _extract_import_candidates(
        ip: str,
        files: list[Path],
    ) -> tuple[list[dict[str, Any]], dict[str, str], dict[str, list[dict[str, str]]]]:
        patterns = {
            "bus_interface": r"\b(APB|APB4|AXI|AXI4|AXI4[- ]?Lite|AHB|Wishbone|I2C|I3C|SMBus|SPI|UART|PCIe|VDM)\b",
            "register_map": r"\b(register|csr|offset|address|addr|0x[0-9a-f]+|CTRL|STATUS|DATA|CMD|PRESCALE|IRQ|W1C|RO|RW)\b",
            "clock_reset": r"\b(clock|clk|reset|rst|rst_n|resetn|frequency|MHz|active[- ]?(low|high))\b",
            "interrupt": r"\b(interrupt|irq|int_|level|pulse|w1c|done|error)\b",
            "memory_map": r"\b(memory map|base address|base|range|window|SRAM|RAM|FIFO|buffer|address map)\b",
            "parameters": r"\b(parameter|localparam|define|configurable|width|depth|DATA_WIDTH|ADDR_WIDTH|FIFO_DEPTH|default)\b",
            "submodule_structure": r"\b(submodule|hierarchy|block|fsm|module|parser|core|regs|fifo|engine|controller)\b",
            "test_expectation": r"\b(test|verify|verification|coverage|scenario|assert|scoreboard|cocotb|uvm|regression|acceptance)\b",
        }
        snippets: dict[str, list[str]] = {key: [] for key, _ in _SSOT_REQUIRED_DECISIONS}
        sources: dict[str, list[dict[str, str]]] = {key: [] for key, _ in _SSOT_REQUIRED_DECISIONS}
        artifacts: list[dict[str, Any]] = []

        for path in files:
            try:
                raw = path.read_text(encoding="utf-8", errors="ignore")
            except Exception as exc:
                artifacts.append({
                    "path": _relative_project_path(path),
                    "error": f"read failed: {exc}",
                    "bytes": 0,
                })
                continue
            text = raw[:262_144]
            rel = _relative_project_path(path)
            assets = extract_markdown_import_assets(text, rel)
            artifact: dict[str, Any] = {
                "path": rel,
                "bytes": len(raw.encode("utf-8", errors="ignore")),
                "truncated": len(raw) > len(text),
            }
            if assets.image_paths:
                artifact["image_paths"] = list(assets.image_paths)
            if assets.visual_paths:
                artifact["visual_paths"] = list(assets.visual_paths)
            artifacts.append(artifact)
            for line in _purpose_lines(ip, path, text):
                if line not in snippets["purpose"]:
                    snippets["purpose"].append(line)
                    sources["purpose"].append({"path": rel, "excerpt": line})
            for key, pattern in patterns.items():
                for line in _snippet_lines(text, pattern):
                    if line not in snippets[key]:
                        snippets[key].append(line)
                        sources[key].append({"path": rel, "excerpt": line})

        candidates: dict[str, str] = {}
        for key, _label in _SSOT_REQUIRED_DECISIONS:
            vals = snippets.get(key) or []
            if vals:
                candidates[key] = "; ".join(vals[:8])[:1200]
        return artifacts, candidates, sources

    def _merge_unique_records(
        existing: list[Any],
        incoming: list[dict[str, Any]],
        key_fields: tuple[str, ...],
        *,
        limit: int = 128,
    ) -> list[dict[str, Any]]:
        out: list[dict[str, Any]] = [dict(x) for x in existing if isinstance(x, dict)]
        index: dict[tuple[str, ...], int] = {}
        for idx, item in enumerate(out):
            key = tuple(str(item.get(field) or "") for field in key_fields)
            if any(key):
                index[key] = idx
        for item in incoming:
            key = tuple(str(item.get(field) or "") for field in key_fields)
            if any(key) and key in index:
                out[index[key]].update(item)
            else:
                if any(key):
                    index[key] = len(out)
                out.append(dict(item))
        return out[-limit:]

    def _merge_todos_by_id(existing: Any, incoming: list[dict[str, Any]]) -> list[dict[str, Any]]:
        out: list[dict[str, Any]] = [dict(x) for x in existing if isinstance(x, dict)] if isinstance(existing, list) else []
        index = {str(item.get("id") or ""): idx for idx, item in enumerate(out) if str(item.get("id") or "")}
        for item in incoming:
            tid = str(item.get("id") or "").strip()
            if not tid:
                continue
            if tid in index:
                prior = out[index[tid]]
                status = prior.get("status")
                prior.update(item)
                if status:
                    prior["status"] = status
            else:
                index[tid] = len(out)
                out.append(dict(item))
        return out

    def _import_evidence_rows(
        artifacts: list[dict[str, Any]],
        sources: dict[str, list[dict[str, str]]],
    ) -> list[dict[str, Any]]:
        paths = [
            str(item.get("path") or "").strip()
            for item in artifacts
            if isinstance(item, dict) and str(item.get("path") or "").strip()
        ]
        by_path: dict[str, list[dict[str, str]]] = {path: [] for path in paths}
        for key, entries in sources.items():
            for entry in entries or []:
                if not isinstance(entry, dict):
                    continue
                path = str(entry.get("path") or "").strip()
                excerpt = str(entry.get("excerpt") or "").strip()
                if not path or not excerpt:
                    continue
                by_path.setdefault(path, []).append({"decision_key": key, "excerpt": excerpt[:500]})
        rows: list[dict[str, Any]] = []
        for artifact in artifacts:
            if not isinstance(artifact, dict):
                continue
            path = str(artifact.get("path") or "").strip()
            if not path:
                continue
            rows.append({
                "path": path,
                "bytes": int(artifact.get("bytes") or 0),
                "truncated": bool(artifact.get("truncated")),
                "excerpts": by_path.get(path, [])[:24],
            })
        return rows

    def _import_section_todos(
        ip: str,
        candidates: dict[str, str],
        sources: dict[str, list[dict[str, str]]],
    ) -> list[dict[str, Any]]:
        todos: list[dict[str, Any]] = []
        all_refs = sorted({
            str(item.get("path") or "").strip()
            for entries in sources.values()
            for item in (entries or [])
            if isinstance(item, dict) and str(item.get("path") or "").strip()
        })
        for section, title, keys in _SSOT_IMPORT_SECTION_TODO_SPECS:
            evidence_keys = [key for key in keys if str(candidates.get(key) or "").strip()]
            refs = sorted({
                str(item.get("path") or "").strip()
                for key in keys
                for item in sources.get(key, [])
                if isinstance(item, dict) and str(item.get("path") or "").strip()
            })
            evidence_note = ", ".join(evidence_keys) if evidence_keys else "no direct heuristic hit"
            todos.append({
                "id": f"IMPORT_SSOT_SECTION_{section.upper()}",
                "content": f"Review imported workspace evidence for SSOT `{section}`",
                "detail": (
                    f"Section {title}: inspect the imported workspace inventory and promote only "
                    f"source-backed facts into `{section}`. Current heuristic evidence keys: {evidence_note}. "
                    "If the workspace lacks this information, leave a precise SSOT QA/TBD item instead of "
                    "inventing fixed-template content."
                ),
                "criteria": [
                    f"`{section}` contains only source-backed values or explicit TBD/none markers",
                    "Every promoted field has a source_ref back to imported workspace evidence",
                    "Contradictions are listed under custom.atlas_import_conflicts or SSOT QA",
                    "No fixed-template behavior is added without evidence",
                ],
                "source_refs": refs or all_refs[:24] or ["custom.atlas_workspace_inventory"],
                "section": section,
                "decision_keys": list(keys),
                "evidence_keys": evidence_keys,
                "command": f"/to-ssot {ip}",
                "script": "$ATLAS_WORKFLOW_ROOT/ssot-gen/scripts/verify_ssot.py",
                "run_command": (
                    f"python3 \"$ATLAS_WORKFLOW_ROOT/ssot-gen/scripts/repair_ssot_schema.py\" {ip} --root \"$ATLAS_PROJECT_ROOT\" --mode engineering && "
                    f"python3 \"$ATLAS_WORKFLOW_ROOT/ssot-gen/scripts/verify_ssot.py\" {ip} --root \"$ATLAS_PROJECT_ROOT\" --mode engineering"
                ),
                "instructions": [
                    f"Read `{ip}/req/import_manifest.json`, `{ip}/req/extracted_decisions.json`, `{ip}/wiki/import-evidence.md`, and every `source_refs` file before editing this section.",
                    f"Write only canonical `{section}` fields supported by imported evidence, approved Q&A, or explicit no-feature policy.",
                    "If source evidence is silent or contradictory, record a precise SSOT QA card instead of filling template defaults.",
                    "After the YAML write, run repair_ssot_schema.py and verify_ssot.py with --root pointing at the active project root.",
                ],
                "priority": "high" if evidence_keys else "normal",
                "required": True,
            })
        return todos

    def _import_downstream_todos(ip: str, source_refs: list[str], has_dv: bool) -> dict[str, list[dict[str, Any]]]:
        refs = source_refs[:24] or ["custom.atlas_import_doc_evidence"]
        todos: dict[str, list[dict[str, Any]]] = {
            "rtl-gen": [
                {
                    "id": "IMPORT_RTL_FROM_DOC_EVIDENCE",
                    "content": "Implement RTL only from imported doc-backed SSOT facts",
                    "detail": (
                        "Use custom.atlas_import_doc_evidence, custom.atlas_decisions, and the canonical "
                        "SSOT sections derived from them. If a required RTL behavior is only implied or "
                        "contradictory in the docs, emit an SSOT question instead of filling a template."
                    ),
                    "criteria": [
                        "RTL TODO plan references imported source_refs for doc-derived behavior",
                        "No RTL behavior is implemented from a fixed template when the import lacks evidence",
                        "DUT compile/lint evidence is fresh after doc-derived RTL edits",
                    ],
                    "source_refs": refs,
                    "command": f"/ssot-rtl {ip}",
                    "script": "$ATLAS_WORKFLOW_ROOT/rtl-gen/scripts/derive_rtl_todos.py",
                    "run_command": f"python3 \"$ATLAS_WORKFLOW_ROOT/rtl-gen/scripts/derive_rtl_todos.py\" {ip} --root \"$ATLAS_PROJECT_ROOT\" --audit-rtl",
                    "instructions": [
                        "Start from the dynamic RTL TODO plan derived from the SSOT, not from a fixed IP template.",
                        "Implement only behavior that traces to SSOT refs, import evidence, or approved workflow_todos.",
                        "Rerun derive_rtl_todos.py --audit-rtl after RTL edits and keep every required gate open until real source evidence exists.",
                    ],
                    "owner_module": f"{ip}_core",
                    "owner_file": f"rtl/{ip}.sv",
                    "priority": "high",
                    "required": True,
                }
            ],
            "tb-gen": [
                {
                    "id": "IMPORT_TB_FROM_DOC_EVIDENCE",
                    "content": "Generate cocotb/pyuvm tests from imported verification evidence",
                    "detail": (
                        "Convert imported scenarios, acceptance criteria, protocol timing, and coverage notes "
                        "into SSOT test_requirements and executable cocotb/pyuvm tests. Use Python/pyuvm by default; "
                        "do not create SV tc/tb files unless the SSOT explicitly requests that backend."
                    ),
                    "criteria": [
                        "Every imported scenario has a cocotb/pyuvm test or a precise blocker",
                        "Scoreboard expectations trace to function_model or imported source_refs",
                        "Simulation emits results.xml, scoreboard_events.jsonl, VCD, and coverage evidence",
                    ],
                    "source_refs": refs,
                    "command": f"/ssot-tb {ip}",
                    "script": "$ATLAS_WORKFLOW_ROOT/tb-gen/scripts/derive_tb_todos.py",
                    "run_command": f"python3 \"$ATLAS_WORKFLOW_ROOT/tb-gen/scripts/derive_tb_todos.py\" {ip} --root \"$ATLAS_PROJECT_ROOT\" --audit-tb",
                    "instructions": [
                        "Refresh tb/tb_todo_plan.json before generation and rerun derive_tb_todos.py --audit-tb after TB edits.",
                        "Generate tests and scoreboards from SSOT function_model, cycle_model, scenarios, and import-backed workflow_todos.",
                        "Preserve source_refs in generated TB manifest/evidence so failures can be traced back to the SSOT.",
                        "Run the generated cocotb simulation and keep results.xml plus scoreboard_events.jsonl as proof.",
                    ],
                    "priority": "high" if has_dv else "normal",
                    "required": True,
                }
            ],
            "sim_debug": [
                {
                    "id": "IMPORT_SIM_DEBUG_EVIDENCE_MAP",
                    "content": "Use imported doc evidence to classify simulation failures",
                    "detail": (
                        "When cocotb/pyuvm results or waveforms disagree with expected behavior, classify the "
                        "mismatch against imported source_refs, SSOT function/cycle model, RTL, or TB ownership."
                    ),
                    "criteria": [
                        "Every failure report cites expected/got evidence and an imported or SSOT source_ref",
                        "Waveform/VCD checks cover imported timing, reset, interrupt, and protocol expectations when present",
                        "Escalations name the owning workflow: ssot-gen, rtl-gen, tb-gen, or coverage",
                    ],
                    "source_refs": refs,
                    "command": f"/wf sim_debug",
                    "script": "$ATLAS_WORKFLOW_ROOT/sim_debug/scripts/compare_fl_rtl_results.py",
                    "run_command": f"python3 \"$ATLAS_WORKFLOW_ROOT/sim_debug/scripts/compare_fl_rtl_results.py\" {ip} --root \"$ATLAS_PROJECT_ROOT\"",
                    "instructions": [
                        "Compare FL, RTL, scoreboard, waveform, and coverage evidence against imported source_refs and SSOT refs.",
                        "Classify each mismatch to ssot-gen, rtl-gen, tb-gen, or human decision ownership.",
                        "Do not close a mismatch without expected/got evidence and the source_ref that defines the expected behavior.",
                    ],
                    "priority": "normal",
                    "required": True,
                }
            ],
        }
        return todos

    def _apply_import_yaml_todos(
        ip: str,
        doc: dict[str, Any],
        custom: dict[str, Any],
        artifacts: list[dict[str, Any]],
        candidates: dict[str, str],
        sources: dict[str, list[dict[str, str]]],
    ) -> dict[str, Any]:
        evidence_rows = _import_evidence_rows(artifacts, sources)
        source_refs = [
            str(row.get("path") or "").strip()
            for row in evidence_rows
            if str(row.get("path") or "").strip()
        ]
        section_todos = _import_section_todos(ip, candidates, sources)
        downstream = _import_downstream_todos(ip, source_refs, bool(candidates.get("test_expectation")))

        custom["atlas_import_doc_evidence"] = _merge_unique_records(
            custom.get("atlas_import_doc_evidence") if isinstance(custom.get("atlas_import_doc_evidence"), list) else [],
            evidence_rows,
            ("path",),
            limit=256,
        )
        custom["atlas_workspace_inventory"] = _merge_unique_records(
            custom.get("atlas_workspace_inventory") if isinstance(custom.get("atlas_workspace_inventory"), list) else [],
            evidence_rows,
            ("path",),
            limit=256,
        )
        prior_draft = custom.get("atlas_import_todo_draft") if isinstance(custom.get("atlas_import_todo_draft"), dict) else {}
        merged_section_todos = _merge_todos_by_id(
            prior_draft.get("section_todos") if isinstance(prior_draft, dict) else [],
            section_todos,
        )
        custom["atlas_import_section_todos"] = merged_section_todos
        custom["atlas_import_todo_draft"] = {
            "updated_at": time.time(),
            "source_refs": source_refs[:64],
            "section_todos": merged_section_todos,
            "downstream_todos": {
                stage: _merge_todos_by_id(
                    (prior_draft.get("downstream_todos") or {}).get(stage) if isinstance(prior_draft.get("downstream_todos"), dict) else [],
                    items,
                )
                for stage, items in downstream.items()
            },
        }

        workflow_todos = doc.get("workflow_todos")
        if not isinstance(workflow_todos, dict):
            workflow_todos = {}
            doc["workflow_todos"] = workflow_todos
        workflow_todos["ssot-gen"] = _merge_todos_by_id(workflow_todos.get("ssot-gen"), section_todos)
        for stage, items in downstream.items():
            workflow_todos[stage] = _merge_todos_by_id(workflow_todos.get(stage), items)

        return {
            "evidence_rows": len(evidence_rows),
            "section_todos": len(section_todos),
            "downstream_todos": {stage: len(items) for stage, items in downstream.items()},
        }

    def _merge_import_candidates(
        ip: str,
        kind: str,
        state: dict[str, Any],
        artifacts: list[dict[str, Any]],
        candidates: dict[str, str],
        sources: dict[str, list[dict[str, str]]],
    ) -> tuple[list[str], list[dict[str, Any]]]:
        filled, conflicts = _record_ssot_decisions(ip, kind, candidates, sources)
        doc, custom = _ssot_custom(ip, kind)
        todo_summary = _apply_import_yaml_todos(ip, doc, custom, artifacts, candidates, sources)
        next_action = "/grill-me" if conflicts or _missing_ssot_decisions(ip, state) else "/to-ssot"
        manifest_path = PROJECT_ROOT / ip / "req" / "import_manifest.json"
        extracted_path = PROJECT_ROOT / ip / "req" / "extracted_decisions.json"
        try:
            manifest_path.parent.mkdir(parents=True, exist_ok=True)
            manifest_doc = {
                "schema_version": "ssot_import_manifest.v1",
                "ip": ip,
                "workflow": "ssot-gen",
                "updated_at": time.time(),
                "kind": kind,
                "artifacts": artifacts,
                "candidate_facts": candidates,
                "sources": sources,
                "filled_decisions": filled,
                "conflicts": conflicts,
                "workflow_todo_summary": todo_summary,
                "next": next_action,
            }
            manifest_path.write_text(json.dumps(manifest_doc, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
            extracted_path.write_text(
                json.dumps(
                    {
                        "schema_version": "ssot_extracted_decisions.v1",
                        "ip": ip,
                        "workflow": "ssot-gen",
                        "updated_at": manifest_doc["updated_at"],
                        "decisions": candidates,
                        "sources": sources,
                        "filled_decisions": filled,
                        "conflicts": conflicts,
                        "next": next_action,
                    },
                    ensure_ascii=False,
                    indent=2,
                )
                + "\n",
                encoding="utf-8",
            )
        except Exception:
            pass
        imports = custom.get("atlas_imports")
        if not isinstance(imports, list):
            imports = []
        imports.append({
            "imported_at": time.time(),
            "artifacts": artifacts,
            "filled": filled,
            "conflicts": conflicts,
            "yaml_todos": todo_summary,
        })
        custom["atlas_imports"] = imports
        _save_ssot_draft(ip, doc)
        imported_artifacts = state.get("imported_artifacts")
        if not isinstance(imported_artifacts, list):
            imported_artifacts = []
        imported_artifacts.extend({
            "path": str(a.get("path") or ""),
            "imported_at": time.time(),
        } for a in artifacts if a.get("path"))
        state["imported_artifacts"] = imported_artifacts[-64:]
        state["last_import_yaml_todos"] = todo_summary
        state["last_step"] = "import"
        if conflicts:
            state["last_issue"] = "import_conflicts"
        if conflicts:
            state["approved"] = False
            state["approved_at"] = 0
        state["status"] = "answered" if not _missing_ssot_decisions(ip, state) else "planned"
        _write_import_wiki_index(ip, kind, artifacts, candidates, sources, filled, conflicts, todo_summary, next_action)
        return filled, conflicts

    def _import_defaults_if_available(ip: str, kind: str, state: dict[str, Any]) -> tuple[list[str], list[dict[str, Any]], list[dict[str, Any]], list[str]]:
        """Import the current IP workspace into the SSOT draft when requested."""
        files, errors = _collect_import_files(ip, [])
        if not files:
            return [], [], [], errors
        artifacts, candidates, sources = _extract_import_candidates(ip, files)
        filled, conflicts = _merge_import_candidates(ip, kind, state, artifacts, candidates, sources)
        state.setdefault("ip", ip)
        state.setdefault("kind", kind)
        state["active_session"] = _ssot_session_for_ip(ip)
        _save_ssot_state(ip, state)
        return filled, conflicts, artifacts, errors

    def _auto_approve_if_complete(ip: str, state: dict[str, Any], *, reason: str) -> bool:
        if state.get("approved"):
            return False
        if _missing_ssot_decisions(ip, state):
            return False
        doc = _load_ssot_draft(ip)
        custom = doc.get("custom") if isinstance(doc.get("custom"), dict) else {}
        conflicts = custom.get("atlas_import_conflicts") if isinstance(custom, dict) else []
        if conflicts:
            return False
        state["approved"] = True
        state["approved_at"] = time.time()
        state["status"] = "approved"
        state["last_step"] = reason
        _save_ssot_state(ip, state)
        return True

    def _rtl_blocker_path(ip: str) -> Path:
        return PROJECT_ROOT / ip / "rtl" / "rtl_blocked.json"

    def _load_rtl_blocker(ip: str) -> dict[str, Any]:
        path = _rtl_blocker_path(ip)
        if not path.is_file():
            return {}
        try:
            doc = json.loads(path.read_text(encoding="utf-8"))
            return doc if isinstance(doc, dict) else {}
        except Exception as exc:
            return {"reason": f"rtl_blocked.json parse failed: {exc}", "questions": []}

    def _rtl_module_contract_placeholder(q: dict[str, Any]) -> str:
        missing = q.get("missing_modules") if isinstance(q.get("missing_modules"), list) else []
        if not missing and isinstance(q.get("candidate_modules"), list):
            missing = q.get("candidate_modules") or []
        available = q.get("available_refs") if isinstance(q.get("available_refs"), dict) else {}
        rows: list[str] = []
        orphan_refs = q.get("orphan_refs") if isinstance(q.get("orphan_refs"), list) else []
        if orphan_refs:
            rows.append("# orphan refs needing an RTL owner: " + ", ".join(str(v) for v in orphan_refs[:16]))
        if available:
            for key in ("source_sections", "function_model_refs", "decomposition_refs", "cycle_model_refs", "feature_refs", "dataflow_refs", "register_refs", "fsm_refs", "test_refs", "ports"):
                vals = available.get(key) if isinstance(available.get(key), list) else []
                if vals:
                    rows.append(f"# available {key}: " + ", ".join(str(v) for v in vals[:10]))
        rows.append("module_contracts:")
        for item in missing[:8]:
            if not isinstance(item, dict):
                continue
            name = str(item.get("name") or "").strip()
            file = str(item.get("file") or "").strip()
            rows += [
                f"  - name: {name}",
                f"    file: {file}",
                "    implements:",
                "      - <specific behavior this module owns>",
                "    source_sections: [<ssot section names>]",
                "    function_model_refs: [<function_model paths>]",
                "    decomposition_refs: [<decomposition paths>]",
                "    cycle_model_refs: [<cycle_model paths>]",
                "    feature_refs: [<feature names or paths>]",
                "    dataflow_refs: [<dataflow paths>]",
                "    register_refs: [<register names or paths>]",
                "    fsm_refs: [<fsm paths>]",
                "    ports: [<owned ports or internal interface ports>]",
                "    connections: {<local_port>: <ssot/interface signal>}",
            ]
        return "\n".join(rows)

    _RTL_OWNERSHIP_BLOCKER_IDS = {
        "RTL_DYNAMIC_TODO_OWNERSHIP",
        "RTL_MODULE_CONTRACTS",
        "RTL_MODULE_BEHAVIOR_MATCH",
        "SSOT_BEHAVIOR_OWNERSHIP",
    }
    _RTL_CONNECTION_BLOCKER_IDS = {
        "RTL_RESOLVE_CONNECTION_CONTRACTS",
        "RTL_CONNECTION_CONTRACTS",
        "RTL_MANIFEST_CONNECTION_CONTRACTS",
    }
    _RTL_IMPL_BLOCKER_IDS = {
        "RTL_TODO_PLAN_MISSING",
        "DETERMINISTIC_RTL_ARTIFACT_NOT_APPROVED",
        "LLM_RTL_IMPLEMENTATION_REQUIRED",
        "COMMON_AI_AGENT_RTL_PROVENANCE_REQUIRED",
    }

    def _rtl_blocker_qa_section(qid: str, raw: dict[str, Any]) -> tuple[str, str, str]:
        text = " ".join(
            [
                qid,
                str(raw.get("decision_needed") or ""),
                str(raw.get("evidence") or ""),
                " ".join(str(ref) for ref in raw.get("source_refs") or [] if isinstance(raw.get("source_refs"), list)),
                " ".join(str(field) for field in raw.get("required_fields") or [] if isinstance(raw.get("required_fields"), list)),
            ]
        ).lower()
        if qid == "RTL_TARGET_SCALE_POLICY" or "target_scale" in text or "target scale" in text:
            return "19_workflow_todos", "19. Workflow / Human Gates", "quality_gates.rtl_gen.target_scale"
        if qid in _RTL_CONNECTION_BLOCKER_IDS or "connection_contract" in text or "integration.connections" in text:
            return "17_integration", "17. Integration / Connection Contracts", "integration.connections"
        if qid in _RTL_OWNERSHIP_BLOCKER_IDS or "sub_modules" in text or "ownership" in text:
            return "04_architecture", "4. Architecture / Decomposition", "sub_modules"
        if "interface" in text or "port" in text or "clock" in text or "reset" in text:
            return "03_interface", "3. Interface", "io_list"
        if "coverage" in text or "test_requirements" in text or "verification" in text:
            return "18_verification", "18. Verification / Gates", "test_requirements"
        return "19_workflow_todos", "19. Workflow / Human Gates", "workflow_todos.rtl-gen"

    def _rtl_blocker_cards(blocker: dict[str, Any]) -> list[dict[str, Any]]:
        cards: list[dict[str, Any]] = []
        for q in blocker.get("questions") if isinstance(blocker.get("questions"), list) else []:
            if not isinstance(q, dict):
                continue
            qid = str(q.get("id") or "").strip() or "RTL_BLOCKER"
            if qid in _RTL_IMPL_BLOCKER_IDS:
                continue
            raw_options = q.get("options") if isinstance(q.get("options"), list) else []
            options = [
                {
                    "id": f"{qid}_opt{idx}",
                    "label": str(opt)[:96],
                    "detail": str(opt),
                }
                for idx, opt in enumerate(raw_options, start=1)
                if str(opt).strip()
            ]
            subtitle_parts = [qid]
            if q.get("evidence"):
                subtitle_parts.append(str(q.get("evidence"))[:220])
            if q.get("recommended_default"):
                subtitle_parts.append("Recommended: " + str(q.get("recommended_default"))[:220])
            card = {
                "id": qid,
                "question": str(q.get("decision_needed") or qid),
                "kind": "single" if options else "input",
                "subtitle": " · ".join(subtitle_parts),
                "options": options,
                "blocker": q,
            }
            if qid in _RTL_OWNERSHIP_BLOCKER_IDS:
                card["kind"] = "input"
                card["multiline"] = True
                card["subtitle"] = (
                    str(card.get("subtitle") or "")
                    + " · Paste YAML/JSON with module_contracts; option clicks alone do not approve RTL ownership."
                )[:900]
                card["placeholder"] = _rtl_module_contract_placeholder(q)
            cards.append(card)
        return cards

    def _rtl_blocker_flow_id(blocker: dict[str, Any]) -> str:
        payload = json.dumps(blocker.get("questions") or blocker, ensure_ascii=False, sort_keys=True, default=str)
        digest = hashlib.sha256(payload.encode("utf-8")).hexdigest()[:12]
        return f"rtl_blocker_{digest}"

    def _rtl_blocker_qa_questions(
        ip: str,
        blocker: dict[str, Any],
        cards: list[dict[str, Any]],
        *,
        reason: str,
    ) -> list[dict[str, Any]]:
        source_path = str(_rtl_blocker_path(ip).relative_to(PROJECT_ROOT))
        questions: list[dict[str, Any]] = []
        for card in cards:
            qid = str(card.get("id") or "RTL_BLOCKER").strip() or "RTL_BLOCKER"
            raw = card.get("blocker") if isinstance(card.get("blocker"), dict) else {}
            orphan_refs = raw.get("orphan_refs") if isinstance(raw.get("orphan_refs"), list) else []
            required_fields = raw.get("required_fields") if isinstance(raw.get("required_fields"), list) else []
            candidate_modules = raw.get("candidate_modules") if isinstance(raw.get("candidate_modules"), list) else []
            section_id, section_title, field_path = _rtl_blocker_qa_section(qid, raw)
            criteria = [
                "Resolve the blocker by updating SSOT-owned authority artifacts, not by hand-editing generated RTL.",
                f"Rerun rtl-gen preflight until `{qid}` no longer appears in `{source_path}`.",
            ]
            if qid in _RTL_OWNERSHIP_BLOCKER_IDS:
                criteria.append("Every orphan source_ref must be covered by an exact or dotted-parent ref in one RTL module contract.")
            elif qid == "RTL_TARGET_SCALE_POLICY":
                criteria.append("Lock positive quality_gates.rtl_gen.target_scale minima or approve target_scale_waiver with owner and reason.")
            elif qid in _RTL_CONNECTION_BLOCKER_IDS:
                criteria.append("Answer with machine-readable module/port/signal connection contracts before top integration signoff.")
            source_refs = [source_path]
            source_refs.extend(str(ref) for ref in orphan_refs[:64])
            if raw.get("evidence"):
                source_refs.append(str(raw.get("evidence")))
            questions.append({
                "id": qid,
                "decision_key": qid,
                "decision_label": str(card.get("question") or raw.get("decision_needed") or qid),
                "question": str(card.get("question") or raw.get("decision_needed") or qid),
                "kind": str(card.get("kind") or "input"),
                "subtitle": str(card.get("subtitle") or ""),
                "options": card.get("options") or [],
                "qa_type": "rtl_blocker",
                "source": reason,
                "source_refs": source_refs,
                "field_path": field_path,
                "section_id": section_id,
                "section_title": section_title,
                "content": f"Resolve rtl-gen blocker `{qid}` for `{ip}`.",
                "detail": (
                    f"Evidence: {raw.get('evidence') or blocker.get('reason') or source_path}. "
                    f"Required fields: {', '.join(str(v) for v in required_fields[:12]) or 'SSOT module contract fields'}. "
                    f"Orphan refs: {len(orphan_refs)}. Candidate modules: {len(candidate_modules)}."
                ),
                "criteria": criteria,
                "placeholder": card.get("placeholder") or "",
                "multiline": bool(card.get("multiline")),
            })
        return questions

    def _start_rtl_blocker_qna(
        ip: str,
        *,
        reason: str = "rtl-gen preflight",
        interactive: bool = True,
        client_session: Any | None = None,
    ) -> bool:
        blocker = _load_rtl_blocker(ip)
        cards = _rtl_blocker_cards(blocker)
        if not blocker or not cards:
            _emit_workflow_result(
                f"[RTL BLOCKER Q&A] no rtl_blocked.json questions found for {ip}",
                "rtl-blocker",
            )
            return True

        # Surface RTL blockers as a single grouped agent-turn message and
        # leave rtl_blocked.json on disk as evidence. The user answers
        # them inline and the SSOT-gen workflow picks them up — there is
        # no separate slash command, no env-var gate, and no auto-skip.
        question_lines: list[str] = []
        for idx, card in enumerate(cards, start=1):
            title = str(card.get("question") or card.get("label") or card.get("title") or "(no question)").strip()
            yaml_path = str(card.get("yaml_path") or card.get("source_ref") or "").strip()
            tag = f"  [{idx}/{len(cards)}]"
            if yaml_path:
                question_lines.append(f"{tag} {title}  ← {yaml_path}")
            else:
                question_lines.append(f"{tag} {title}")
        msg = (
            f"[RTL BLOCKER] {len(cards)} SSOT decision(s) needed for {ip} "
            f"(see {_rtl_blocker_path(ip).relative_to(PROJECT_ROOT)}):\n"
            + "\n".join(question_lines[:20])
            + ("\n  ... (additional blockers truncated)" if len(cards) > 20 else "")
            + "\nAnswer inline so the SSOT-gen workflow can incorporate them."
        )
        _append_session_message(_canonical_session_string(ip), "assistant", msg)
        _append_workflow_history("ssot-gen", "assistant", msg)
        _emit_workflow_result(msg, "rtl-blocker")
        return True

    def _sim_human_gate_cards(ip: str, classify_doc: dict[str, Any]) -> list[dict[str, Any]]:
        cards: list[dict[str, Any]] = []
        classifications = classify_doc.get("classifications")
        for item in classifications if isinstance(classifications, list) else []:
            if not isinstance(item, dict) or item.get("llm_loop_allowed") is not False:
                continue
            goal_id = str(item.get("goal_id") or "EQ_HUMAN_GATE").strip() or "EQ_HUMAN_GATE"
            human_question = str(item.get("human_question") or "").strip()
            evidence = item.get("evidence") if isinstance(item.get("evidence"), dict) else {}
            ssot_refs = evidence.get("ssot_refs") if isinstance(evidence.get("ssot_refs"), list) else []
            fl_expected = evidence.get("fl_expected") if isinstance(evidence.get("fl_expected"), dict) else {}
            rtl_observed = evidence.get("rtl_observed") if isinstance(evidence.get("rtl_observed"), dict) else {}
            evidence_bits = [
                f"class={item.get('classification') or 'unknown'}",
                f"owner={item.get('owner') or 'human'}",
            ]
            if ssot_refs:
                evidence_bits.append("SSOT " + ", ".join(str(x) for x in ssot_refs[:3]))
            if fl_expected:
                evidence_bits.append("FL " + json.dumps(fl_expected, sort_keys=True)[:180])
            if rtl_observed:
                evidence_bits.append("RTL " + json.dumps(rtl_observed, sort_keys=True)[:180])
            cards.append({
                "id": goal_id,
                "question": f"Decision needed for {goal_id}: define expected behavior or approve waiver",
                "kind": "single",
                "subtitle": " · ".join(evidence_bits)[:500],
                "options": [
                    {
                        "id": f"{goal_id}_update_ssot",
                        "label": "Update SSOT",
                        "detail": "Record the intended behavior in SSOT/requirements, then regenerate FL, equivalence goals, TB, and coverage.",
                    },
                    {
                        "id": f"{goal_id}_waiver",
                        "label": "Waiver",
                        "detail": "Record an explicit waiver/rationale; signoff remains human-owned until the waiver is reviewed.",
                    },
                ],
                "human_question": human_question,
                "classification": item,
            })
        return cards

    def _persist_sim_human_gate_answers(
        ip: str,
        classify_doc: dict[str, Any],
        cards: list[dict[str, Any]],
        answer_entries: list[dict[str, Any]],
    ) -> Path:
        sim_dir = PROJECT_ROOT / ip / "sim"
        sim_dir.mkdir(parents=True, exist_ok=True)
        out_path = sim_dir / "human_gate_answers.json"
        doc = {
            "schema_version": 1,
            "type": "fl_rtl_human_gate_answers",
            "ip": ip,
            "captured_at": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime()),
            "source": str((PROJECT_ROOT / ip / "sim" / "mismatch_classification.json").relative_to(PROJECT_ROOT)),
            "answers": answer_entries,
        }
        out_path.write_text(json.dumps(doc, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")

        state = _load_ssot_state(ip)
        state.setdefault("ip", ip)
        state.setdefault("kind", "equivalence human gate")
        state["status"] = "equivalence_human_gate_answered"
        state["equivalence_human_gate_source"] = doc["source"]
        state["equivalence_human_gate_answers"] = answer_entries
        state["equivalence_human_gate_classifications"] = [
            card.get("classification") for card in cards if isinstance(card.get("classification"), dict)
        ]
        _save_ssot_state(ip, state)

        return out_path

    def _start_sim_human_gate_qna(ip: str, classify_doc: dict[str, Any], *, reason: str = "sim-debug", client_session: Any | None = None) -> bool:
        cards = _sim_human_gate_cards(ip, classify_doc)
        if not cards:
            return False

        def _worker() -> None:
            import uuid as _uuid

            flow_id = "simq_" + _uuid.uuid4().hex[:10]
            bridge.open_question(flow_id)
            bridge.emit("ask_user", flow_id=flow_id, questions=cards)
            bridge.emit("agent_state", running=True)
            try:
                ans = bridge.wait_answer(flow_id, timeout=900)
            finally:
                bridge.close_question(flow_id)
            if not isinstance(ans, dict) or not isinstance(ans.get("answers"), list):
                msg = (
                    f"[SIM HUMAN GATE] {ip}: no answer received; mismatch remains human-gated.\n"
                    f"source: {ip}/sim/mismatch_classification.json"
                )
                _append_session_message(f"{ip}/sim_debug", "assistant", msg)
                _append_workflow_history("sim_debug", "assistant", msg)
                _append_active_history("assistant", "```\n" + msg + "\n```")
                _emit_workflow_result(msg, "sim-debug")
                return

            answer_entries: list[dict[str, Any]] = []
            for card, raw_answer in zip(cards, ans.get("answers") or []):
                qa = raw_answer if isinstance(raw_answer, dict) else {}
                classification = card.get("classification") if isinstance(card.get("classification"), dict) else {}
                answer_entries.append({
                    "goal_id": card.get("id"),
                    "decision_needed": card.get("question"),
                    "answer": _answer_text(qa, card),
                    "selected": qa.get("selected") or [],
                    "custom": str(qa.get("custom") or "").strip(),
                    "source": reason,
                    "classification": classification.get("classification"),
                    "owner": classification.get("owner"),
                    "evidence": classification.get("evidence") if isinstance(classification.get("evidence"), dict) else {},
                    "human_question": card.get("human_question") or "",
                })
            try:
                out_path = _persist_sim_human_gate_answers(ip, classify_doc, cards, answer_entries)
                msg = (
                    f"[SIM HUMAN GATE] captured {len(answer_entries)} answer(s) for {ip}\n"
                    f"answers: {out_path.relative_to(PROJECT_ROOT)}\n"
                    "next: rerun SSOT/FL/equivalence generation if behavior changed, or keep signoff human-owned for waivers"
                )
            except Exception as exc:
                msg = f"[SIM HUMAN GATE] {ip}: failed to persist answers: {exc}"
            _append_session_message(f"{ip}/sim_debug", "assistant", msg)
            _append_workflow_history("sim_debug", "assistant", msg)
            _append_active_history("assistant", "```\n" + msg + "\n```")
            _emit_workflow_result(msg, "sim-debug")

        ctx = contextvars.copy_context()
        threading.Thread(target=ctx.run, args=(_worker,), daemon=True).start()
        return True

    app.state.atlas_bridge = bridge
    app.state.start_rtl_blocker_qna = _start_rtl_blocker_qna
    app.state.active_ssot_qa_context = _active_ssot_qa_context
    app.state.valid_ip_name = _valid_ip_name
    app.state.ssot_q_pairs_from_questions = _ssot_q_pairs_from_questions
    app.state.load_ssot_state = _load_ssot_state
    app.state.upsert_ssot_qa_items = _upsert_ssot_qa_items
    app.state.status_group = _status_group

    def _emit_workflow_result(text: str, tool: str = "workflow") -> None:
        body = (text or "").strip() or "(no output)"
        payload = "```\n" + body + "\n```"
        bridge.emit("tool_result", text=payload, tool=tool, truncated=False)
        bridge.emit("slash_output", text=payload)
        bridge.emit("flush")
        bridge.emit("commands_changed")
        bridge.emit("agent_state", running=False)


    def _safe_import_upload_name(name: str) -> str:
        raw = Path(str(name or "import.txt")).name
        safe = re.sub(r"[^A-Za-z0-9._-]+", "_", raw).strip("._")
        if not safe:
            safe = "import.txt"
        return safe[:120]

    def _markitdown_convert(src_path: Path) -> tuple[str, str]:
        """Run markitdown on src_path under an external Python ≥3.10.

        Returns (md_text, error). An empty error string means success.

        Probe order:
          1. ATLAS_MARKITDOWN_PYTHON env (explicit override).
          2. Windows: `python` / `python.exe`, then `python3.12`,
             then `py -3.12`.
          3. macOS/Linux: `python3.12`, `python3.11`, `python3.10` (PATH),
             then `/opt/homebrew/bin/python3.10` as a last resort.
        """
        import shutil as _shutil
        import subprocess as _subprocess

        is_windows = sys.platform.startswith("win")

        def _candidates() -> list[list[str]]:
            cands: list[list[str]] = []
            if _SSOT_MARKITDOWN_PY:
                cands.append([_SSOT_MARKITDOWN_PY])
            if is_windows:
                # Plain `python` (= python.exe on PATH) is the user's
                # preferred entry point on Windows; only fall back to
                # the py launcher / versioned aliases when it isn't
                # installed.
                for name in ("python", "python.exe", "python3.12"):
                    found = _shutil.which(name)
                    if found:
                        cands.append([found])
                py_launcher = _shutil.which("py")
                if py_launcher:
                    cands.append([py_launcher, "-3.12"])
            else:
                for name in ("python3.12", "python3.11", "python3.10"):
                    found = _shutil.which(name)
                    if found:
                        cands.append([found])
                brew = "/opt/homebrew/bin/python3.10"
                if Path(brew).exists():
                    cands.append([brew])
            return cands

        candidates = _candidates()
        if not candidates:
            return "", (
                "markitdown needs Python 3.10+: no candidate interpreter found. "
                "Set ATLAS_MARKITDOWN_PYTHON or install Python "
                + ("(`python -m pip install markitdown`)" if is_windows else "(python3.12 with markitdown)")
            )

        last_err = ""
        for cmd_prefix in candidates:
            try:
                # text=True without an explicit encoding uses the OS
                # locale codec — Korean Windows defaults to cp949, which
                # blows up on the 0xf0 leading bytes of any 4-byte UTF-8
                # sequence (emoji, CJK extension B+, etc.) that
                # markitdown emits. Pin utf-8 explicitly.
                result = _subprocess.run(
                    [*cmd_prefix, "-m", "markitdown", str(src_path)],
                    capture_output=True, text=True,
                    encoding="utf-8", errors="replace",
                    timeout=60,
                )
            except _subprocess.TimeoutExpired:
                last_err = "markitdown timed out (60s)"
                continue
            except FileNotFoundError as exc:
                last_err = f"{' '.join(cmd_prefix)} not found: {exc}"
                continue
            except Exception as exc:
                last_err = f"markitdown error via {' '.join(cmd_prefix)}: {exc}"
                continue
            if result.returncode == 0:
                return result.stdout, ""
            last_err = (
                f"markitdown failed via {' '.join(cmd_prefix)} "
                f"(rc={result.returncode}): {result.stderr[:300]}"
            )
        return "", last_err or "markitdown not runnable"

    _DATA_IMAGE_BASE64_RE = re.compile(
        r"data:(image/[A-Za-z0-9.+-]+)\s*[:;]\s*base64\s*,",
        re.IGNORECASE,
    )

    def _normalize_markdown_data_image_uris(text: str) -> str:
        """Normalize malformed doc-converter image data URIs in markdown."""
        return _DATA_IMAGE_BASE64_RE.sub(r"data:\1;base64,", str(text or ""))

    def _describe_image(img: Path) -> str:
        """Describe an image for SSOT import evidence.

        Uses core.tools.read_image, which sends the image through the active
        multimodal LLM client. Import should fail visibly when image reading
        is not configured, so development does not hide missing vision support.
        """
        prompt_text = (
            "Describe this image for SSOT import evidence in 2-3 sentences. "
            "Include visible text, diagrams, charts, key data, interfaces, "
            "signals, states, registers, or requirements if present."
        )

        from core.tools import read_image  # type: ignore

        desc = str(read_image(path=str(img), prompt=prompt_text) or "").strip()
        if not desc:
            raise RuntimeError(f"read_image returned an empty description for {img}")
        return desc

    def _import_visual_evidence_enabled() -> bool:
        raw = str(os.environ.get("ATLAS_IMPORT_VISUAL_EVIDENCE", "true") or "true").strip().lower()
        return raw not in {"0", "false", "no", "off"}

    def _import_visual_evidence_max_pages() -> int:
        raw = str(os.environ.get("ATLAS_IMPORT_VISUAL_EVIDENCE_MAX_PAGES", "100") or "100").strip()
        try:
            value = int(raw)
        except ValueError:
            return 100
        return max(0, min(value, 100))

    def _import_visual_evidence_scale() -> float:
        raw = str(os.environ.get("ATLAS_IMPORT_VISUAL_EVIDENCE_SCALE", "2.0") or "2.0").strip()
        try:
            value = float(raw)
        except ValueError:
            return 2.0
        return max(1.0, min(value, 4.0))

    def _describe_visual_surface(surface: VisualSurface) -> str:
        from core.tools import read_image  # type: ignore

        desc = str(read_image(path=str(surface.image_path), prompt=visual_evidence_prompt(surface)) or "").strip()
        if not desc:
            raise RuntimeError(f"read_image returned an empty visual evidence description for {surface.image_path}")
        return desc

    def _convert_upload_to_markdown(
        original_path: Path,
        suffix: str,
        dest_dir: Path,
        images_dir: Path,
        visual_dir: Path,
        stamp: int,
        idx: int,
        basename: str,
    ) -> tuple[Optional[Path], list[Path], list[Path], str]:
        """Convert an uploaded document to Markdown plus extracted images.

        Strategy: use markitdown for document-to-markdown conversion, and
        return its error directly when it is unavailable. Image extraction
        still runs the per-format path (PyMuPDF / python-pptx / python-docx)
        since markitdown does not emit clean image files. Image descriptions
        are appended at the bottom of the Markdown under '## Extracted Images'.
        PDF figure pages are rendered and appended under '## Visual Evidence'
        so vector drawings and page-layout diagrams are visible to image read.

        Returns (md_path, image_paths, visual_paths, error). md_path is None
        when conversion fails; the caller keeps the saved original regardless.
        """
        md_target = dest_dir / f"{stamp}_{idx}_{basename}.md"
        image_paths: list[Path] = []
        visual_paths: list[Path] = []
        visual_surfaces: list[VisualSurface] = []
        image_seq = 0

        def _low_information_import_image_reason(blob: bytes, ext: str) -> str:
            normalized_ext = (ext or "png").lower().lstrip(".")
            if normalized_ext in {"svg", "svg+xml"}:
                return ""
            if len(blob or b"") > 2048:
                return ""
            try:
                from PIL import Image, ImageStat  # type: ignore

                img = Image.open(io.BytesIO(blob)).convert("RGBA")
                width, height = img.size
                if width <= 1 or height <= 1:
                    return "degenerate extracted image"
                if width * height > 25000:
                    return ""
                thumb = img.copy()
                thumb.thumbnail((32, 32))
                colors = thumb.getcolors(maxcolors=4096) or []
                extrema = ImageStat.Stat(thumb.convert("RGB")).extrema
                spread = max((hi - lo) for lo, hi in extrema) if extrema else 255
                if len(colors) <= 2:
                    return "flat tiny extracted image"
                if len(colors) <= 4 and spread <= 8:
                    return "near-flat tiny extracted image"
            except Exception:
                return ""
            return ""

        def _write_image(
            blob: bytes,
            ext: str,
            n: int | None = None,
            *,
            filter_noise: bool = False,
        ) -> Optional[Path]:
            nonlocal image_seq
            if n is None:
                image_seq += 1
                n = image_seq
            else:
                image_seq = max(image_seq, int(n or 0))
            ext = (ext or "png").lower().lstrip(".") or "png"
            if ext in {"jpeg", "jpg"}:
                ext = "jpg"
            elif ext in {"svg+xml", "svg"}:
                ext = "svg"
            elif not re.match(r"^[A-Za-z0-9]+$", ext):
                ext = "png"
            if filter_noise and _low_information_import_image_reason(blob, ext):
                return None
            img_path = images_dir / f"{stamp}_{idx}_{n}.{ext}"
            img_path.write_bytes(blob)
            image_paths.append(img_path)
            return img_path

        def _inline_image_rel(path: Path) -> str:
            try:
                return path.relative_to(dest_dir).as_posix()
            except ValueError:
                try:
                    return path.relative_to(PROJECT_ROOT).as_posix()
                except ValueError:
                    return path.as_posix()

        def _materialize_data_image_uris(text: str) -> str:
            """Replace inline base64 image data URIs with saved image files."""
            import base64 as _base64

            def _save_data_image(mime: str, encoded: str) -> str:
                payload = re.sub(r"\s+", "", encoded or "")
                if not payload:
                    return ""
                try:
                    blob = _base64.b64decode(payload, validate=True)
                except Exception:
                    return ""
                ext = str(mime or "image/png").split("/", 1)[-1]
                saved = _write_image(blob, ext)
                return _inline_image_rel(saved) if saved else ""

            md_re = re.compile(
                r"!\[([^\]]*)\]\(\s*data:(image/[A-Za-z0-9.+-]+)\s*[:;]\s*base64\s*,"
                r"([A-Za-z0-9+/=_\-\s\r\n]+)\s*(?:[\"'][^)]*[\"'])?\)",
                re.IGNORECASE,
            )
            html_re = re.compile(
                r"(\bsrc\s*=\s*[\"'])data:(image/[A-Za-z0-9.+-]+)\s*[:;]\s*base64\s*,"
                r"([^\"']+)([\"'])",
                re.IGNORECASE,
            )

            def _replace_md(match: re.Match[str]) -> str:
                rel = _save_data_image(match.group(2), match.group(3))
                if not rel:
                    return match.group(0)
                return f"![{match.group(1)}]({rel})"

            def _replace_html(match: re.Match[str]) -> str:
                rel = _save_data_image(match.group(2), match.group(3))
                if not rel:
                    return match.group(0)
                return f"{match.group(1)}{rel}{match.group(4)}"

            return html_re.sub(_replace_html, md_re.sub(_replace_md, str(text or "")))

        def _write_markdown(text: str) -> None:
            md_target.write_text(
                _normalize_markdown_data_image_uris(_materialize_data_image_uris(text)),
                encoding="utf-8",
            )

        def _add_office_visual_surfaces() -> None:
            if not _import_visual_evidence_enabled():
                return
            surfaces = render_office_visual_surfaces(
                original_path,
                visual_dir,
                stamp,
                idx,
                max_pages=_import_visual_evidence_max_pages(),
                scale=_import_visual_evidence_scale(),
            )
            visual_surfaces.extend(surfaces)
            visual_paths.extend(surface.image_path for surface in surfaces)

        try:
            if suffix in _SSOT_IMPORT_PASSTHROUGH:
                text = original_path.read_bytes().decode("utf-8", errors="replace")
                _write_markdown(text)

            elif suffix == ".doc":
                return None, [], [], "legacy .doc not supported (save as .docx)"

            elif suffix in _SSOT_IMPORT_IMAGE_EXTENSIONS:
                desc = _describe_image(original_path)
                try:
                    rel = original_path.relative_to(PROJECT_ROOT).as_posix()
                except ValueError:
                    rel = original_path.as_posix()
                md_target.write_text(
                    "\n".join([
                        f"# Image Import: {original_path.name}",
                        "",
                        f"Source: `{rel}`",
                        "",
                        "## Image Description",
                        "",
                        desc if desc else "_(no image description)_",
                        "",
                    ]),
                    encoding="utf-8",
                )
                err = desc if desc.startswith("Error:") else ""
                return md_target, [original_path], [], err

            else:
                md_text, mk_err = _markitdown_convert(original_path)
                if mk_err:
                    return None, image_paths, visual_paths, mk_err
                if not md_text:
                    return None, image_paths, visual_paths, "markitdown produced no markdown output"
                _write_markdown(md_text)

            if suffix == ".pdf":
                import fitz  # type: ignore
                doc = fitz.open(str(original_path))
                n = 0
                for page in doc:
                    for img in page.get_images(full=True):
                        xref = img[0]
                        extracted = doc.extract_image(xref)
                        if not extracted or not extracted.get("image"):
                            continue
                        n += 1
                        _write_image(
                            extracted["image"],
                            extracted.get("ext") or "png",
                            n,
                            filter_noise=True,
                        )
                doc.close()
                if _import_visual_evidence_enabled():
                    surfaces = render_pdf_visual_surfaces(
                        original_path,
                        visual_dir,
                        stamp,
                        idx,
                        max_pages=_import_visual_evidence_max_pages(),
                        scale=_import_visual_evidence_scale(),
                    )
                    visual_surfaces.extend(surfaces)
                    visual_paths.extend(surface.image_path for surface in surfaces)

            elif suffix == ".pptx":
                from pptx import Presentation  # type: ignore
                from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore

                prs = Presentation(str(original_path))
                n = 0
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            image = shape.image
                            n += 1
                            _write_image(image.blob, image.ext or "png", n, filter_noise=True)
                _add_office_visual_surfaces()

            elif suffix == ".docx":
                from docx import Document  # type: ignore

                doc = Document(str(original_path))

                def _docx_tables_as_markdown(d) -> str:
                    """Render every <w:tbl> in the docx as a Markdown table.

                    markitdown frequently flattens .docx tables (loses the
                    separator row or collapses cells). Re-emit tables from
                    `Document.tables` so the SSOT importer keeps register
                    maps, pin lists, and similar structured evidence.
                    """
                    blocks: list[str] = []
                    for tbl in getattr(d, "tables", []) or []:
                        rows = []
                        for row in tbl.rows:
                            cells = [
                                " ".join((c.text or "").split()).replace("|", "\\|") or " "
                                for c in row.cells
                            ]
                            rows.append(cells)
                        if not rows:
                            continue
                        width = max(len(r) for r in rows)
                        norm = [r + [" "] * (width - len(r)) for r in rows]
                        header = "| " + " | ".join(norm[0]) + " |"
                        sep = "| " + " | ".join(["---"] * width) + " |"
                        body = "\n".join("| " + " | ".join(r) + " |" for r in norm[1:])
                        blocks.append(header + "\n" + sep + ("\n" + body if body else ""))
                    return "\n\n".join(blocks)

                docx_tables_md = _docx_tables_as_markdown(doc)

                # Always append the python-docx-rendered tables after the
                # main body (whether markitdown wrote it or the
                # original document did). They sit under a "## Source Tables"
                # heading so downstream consumers can locate them
                # deterministically even when markitdown flattened the
                # originals.
                if docx_tables_md:
                    existing = md_target.read_text(encoding="utf-8", errors="replace")
                    if "## Source Tables" not in existing:
                        with md_target.open("a", encoding="utf-8") as _fh:
                            if existing and not existing.endswith("\n"):
                                _fh.write("\n")
                            _fh.write("\n## Source Tables\n\n")
                            _fh.write(docx_tables_md)
                            _fh.write("\n")

                n = 0
                for rel in doc.part.rels.values():
                    if "image" in (rel.reltype or ""):
                        target = rel.target_part
                        blob = getattr(target, "blob", None)
                        if not blob:
                            continue
                        content_type = getattr(target, "content_type", "") or ""
                        ext = content_type.split("/")[-1] or "png"
                        n += 1
                        _write_image(blob, ext, n, filter_noise=True)
                _add_office_visual_surfaces()

            if image_paths:
                desc_lines = ["", "", "## Extracted Images", ""]
                from concurrent.futures import ThreadPoolExecutor
                with ThreadPoolExecutor(max_workers=min(8, len(image_paths))) as pool:
                    descs = list(pool.map(_describe_image, image_paths))
                for image_index, img_path in enumerate(image_paths):
                    try:
                        rel = img_path.relative_to(PROJECT_ROOT).as_posix()
                    except ValueError:
                        rel = img_path.as_posix()
                    desc_lines.append(f"### `{rel}`")
                    desc_text = str(descs[image_index]).strip()
                    if not desc_text:
                        raise RuntimeError(f"read_image returned an empty description for {img_path}")
                    desc_lines.append(desc_text)
                    desc_lines.append("")
                with open(md_target, "a", encoding="utf-8") as fh:
                    fh.write("\n".join(desc_lines))

            if visual_surfaces:
                desc_lines = ["", "", "## Visual Evidence", ""]
                from concurrent.futures import ThreadPoolExecutor
                with ThreadPoolExecutor(max_workers=min(4, len(visual_surfaces))) as pool:
                    descs = list(pool.map(_describe_visual_surface, visual_surfaces))
                for surface_index, surface in enumerate(visual_surfaces):
                    try:
                        source_rel = surface.image_path.relative_to(PROJECT_ROOT).as_posix()
                    except ValueError:
                        source_rel = surface.image_path.as_posix()
                    desc_text = str(descs[surface_index]).strip()
                    if not desc_text:
                        raise RuntimeError(
                            f"read_image returned an empty visual evidence description for {surface.image_path}"
                    )
                    title = surface.title or "rendered document page"
                    labels = {
                        "pdf_page": "PDF page",
                        "docx_page": "DOCX page",
                        "pptx_slide": "PPTX slide",
                    }
                    surface_label = labels.get(surface.source_kind, surface.source_kind.replace("_", " "))
                    desc_lines.append(f"### {surface_label} {surface.page_number} — {title}")
                    desc_lines.append("")
                    desc_lines.append(f"Source: `{source_rel}`")
                    desc_lines.append(f"Reason: {surface.reason}")
                    desc_lines.append(f"Rendered: ![]({_inline_image_rel(surface.image_path)})")
                    if surface.text_hint:
                        desc_lines.append("")
                        desc_lines.append(f"Text hint: {surface.text_hint}")
                    desc_lines.append("")
                    desc_lines.append(desc_text)
                    desc_lines.append("")
                with open(md_target, "a", encoding="utf-8") as fh:
                    fh.write("\n".join(desc_lines))

            return md_target, image_paths, visual_paths, ""
        except Exception as exc:
            return None, image_paths, visual_paths, f"convert failed: {exc}"

    _SSOT_IMPORT_CONVERTERS = ("markitdown",)

    @app.get("/api/ssot/import/converter")
    async def api_ssot_import_converter_get():
        """Return the active document→markdown converter preference."""
        current = (os.environ.get("ATLAS_IMPORT_CONVERTER", "markitdown") or "markitdown").strip().lower()
        if current not in _SSOT_IMPORT_CONVERTERS:
            return JSONResponse({
                "ok": False,
                "error": f"unsupported import converter {current!r}; expected one of {_SSOT_IMPORT_CONVERTERS}",
                "options": list(_SSOT_IMPORT_CONVERTERS),
            }, status_code=500)
        return JSONResponse({
            "ok": True,
            "converter": current,
            "options": list(_SSOT_IMPORT_CONVERTERS),
        })

    @app.post("/api/ssot/import/converter")
    async def api_ssot_import_converter_set(request: Request):
        """Set the active document→markdown converter preference."""
        try:
            body = await request.json()
        except Exception:
            body = {}
        choice = str((body or {}).get("converter") or "").strip().lower()
        if choice not in _SSOT_IMPORT_CONVERTERS:
            return JSONResponse({
                "ok": False,
                "error": f"converter must be one of {_SSOT_IMPORT_CONVERTERS}",
            }, status_code=400)
        os.environ["ATLAS_IMPORT_CONVERTER"] = choice
        try:
            _persist_env_values({"ATLAS_IMPORT_CONVERTER": choice})
        except Exception:
            pass
        return JSONResponse({"ok": True, "converter": choice})

    @app.post("/api/ssot/import/upload")
    async def api_ssot_import_upload(request: Request):
        """Upload requirement/source evidence into <ip>/req/imports/.

        The route only stores the attachment and returns the exact /import
        command to run. That keeps import semantics in the existing slash
        command path, so Web UI, Textual UI, and headless command handling
        stay aligned.
        """
        content_type = (request.headers.get("content-type") or "").lower()
        upload_entries: list[tuple[str, bytes]] = []
        ip = ""
        if "json" in content_type:
            try:
                body = await request.json()
            except Exception as exc:
                return JSONResponse({"error": f"invalid json body: {exc}"}, status_code=400)
            ip = str((body or {}).get("ip") or _active_ssot_ip() or "").strip()
            import base64 as _base64_import

            for idx, item in enumerate((body or {}).get("files") or []):
                if not isinstance(item, dict):
                    continue
                filename = str(item.get("name") or f"import_{idx}.txt")
                encoded = str(item.get("content_b64") or "")
                if not encoded:
                    continue
                try:
                    upload_entries.append((filename, _base64_import.b64decode(encoded, validate=True)))
                except Exception as exc:
                    return JSONResponse({"error": f"{filename}: invalid base64 payload: {exc}"}, status_code=400)
        else:
            try:
                form = await request.form()
            except Exception as exc:
                return JSONResponse({"error": f"invalid multipart form: {exc}"}, status_code=400)
            ip = str(form.get("ip") or _active_ssot_ip() or "").strip()
            uploads = []
            try:
                uploads.extend(form.getlist("files"))  # type: ignore[attr-defined]
                uploads.extend(form.getlist("file"))  # type: ignore[attr-defined]
            except Exception:
                one = form.get("files") or form.get("file")
                if one is not None:
                    uploads.append(one)
            uploads = [up for up in uploads if hasattr(up, "read")]
            for idx, upload in enumerate(uploads[:16]):
                filename = getattr(upload, "filename", "") or f"import_{idx}.txt"
                try:
                    raw = await upload.read()
                except Exception as exc:
                    return JSONResponse({"error": f"{filename}: read failed: {exc}"}, status_code=400)
                if not isinstance(raw, (bytes, bytearray)):
                    raw = str(raw).encode("utf-8", errors="replace")
                upload_entries.append((filename, bytes(raw)))

        if not _valid_ip_name(ip):
            return JSONResponse({"error": f"invalid ip {ip!r}"}, status_code=400)
        if not upload_entries:
            return JSONResponse({"error": "missing file upload"}, status_code=400)

        dest_dir = _ip_root(ip) / "req" / "imports"
        originals_dir = dest_dir / "originals"
        images_dir = dest_dir / "images"
        visual_dir = dest_dir / "visual"
        try:
            dest_dir.mkdir(parents=True, exist_ok=True)
            originals_dir.mkdir(parents=True, exist_ok=True)
            images_dir.mkdir(parents=True, exist_ok=True)
            visual_dir.mkdir(parents=True, exist_ok=True)
        except OSError as exc:
            return JSONResponse({"error": f"cannot create import dir: {exc}"}, status_code=500)

        saved: list[dict[str, Any]] = []
        errors: list[str] = []
        max_bytes = 32 * 1024 * 1024
        for idx, (raw_name, raw) in enumerate(upload_entries[:16]):
            filename = _safe_import_upload_name(raw_name or f"import_{idx}.txt")
            suffix = Path(filename).suffix.lower()
            if suffix not in _SSOT_IMPORT_EXTENSIONS:
                errors.append(f"{filename}: unsupported extension {suffix or '(none)'}")
                continue
            if len(raw) > max_bytes:
                errors.append(f"{filename}: file too large ({len(raw)} bytes > {max_bytes})")
                continue
            stamp = int(time.time() * 1000)
            basename = Path(filename).stem or f"import_{idx}"
            original_target = originals_dir / f"{stamp}_{idx}_{filename}"
            try:
                original_target.write_bytes(bytes(raw))
            except OSError as exc:
                errors.append(f"{filename}: write failed: {exc}")
                continue

            md_path, image_paths, visual_paths, conv_err = await asyncio.to_thread(
                _convert_upload_to_markdown,
                original_target, suffix, dest_dir, images_dir, visual_dir, stamp, idx, basename,
            )
            entry: dict[str, Any] = {
                "name": filename,
                "bytes": len(raw),
                "original_path": original_target.relative_to(PROJECT_ROOT).as_posix(),
                "md_path": md_path.relative_to(PROJECT_ROOT).as_posix() if md_path else None,
                "image_paths": [
                    p.relative_to(PROJECT_ROOT).as_posix() for p in image_paths
                ],
                "visual_paths": [
                    p.relative_to(PROJECT_ROOT).as_posix() for p in visual_paths
                ],
                "image_count": len(image_paths),
                "visual_count": len(visual_paths),
            }
            if conv_err:
                entry["convert_error"] = conv_err
                errors.append(f"{filename}: {conv_err}")
            entry["path"] = entry["md_path"] or entry["original_path"]
            saved.append(entry)

        if not saved:
            return JSONResponse({"error": "no files saved", "errors": errors}, status_code=400)
        paths = [item["path"] for item in saved]
        command = f"/import --ip {ip} " + " ".join(f"@{path}" for path in paths)
        return JSONResponse({
            "ok": True,
            "ip": ip,
            "saved": saved,
            "paths": paths,
            "errors": errors,
            "command": command,
        })

    def _ssot_context_for_session(session_name: str) -> tuple[AtlasContext | None, JSONResponse | None]:
        raw = normalize_session_name(str(session_name or ""))
        if not raw:
            return None, None
        try:
            return AtlasContext.from_session_key(
                raw,
                atlas_root=os.environ.get("ATLAS_ROOT") or str(PROJECT_ROOT),
            ), None
        except ValueError as exc:
            return None, JSONResponse({"ok": False, "error": str(exc)}, status_code=400)

    def _ssot_base_root_for_context(context: AtlasContext | None) -> Path | None:
        if context is not None and not context.legacy:
            return context.workspace_root
        return None

    def _ssot_ip_root_for_context(ip: str, context: AtlasContext | None) -> Path:
        if context is not None and not context.legacy:
            return context.ip_root
        return _ip_root(ip)

    def _ssot_relative_path(path: Path, context: AtlasContext | None) -> str:
        base = context.workspace_root if context is not None and not context.legacy else PROJECT_ROOT
        try:
            return str(path.relative_to(base))
        except Exception:
            return str(path)

    def _ssot_context_denied(request: Request, context: AtlasContext | None, ip: str):
        if context is None or context.legacy:
            return None
        if context.ip_name and context.ip_name != "default" and ip != context.ip_name:
            return JSONResponse({"ok": False, "error": "session ip mismatch"}, status_code=400)
        uid, uname, is_admin = _authz_identity(request)
        if not uid or uid == "default":
            return JSONResponse({"ok": False, "error": "login required"}, status_code=401)
        if is_admin or uname == context.user_name:
            return None
        return JSONResponse({"ok": False, "error": "session owner mismatch"}, status_code=403)

    def _ssot_html_doc_url(ip: str, session_name: str = "") -> str:
        from urllib.parse import urlencode

        params = {"ip": ip, "format": "html", "inline": "1"}
        if session_name:
            params["session_id"] = session_name
        return f"/api/ssot/export?{urlencode(params)}"

    def _load_ssot_yaml_for_export(ip: str, base_root: Path | None = None) -> dict[str, Any]:
        import yaml as _yaml  # type: ignore

        path = _ssot_yaml_path(ip, base_root=base_root)
        if not path.is_file():
            raise FileNotFoundError(str(path))
        try:
            data = _yaml.safe_load(path.read_text(encoding="utf-8"))
        except Exception as exc:
            raise ValueError(f"invalid yaml: {exc}") from exc
        if data is None:
            return {}
        if not isinstance(data, dict):
            raise ValueError("invalid yaml: top-level must be a mapping")
        return data

    @app.post("/api/ssot/doc-feedback")
    async def api_ssot_doc_feedback(request: Request, session_id: str = "", session: str = ""):
        """Apply DOC-tab feedback to the SSOT draft and anchored DOC export."""
        try:
            body = await request.json()
        except Exception:
            body = {}
        body = body if isinstance(body, dict) else {}
        ip = str(body.get("ip") or "").strip()
        ip = ip if _valid_ip_name(ip) else _active_ssot_ip()
        if not _valid_ip_name(ip):
            return JSONResponse({"ok": False, "error": "no active IP found"}, status_code=400)

        active_session = str(body.get("session_id") or body.get("session") or session_id or session or "").strip()
        context, context_error = _ssot_context_for_session(active_session)
        if context_error is not None:
            return context_error
        denied = _ssot_context_denied(request, context, ip)
        if denied is not None:
            return denied
        base_root = _ssot_base_root_for_context(context)
        path = _ssot_yaml_path(ip, base_root=base_root)
        if not path.is_file():
            return JSONResponse({"ok": False, "error": f"ssot yaml not found for ip {ip!r}"}, status_code=404)
        doc = _load_ssot_draft(ip, base_root=base_root)
        if not isinstance(doc, dict) or not doc:
            return JSONResponse({"ok": False, "error": f"ssot yaml could not be parsed for ip {ip!r}"}, status_code=400)

        section = _ssot_feedback_section(str(body.get("section") or "custom"))
        yaml_path = str(body.get("path") or body.get("yaml_path") or "").strip()
        raw_yaml_path = yaml_path
        field = str(body.get("field") or "").strip()
        value = body.get("value")
        comment = str(body.get("comment") or "").strip()
        value_text = "" if value is None else str(value).strip()
        if not comment and not value_text and not field and not yaml_path:
            return JSONResponse({"ok": False, "error": "feedback, value, field, or path is required"}, status_code=400)

        if not yaml_path and field:
            yaml_path = f"{section}.{field}"
        try:
            tokens = _parse_ssot_feedback_path(yaml_path) if yaml_path else []
            if tokens and (value_text or field):
                write_tokens = list(tokens)
                if raw_yaml_path and field:
                    field_tokens = _parse_ssot_feedback_path(field)
                    if field_tokens and write_tokens[-len(field_tokens):] != field_tokens:
                        write_tokens.extend(field_tokens)
                elif value_text:
                    try:
                        current_value = _get_ssot_feedback_path(doc, write_tokens)
                    except KeyError:
                        current_value = None
                    if isinstance(current_value, dict):
                        write_tokens.append("review_note")
                    elif isinstance(current_value, list):
                        write_tokens = []
                target_value = value if value_text else comment
                if write_tokens:
                    _set_ssot_feedback_path(doc, write_tokens, target_value)
                    yaml_path = _ssot_feedback_tokens_to_path(write_tokens)
                    if isinstance(write_tokens[0], str):
                        section = _ssot_feedback_section(write_tokens[0], section)
        except ValueError as exc:
            return JSONResponse({"ok": False, "error": str(exc)}, status_code=400)

        custom = doc.get("custom")
        if not isinstance(custom, dict):
            custom = {}
            doc["custom"] = custom
        feedback_rows = custom.get("atlas_doc_feedback")
        if not isinstance(feedback_rows, list):
            feedback_rows = []
            custom["atlas_doc_feedback"] = feedback_rows

        ts_ms = int(time.time() * 1000)
        feedback_id = f"doc_feedback_{ts_ms}_{_ssot_feedback_slug(yaml_path or field or section)}"
        record = {
            "id": feedback_id,
            "section": section,
            "path": yaml_path,
            "field": field,
            "value": value if value is not None else "",
            "comment": comment,
            "source": "ssot-doc",
            "created_at_ms": ts_ms,
        }
        feedback_rows.append(record)

        blocks = doc.get("custom_blocks")
        if not isinstance(blocks, list):
            blocks = []
            doc["custom_blocks"] = blocks
        title_target = field or yaml_path or section
        inline_lines = []
        if yaml_path:
            inline_lines.append(f"- path: `{yaml_path}`")
        if field and field not in yaml_path:
            inline_lines.append(f"- field: `{field}`")
        if value_text:
            inline_lines.append(f"- value: {value_text}")
        if comment:
            if inline_lines:
                inline_lines.append("")
            inline_lines.append(comment)
        blocks.append({
            "after": section,
            "title": f"Doc Feedback: {title_target}",
            "type": "markdown",
            "inline": "\n".join(inline_lines).strip() or "Feedback recorded.",
            "source": "ssot-doc",
            "id": feedback_id,
        })

        try:
            _save_ssot_draft(ip, doc, base_root=base_root)
        except Exception as exc:
            return JSONResponse({"ok": False, "error": f"failed to save ssot: {exc}"}, status_code=500)

        rel_path = _ssot_relative_path(path, context)
        return JSONResponse({
            "ok": True,
            "ip": ip,
            "section": section,
            "path": yaml_path,
            "field": field,
            "feedback_id": feedback_id,
            "feedback_count": len(feedback_rows),
            "ssot_path": rel_path,
            "doc_url": _ssot_html_doc_url(ip, active_session),
        })

    @app.get("/api/ssot/doc-source")
    async def api_ssot_doc_source(
        request: Request,
        ip: str = "",
        path: str = "",
        session_id: str = "",
        session: str = "",
    ):
        yaml_path = str(path or "").strip()
        if not yaml_path:
            return JSONResponse({"ok": False, "error": "path is required"}, status_code=400)
        clean_ip = str(ip or "").strip()
        if not _valid_ip_name(clean_ip):
            return JSONResponse({"ok": False, "error": f"invalid ip {ip!r}"}, status_code=400)
        context, context_error = _ssot_context_for_session(session_id or session)
        if context_error is not None:
            return context_error
        denied = _ssot_context_denied(request, context, clean_ip)
        if denied is not None:
            return denied
        base_root = _ssot_base_root_for_context(context)
        ssot_path = _ssot_yaml_path(clean_ip, base_root=base_root)
        if not ssot_path.is_file():
            return JSONResponse({"ok": False, "error": f"ssot yaml not found for ip {clean_ip!r}"}, status_code=404)
        try:
            tokens = _parse_ssot_feedback_path(yaml_path)
        except ValueError as exc:
            return JSONResponse({"ok": False, "error": str(exc)}, status_code=400)
        doc = _load_ssot_draft(clean_ip, base_root=base_root)
        if not isinstance(doc, dict) or not doc:
            return JSONResponse({"ok": False, "error": f"ssot yaml could not be parsed for ip {clean_ip!r}"}, status_code=400)
        try:
            from src.atlas_ssot_doc_map import lookup_ssot_doc_source

            rel_path = _ssot_relative_path(ssot_path, context)
            payload = lookup_ssot_doc_source(doc, path=yaml_path, tokens=tokens, ip=clean_ip, ssot_path=rel_path)
        except KeyError:
            return JSONResponse({"ok": False, "error": "path not found"}, status_code=404)
        except Exception as exc:
            return JSONResponse({"ok": False, "error": f"failed to resolve ssot source: {exc}"}, status_code=500)
        return JSONResponse(payload)

    @app.get("/api/ssot/export")
    async def api_ssot_export(
        request: Request,
        ip: str,
        format: str = "md",
        inline: bool = False,
        session_id: str = "",
        session: str = "",
    ):
        """Export the canonical ssot yaml as md/docx/html for human review.

        Reverse direction of /api/ssot/import/upload. Writes
        <ip>/doc/<ip>_ssot.<ext> deterministically (yaml-walker, no LLM)
        and streams it back via FileResponse with the right Content-Type
        and a download filename.
        """
        fmt = (format or "md").strip().lower()
        if fmt not in {"md", "docx", "html"}:
            return JSONResponse({"error": f"invalid format {format!r}"}, status_code=400)
        if not _valid_ip_name(ip):
            return JSONResponse({"error": f"invalid ip {ip!r}"}, status_code=400)
        context, context_error = _ssot_context_for_session(session_id or session)
        if context_error is not None:
            return context_error
        denied = _ssot_context_denied(request, context, ip)
        if denied is not None:
            return denied
        base_root = _ssot_base_root_for_context(context)
        try:
            data = _load_ssot_yaml_for_export(ip, base_root=base_root)
        except FileNotFoundError:
            return JSONResponse(
                {"error": f"ssot yaml not found for ip {ip!r}"},
                status_code=404,
            )
        except ValueError as exc:
            return JSONResponse({"error": str(exc)}, status_code=400)

        out_dir = _ssot_ip_root_for_context(ip, context) / "doc"
        try:
            out_dir.mkdir(parents=True, exist_ok=True)
        except OSError as exc:
            return JSONResponse(
                {"error": f"cannot create doc dir: {exc}"},
                status_code=500,
            )
        out_path = out_dir / f"{ip}_ssot.{fmt}"

        try:
            if fmt == "md":
                md_text = _ssot_to_markdown(data, ip)
                out_path.write_text(md_text, encoding="utf-8")
                media = "text/markdown; charset=utf-8"
            elif fmt == "html":
                md_text = _ssot_to_markdown(data, ip)
                html_text = _ssot_to_html(md_text, ip, data)
                out_path.write_text(html_text, encoding="utf-8")
                media = "text/html; charset=utf-8"
            else:
                _ssot_to_docx(data, ip, out_path)
                media = (
                    "application/vnd.openxmlformats-officedocument."
                    "wordprocessingml.document"
                )
        except Exception as exc:
            return JSONResponse(
                {"error": f"render failed: {exc}"},
                status_code=500,
            )

        filename = f"{ip}_ssot.{fmt}"
        disposition = "inline" if inline and fmt == "html" else "attachment"
        return FileResponse(
            str(out_path),
            media_type=media,
            filename=filename,
            headers={"Content-Disposition": f'{disposition}; filename="{filename}"'},
        )

    @app.get("/api/req/export")
    async def api_req_export(
        request: Request,
        ip: str,
        format: str = "html",
        variant: str = "full",
        inline: bool = False,
        session_id: str = "",
        session: str = "",
    ):
        """Export the per-IP REQ bundle as one human-reviewable HTML document.

        Companion to /api/ssot/export (the DOC tab). Aggregates the locked-truth
        bundle — requirements, obligations, contract, evidence, approval/lock —
        read-only and streams the rendered HTML into the REQ tab iframe. The
        rendered file is written to <ip>/doc/<ip>_req.html (same convention as
        the SSOT export); no IP source artifacts are modified.
        """
        fmt = (format or "html").strip().lower()
        if fmt != "html":
            return JSONResponse({"error": f"invalid format {format!r}"}, status_code=400)
        if variant not in {"full", "core4"}:
            return JSONResponse({"error": f"invalid variant {variant!r}"}, status_code=400)
        if not _valid_ip_name(ip):
            return JSONResponse({"error": f"invalid ip {ip!r}"}, status_code=400)
        context, context_error = _ssot_context_for_session(session_id or session)
        if context_error is not None:
            return context_error
        denied = _ssot_context_denied(request, context, ip)
        if denied is not None:
            return denied
        ip_root = _ssot_ip_root_for_context(ip, context)
        if not ip_root.is_dir():
            return JSONResponse({"error": f"ip root not found for {ip!r}"}, status_code=404)
        out_dir = ip_root / "doc"
        try:
            out_dir.mkdir(parents=True, exist_ok=True)
        except OSError as exc:
            return JSONResponse({"error": f"cannot create doc dir: {exc}"}, status_code=500)
        out_path = out_dir / f"{ip}_req.html"
        try:
            from src.atlas_req_export import req_html_for_ip

            html_text = req_html_for_ip(ip_root, ip, variant=variant)
            out_path.write_text(html_text, encoding="utf-8")
        except Exception as exc:
            return JSONResponse({"error": f"render failed: {exc}"}, status_code=500)

        filename = f"{ip}_req.html"
        disposition = "inline" if inline else "attachment"
        return FileResponse(
            str(out_path),
            media_type="text/html; charset=utf-8",
            filename=filename,
            headers={"Content-Disposition": f'{disposition}; filename="{filename}"'},
        )

    @app.post("/api/ssot/validate")
    async def api_ssot_validate(request: Request):
        """Run the disk-truth SSOT validator for one IP.

        Use the Python verifier as the UI entry point so ATLAS_PROJECT_ROOT /
        --root stays explicit. The verifier still calls check_ssot_disk.sh
        when that shell wrapper is available, but can report/fallback cleanly
        when bash or the wrapper path is unavailable.
        """
        try:
            body = await request.json()
        except Exception:
            body = {}
        ip = str((body or {}).get("ip") or "").strip()
        ip = ip if _valid_ip_name(ip) else _active_ssot_ip()
        mode = str((body or {}).get("mode") or "engineering").strip().lower().replace("_", "-")
        if mode == "eng":
            mode = "engineering"
        if mode == "sign-off":
            mode = "signoff"
        if mode not in {"starter", "engineering", "signoff"}:
            return JSONResponse({"error": f"invalid mode {mode!r}"}, status_code=400)
        if not _valid_ip_name(ip):
            return JSONResponse({"error": "no active IP found"}, status_code=400)

        script = WORKFLOW_ROOT / "ssot-gen" / "scripts" / "verify_ssot.py"
        if not script.is_file():
            return JSONResponse({"error": f"validator script not found: {script}"}, status_code=404)

        env = os.environ.copy()
        env["IP_NAME"] = ip
        env["ATLAS_RUN_MODE"] = mode
        env["ATLAS_PROJECT_ROOT"] = str(PROJECT_ROOT)
        cmd = [
            "python",
            str(script),
            ip,
            "--root",
            str(PROJECT_ROOT),
            "--mode",
            mode,
            "--preview",
            "strict",
        ]
        started = time.time()
        command_text = " ".join(shlex.quote(part) for part in cmd)

        def _run_verify_inprocess() -> tuple[int, str, str]:
            spec = importlib.util.spec_from_file_location("_atlas_verify_ssot_runtime", script)
            if spec is None or spec.loader is None:
                raise RuntimeError(f"cannot load verifier module: {script}")
            module = importlib.util.module_from_spec(spec)
            spec.loader.exec_module(module)
            old_argv = sys.argv[:]
            stdout = io.StringIO()
            stderr = io.StringIO()
            try:
                sys.argv = cmd[1:]
                with contextlib.redirect_stdout(stdout), contextlib.redirect_stderr(stderr):
                    rc = int(module.main())
            finally:
                sys.argv = old_argv
            return rc, stdout.getvalue(), stderr.getvalue()

        try:
            returncode, stdout, stderr = await asyncio.to_thread(_run_verify_inprocess)
        except subprocess.TimeoutExpired as exc:
            return JSONResponse({
                "ok": False,
                "ip": ip,
                "mode": mode,
                "returncode": 124,
                "stdout": exc.stdout or "",
                "stderr": exc.stderr or "validation timed out",
                "elapsed_ms": int((time.time() - started) * 1000),
                "command": command_text,
            }, status_code=504)
        except Exception as exc:
            return JSONResponse({"error": f"validation failed to launch: {exc}"}, status_code=500)

        return JSONResponse({
            "ok": returncode == 0,
            "ip": ip,
            "mode": mode,
            "returncode": returncode,
            "stdout": stdout or "",
            "stderr": stderr or "",
            "elapsed_ms": int((time.time() - started) * 1000),
            "command": command_text,
        })













    @app.get("/api/catalog/models")
    async def api_catalog_models():
        """Return available IP models/templates discoverable in the project.

        Catalog models are reusable sources that can be instantiated.
        /api/soc remains the actual placed instance hierarchy.
        """
        try:
            import yaml as _yaml
        except ImportError:
            return JSONResponse({"models": [], "count": 0,
                                 "error": "PyYAML not installed"})
        models = []
        seen = set()
        def _catalog_kind(name: str) -> str:
            n = (name or "").lower()
            if any(x in n for x in ["cpu", "core", "cortex", "riscv"]): return "cpu"
            if any(x in n for x in ["noc", "axi", "crossbar", "interconnect", "cci"]): return "bus"
            if any(x in n for x in ["ddr", "sram", "mem", "dram"]): return "mem"
            if any(x in n for x in ["pll", "adc", "dac", "phy"]): return "analog"
            return "periph"
        for p in sorted(PROJECT_ROOT.glob("*/yaml/*.ssot.yaml")):
            try:
                doc = _yaml.safe_load(p.read_text(encoding="utf-8", errors="replace")) or {}
            except Exception:
                doc = {}
            if not isinstance(doc, dict):
                doc = {}
            ip_dir = p.parents[1]
            top = doc.get("top_module")
            name = top if isinstance(top, str) and top.strip() else ip_dir.name
            if name in seen:
                continue
            seen.add(name)
            ports = []
            for bi in (doc.get("busInterfaces") or []):
                if isinstance(bi, dict):
                    ports.append({
                        "name": bi.get("name"),
                        "proto": bi.get("proto"),
                        "role": bi.get("role"),
                        "side": bi.get("side"),
                    })
            models.append({
                "name": name,
                "id": ip_dir.name,
                "kind": _catalog_kind(name),
                "source": "project",
                "ssot_path": p.relative_to(PROJECT_ROOT).as_posix(),
                "ports": ports,
            })
        return JSONResponse({"models": models, "count": len(models)})

    @app.get("/api/workspace/tree")
    async def api_workspace_tree(depth: int = 2):
        """Return the real project directory hierarchy for Architect.

        This is separate from /api/catalog/models: catalog is reusable IP
        models; workspace tree is the on-disk project shape.
        """
        max_depth = max(1, min(int(depth or 2), 4))
        skip = {
            ".git", "__pycache__", ".pytest_cache", ".mypy_cache",
            ".ruff_cache", "node_modules", ".venv", "venv", "vendor",
            ".session", ".rag", ".claude", ".omc", ".benchmark",
            ".benchmarks", ".common_ai_agent", ".session_debug",
        }
        ip_artifacts = {"yaml", "rtl", "tb", "sim", "lint", "syn", "sta", "sta-post", "pnr", "dft", "doc", "req", "list"}

        def _meta(p: Path) -> dict:
            child_names = set()
            try:
                child_names = {c.name for c in p.iterdir() if c.is_dir()}
            except OSError:
                pass
            ssot_count = 0
            try:
                ssot_count = len(list((p / "yaml").glob("*.ssot.yaml"))) if (p / "yaml").is_dir() else 0
            except OSError:
                ssot_count = 0
            artifacts = sorted(child_names & ip_artifacts)
            return {
                "is_ip": ssot_count > 0,
                "ssot_count": ssot_count,
                "artifacts": artifacts,
            }

        def _node(p: Path, level: int) -> dict | None:
            name = p.name or str(p)
            if name in skip or name.startswith("."):
                return None
            node = {
                "name": name,
                "path": p.relative_to(PROJECT_ROOT).as_posix(),
                "kind": "dir",
                **_meta(p),
                "children": [],
            }
            if level >= max_depth:
                return node
            try:
                dirs = sorted([c for c in p.iterdir() if c.is_dir()],
                              key=lambda c: (not ((c / "yaml").is_dir()), c.name.lower()))
            except OSError:
                dirs = []
            for child in dirs:
                child_node = _node(child, level + 1)
                if child_node is not None:
                    node["children"].append(child_node)
            return node

        root = {
            "name": PROJECT_ROOT.name,
            "path": ".",
            "kind": "root",
            "children": [],
        }
        try:
            top_dirs = sorted([p for p in PROJECT_ROOT.iterdir() if p.is_dir()],
                              key=lambda p: (not ((p / "yaml").is_dir()), p.name.lower()))
        except OSError:
            top_dirs = []
        for p in top_dirs:
            n = _node(p, 1)
            if n is not None:
                root["children"].append(n)
        return JSONResponse({"root": root, "count": len(root["children"]),
                             "project_root": str(PROJECT_ROOT)})

    @app.post("/api/soc/layout")
    async def api_soc_layout(request: Request):
        """Persist user-dragged block positions back into soc.ssot.yaml.

        Body (JSON): `{"layout": {"<cluster>/<inst>": {"x": <num>, "y": <num>}, …}}`

        For each entry, find the matching `instances[].id` (`<cluster>/<inst>`
        is split on `/`; we use just the inst id since SoC SSOT instance
        ids are unique) and set its `x:` / `y:` keys. Other fields are
        left untouched. The file is rewritten in-place with
        `yaml.safe_dump`. Empty layout `{}` clears all x/y from every
        instance (paired with the frontend's [reset] button).

        Architect screen reads these on the next /api/soc fetch and
        uses them as the default block position (overriding the
        auto-grid). LocalStorage layout still wins as the most-local
        cache, so a user can preview drag-arounds before committing.
        """
        try:
            body = await request.json()
        except Exception as e:
            return JSONResponse({"error": f"bad json: {e}"}, status_code=400)
        layout = body.get("layout") if isinstance(body, dict) else None
        if not isinstance(layout, dict):
            return JSONResponse({"error": "missing 'layout' object"}, status_code=400)
        try:
            import yaml as _yaml
        except ImportError:
            return JSONResponse({"error": "PyYAML not installed"}, status_code=500)

        soc_path = PROJECT_ROOT / "soc.ssot.yaml"
        if not soc_path.is_file():
            return JSONResponse({"error": "soc.ssot.yaml not found at project root"},
                                 status_code=404)
        try:
            doc = _yaml.safe_load(soc_path.read_text(encoding="utf-8")) or {}
        except Exception as e:
            return JSONResponse({"error": f"parse: {e}"}, status_code=500)
        if not isinstance(doc, dict): doc = {}
        instances = doc.get("instances")
        if not isinstance(instances, list):
            return JSONResponse({"error": "soc.ssot.yaml has no instances[]"},
                                 status_code=400)

        # Build lookup `inst_id → ref` from layout keys.
        ref_for_inst = {}
        for ref in layout.keys():
            if not isinstance(ref, str) or "/" not in ref: continue
            inst_id = ref.split("/", 1)[1]
            ref_for_inst[inst_id] = ref

        touched = 0
        cleared = 0
        for inst in instances:
            if not isinstance(inst, dict): continue
            iid = inst.get("id")
            if not iid: continue
            ref = ref_for_inst.get(iid)
            if ref is None:
                # Instance not in incoming layout: if it has stale x/y
                # AND the incoming layout was empty `{}`, clear them.
                if not layout:
                    if "x" in inst: inst.pop("x"); cleared += 1
                    if "y" in inst: inst.pop("y"); cleared += 1
                    if "top_x" in inst: inst.pop("top_x"); cleared += 1
                    if "top_y" in inst: inst.pop("top_y"); cleared += 1
                continue
            pos = layout.get(ref)
            if isinstance(pos, dict) and isinstance(pos.get("x"), (int, float)) \
               and isinstance(pos.get("y"), (int, float)):
                if ref.startswith("top:"):
                    inst["top_x"] = round(float(pos["x"]), 1)
                    inst["top_y"] = round(float(pos["y"]), 1)
                else:
                    inst["x"] = round(float(pos["x"]), 1)
                    inst["y"] = round(float(pos["y"]), 1)
                touched += 1

        # Preserve hex formatting on address fields. PyYAML parses
        # `0x4000_2000` to int 1073750016 on load; safe_dump would write
        # it back as decimal. Walk the doc and stringify any int in
        # known address slots so the rewritten file keeps the canonical
        # `0x4000_2000` notation the user wrote.
        def _hex8(n):
            h = f"{n:x}"
            target = max(8, ((len(h) - 1) // 4 + 1) * 4)
            h = h.zfill(target)
            if len(h) > 4:
                rev = h[::-1]
                groups = [rev[i:i+4] for i in range(0, len(rev), 4)]
                h = "_".join(groups)[::-1]
            return f"0x{h}"
        for inst in (doc.get("instances") or []):
            if isinstance(inst, dict) and isinstance(inst.get("addr"), int):
                inst["addr"] = _hex8(inst["addr"])
        for e in (doc.get("addrMap") or []):
            if isinstance(e, dict):
                if isinstance(e.get("base"), int):  e["base"]  = _hex8(e["base"])
                if isinstance(e.get("range"), int): e["range"] = _hex8(e["range"])

        try:
            with open(soc_path, "w", encoding="utf-8") as f:
                _yaml.safe_dump(doc, f, sort_keys=False,
                                default_flow_style=False, allow_unicode=True)
        except OSError as e:
            return JSONResponse({"error": f"write: {e}"}, status_code=500)
        return JSONResponse({"ok": True, "touched": touched, "cleared": cleared,
                              "path": soc_path.relative_to(PROJECT_ROOT).as_posix()})

    @app.post("/api/soc/connect")
    async def api_soc_connect(request: Request):
        """Append a port-to-port connection to soc.ssot.yaml.

        Body: {"from": "ip/PORT", "to": "ip/PORT", "proto": "AXI4"}
        """
        try:
            body = await request.json()
        except Exception as e:
            return JSONResponse({"error": f"bad json: {e}"}, status_code=400)
        if not isinstance(body, dict):
            return JSONResponse({"error": "expected json object"}, status_code=400)
        src = str(body.get("from") or "").strip()
        dst = str(body.get("to") or "").strip()
        proto = str(body.get("proto") or "").strip().upper()
        if "/" not in src or "/" not in dst:
            return JSONResponse({"error": "from/to must look like ip/PORT"},
                                status_code=400)
        if src == dst:
            return JSONResponse({"error": "cannot connect a port to itself"},
                                status_code=400)
        try:
            import yaml as _yaml
        except ImportError:
            return JSONResponse({"error": "PyYAML not installed"}, status_code=500)

        soc_path = PROJECT_ROOT / "soc.ssot.yaml"
        if not soc_path.is_file():
            return JSONResponse({"error": "soc.ssot.yaml not found at project root"},
                                status_code=404)
        try:
            doc = _yaml.safe_load(soc_path.read_text(encoding="utf-8")) or {}
        except Exception as e:
            return JSONResponse({"error": f"parse: {e}"}, status_code=500)
        if not isinstance(doc, dict):
            doc = {}
        conns = doc.setdefault("connections", [])
        if not isinstance(conns, list):
            return JSONResponse({"error": "soc.ssot.yaml connections is not a list"},
                                status_code=400)
        for c in conns:
            if isinstance(c, dict) and c.get("from") == src and c.get("to") == dst:
                return JSONResponse({"ok": True, "duplicate": True,
                                     "connection": c,
                                     "path": soc_path.relative_to(PROJECT_ROOT).as_posix()})
        entry = {"from": src, "to": dst}
        if proto:
            entry["proto"] = proto
        conns.append(entry)

        def _hex8(n):
            h = f"{n:x}"
            target = max(8, ((len(h) - 1) // 4 + 1) * 4)
            h = h.zfill(target)
            if len(h) > 4:
                rev = h[::-1]
                groups = [rev[i:i+4] for i in range(0, len(rev), 4)]
                h = "_".join(groups)[::-1]
            return f"0x{h}"
        for inst in (doc.get("instances") or []):
            if isinstance(inst, dict) and isinstance(inst.get("addr"), int):
                inst["addr"] = _hex8(inst["addr"])
        for e in (doc.get("addrMap") or []):
            if isinstance(e, dict):
                if isinstance(e.get("base"), int):  e["base"]  = _hex8(e["base"])
                if isinstance(e.get("range"), int): e["range"] = _hex8(e["range"])
        try:
            with open(soc_path, "w", encoding="utf-8") as f:
                _yaml.safe_dump(doc, f, sort_keys=False,
                                default_flow_style=False, allow_unicode=True)
        except OSError as e:
            return JSONResponse({"error": f"write: {e}"}, status_code=500)
        return JSONResponse({"ok": True, "connection": entry,
                             "path": soc_path.relative_to(PROJECT_ROOT).as_posix()})

    @app.post("/api/soc/instance/add")
    async def api_soc_instance_add(request: Request):
        """Instantiate a catalog model into soc.ssot.yaml.

        Body: {"model":"spi_master", "id":"spi_master_0",
               "cluster":"periph_ss", "addr":"0x4000_3000"}
        """
        try:
            body = await request.json()
        except Exception as e:
            return JSONResponse({"error": f"bad json: {e}"}, status_code=400)
        if not isinstance(body, dict):
            return JSONResponse({"error": "expected json object"}, status_code=400)
        model = str(body.get("model") or body.get("name") or "").strip()
        req_id = str(body.get("id") or "").strip()
        cluster_id = str(body.get("cluster") or "").strip()
        addr = body.get("addr")
        if not model:
            return JSONResponse({"error": "missing model"}, status_code=400)
        try:
            import yaml as _yaml
        except ImportError:
            return JSONResponse({"error": "PyYAML not installed"}, status_code=500)
        soc_path = PROJECT_ROOT / "soc.ssot.yaml"
        if not soc_path.is_file():
            return JSONResponse({"error": "soc.ssot.yaml not found at project root"},
                                status_code=404)

        catalog = []
        for p in sorted(PROJECT_ROOT.glob("*/yaml/*.ssot.yaml")):
            try:
                d = _yaml.safe_load(p.read_text(encoding="utf-8", errors="replace")) or {}
            except Exception:
                d = {}
            if not isinstance(d, dict):
                d = {}
            ip_dir = p.parents[1]
            top = d.get("top_module")
            name = top if isinstance(top, str) and top.strip() else ip_dir.name
            catalog.append({"name": name, "id": ip_dir.name,
                            "ssot": p.relative_to(PROJECT_ROOT).as_posix()})
        found = next((m for m in catalog
                      if m["name"] == model or m["id"] == model), None)
        if not found:
            return JSONResponse({"error": f"model not found: {model}"}, status_code=404)

        try:
            doc = _yaml.safe_load(soc_path.read_text(encoding="utf-8")) or {}
        except Exception as e:
            return JSONResponse({"error": f"soc parse: {e}"}, status_code=500)
        if not isinstance(doc, dict):
            doc = {}
        instances = doc.setdefault("instances", [])
        clusters = doc.setdefault("clusters", [])
        if not isinstance(instances, list) or not isinstance(clusters, list):
            return JSONResponse({"error": "soc.ssot.yaml instances/clusters must be lists"},
                                status_code=400)
        existing_ids = {str(i.get("id")) for i in instances if isinstance(i, dict) and i.get("id")}
        base = req_id or found["name"]
        inst_id = base
        if inst_id in existing_ids:
            i = 0
            while f"{base}_{i}" in existing_ids:
                i += 1
            inst_id = f"{base}_{i}"

        def _role_for_name(name):
            n = (name or "").lower()
            if any(x in n for x in ["cpu", "core", "cortex", "riscv"]): return ("cpu_ss", "CPU")
            if any(x in n for x in ["noc", "axi", "crossbar", "interconnect", "cci"]): return ("noc", "BUS")
            if any(x in n for x in ["ddr", "sram", "mem", "dram"]): return ("mem_ss", "MEM")
            return ("periph_ss", "PERIPH")

        default_cluster, default_role = _role_for_name(found["name"])
        cluster_id = cluster_id or default_cluster
        cluster = next((c for c in clusters
                        if isinstance(c, dict) and (c.get("id") or c.get("name")) == cluster_id), None)
        if cluster is None:
            cluster = {"id": cluster_id, "role": default_role, "members": []}
            clusters.append(cluster)
        members = cluster.setdefault("members", [])
        if isinstance(members, list) and inst_id not in members:
            members.append(inst_id)

        inst = {"id": inst_id, "ssot": found["ssot"]}
        if addr not in (None, ""):
            inst["addr"] = str(addr)
        if isinstance(body.get("x"), (int, float)): inst["top_x"] = round(float(body["x"]), 1)
        if isinstance(body.get("y"), (int, float)): inst["top_y"] = round(float(body["y"]), 1)
        instances.append(inst)

        try:
            with open(soc_path, "w", encoding="utf-8") as f:
                _yaml.safe_dump(doc, f, sort_keys=False,
                                default_flow_style=False, allow_unicode=True)
        except OSError as e:
            return JSONResponse({"error": f"write: {e}"}, status_code=500)
        return JSONResponse({"ok": True, "instance": inst, "cluster": cluster_id,
                             "model": found, "path": soc_path.relative_to(PROJECT_ROOT).as_posix()})

    @app.post("/api/soc/instance/delete")
    async def api_soc_instance_delete(request: Request):
        """Remove an instance from the SoC hierarchy without deleting model files.

        Body: {"id":"counter"}
        Removes:
          - instances[] entry
          - clusters[].members[] reference
          - connections touching <id>/*
          - addrMap entry matching id/name
        """
        try:
            body = await request.json()
        except Exception as e:
            return JSONResponse({"error": f"bad json: {e}"}, status_code=400)
        inst_id = str((body or {}).get("id") or "").strip()
        if not inst_id:
            return JSONResponse({"error": "missing id"}, status_code=400)
        try:
            import yaml as _yaml
        except ImportError:
            return JSONResponse({"error": "PyYAML not installed"}, status_code=500)
        soc_path = PROJECT_ROOT / "soc.ssot.yaml"
        if not soc_path.is_file():
            return JSONResponse({"error": "soc.ssot.yaml not found at project root"},
                                status_code=404)
        try:
            doc = _yaml.safe_load(soc_path.read_text(encoding="utf-8")) or {}
        except Exception as e:
            return JSONResponse({"error": f"soc parse: {e}"}, status_code=500)
        if not isinstance(doc, dict):
            doc = {}
        removed = {"instances": 0, "members": 0, "connections": 0, "addrMap": 0}
        instances = doc.get("instances") or []
        if isinstance(instances, list):
            kept = []
            for inst in instances:
                if isinstance(inst, dict) and str(inst.get("id") or "") == inst_id:
                    removed["instances"] += 1
                else:
                    kept.append(inst)
            doc["instances"] = kept
        for c in (doc.get("clusters") or []):
            if not isinstance(c, dict) or not isinstance(c.get("members"), list):
                continue
            before = len(c["members"])
            c["members"] = [m for m in c["members"] if str(m) != inst_id]
            removed["members"] += before - len(c["members"])
        conns = doc.get("connections") or []
        if isinstance(conns, list):
            kept = []
            prefix = f"{inst_id}/"
            for conn in conns:
                if isinstance(conn, dict) and (
                    str(conn.get("from") or "").startswith(prefix) or
                    str(conn.get("to") or "").startswith(prefix)
                ):
                    removed["connections"] += 1
                else:
                    kept.append(conn)
            doc["connections"] = kept
        amap = doc.get("addrMap") or []
        if isinstance(amap, list):
            kept = []
            for ent in amap:
                if isinstance(ent, dict) and str(ent.get("name") or "") == inst_id:
                    removed["addrMap"] += 1
                else:
                    kept.append(ent)
            doc["addrMap"] = kept
        if removed["instances"] == 0:
            return JSONResponse({"error": f"instance not found: {inst_id}",
                                 "removed": removed}, status_code=404)

        def _hex8(n):
            if isinstance(n, int):
                return "0x" + format(n, "08x")
            s = str(n)
            if s.startswith("0x"):
                return s
            return n

        for inst in (doc.get("instances") or []):
            if isinstance(inst, dict) and isinstance(inst.get("addr"), int):
                inst["addr"] = _hex8(inst["addr"])
        for e in (doc.get("addrMap") or []):
            if isinstance(e, dict):
                if isinstance(e.get("base"), int): e["base"] = _hex8(e["base"])
                if isinstance(e.get("range"), int): e["range"] = _hex8(e["range"])
        try:
            with open(soc_path, "w", encoding="utf-8") as f:
                _yaml.safe_dump(doc, f, sort_keys=False,
                                default_flow_style=False, allow_unicode=True)
        except OSError as e:
            return JSONResponse({"error": f"write: {e}"}, status_code=500)
        return JSONResponse({"ok": True, "id": inst_id, "removed": removed,
                             "path": soc_path.relative_to(PROJECT_ROOT).as_posix()})

    # Phase 16: /api/diagram/plan extracted to atlas_api_diagram_plan.py.
    from src.atlas_api_diagram_plan import register_diagram_plan_routes as _register_diagram_plan_routes
    _register_diagram_plan_routes(app, PROJECT_ROOT=PROJECT_ROOT)

    @app.post("/api/ipxact/import")
    async def api_ipxact_import(request: Request):
        """Import an IP-XACT XML payload into the project as a new IP.

        Accepts either:
          • multipart/form-data with a `xml` file part + optional `name`
          • application/json: {"xml": "<XML…>", "name": "spi_master"}
          • application/xml or text/xml body + ?name=<ip_name> query

        Writes <project_root>/<name>/yaml/<name>.ssot.yaml and scaffolds
        the surrounding IP layout. Returns the parsed SSOT + path.
        """
        try:
            ct = (request.headers.get("content-type") or "").lower()
            xml_text: str = ""
            ip_name: str = ""
            if ct.startswith("multipart/form-data"):
                form = await request.form()
                up = form.get("xml")
                if up is None:
                    return JSONResponse({"error": "missing 'xml' file part"}, status_code=400)
                if hasattr(up, "read"):
                    raw = await up.read()
                    xml_text = raw.decode("utf-8", errors="replace") if isinstance(raw, (bytes, bytearray)) else str(raw)
                else:
                    xml_text = str(up)
                ip_name = (form.get("name") or "").strip()
            elif "json" in ct:
                body = await request.json()
                xml_text = body.get("xml", "") or ""
                ip_name = (body.get("name") or "").strip()
            else:
                raw = await request.body()
                xml_text = raw.decode("utf-8", errors="replace") if isinstance(raw, (bytes, bytearray)) else str(raw)
                ip_name = (request.query_params.get("name") or "").strip()
            if not xml_text.strip():
                return JSONResponse({"error": "empty XML payload"}, status_code=400)

            try:
                from core.ipxact_import import import_ipxact as _conv
            except Exception:
                try: from ipxact_import import import_ipxact as _conv  # type: ignore
                except Exception as e:
                    return JSONResponse({"error": f"importer unavailable: {e}"}, status_code=500)
            try:
                ssot = _conv(xml_text, ip_name=ip_name or None)
            except Exception as e:
                return JSONResponse({"error": f"parse error: {e}"}, status_code=400)
            name = (ip_name or ssot.get("top_module") or "").strip()
            if not name or not re.match(r"^[A-Za-z][A-Za-z0-9_]*$", name):
                return JSONResponse({"error": f"invalid ip name {name!r}"}, status_code=400)

            # Write into <project_root>/<name>/yaml/<name>.ssot.yaml.
            ip_dir = PROJECT_ROOT / name
            (ip_dir / "yaml").mkdir(parents=True, exist_ok=True)
            for sub in ("rtl", "sim", "tb", "list", "lint", "doc"):
                (ip_dir / sub).mkdir(parents=True, exist_ok=True)
            yaml_path = ip_dir / "yaml" / f"{name}.ssot.yaml"
            try:
                import yaml as _yaml
                with open(yaml_path, "w", encoding="utf-8") as f:
                    f.write("# Auto-imported from IP-XACT — review and edit as needed.\n")
                    _yaml.safe_dump(ssot, f, sort_keys=False, default_flow_style=False, allow_unicode=True)
            except ImportError:
                # No PyYAML → write a JSON sidecar so the import still
                # produces something usable.
                with open(yaml_path.with_suffix(".json"), "w", encoding="utf-8") as f:
                    json.dump(ssot, f, indent=2)
                yaml_path = yaml_path.with_suffix(".json")
            except OSError as e:
                return JSONResponse({"error": f"write error: {e}"}, status_code=500)

            # Auto-register into soc.ssot.yaml when Tier-1 is active so
            # the new IP appears in the architect tree on the next
            # /api/soc fetch (was P1 bug — yaml on disk but tree empty).
            registered = False
            try:
                soc_path = PROJECT_ROOT / "soc.ssot.yaml"
                if soc_path.is_file():
                    import yaml as _y
                    sd = _y.safe_load(soc_path.read_text(encoding="utf-8")) or {}
                    if isinstance(sd, dict):
                        instances = sd.setdefault("instances", [])
                        # Skip if an instance with this id already exists.
                        existing = next((i for i in instances
                                          if isinstance(i, dict) and i.get("id") == name), None)
                        if existing is None:
                            new_inst = {
                                "id": name,
                                "ssot": yaml_path.relative_to(PROJECT_ROOT).as_posix(),
                            }
                            # Pull addr from the imported IP's memoryMap
                            # so addrmap_check can validate it.
                            mm = (ssot or {}).get("memoryMap") or []
                            if isinstance(mm, list) and mm and isinstance(mm[0], dict):
                                base = mm[0].get("base")
                                if base: new_inst["addr"] = base
                            instances.append(new_inst)
                            # Drop into a synthetic "uncategorized" cluster
                            # if nothing else claims it (clusters[].members
                            # is the source of truth — auto-add a stub).
                            clusters = sd.setdefault("clusters", [])
                            uncat = next((c for c in clusters
                                          if isinstance(c, dict) and c.get("id") == "uncategorized"),
                                          None)
                            if uncat is None:
                                clusters.append({
                                    "id": "uncategorized",
                                    "role": "PERIPH",
                                    "label": "Uncategorized (auto-imported)",
                                    "members": [name],
                                })
                            else:
                                members = uncat.setdefault("members", [])
                                if name not in members: members.append(name)
                            with open(soc_path, "w", encoding="utf-8") as f:
                                _y.safe_dump(sd, f, sort_keys=False,
                                             default_flow_style=False, allow_unicode=True)
                            registered = True
            except Exception as e:
                # Non-fatal — IP file is on disk, just couldn't auto-register.
                # Frontend will still see it via Tier-2 fallback.
                pass

            return JSONResponse({
                "ok": True,
                "name": name,
                "path": yaml_path.relative_to(PROJECT_ROOT).as_posix(),
                "registered_in_soc": registered,
                "ssot": ssot,
            })
        except Exception as e:
            return JSONResponse({"error": str(e)}, status_code=500)

    @app.get("/api/conversation")
    async def api_conversation(limit: int = 200):
        """Return the last N messages from the active workspace's
        conversation.json. Used by the Atlas frontend to hydrate the
        chat feed when the user switches workflow (/wf <name>) — without
        this the Atlas chat is browser-session-only and prior context
        from the workflow is invisible.

        config.HISTORY_FILE is already redirected per workspace by
        session_setup.setup_session, so we just read whichever file
        the live session points at right now.
        """
        try:
            try: import src.config as _cfg  # type: ignore
            except Exception:
                try: import config as _cfg  # type: ignore
                except Exception: _cfg = None
            if _cfg is None:
                return JSONResponse({"messages": [], "error": "config unavailable"})
            hpath = Path(getattr(_cfg, "HISTORY_FILE", "") or "")

            def _read_history(path: Path) -> list[dict[str, Any]]:
                if not path.is_file():
                    return []
                raw = json.loads(path.read_text(encoding="utf-8"))
                return raw if isinstance(raw, list) else []

            fallback_session = ""
            try:
                msgs = _read_history(hpath)
            except Exception as e:
                return JSONResponse({"messages": [], "path": str(hpath),
                                       "error": f"parse: {e}"})
            non_system = [m for m in msgs if isinstance(m, dict) and m.get("role") != "system"]
            if not non_system:
                root = PROJECT_ROOT / ".session"
                candidates = []
                if root.is_dir():
                    for p in root.rglob("conversation.json"):
                        try:
                            candidates.append((p.stat().st_mtime, p))
                        except Exception:
                            pass
                for _, p in sorted(candidates, reverse=True):
                    try:
                        alt = _read_history(p)
                    except Exception:
                        continue
                    alt_non_system = [
                        m for m in alt
                        if isinstance(m, dict) and m.get("role") != "system"
                    ]
                    if alt_non_system:
                        msgs = alt
                        hpath = p
                        try:
                            fallback_session = p.parent.relative_to(root).as_posix()
                        except Exception:
                            fallback_session = str(p.parent)
                        break
            # Drop system prompts (huge, useless in chat replay) and
            # keep only the last `limit` items.
            msgs = [m for m in msgs if isinstance(m, dict) and m.get("role") != "system"]
            if len(msgs) > limit:
                msgs = msgs[-limit:]
            return JSONResponse({"messages": msgs, "path": str(hpath),
                                  "fallback_session": fallback_session,
                                  "truncated_to": limit})
        except Exception as e:
            return JSONResponse({"messages": [], "error": str(e)},
                                 status_code=500)

    # /api/session/history, /api/session/state, /api/session/list registered
    # via register_sessions_routes() (see atlas_api_sessions.py).

    # Jobs API (/api/job/dispatch, /api/jobs, /api/job/{id}/log, /api/job/{id}/cancel,
    # /api/jobs/dispatch_many, /api/jobs/clear, /api/pipeline/*) lives in
    # src/atlas_api_jobs.py. Inject runtime callables so routes see PROJECT_ROOT
    # changes from --root and the live normalize_session_name function.
    from atlas_api_jobs import register_jobs_routes, get_jobs_state as _get_jobs_state  # noqa: WPS433
    register_jobs_routes(
        app,
        project_root=lambda: PROJECT_ROOT,
        normalize_session_name=normalize_session_name,
        persist_config_values=_persist_config_values,
    )

    # ── Git API — status / diff / commit / push ─────────────────
    # All git commands run inside PROJECT_ROOT (the user's cwd at
    # launch). Read-only ops stream back; commit + push run sync
    # and return their stdout/stderr. Push includes an explicit
    # confirm flag because it's destructive (remote-visible).
    import subprocess as _sp_git
    # Git API (status / log / show / diff / commit / push) lives in
    # src/atlas_api_git.py. Inject runtime callables so the routes see
    # PROJECT_ROOT changes from --root and the live active-IP value.
    from atlas_api_git import register_git_routes  # noqa: WPS433
    register_git_routes(
        app,
        project_root=lambda: PROJECT_ROOT,
        active_ip_value=_active_ip_value,
        valid_ip_name=_valid_ip_name,
    )
    from atlas_api_vcd import register_vcd_routes  # noqa: WPS433
    register_vcd_routes(
        app,
        project_root=lambda: PROJECT_ROOT,
        safe_path=_safe,
        skip_dirs=SKIP_DIRS,
        max_vcd_bytes=MAX_VCD_BYTES,
        fs_authz=_fs_authz,
    )
    # Workspaces API (list workflow definitions + download.zip) lives in
    # src/atlas_api_workspaces.py. Inject runtime callables so routes see
    # PROJECT_ROOT changes from --root.
    from atlas_api_workspaces import register_workspaces_routes  # noqa: WPS433
    register_workspaces_routes(
        app,
        project_root=lambda: PROJECT_ROOT,
        source_root=SOURCE_ROOT,
        safe_path=_safe,
    )
    # Orchestrator Chat API (per-IP rooms + _global). Routes share the
    # AtlasDB used everywhere else and route chat events through the
    # multi-user bridge's broadcast_all so cross-session clients see
    # them in real time.
    from atlas_api_chat import register_chat_routes  # noqa: WPS433
    from core.atlas_db import AtlasDB as _ChatAtlasDB  # noqa: WPS433
    from core.atlas_permissions import PermissionPolicy as _ChatPermissionPolicy  # noqa: WPS433
    _chat_db = _ChatAtlasDB()
    register_chat_routes(
        app,
        db=_chat_db,
        bridge=bridge,
        permissions=_ChatPermissionPolicy(_chat_db),
    )
    # SSOT API (/api/ssot, /api/ssot/qa, /api/ssot/qa/sessions,
    # /api/ssot/qa/answer) lives in src/atlas_api_ssot.py.
    from atlas_api_ssot import register_ssot_routes  # noqa: WPS433
    register_ssot_routes(
        app,
        project_root=lambda: PROJECT_ROOT,
        safe_path=_safe,
        skip_dirs=SKIP_DIRS,
        max_read_bytes=MAX_READ_BYTES,
        valid_ip_name=_valid_ip_name,
        active_ssot_ip=_active_ssot_ip,
        ssot_qa_view=_ssot_qa_view,
        ssot_qa_sessions_view=_ssot_qa_sessions_view,
        ssot_qa_path=_ssot_qa_path,
        qa_slug=_qa_slug,
        upsert_ssot_qa_items=_upsert_ssot_qa_items,
        load_ssot_state=_load_ssot_state,
        canonical_session_string=_canonical_session_string,
        normalize_session_name=normalize_session_name,
        append_session_message=_append_session_message,
        bridge=bridge,
        fs_authz=_fs_authz,
    )

    # ── Sessions API (/api/session*, /api/sessions*) ────────────────
    # Routes live in src/atlas_api_sessions.py. Inject runtime callables
    # so routes see PROJECT_ROOT changes from --root and the live bridge.
    from atlas_api_sessions import register_sessions_routes  # noqa: WPS433

    # Lazy proxy for main._setup_workspace — main is imported later by
    # run_atlas_ui, so this wrapper does the import on first call.
    # Without it, /api/session/activate could only mirror env vars and
    # the live workspace stayed pinned to whatever was last loaded by
    # main.py's chat_loop /wf handler, producing the UI(tb-gen)/backend
    # (fl-model-gen) desync the user reported.
    def _setup_workspace_proxy(name: str) -> None:
        try:
            import main as _main_mod  # type: ignore
        except ImportError:
            try:
                from src import main as _main_mod  # type: ignore
            except ImportError:
                return
        fn = getattr(_main_mod, "_setup_workspace", None)
        if callable(fn):
            fn(name)

    def _setup_session_proxy(
        session_id: str,
        *,
        mirror_env: bool = True,
        apply_main: bool = True,
    ) -> None:
        session = normalize_session_name(str(session_id or ""))
        if not session:
            return
        parts = [part for part in session.split("/") if part]
        context: AtlasContext | None = None
        if len(parts) >= 3:
            try:
                context = AtlasContext.from_session_key(
                    session,
                    atlas_root=os.environ.get("ATLAS_ROOT") or str(PROJECT_ROOT),
                )
            except ValueError:
                context = None
        owner = context.user_name if context is not None else (parts[0].strip() if len(parts) >= 1 else "")
        ip = context.ip_name if context is not None else (parts[1].strip() if len(parts) >= 2 else "")
        workflow = context.workflow if context is not None else (parts[2].strip() if len(parts) >= 3 else "")
        _atlas_active_session_cv.set(session)
        if ip:
            _atlas_active_ip_cv.set(ip)
        if mirror_env:
            if context is not None:
                os.environ.update(context.export_env())
            else:
                os.environ["ATLAS_ACTIVE_SESSION"] = session
                if owner:
                    os.environ["ATLAS_MEMORY_USER"] = owner
                    os.environ["ATLAS_DEFAULT_SESSION_ID"] = owner
                if ip:
                    os.environ["ATLAS_ACTIVE_IP"] = ip
                if workflow:
                    os.environ["ATLAS_DEFAULT_WORKFLOW"] = workflow
                    os.environ["ACTIVE_WORKSPACE"] = workflow
        if not apply_main:
            return
        try:
            import main as _main_mod  # type: ignore
        except ImportError:
            try:
                from src import main as _main_mod  # type: ignore
            except ImportError:
                _main_mod = None
        fn = getattr(_main_mod, "_setup_session", None) if _main_mod is not None else None
        if callable(fn):
            fn(session)
        else:
            from core.session_setup import setup_session as _shared_setup_session
            _shared_setup_session(session)
        if mirror_env:
            os.environ["ATLAS_SESSION_APPLIED"] = session

    # Build the admin-check BEFORE registering session routes so the interactive
    # worker status endpoint can gate its all-owner view behind it (Wave-3 H8).
    # The richer _admin_required(request) closure defined below for /api/admin/*
    # reuses the same core.atlas_auth predicates so both stay in lock-step.
    from core.atlas_auth import (
        is_admin_user as _is_admin_user_for_sessions,
        is_local_admin_mode as _is_local_admin_mode_for_sessions,
    )

    def _sessions_admin_check(request: Request) -> bool:
        if _is_local_admin_mode_for_sessions():
            return True
        return bool(_is_admin_user_for_sessions(request.scope.get("user") or {}))

    register_sessions_routes(
        app,
        project_root=lambda: PROJECT_ROOT,
        normalize_session_name=normalize_session_name,
        active_session_value=_active_session_value,
        atlas_active_session_cv=_atlas_active_session_cv,
        atlas_active_ip_cv=_atlas_active_ip_cv,
        bridge=bridge,
        get_jobs_state=_get_jobs_state,
        atlas_db_factory=AtlasDB,
        setup_session=_setup_session_proxy,
        setup_workspace=_setup_workspace_proxy,
        admin_check=_sessions_admin_check,
    )

    # ── Admin endpoints ──────────────────────────────────────────
    from core.atlas_auth import admin_auth_status, is_admin_user, is_local_admin_mode

    def _admin_required(request: Request) -> Optional[dict]:
        user = request.scope.get("user") or {}
        if is_local_admin_mode():
            return {
                "id": user.get("id") or "local-admin",
                "username": user.get("username") or "local-admin",
                "role": "admin",
            }
        return user if is_admin_user(user) else None

    def _admin_denied(request: Request) -> JSONResponse:
        if request.scope.get("user"):
            return JSONResponse({"error": "Admin role required"}, status_code=403)
        return JSONResponse({"error": "Admin login required"}, status_code=401)

    def _admin_runtime_payload() -> dict[str, Any]:
        try:
            from atlas_api_jobs import worker_runtime_snapshot  # noqa: WPS433
        except ImportError:
            from src.atlas_api_jobs import worker_runtime_snapshot  # type: ignore  # noqa: WPS433

        override_ref = _scm_ui_override_ref()
        override_local = _scm_ui_override_local_path()
        scm_override: dict[str, Any] = {
            "enabled": bool(override_ref),
            "kind": "remote" if _scm_ui_override_is_url(override_ref) else ("local" if override_ref else ""),
            "ref": override_ref,
        }
        if override_local is not None:
            scm_override.update({
                "path": str(override_local),
                "exists": override_local.is_file(),
                "mtime": override_local.stat().st_mtime if override_local.is_file() else None,
            })
        return {
            "worker_runtime": worker_runtime_snapshot(PROJECT_ROOT),
            "scm": {
                "provider": configured_scm_provider(),
                "ui_override": scm_override,
            },
            "atlas": {
                "run_mode": os.environ.get("ATLAS_RUN_MODE", "engineering"),
                "exec_mode": _current_atlas_exec_mode(),
                "multi_user": os.environ.get("ATLAS_MULTI_USER", "1"),
                "multi_user_proc": os.environ.get("ATLAS_MULTI_USER_PROC", "1"),
            },
        }

    @app.get("/api/admin/auth/status")
    async def api_admin_auth_status(request: Request):
        try:
            with AtlasDB() as db:
                return JSONResponse(admin_auth_status(db, request.scope.get("user")))
        except Exception as exc:
            return JSONResponse({"error": str(exc)}, status_code=500)

    @app.get("/api/admin/users")
    async def api_admin_users(request: Request):
        if _admin_required(request) is None:
            return _admin_denied(request)
        try:
            with AtlasDB() as db:
                users = db.list_all_users()
                counts = db.count_sessions_by_user()
                active_by_user: dict[str, dict] = {}
                for session in db.list_all_sessions():
                    user_id = str(session.get("user_id") or "")
                    if not user_id or user_id in active_by_user:
                        continue
                    if str(session.get("status") or "").lower() != "active":
                        continue
                    active_by_user[user_id] = session
                for u in users:
                    u["session_count"] = counts.get(u["id"], 0)
                    active = active_by_user.get(str(u["id"] or ""))
                    u["active_session_id"] = (active or {}).get("id") or ""
                    u["active_ip"] = (active or {}).get("ip") or ""
                    u["active_workflow"] = (active or {}).get("workflow") or ""
                    u["active_session_updated_at"] = (active or {}).get("updated_at")
                    u["active_workflow_status"] = (
                        (active or {}).get("latest_workflow_status") or ""
                    )
                return JSONResponse({"users": users})
        except Exception as e:
            print(f"api_admin_users error: {e}")
            return JSONResponse({"error": str(e)}, status_code=500)

    @app.get("/api/admin/sessions")
    async def api_admin_sessions(request: Request):
        if _admin_required(request) is None:
            return _admin_denied(request)
        try:
            with AtlasDB() as db:
                sessions = db.list_all_sessions()
                return JSONResponse({"sessions": sessions})
        except Exception as e:
            print(f"api_admin_sessions error: {e}")
            return JSONResponse({"error": str(e)}, status_code=500)

    @app.get("/api/admin/ips")
    async def api_admin_ips(request: Request):
        if _admin_required(request) is None:
            return _admin_denied(request)
        try:
            with AtlasDB() as db:
                return JSONResponse({"ips": db.list_all_ip_pointers()})
        except Exception as e:
            print(f"api_admin_ips error: {e}")
            return JSONResponse({"error": str(e)}, status_code=500)

    @app.delete("/api/admin/ips/{ip_id}")
    async def api_admin_delete_ip_pointer(ip_id: str, request: Request):
        if _admin_required(request) is None:
            return _admin_denied(request)
        try:
            with AtlasDB() as db:
                if db.get_ip_block(ip_id) is None:
                    return JSONResponse({"error": "ip pointer not found"}, status_code=404)
                result = db.delete_ip_pointer(ip_id)
            return JSONResponse({**result, "filesystem_deleted": False})
        except Exception as e:
            print(f"api_admin_delete_ip_pointer error: {e}")
            return JSONResponse({"error": str(e)}, status_code=500)

    @app.delete("/api/admin/users/{user_id}")
    async def api_admin_delete_user_pointer(user_id: str, request: Request):
        admin = _admin_required(request)
        if admin is None:
            return _admin_denied(request)
        if str(admin.get("id") or "") == str(user_id or ""):
            return JSONResponse({"error": "cannot delete the signed-in admin user"}, status_code=400)
        try:
            with AtlasDB() as db:
                user = db.get_user(user_id)
                if user is None:
                    return JSONResponse({"error": "user pointer not found"}, status_code=404)
                if str(user.get("role") or "").lower() == "admin":
                    remaining = db._fetchone(
                        "SELECT COUNT(*) AS cnt FROM users WHERE role = 'admin' AND id != ?",
                        (user_id,),
                    )
                    if int(remaining["cnt"] if remaining is not None else 0) <= 0:
                        return JSONResponse({"error": "cannot delete the last admin user"}, status_code=400)
                result = db.delete_user_pointer(user_id)
            return JSONResponse({**result, "filesystem_deleted": False})
        except Exception as e:
            print(f"api_admin_delete_user_pointer error: {e}")
            return JSONResponse({"error": str(e)}, status_code=500)

    @app.get("/api/admin/usage")
    async def api_admin_usage(request: Request):
        """Per-user usage aggregation: tokens, cost, message count, model
        distribution. Joined across users / sessions / messages."""
        if _admin_required(request) is None:
            return _admin_denied(request)
        try:
            with AtlasDB() as db:
                from core.atlas_admin_usage import build_admin_usage_payload
                return JSONResponse(build_admin_usage_payload(db))
        except Exception as e:
            print(f"api_admin_usage error: {e}")
            return JSONResponse({"error": str(e)}, status_code=500)

    @app.get("/api/admin/session-flow")
    async def api_admin_session_flow(request: Request):
        """Read-only Session Flow read-model: per-session flow state, risk, input/
        LLM/worker/artifact counters, IP flow, and attribution gaps. Admin-only."""
        if _admin_required(request) is None:
            return _admin_denied(request)
        try:
            with AtlasDB() as db:
                from core.session_flow_usage import build_session_flow_payload
                qp = request.query_params
                payload = build_session_flow_payload(db, {
                    "range": (qp.get("range") or "7d").strip(),
                    "lens": (qp.get("lens") or "team_lead").strip(),
                    "risk": (qp.get("risk") or "all").strip(),
                    "ip_id": (qp.get("ip_id") or "").strip() or None,
                    "workflow": (qp.get("workflow") or "").strip() or None,
                    "user_id": (qp.get("user_id") or "").strip() or None,
                    "session_id": (qp.get("session_id") or "").strip() or None,
                    "limit": qp.get("limit"),
                    "offset": qp.get("offset"),
                })
                # Plan response shape exposes the page window as `pagination`.
                payload["pagination"] = payload.pop("limits", {})
                return JSONResponse(payload)
        except Exception as e:
            print(f"api_admin_session_flow error: {e}")
            return JSONResponse({"error": str(e)}, status_code=500)

    @app.get("/api/admin/runtime")
    async def api_admin_runtime(request: Request):
        if _admin_required(request) is None:
            return _admin_denied(request)
        try:
            return JSONResponse(_admin_runtime_payload())
        except Exception as e:
            print(f"api_admin_runtime error: {e}")
            return JSONResponse({"error": str(e)}, status_code=500)

    @app.post("/api/admin/chat")
    async def api_admin_chat(request: Request):
        """DB-backed admin Q&A for usage, memory, feedback, and user inputs."""
        if _admin_required(request) is None:
            return _admin_denied(request)
        try:
            body = await request.json()
        except Exception:
            body = {}
        question = str((body or {}).get("message") or (body or {}).get("question") or "").strip()
        if not question:
            return JSONResponse({"error": "message required"}, status_code=400)
        if len(question) > 2000:
            return JSONResponse({"error": "message too long"}, status_code=413)
        try:
            with AtlasDB() as db:
                from core.atlas_admin_chat import answer_admin_question
                return JSONResponse(answer_admin_question(db, question))
        except Exception as e:
            print(f"api_admin_chat error: {e}")
            return JSONResponse({"error": str(e)}, status_code=500)

    @app.get("/api/user/dashboard")
    async def api_user_dashboard(request: Request):
        user = request.scope.get("user")
        if not user:
            return JSONResponse({"error": "login required"}, status_code=401)
        try:
            with AtlasDB() as db:
                from core.atlas_user_dashboard import build_user_dashboard_payload
                return JSONResponse(build_user_dashboard_payload(
                    db,
                    user,
                    run_mode=os.environ.get("ATLAS_RUN_MODE", ""),
                    exec_mode=os.environ.get("ATLAS_EXEC_MODE")
                    or os.environ.get("ATLAS_DEFAULT_EXEC_MODE", ""),
                ))
        except Exception as e:
            print(f"api_user_dashboard error: {e}")
            return JSONResponse({"error": str(e)}, status_code=500)

    @app.post("/api/feedback")
    async def api_feedback_submit(request: Request):
        """Any logged-in user can drop a feedback message via /feedback
        slash. Stored in the `feedback` table and surfaced in the admin
        dashboard's Feedback tab."""
        user = request.scope.get("user")
        cookie_user = auth.get_user_from_cookie(request)
        if cookie_user is not None:
            user = cookie_user
        elif user and str(user.get("id") or "") == "local-admin":
            # Local admin mode synthesizes this identity for admin surfaces.
            # Feedback rows need a real submitter so admin attribution stays useful.
            user = None
        if not user:
            return JSONResponse({"error": "login required"}, status_code=401)
        try:
            body = await request.json()
        except Exception:
            body = {}
        content = str((body or {}).get("content") or "").strip()
        if not content:
            return JSONResponse({"error": "content required"}, status_code=400)
        if len(content) > 4000:
            return JSONResponse({"error": "content too long (max 4000 chars)"},
                                status_code=413)
        try:
            import uuid as _uuid
            email_sent = False
            with AtlasDB() as db:
                fid = _uuid.uuid4().hex
                db._execute(
                    "INSERT INTO feedback (id, user_id, content, status, created_at) "
                    "VALUES (?, ?, ?, 'open', ?)",
                    (fid, user["id"], content, time.time()),
                )
                try:
                    from core.atlas_auth import send_feedback_email
                    email_sent = send_feedback_email(db, user, fid, content)
                except Exception:
                    email_sent = False
            return JSONResponse({"ok": True, "id": fid, "email_sent": email_sent})
        except Exception as exc:
            return JSONResponse({"error": str(exc)}, status_code=500)

    @app.get("/api/admin/feedback")
    async def api_admin_feedback(request: Request):
        """Admin view of every feedback row, joined to the submitter's
        username for readability."""
        if _admin_required(request) is None:
            return _admin_denied(request)
        try:
            with AtlasDB() as db:
                rows = db._fetchall(
                    "SELECT f.id, f.user_id, u.username, f.content, f.status, "
                    "       f.created_at, f.resolved_at, f.resolved_by, f.notes "
                    "  FROM feedback f "
                    "  LEFT JOIN users u ON u.id = f.user_id "
                    " ORDER BY f.created_at DESC"
                )
                items = [dict(r) for r in rows]
            return JSONResponse({"feedback": items})
        except Exception as exc:
            return JSONResponse({"error": str(exc)}, status_code=500)

    @app.post("/api/admin/feedback/{fid}/resolve")
    async def api_admin_feedback_resolve(fid: str, request: Request):
        """Mark a feedback item resolved. Body: {notes: str (optional)}."""
        admin = _admin_required(request)
        if admin is None:
            return _admin_denied(request)
        try:
            body = await request.json()
        except Exception:
            body = {}
        notes = str((body or {}).get("notes") or "").strip()
        try:
            with AtlasDB() as db:
                db._execute(
                    "UPDATE feedback SET status = 'resolved', resolved_at = ?, "
                    "       resolved_by = ?, notes = ? WHERE id = ?",
                    (time.time(), admin.get("username", ""), notes, fid),
                )
            return JSONResponse({"ok": True})
        except Exception as exc:
            return JSONResponse({"error": str(exc)}, status_code=500)

    @app.delete("/api/admin/sessions/{session_id}")
    async def api_admin_delete_session(session_id: str, request: Request):
        if _admin_required(request) is None:
            return _admin_denied(request)
        # Propagate the runtime delete outcome (review #2 gap1): pending runtime
        # queue without ?force=1 -> 409 / deleted=False (no orphaned runtime file).
        from atlas_session_delete import force_delete_requested, session_delete_response
        try:
            with AtlasDB() as db:
                if db.get_session(session_id) is None:
                    return JSONResponse({"error": "session not found"}, status_code=404)
                result = db.delete_session(session_id, force=force_delete_requested(request))
                return session_delete_response(result)
        except Exception as e:
            print(f"api_admin_delete_session error: {e}")
            return JSONResponse({"error": str(e)}, status_code=500)

    @app.get("/api/admin/permissions")
    async def api_admin_permissions(request: Request):
        if _admin_required(request) is None:
            return _admin_denied(request)
        try:
            with AtlasDB() as db:
                rows = db._fetchall(
                    "SELECT p.id, p.ip_id, i.ip_name, p.grantee_user_id, u.username, "
                    "       p.granted_by_user_id, gu.username AS granted_by_username, "
                    "       p.permission, p.created_at, p.expires_at "
                    "  FROM ip_permissions p "
                    "  LEFT JOIN ip_blocks i ON i.id = p.ip_id "
                    "  LEFT JOIN users u ON u.id = p.grantee_user_id "
                    "  LEFT JOIN users gu ON gu.id = p.granted_by_user_id "
                    " ORDER BY p.created_at DESC"
                )
            return JSONResponse({"permissions": [dict(r) for r in rows]})
        except Exception as exc:
            return JSONResponse({"error": str(exc)}, status_code=500)

    @app.get("/api/admin/permissions/options")
    async def api_admin_permissions_options(request: Request):
        if _admin_required(request) is None:
            return _admin_denied(request)
        try:
            with AtlasDB() as db:
                ip_rows = db._fetchall(
                    "SELECT i.id, i.ip_name, w.name AS workspace_name, w.owner_user_id, "
                    "       ou.username AS owner_username "
                    "  FROM ip_blocks i "
                    "  LEFT JOIN workspaces w ON w.id = i.workspace_id "
                    "  LEFT JOIN users ou ON ou.id = w.owner_user_id "
                    " ORDER BY i.ip_name"
                )
                user_rows = db._fetchall(
                    "SELECT id, username, display_name, role FROM users ORDER BY username"
                )
            return JSONResponse({
                "ips": [dict(r) for r in ip_rows],
                "users": [dict(r) for r in user_rows],
                "levels": ["view", "import", "write", "admin"],
            })
        except Exception as exc:
            return JSONResponse({"error": str(exc)}, status_code=500)

    @app.post("/api/admin/permissions")
    async def api_admin_permissions_grant(request: Request):
        admin = _admin_required(request)
        if admin is None:
            return _admin_denied(request)
        try:
            body = await request.json()
        except Exception:
            body = {}
        ip_id = str((body or {}).get("ip_id") or "").strip()
        grantee = str((body or {}).get("grantee_user_id") or "").strip()
        permission = str((body or {}).get("permission") or "").strip().lower()
        expires_at = (body or {}).get("expires_at")
        if not ip_id or not grantee or permission not in {"view", "import", "write", "admin"}:
            return JSONResponse({"error": "ip_id, grantee_user_id, permission required"}, status_code=400)
        try:
            with AtlasDB() as db:
                row = db.grant_ip_permission(
                    ip_id=ip_id,
                    grantee_user_id=grantee,
                    permission=permission,
                    granted_by_user_id=admin.get("id") or "",
                    expires_at=expires_at if expires_at else None,
                )
            return JSONResponse({"permission": row})
        except ValueError as exc:
            return JSONResponse({"error": str(exc)}, status_code=400)
        except Exception as exc:
            return JSONResponse({"error": str(exc)}, status_code=500)

    @app.delete("/api/admin/permissions")
    async def api_admin_permissions_revoke(request: Request):
        if _admin_required(request) is None:
            return _admin_denied(request)
        try:
            body = await request.json()
        except Exception:
            body = {}
        ip_id = str((body or {}).get("ip_id") or "").strip()
        grantee = str((body or {}).get("grantee_user_id") or "").strip()
        permission = (body or {}).get("permission")
        if not ip_id or not grantee:
            return JSONResponse({"error": "ip_id and grantee_user_id required"}, status_code=400)
        try:
            with AtlasDB() as db:
                removed = db.revoke_ip_permission(ip_id, grantee, permission)
            return JSONResponse({"revoked": removed})
        except Exception as exc:
            return JSONResponse({"error": str(exc)}, status_code=500)

    @app.get("/api/admin/db/tables")
    async def api_admin_db_tables(request: Request):
        if _admin_required(request) is None:
            return _admin_denied(request)
        from core.atlas_admin_db import list_tables
        try:
            with AtlasDB() as db:
                return JSONResponse(list_tables(db))
        except Exception as exc:
            return JSONResponse({"error": str(exc)}, status_code=500)

    @app.get("/api/admin/db/preview")
    async def api_admin_db_preview(request: Request, per_table: int = 3):
        if _admin_required(request) is None:
            return _admin_denied(request)
        from core.atlas_admin_db import preview_all
        try:
            with AtlasDB() as db:
                return JSONResponse(preview_all(db, per_table=per_table))
        except Exception as exc:
            return JSONResponse({"error": str(exc)}, status_code=500)

    @app.get("/api/admin/db/table/{name}")
    async def api_admin_db_table(name: str, request: Request,
                                 limit: int = 50, offset: int = 0,
                                 order: str = "desc"):
        if _admin_required(request) is None:
            return _admin_denied(request)
        from core.atlas_admin_db import read_table
        try:
            with AtlasDB() as db:
                payload, err = read_table(db, name, limit=limit, offset=offset, order=order)
        except Exception as exc:
            return JSONResponse({"error": str(exc)}, status_code=500)
        if err:
            return JSONResponse({"error": err}, status_code=400)
        return JSONResponse(payload)

    @app.get("/api/admin/db/runtime/{session_uid}")
    async def api_admin_db_runtime(session_uid: str, request: Request,
                                   table: str = "", limit: int = 50,
                                   offset: int = 0, order: str = "desc"):
        """Inspect ONE session's per-session runtime DB (mirror of atlas_admin).

        Uses the SAME shared hardened resolver as the standalone admin server so
        the security policy cannot drift between the two copies (plan §2.11 / R24):
        session_uid ONLY (path-like rejected), manifest-resolved + containment-
        checked under ATLAS_RUNTIME_DB_ROOT, ownership enforced even under
        local-admin bypass, and NO filesystem path returned.
        """
        admin = _admin_required(request)
        if admin is None:
            return _admin_denied(request)
        from core.atlas_admin_db import inspect_runtime_table, RuntimeInspectError
        user = request.scope.get("user") or {}
        try:
            with AtlasDB() as db:
                payload = inspect_runtime_table(
                    db,
                    session_uid=session_uid,
                    requesting_user_id=str(user.get("id") or admin.get("id") or ""),
                    requesting_username=str(user.get("username") or admin.get("username") or ""),
                    is_admin=True,
                    table=(table or None),
                    limit=limit,
                    offset=offset,
                    order=order,
                )
            return JSONResponse(payload)
        except RuntimeInspectError as exc:
            return JSONResponse({"error": exc.message}, status_code=exc.status)
        except Exception:
            # Never echo str(exc): an unexpected error message could leak a
            # filesystem path. Path-free generic 500 (mirrors atlas_admin, R24).
            return JSONResponse({"error": "runtime inspect failed"}, status_code=500)

    @app.get("/admin")
    async def admin_page(request: Request):
        html = (FRONTEND / "admin.html").read_text(encoding="utf-8")

        def _inline_script(match):
            attrs = match.group("attrs")
            src = match.group("src").split("?", 1)[0]
            if not src.endswith((".jsx", ".js")):
                return match.group(0)
            path = (FRONTEND / src).resolve()
            try:
                path.relative_to(FRONTEND.resolve())
            except Exception:
                return match.group(0)
            if not path.is_file():
                return match.group(0)
            code = path.read_text(encoding="utf-8")
            if "data-filename" not in attrs:
                attrs = f'data-filename="{html_lib.escape(src, quote=True)}" {attrs}'
            return f'<script type="text/babel" {attrs}>{code.rstrip()}\n//# sourceURL={src}</script>'

        html = re.sub(
            r'<script\s+type="text/babel"\s+(?P<attrs>[^>]*?)src="(?P<src>[^"]+)"[^>]*>\s*</script>',
            _inline_script,
            html,
        )
        return HTMLResponse(html)

    # NOTE: WebSocket endpoint is registered via Starlette's WebSocketRoute
    # (added to app.router.routes below) instead of the @app.websocket
    # decorator. The decorator routes through FastAPI's dependency-injection
    # layer, which can't resolve the `websocket: WebSocket` annotation when
    # `from __future__ import annotations` is active (PEP 563 turns all
    # annotations into strings) and rejects the handshake with HTTP 403.
    # Starlette's WebSocketRoute talks to the function directly and ignores
    # parameter annotations entirely.
    async def ws_agent(websocket: WebSocket):
        try:
            await websocket.accept()
        except Exception as exc:
            if _is_websocket_disconnect(exc):
                return
            raise
        session_id = websocket.query_params.get("session_id", "")
        _multi_user = _multi_user_enabled()

        class _WebSocketCookieRequest:
            def __init__(self, cookies: dict):
                self.cookies = cookies

        cookies = getattr(websocket, "cookies", None)
        if cookies is None:
            cookies = websocket.scope.get("cookies") or {}
        user = auth.get_user_from_cookie(_WebSocketCookieRequest(cookies))
        if user is None:
            try:
                from core.atlas_auth import is_local_admin_mode, local_admin_user
            except Exception:
                try:
                    from atlas_auth import is_local_admin_mode, local_admin_user  # type: ignore
                except Exception:
                    is_local_admin_mode = None  # type: ignore
                    local_admin_user = None  # type: ignore
            if callable(is_local_admin_mode) and is_local_admin_mode() and callable(local_admin_user):
                user = local_admin_user()
        if user is None:
            await _close_websocket_quietly(websocket, code=1008, reason="unauthenticated")
            return

        username = normalize_session_name(str(user.get("username") or ""))

        def _authorize_ws_session(raw_session: str) -> str | None:
            normalized = normalize_session_name(str(raw_session or ""))
            if not normalized or normalized == "default":
                normalized = f"{username}/default/default/default" if username else "default/default/default/default"
            elif username and normalized == username:
                normalized = f"{username}/default/default/default"
            else:
                parts = [part for part in normalized.split("/") if part]
                if len(parts) == 2 and username and parts[0] == username:
                    normalized = f"{parts[0]}/default/{parts[1]}/default"
                elif len(parts) == 3 and username and parts[0] == username:
                    normalized = f"{parts[0]}/default/{parts[1]}/{parts[2]}"
            owner = normalized.split("/", 1)[0]
            if _multi_user and username and owner != username:
                with AtlasDB() as db:
                    owned = db.get_session(normalized)
                if not (owned and owned.get("user_id") == user["id"]):
                    return None
            return normalized

        def _prompt_target_session(raw_session: str, msg: dict) -> str | None:
            """Canonicalize prompt routing from the visible browser target.

            The browser can briefly hold a stale ACTIVE_SESSION while the top
            IP/workflow controls already show the new target.  The prompt
            message carries explicit `ip` and `workflow`; prefer those so we
            do not spawn queues like `<user>/<ip>/default` when the user is
            visibly on `<user>/<ip>/ssot-gen`.
            """
            base = _authorize_ws_session(raw_session)
            if not base:
                return None
            parts = [part for part in normalize_session_name(base).split("/") if part]
            while len(parts) < 3:
                parts.append("default")
            ip = normalize_session_name(str(
                msg.get("ip")
                or msg.get("scope")
                or msg.get("ip_id")
                or ""
            ))
            workflow = normalize_session_name(str(
                msg.get("workflow")
                or msg.get("workspace")
                or msg.get("active_workflow")
                or ""
            ))
            if ip and ip not in {"default", "soc", "user"}:
                parts[2 if len(parts) >= 4 else 1] = ip
            if workflow and workflow not in {"user", "soc"}:
                parts[3 if len(parts) >= 4 else 2] = workflow
            return _authorize_ws_session("/".join(parts[:4] if len(parts) >= 4 else parts[:3]))

        def _prompt_images_for_bridge(msg: dict) -> list[dict[str, str]]:
            return [
                {"image_url": image.image_url, "detail": image.detail}
                for image in normalize_prompt_images(msg.get("images"))
            ]

        # Identity-driven default: empty / legacy "default" session_id
        # collapses to the user's default namespace. Full
        # <user>/<ip>/<workflow> namespaces are allowed for that user, so
        # two browser tabs on different IP/workflow views do not receive
        # each other's backend stream.
        session_id = _authorize_ws_session(session_id)
        if session_id is None:
            await _close_websocket_quietly(websocket, code=1008, reason="forbidden")
            return
        bridge.bind_client(websocket, session_id)
        try:
            if bridge._using_processes():
                _setup_session_proxy(session_id, mirror_env=False, apply_main=False)
            else:
                _setup_session_proxy(session_id)
        except Exception:
            pass
        _ensure_broadcaster()
        # Greeting — surface user-tunable layout settings so the frontend
        # can pick its center-column shape (classic vs tabbed Chat/Preview/Q&A).
        _center_layout = "classic"
        _chat_feed_summary = True
        try:
            import src.config as _cfg_hello
            # /healthz already calls reload_env on every poll, so the
            # WS-connect handshake doesn't need to repeat the sync .env
            # read. Reading the cached attribute values is plenty.
            _center_layout = getattr(_cfg_hello, "ATLAS_CENTER_LAYOUT", "classic")
            _chat_feed_summary = bool(getattr(_cfg_hello, "ATLAS_CHAT_FEED_SUMMARY", True))
        except Exception:
            pass
        try:
            try:
                _session_running = bool(getattr(bridge.get_session(session_id), "agent_running", False))
            except Exception:
                _session_running = bool(bridge.agent_running)
            await websocket.send_json({"type": "hello", "frontend": "atlas",
                                        "running": _session_running,
                                        "center_layout": _center_layout,
                                        "chat_feed_summary": _chat_feed_summary})
            for pending_event in bridge.session_pending_ask_user_events(session_id):
                await websocket.send_json(pending_event)
        except Exception as exc:
            if not _is_websocket_disconnect(exc):
                raise
            bridge.unbind_client(websocket)
            return
        try:
            while True:
                data = await websocket.receive_text()
                try:
                    msg = json.loads(data)
                except Exception:
                    continue
                session = bridge.get_client_session(websocket)
                if session is None:
                    continue
                set_atlas_bridge_session_id(session.session_id)
                t = msg.get("type")
                if t in ("session_switch", "client_session_switch"):
                    _session_raw = str(
                        msg.get("session_id")
                        or msg.get("session")
                        or msg.get("namespace")
                        or ""
                    ).strip()
                    _session = _authorize_ws_session(_session_raw)
                    if not _session:
                        await websocket.send_json({
                            "type": "error",
                            "message": f"invalid or forbidden session: {_session_raw!r}",
                        })
                        continue
                    if _session != session.session_id:
                        bridge.bind_client(websocket, _session)
                        session = bridge.get_client_session(websocket)
                        if session is None:
                            continue
                        set_atlas_bridge_session_id(session.session_id)
                        try:
                            if bridge._using_processes():
                                _setup_session_proxy(_session, mirror_env=False, apply_main=False)
                            else:
                                _setup_session_proxy(_session)
                        except Exception as exc:
                            await websocket.send_json({
                                "type": "error",
                                "message": f"session setup failed: {exc}",
                            })
                            continue
                    await websocket.send_json({
                        "type": "session_switched",
                        "session_id": session.session_id,
                    })
                    continue
                if t in ("prompt", "send"):
                    _images = _prompt_images_for_bridge(msg)
                    if not (msg.get("text") or _images):
                        continue
                    _txt = str(msg.get("text") or "").strip()
                    _msg_id = str(msg.get("msg_id") or "").strip()
                    _txt_preview = (
                        str(msg.get("text") or "")[:80].replace("\n", " ")
                        or ("[image]" if _images else "")
                    )
                    _session_raw = str(msg.get("session") or "").strip()
                    _session = session.session_id
                    if _session_raw:
                        _session = _prompt_target_session(_session_raw, msg)
                        if not _session:
                            for frame in _prompt_ack_frames(
                                msg_id=_msg_id,
                                text_preview=_txt_preview,
                                session_id=session.session_id,
                                ok=False,
                                error=f"invalid or forbidden session: {_session_raw!r}",
                            ):
                                await websocket.send_json(frame)
                            continue
                        if _session != session.session_id:
                            bridge.bind_client(websocket, _session)
                            session = bridge.get_client_session(websocket)
                            if session is None:
                                continue
                            set_atlas_bridge_session_id(session.session_id)
                    # Idempotent submit + ack:
                    # The frontend retransmits a prompt with the same
                    # msg_id if it doesn't see an `agent_received`
                    # ack within ~3s. Send the ack directly on this
                    # websocket before any session setup/spawn work so
                    # slow worker startup cannot trigger a duplicate send.

                    async def _send_prompt_acceptance(
                        *,
                        ok: bool,
                        queued: bool = False,
                        handled: str = "",
                        duplicate: bool = False,
                        error: str = "",
                    ) -> None:
                        # AC-2: emit agent_received (transport ack) ONLY on a
                        # genuine accept (ok and not duplicate), immediately
                        # before agent_accepted (delivery ack). A dropped
                        # prompt (ok:false) sends NO agent_received, so the
                        # frontend transport-retry timer fires and the prompt
                        # auto-lands once the worker is warm. _prompt_ack_frames
                        # encodes this ordering/conditionality contract.
                        for frame in _prompt_ack_frames(
                            msg_id=_msg_id,
                            text_preview=_txt_preview,
                            session_id=session.session_id,
                            ok=ok,
                            queued=queued,
                            handled=handled,
                            duplicate=duplicate,
                            error=error,
                        ):
                            await websocket.send_json(frame)

                    async def _accept_handled(kind: str) -> None:
                        if _msg_id:
                            session.mark_msg_id_seen(_msg_id)
                        try:
                            await _send_prompt_acceptance(ok=True, handled=kind)
                        except Exception as exc:
                            if not _is_websocket_disconnect(exc):
                                session.emit("error", message=f"acceptance ack failed: {exc}")

                    async def _accept_queued(kind: str = "") -> None:
                        import time as _t
                        _t0 = _t.monotonic()
                        # Wave-3 C2/H6: use the parallel result variant so a
                        # capacity refusal surfaces as agent_accepted{ok:false,
                        # error:"capacity_wait"} with NO agent_received — keeping
                        # the frontend transport-retry/hold path armed. The bool
                        # submit_prompt_for_session is intentionally left intact
                        # for non-WS callers/tests.
                        _result = bridge.submit_prompt_result_for_session(
                            session.session_id, _txt, images=_images
                        )
                        delivered = bool(getattr(_result, "ok", False))
                        _submit_ms = (_t.monotonic() - _t0) * 1000.0
                        if delivered and _msg_id:
                            session.mark_msg_id_seen(_msg_id)
                        _t1 = _t.monotonic()
                        try:
                            await _send_prompt_acceptance(
                                ok=delivered,
                                queued=delivered,
                                handled=kind,
                                error="" if delivered else (getattr(_result, "error", "") or "input was not delivered to the agent worker"),
                            )
                        except Exception as exc:
                            if not _is_websocket_disconnect(exc):
                                session.emit("error", message=f"acceptance ack failed: {exc}")
                        _ack_ms = (_t.monotonic() - _t1) * 1000.0
                        # (a) diagnostic: the frontend declares the prompt failed
                        # if agent_accepted does not arrive within ~7s. Log when
                        # the submit (spawn + DB enqueue) or the ack send is slow
                        # so a production "Input not confirmed" is traceable to
                        # its real cause (DB write-lock contention vs WS
                        # backpressure) instead of being guessed at.
                        if _submit_ms + _ack_ms > 1000.0:
                            print(
                                f"[prompt-latency] session={session.session_id} "
                                f"msg_id={_msg_id} submit={_submit_ms:.0f}ms "
                                f"ack_send={_ack_ms:.0f}ms delivered={bool(delivered)}",
                                flush=True,
                            )

                    # AC-2: the up-front, UNCONDITIONAL agent_received that
                    # used to be sent here lied about delivery and disarmed the
                    # frontend transport-retry before a live worker had accepted
                    # the prompt — the root cause of warmup/crash silent loss.
                    # agent_received is now emitted by _send_prompt_acceptance
                    # (via _prompt_ack_frames) ONLY on a real accept. The
                    # duplicate guard below still short-circuits a re-sent
                    # msg_id, and msg_id dedup keeps a raced retransmit from
                    # double-running the worker.
                    if _msg_id and session.has_msg_id(_msg_id):
                        try:
                            await _send_prompt_acceptance(
                                ok=True,
                                handled="duplicate",
                                duplicate=True,
                            )
                        except Exception as exc:
                            if not _is_websocket_disconnect(exc):
                                session.emit("error", message=f"acceptance ack failed: {exc}")
                        continue
                    if _session_raw:
                        try:
                            if bridge._using_processes():
                                _setup_session_proxy(_session, mirror_env=False, apply_main=False)
                            else:
                                _setup_session_proxy(_session)
                        except Exception as exc:
                            try:
                                await _send_prompt_acceptance(
                                    ok=False,
                                    error=f"session setup failed: {exc}",
                                )
                            except Exception as ack_exc:
                                if not _is_websocket_disconnect(ack_exc):
                                    session.emit("error", message=f"acceptance ack failed: {ack_exc}")
                            session.emit("error", message=f"session setup failed: {exc}")
                            continue
                    _txt, _ = _apply_locked_truth_draft_overlay(
                        PROJECT_ROOT,
                        session.session_id,
                        msg,
                        _txt,
                    )
                    import os as _os
                    # ── Mode-flip slashes need to apply mid-loop ──
                    # `/mode normal` and `/plan` typed while the agent is
                    # running normally land in the _interrupts queue,
                    # which feeds the agent as conversational text — the
                    # slash dispatcher only runs against _inbox between
                    # turns. So agent_mode never actually flips and the
                    # agent stays trapped in plan_q forever, hitting
                    # `[Plan Mode] blocked` on every write_file. We
                    # intercept those four canonical forms here and set
                    # AGENT_MODE_OVERRIDE in the environment; react_loop
                    # reads it at the top of each iteration.
                    _low = _txt.lower()
                    if _handle_bang_shell_command(_txt, client_session=session):
                        await _accept_handled("bang")
                        continue
                    if _txt.startswith("/"):
                        if _handle_new_ip_command(_txt, client_session=session):
                            await _accept_handled("new_ip")
                            continue
                        if _handle_ip_command(_txt, client_session=session):
                            await _accept_handled("ip")
                            continue
                        if _handle_session_command(_txt, client_session=session):
                            await _accept_handled("session")
                            continue
                        if _handle_import_command(_txt, client_session=session):
                            await _accept_handled("import")
                            continue
                        if _handle_grill_me_command(_txt, client_session=session):
                            await _accept_handled("grill")
                            continue
                        if _handle_approval_command(_txt, client_session=session):
                            await _accept_handled("approval")
                            continue
                        if _handle_verify_ssot_command(_txt, client_session=session):
                            await _accept_handled("verify_ssot")
                            continue
                        if _handle_repair_ssot_command(_txt, client_session=session):
                            await _accept_handled("repair_ssot")
                            continue
                        if _handle_repair_rtl_command(_txt, client_session=session):
                            await _accept_handled("repair_rtl")
                            continue
                        if _handle_repair_equiv_command(_txt, client_session=session):
                            await _accept_handled("repair_equiv")
                            continue
                        if _handle_to_ssot_gate(_txt, client_session=session):
                            await _accept_handled("to_ssot")
                            continue
                        if _run_stage_command(_txt, client_session=session):
                            await _accept_handled("stage")
                            continue
                        if _handle_refresh_wiki_command(_txt, client_session=session):
                            await _accept_handled("refresh_wiki")
                            continue
                        if _execute_generic_slash_command(_txt, session):
                            await _accept_handled("slash")
                            continue
                        if _low in ("/plan", "/mode plan", "/mode normal", "/normal"):
                            is_plan = _low in ("/plan", "/mode plan")
                            if is_plan:
                                _agent_mode_override_cv.set("plan_q")
                                _plan_mode_cv.set("true")
                                os.environ["AGENT_MODE_OVERRIDE"] = "plan_q"
                                os.environ["PLAN_MODE"] = "true"
                            else:
                                _agent_mode_override_cv.set("normal")
                                _plan_mode_cv.set("false")
                                os.environ["AGENT_MODE_OVERRIDE"] = "normal"
                                os.environ["PLAN_MODE"] = "false"
                                _os.environ.pop("_PLAN_TODO_WRITE_COUNT", None)
                            # Immediate UI feedback so the chat reflects the
                            # mode flip the moment the user clicks the
                            # NORMAL/PLAN pill (or types /plan), instead of
                            # waiting for the next turn boundary. Previously
                            # we deferred the banner to main.py's dispatcher
                            # which only runs between turns — clicking PLAN
                            # while the agent was idle left the chat silent
                            # until the user typed a message.
                            try:
                                if is_plan:
                                    session.emit("agent", text=(
                                        "✅ Plan mode — read-only. "
                                        "The agent will analyze and propose without mutating tools. "
                                        "Type `apply` or click NORMAL to execute."
                                    ))
                                else:
                                    session.emit("agent", text=(
                                        "✅ Normal mode — tools enabled."
                                    ))
                                session.emit("flush")
                            except Exception:
                                pass
                            # Two-pronged dispatch:
                            # (1) AGENT_MODE_OVERRIDE handles the MID-LOOP case
                            #     (agent currently iterating; react_loop top
                            #     pops it on the next pass and flips local
                            #     agent_mode for parallel_executor).
                            # (2) submit_prompt forwards the slash so main.py's
                            #     dispatcher can fire AGENT_MODE:normal/plan
                            #     when the loop is IDLE — that path is what
                            #     keeps main.py's local agent_mode + the
                            #     system prompt in messages[0] consistent
                            #     across turns. Without this submit, the
                            #     UI's "● NORMAL" pill could click without
                            #     ever telling main.py to flip — desync.
                            await _accept_queued("mode")
                            continue
                    await _accept_queued()
                elif t == "interrupt":
                    bridge.submit_interrupt_for_session(session.session_id, msg.get("text", ""))
                elif t == "answer" and msg.get("flow_id"):
                    accepted = bridge.submit_answer_for_session(session.session_id, msg["flow_id"], msg)
                    if accepted:
                        session.emit("agent_state", running=True)
                    else:
                        session.emit(
                            "error",
                            message=(
                                "answer rejected: no pending ask_user flow "
                                f"for {msg['flow_id']}"
                            ),
                        )
                elif t == "stop":
                    # Esc from the UI — abort the current iteration.
                    # Surface the stop on the chat feed so the user sees
                    # the backend acknowledged the ESC immediately,
                    # before the agent thread gets to its next poll.
                    session.emit("token", text="\n⏹  stop received\n")
                    session.emit("flush")
                    bridge.request_stop_for_session(session.session_id)
                    session.emit("agent_state", running=False)
                elif t == "shutdown":
                    # Exit button — terminate only this session's worker.
                    # Atlas UI is the shared backend server, so keep it alive.
                    session.emit("token", text="\n⏹  worker exit requested\n")
                    session.emit("flush")
                    bridge.exit_session(session.session_id)
                # Other types (e.g. run_stage, tool_call) can be wired later
        except Exception as exc:
            if not _is_websocket_disconnect(exc):
                raise
        finally:
            bridge.unbind_client(websocket)

    from core.atlas_auth import GuestAuth, AuthMiddleware, create_auth_endpoints
    auth = GuestAuth(AtlasDB())
    app.state.auth = auth
    app.add_middleware(AuthMiddleware, auth=auth)
    create_auth_endpoints(app, auth)

    # Register the WebSocket endpoint via Starlette so we don't go through
    # FastAPI's DI layer (see the long comment above the ws_agent definition).
    app.router.routes.append(WebSocketRoute("/ws/agent", ws_agent))

    # Static assets — jsx, css, js, fonts (registered LAST so it doesn't
    # shadow the explicit routes above). Disable client-side caching so
    # a normal page refresh always picks up new JSX/CSS.
    import mimetypes as _mimetypes
    _JS_CONTENT_TYPE = "application/javascript; charset=utf-8"
    _MIME_OVERRIDES: dict[str, str] = {
        ".js": _JS_CONTENT_TYPE,
        ".jsx": _JS_CONTENT_TYPE,
        ".mjs": _JS_CONTENT_TYPE,
    }

    class _NoCacheStatic(StaticFiles):
        async def get_response(self, path, scope):
            resp = await super().get_response(path, scope)
            resp.headers["Cache-Control"] = "no-store, max-age=0"
            p = str(path)
            ext = p[p.rfind("."):] if "." in p else ""
            if ext in _MIME_OVERRIDES:
                resp.headers["Content-Type"] = _MIME_OVERRIDES[ext]
            else:
                guessed = _mimetypes.guess_type(p)[0]
                if guessed and resp.headers.get("Content-Type", "").startswith("text/plain"):
                    resp.headers["Content-Type"] = guessed
            return resp

    # Vite build assets (ATLAS_FRONTEND_MODE=vite). Guarded so it only
    # mounts when a real built dist/assets dir exists; harmless no-op in the
    # default legacy path. Mounted BEFORE the catch-all "/" static mount so
    # /assets/* resolves to the built bundle. No existing /assets route, so
    # this does not shadow anything.
    _vite_assets_dir = FRONTEND / "dist" / "assets"
    if _vite_assets_dir.is_dir():
        app.mount("/assets", _NoCacheStatic(directory=str(_vite_assets_dir), html=False),
                  name="vite-assets")

    app.mount("/", _NoCacheStatic(directory=str(FRONTEND), html=False),
              name="atlas-static")

    app.state.bridge = bridge

    # Expose SSOT-QA helpers so run_atlas_ui's nested callbacks
    # (_ask_user_cb, _record_ssot_qa_cb) can reach them across function
    # scopes — these helpers live in create_app's local closure and were
    # otherwise invisible from run_atlas_ui, causing NameError at the
    # first ask_user / record_ssot_qa invocation.
    app.state.active_ssot_qa_context = _active_ssot_qa_context
    app.state.ssot_q_pairs_from_questions = _ssot_q_pairs_from_questions
    app.state.upsert_ssot_qa_items = _upsert_ssot_qa_items
    app.state.load_ssot_state = _load_ssot_state
    app.state.valid_ip_name = _valid_ip_name
    app.state.status_group = _status_group
    app.state.answer_text = _answer_text
    return app


# ── Entry point ────────────────────────────────────────────────────
def _local_ipv4_addresses() -> list[str]:
    """Return local IPv4 bind candidates for clearer startup errors."""

    addrs: set[str] = {"127.0.0.1"}
    try:
        for info in socket.getaddrinfo(socket.gethostname(), None, socket.AF_INET):
            ip = str(info[4][0] or "").strip()
            if ip:
                addrs.add(ip)
    except Exception:
        pass
    try:
        proc = subprocess.run(
            ["ifconfig"],
            check=False,
            capture_output=True,
            text=True,
            timeout=1,
        )
        for ip in re.findall(r"\binet\s+(\d+\.\d+\.\d+\.\d+)\b", proc.stdout or ""):
            addrs.add(ip)
    except Exception:
        pass
    return sorted(addrs, key=lambda ip: (ip.startswith("127."), ip))


def _lan_ipv4_addresses() -> list[str]:
    return [ip for ip in _local_ipv4_addresses() if not ip.startswith("127.")]


def _bind_help(host: str, port: int) -> str:
    ips = _local_ipv4_addresses()
    lan = _lan_ipv4_addresses()
    options = ["127.0.0.1", "0.0.0.0", *lan]
    option_text = ", ".join(dict.fromkeys(options))
    return (
        f"Cannot bind {host}:{port} because that address is not assigned to this Mac.\n"
        f"Current local IPv4 addresses: {', '.join(ips) or '(none)'}.\n"
        f"Use one of: --host {option_text}."
    )


def _assert_bind_target_available(host: str, port: int, label: str) -> None:
    bind_host = str(host or "127.0.0.1").strip() or "127.0.0.1"
    family = socket.AF_INET6 if ":" in bind_host and bind_host != "0.0.0.0" else socket.AF_INET
    sock = socket.socket(family, socket.SOCK_STREAM)
    try:
        sock.setsockopt(socket.SOL_SOCKET, socket.SO_REUSEADDR, 1)
        sock.bind((bind_host, int(port)))
    except OSError as exc:
        err = getattr(exc, "errno", None)
        if err in {errno.EADDRNOTAVAIL, 49, 99}:
            sys.exit(f"{label}: {_bind_help(bind_host, int(port))}")
        if err in {errno.EADDRINUSE, 48, 98}:
            sys.exit(f"{label}: port {port} is already in use on {bind_host}.")
        sys.exit(f"{label}: cannot bind {bind_host}:{port}: {exc}")
    finally:
        sock.close()


def _access_url(host: str, port: int, path: str = "") -> str:
    display_host = str(host or "127.0.0.1").strip() or "127.0.0.1"
    if display_host in {"0.0.0.0", "::"}:
        lan = _lan_ipv4_addresses()
        if lan:
            display_host = lan[0]
    return f"http://{display_host}:{port}{path}"


# Canonical launcher entry. Phase 4b moved `main()` into atlas_runtime
# (and atlas_ui.py:415 re-imports it for backward compat), but the
# original `if __name__ == "__main__": main()` guard never came back
# to atlas_ui.py. Result: `python3 src/atlas_ui.py …` ran the module
# top-to-bottom (imports + class/function defs + module-level constants)
# and exited without ever calling main() — no uvicorn, no listener.
# The README at frontend/atlas/README.md:37 still points users at this
# invocation, so the guard belongs here for the canonical CLI path.
if __name__ == "__main__":
    main()
