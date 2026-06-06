import base64
import importlib
import io
import sys
from pathlib import Path
from typing import Optional

import fitz
import pytest
from fastapi.testclient import TestClient


ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))
if str(ROOT / "src") not in sys.path:
    sys.path.insert(0, str(ROOT / "src"))


def _register(client: TestClient) -> None:
    response = client.post(
        "/api/auth/register",
        json={"username": "visual_tester", "password": "pw"},
    )
    assert response.status_code == 200, response.text


def _figure_pdf_bytes() -> bytes:
    doc = fitz.open()
    toc = doc.new_page(width=612, height=792)
    toc.insert_text((72, 96), "List of Figures")
    toc.insert_text((72, 128), "FIGURE 1. TEST BLOCK DIAGRAM ................................ 2")

    page = doc.new_page(width=612, height=792)
    page.insert_text((72, 96), "1.2. Block Diagram")
    page.insert_text((72, 132), "Figure 1. Test Block Diagram")
    page.draw_rect(fitz.Rect(120, 200, 260, 280))
    page.insert_text((145, 240), "AHB Slave")
    page.draw_rect(fitz.Rect(340, 200, 480, 280))
    page.insert_text((365, 240), "DMA Core")
    page.draw_line((260, 240), (340, 240))

    quiet = doc.new_page(width=612, height=792)
    quiet.insert_text((72, 96), "Appendix")
    raw = doc.write()
    doc.close()
    return raw


def _office_bytes(kind: str) -> bytes:
    buffer = io.BytesIO()
    if kind == "docx":
        from docx import Document

        doc = Document()
        doc.add_heading("Interface Timing", level=1)
        doc.add_paragraph("The visual page contains valid/ready timing evidence.")
        doc.save(buffer)
        return buffer.getvalue()
    if kind == "pptx":
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(7), Inches(1))
        box.text = "Interface Timing Slide"
        prs.save(buffer)
        return buffer.getvalue()
    raise AssertionError(f"unexpected office kind {kind}")


def test_ssot_pdf_import_adds_visual_evidence_for_rendered_figure_pages(tmp_path, monkeypatch):
    import src.atlas_ui as atlas_ui

    core_tools = importlib.import_module("core.tools")
    described: list[str] = []

    def fake_read_image(path: Optional[str] = None, prompt: str = "") -> str:
        assert path is not None
        assert "rendered document page" in prompt
        described.append(Path(path).name)
        return "visual analysis: AHB Slave connects to DMA Core."

    fake_markitdown = tmp_path / "markitdown.py"
    fake_markitdown.write_text(
        "print('# Converted PDF\\n\\nFigure 1. Test Block Diagram')\n",
        encoding="utf-8",
    )

    monkeypatch.setenv("HOME", str(tmp_path / "home"))
    monkeypatch.setenv("ATLAS_MULTI_USER", "0")
    monkeypatch.setenv("ATLAS_MULTI_USER_PROC", "0")
    monkeypatch.setenv("ATLAS_IMPORT_CONVERTER", "markitdown")
    monkeypatch.setenv("ATLAS_MARKITDOWN_PYTHON", sys.executable)
    monkeypatch.chdir(tmp_path)
    monkeypatch.setattr(atlas_ui, "PROJECT_ROOT", tmp_path)
    monkeypatch.setattr(core_tools, "read_image", fake_read_image)

    client = TestClient(atlas_ui.create_app())
    _register(client)
    response = client.post(
        "/api/ssot/import/upload",
        json={
            "ip": "mctp_assembler",
            "files": [
                {
                    "name": "figure.pdf",
                    "content_b64": base64.b64encode(_figure_pdf_bytes()).decode("ascii"),
                }
            ],
        },
    )

    assert response.status_code == 200, response.text
    payload = response.json()
    saved = payload["saved"][0]
    saved_path = tmp_path / payload["paths"][0]
    saved_text = saved_path.read_text(encoding="utf-8")

    assert payload["errors"] == []
    assert len(described) == 1
    assert saved["visual_paths"] == [f"mctp_assembler/req/imports/visual/{described[0]}"]
    assert saved["visual_count"] == 1
    assert "## Visual Evidence" in saved_text
    assert "PDF page 2" in saved_text
    assert "Figure 1. Test Block Diagram" in saved_text
    assert "visual analysis: AHB Slave connects to DMA Core." in saved_text
    assert "PDF page 1" not in saved_text


@pytest.mark.parametrize(
    ("kind", "expected_heading", "expected_prompt"),
    [
        ("docx", "DOCX page 1", "Surface: docx_page page 1"),
        ("pptx", "PPTX slide 1", "Surface: pptx_slide page 1"),
    ],
)
def test_ssot_office_import_adds_visual_evidence_after_pdf_normalization(
    tmp_path,
    monkeypatch,
    kind: str,
    expected_heading: str,
    expected_prompt: str,
):
    import src.atlas_ui as atlas_ui

    core_tools = importlib.import_module("core.tools")
    visual_evidence = importlib.import_module("core.visual_evidence")
    described: list[str] = []

    def fake_convert_office_document_to_pdf(
        document_path: Path,
        output_dir: Path,
        timeout_seconds: int = 180,
    ) -> Path:
        output_dir.mkdir(parents=True, exist_ok=True)
        pdf_path = output_dir / f"{document_path.stem}.pdf"
        doc = fitz.open()
        page = doc.new_page(width=612, height=792)
        page.insert_text((72, 96), f"{kind.upper()} visual timing surface")
        page.draw_rect(fitz.Rect(120, 200, 260, 280))
        page.insert_text((135, 240), "valid/ready")
        pdf_path.write_bytes(doc.write())
        doc.close()
        return pdf_path

    def fake_read_image(path: Optional[str] = None, prompt: str = "") -> str:
        assert path is not None
        assert expected_prompt in prompt
        described.append(Path(path).name)
        return f"visual analysis for {kind}: valid/ready timing is visible."

    fake_markitdown = tmp_path / "markitdown.py"
    fake_markitdown.write_text(
        "print('# Converted Office\\n\\nInterface timing evidence')\n",
        encoding="utf-8",
    )

    monkeypatch.setenv("HOME", str(tmp_path / "home"))
    monkeypatch.setenv("ATLAS_MULTI_USER", "0")
    monkeypatch.setenv("ATLAS_MULTI_USER_PROC", "0")
    monkeypatch.setenv("ATLAS_IMPORT_CONVERTER", "markitdown")
    monkeypatch.setenv("ATLAS_MARKITDOWN_PYTHON", sys.executable)
    monkeypatch.chdir(tmp_path)
    monkeypatch.setattr(atlas_ui, "PROJECT_ROOT", tmp_path)
    monkeypatch.setattr(core_tools, "read_image", fake_read_image)
    monkeypatch.setattr(visual_evidence, "convert_office_document_to_pdf", fake_convert_office_document_to_pdf)

    client = TestClient(atlas_ui.create_app())
    _register(client)
    response = client.post(
        "/api/ssot/import/upload",
        json={
            "ip": "mctp_assembler",
            "files": [
                {
                    "name": f"timing.{kind}",
                    "content_b64": base64.b64encode(_office_bytes(kind)).decode("ascii"),
                }
            ],
        },
    )

    assert response.status_code == 200, response.text
    payload = response.json()
    saved = payload["saved"][0]
    saved_text = (tmp_path / payload["paths"][0]).read_text(encoding="utf-8")

    assert payload["errors"] == []
    assert len(described) == 1
    assert saved["visual_count"] == 1
    assert saved["visual_paths"] == [f"mctp_assembler/req/imports/visual/{described[0]}"]
    assert f"_{kind}_{'page' if kind == 'docx' else 'slide'}_001.png" in described[0]
    assert "## Visual Evidence" in saved_text
    assert expected_heading in saved_text
    assert f"visual analysis for {kind}: valid/ready timing is visible." in saved_text
