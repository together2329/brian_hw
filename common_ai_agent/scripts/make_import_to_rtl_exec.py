#!/usr/bin/env python3
"""Executive one-pager: purpose + flow + expected impact.
Import → Deep Interview → SSOT → RTL.  Run: python3 scripts/make_import_to_rtl_exec.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY  = RGBColor(0x14, 0x1E, 0x3A)
VIOLET= RGBColor(0x6A, 0x4C, 0xC4)
BLUE  = RGBColor(0x2D, 0x6C, 0xDF)
TEAL  = RGBColor(0x17, 0x9E, 0x8A)
AMBER = RGBColor(0xE0, 0x8A, 0x1E)
GREEN = RGBColor(0x2E, 0x9E, 0x5B)
GREY  = RGBColor(0x5B, 0x63, 0x72)
LGREY = RGBColor(0x8A, 0x93, 0xA3)
LIGHT = RGBColor(0xF1, 0xF4, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK  = RGBColor(0x1E, 0x24, 0x30)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
s = prs.slides.add_slide(prs.slide_layouts[6])


def box(x, y, w, h, *, fill=None, line=None, line_w=1.0, rounded=True):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shp.shadow.inherit = False
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    return shp


def text(x, y, w, h, lines, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, (t, sz, b, c) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; r = p.add_run(); r.text = t
        r.font.size = Pt(sz); r.font.bold = b; r.font.color.rgb = c; r.font.name = "Calibri"
        p.space_after = Pt(2)
    return tb


def chevron(x, y, c=GREY):
    a = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(0.4), Inches(0.62))
    a.fill.solid(); a.fill.fore_color.rgb = c; a.line.fill.background(); a.shadow.inherit = False


# ── header: purpose-revealing title ──
box(0, 0, 13.333, 1.5, fill=NAVY, rounded=False)
box(0, 1.5, 13.333, 0.07, fill=TEAL, rounded=False)
text(0.55, 0.2, 12.3, 0.4, [("AI 반도체 IP 설계 자동화", 12, True, TEAL)])
text(0.53, 0.46, 12.3, 0.7, [("스펙 문서 → 검증된 RTL, 자동 설계 파이프라인", 27, True, WHITE)])
text(0.55, 1.08, 12.3, 0.4,
     [("목적: 사람이 쓰던 설계 명세·코드 작성을 AI가 자동화 — 문서 한 건에서 합성가능 RTL까지", 12.5, False, RGBColor(0xB9,0xC6,0xDE))])

# ── flow band (top half) ──
stages = [
    ("1", "IMPORT", VIOLET, "문서 반입", "스펙·PDF에서 요구 추출"),
    ("2", "DEEP INTERVIEW", BLUE, "요구사항 확정", "핵심 의사결정만 사람과 Q&A"),
    ("3", "SSOT", TEAL, "단일 명세 생성", "AI가 표준 명세 자동 작성"),
    ("4", "RTL", AMBER, "구현 생성", "AI가 검증가능 RTL 자동 작성"),
]
n = len(stages); cw = 2.85; gap = 0.36
total = n * cw + (n - 1) * gap
x0 = (13.333 - total) / 2
y = 1.95; ch = 2.15
for i, (num, name, c, what, sub) in enumerate(stages):
    box(x0 + i*(cw+gap), y, cw, ch, fill=WHITE, line=c, line_w=2.25)
    x = x0 + i*(cw+gap)
    box(x, y, cw, 0.62, fill=c)
    text(x+0.18, y, 0.5, 0.62, [(num, 16, True, WHITE)], anchor=MSO_ANCHOR.MIDDLE)
    text(x+0.7, y, cw-0.8, 0.62, [(name, 14, True, WHITE)], anchor=MSO_ANCHOR.MIDDLE)
    text(x+0.22, y+0.78, cw-0.44, 0.5, [(what, 14, True, DARK)])
    text(x+0.22, y+1.28, cw-0.44, 0.7, [(sub, 11, False, GREY)])
    if i < n-1:
        chevron(x+cw-0.02, y+ch/2-0.31, c=c)

# ── expected impact band (bottom half) ──
text(0.55, 4.35, 12.3, 0.4, [("기대효과", 17, True, NAVY)])
box(0.55, 4.78, 2.0, 0.05, fill=AMBER)
impacts = [
    ("설계 기간 단축", "문서→명세→RTL을 AI가\n자동 생성, 수작업 최소화", TEAL),
    ("일관성·추적성", "단일 명세(SSOT)에서\n모든 산출물 파생 → 불일치 제거", BLUE),
    ("품질 자동 검증", "단계별 자동 게이트로\n결함을 조기 차단", GREEN),
    ("사람은 핵심에 집중", "반복작업은 AI가,\n사람은 의사결정에 집중", AMBER),
]
iw = 2.98; ig = 0.25; ix0 = (13.333 - (4*iw + 3*ig)) / 2
iy = 5.0; ih = 1.75
for i, (title_i, body, c) in enumerate(impacts):
    x = ix0 + i*(iw+ig)
    box(x, iy, iw, ih, fill=LIGHT)
    box(x, iy, iw, 0.12, fill=c)
    text(x+0.22, iy+0.3, iw-0.44, 0.5, [(title_i, 13.5, True, c)])
    text(x+0.22, iy+0.85, iw-0.44, 0.85, [(body, 11, False, DARK)])

text(0.55, 7.0, 12.3, 0.35,
     [("ATLAS · Import → Deep Interview → SSOT → RTL  |  단계별 검증 게이트 통과 후 진행", 9.5, False, LGREY)])

import os
out = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "doc", "import_to_rtl_exec.pptx")
prs.save(out); print("saved:", out)
