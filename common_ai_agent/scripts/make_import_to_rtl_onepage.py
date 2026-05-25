#!/usr/bin/env python3
"""One-page flow: Import → Deep Interview → SSOT → RTL (ATLAS).

Grounded: import_document (tools.py:1158, PDF/doc → req/ extraction),
requirement elicitation / interactive ssot-gen, ssot-gen → yaml/<ip>.ssot.yaml,
rtl-gen → rtl/<ip>.sv + gate. Run: python3 scripts/make_import_to_rtl_onepage.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY  = RGBColor(0x16, 0x21, 0x3D)
VIOLET= RGBColor(0x6A, 0x4C, 0xC4)   # Import
BLUE  = RGBColor(0x2D, 0x6C, 0xDF)   # Deep Interview
TEAL  = RGBColor(0x17, 0x9E, 0x8A)   # SSOT
AMBER = RGBColor(0xE0, 0x8A, 0x1E)   # RTL
GREY  = RGBColor(0x5B, 0x63, 0x72)
LIGHT = RGBColor(0xF1, 0xF4, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK  = RGBColor(0x1E, 0x24, 0x30)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
s = prs.slides.add_slide(prs.slide_layouts[6])


def box(x, y, w, h, *, fill=None, line=None, line_w=1.0, rounded=True):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shp.shadow.inherit = False
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    return shp


def text(x, y, w, h, lines, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, (t, sz, b, c) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; r = p.add_run(); r.text = t
        r.font.size = Pt(sz); r.font.bold = b; r.font.color.rgb = c; r.font.name = "Calibri"
        p.space_after = Pt(3)
    return tb


def chevron(x, y):
    a = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(0.45), Inches(0.7))
    a.fill.solid(); a.fill.fore_color.rgb = GREY; a.line.fill.background(); a.shadow.inherit = False


# ── header band ──
box(0, 0, 13.333, 1.2, fill=NAVY, rounded=False)
box(0, 1.2, 13.333, 0.06, fill=TEAL, rounded=False)
text(0.55, 0.18, 12.3, 0.4, [("END-TO-END FLOW", 12, True, TEAL)])
text(0.53, 0.44, 12.3, 0.7, [("Import → Deep Interview → SSOT → RTL", 28, True, WHITE)])

# ── 4 stage cards ──
stages = [
    ("1", "IMPORT", VIOLET, "소스 문서 반입",
     ["PDF·스펙·문서를 가져와", "req/ 로 텍스트 추출", "(import_document)"],
     "IN:  스펙 문서", "OUT: req/ 추출물"),
    ("2", "DEEP INTERVIEW", BLUE, "요구사항 정제",
     ["모호점을 사람과 Q&A로", "확정 — 아키텍처 결정·", "요구 elicitation"],
     "IN:  추출물 + 질문", "OUT: 확정 요구사항"),
    ("3", "SSOT", TEAL, "단일 명세 생성",
     ["ssot-gen(LLM)이 명세 작성", "yaml/<ip>.ssot.yaml", "(20개 섹션)"],
     "IN:  확정 요구사항", "OUT: ssot.yaml ✓"),
    ("4", "RTL", AMBER, "구현 생성",
     ["rtl-gen(LLM)이 SV 작성", "rtl/<ip>.sv + list/<ip>.f", "compile + gate 검증"],
     "IN:  ssot.yaml", "OUT: <ip>.sv ✓"),
]
n = len(stages); cw = 2.78; gap = 0.42
total = n * cw + (n - 1) * gap
x0 = (13.333 - total) / 2
y = 1.75; ch = 4.0
for i, (num, name, c, what, body, ti, to) in enumerate(stages):
    x = x0 + i * (cw + gap)
    box(x, y, cw, ch, fill=WHITE, line=c, line_w=2.25)
    # number badge
    box(x + 0.18, y + 0.2, 0.55, 0.55, fill=c)
    text(x + 0.18, y + 0.2, 0.55, 0.55, [(num, 18, True, WHITE)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(x + 0.85, y + 0.22, cw - 0.95, 0.6, [(name, 15.5, True, c)], anchor=MSO_ANCHOR.MIDDLE)
    text(x + 0.22, y + 0.95, cw - 0.44, 0.4, [(what, 12.5, True, DARK)])
    text(x + 0.22, y + 1.4, cw - 0.44, 1.5, [(b, 10.5, False, GREY) for b in body])
    # io chips
    box(x + 0.18, y + 2.95, cw - 0.36, 0.45, fill=LIGHT)
    text(x + 0.3, y + 2.97, cw - 0.5, 0.42, [(ti, 9.5, False, DARK)], anchor=MSO_ANCHOR.MIDDLE)
    box(x + 0.18, y + 3.45, cw - 0.36, 0.45, fill=c)
    text(x + 0.3, y + 3.47, cw - 0.5, 0.42, [(to, 9.5, True, WHITE)], anchor=MSO_ANCHOR.MIDDLE)
    if i < n - 1:
        chevron(x + cw + 0.0, y + ch/2 - 0.35)

# ── bottom engine note ──
box(x0, 6.05, total, 1.05, fill=LIGHT)
box(x0, 6.05, 0.14, 1.05, fill=NAVY)
text(x0 + 0.35, 6.18, total - 0.6, 0.85, [
    ("조율 엔진: Orchestrator가 각 단계를 worker에 dispatch → yield_run으로 대기 → 완료 시 깨어나 다음 단계.", 11.5, True, NAVY),
    ("각 단계는 게이트 통과(passed) 후 다음으로 진행. 상태는 pipeline state(idle/running/passed/failed)로 가시화.", 10.5, False, GREY),
])

import os
out = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "doc", "import_to_rtl_onepage.pptx")
prs.save(out)
print("saved:", out)
