#!/usr/bin/env python3
"""Generate a presentation on the ATLAS SSOT -> RTL generation process.

Content is grounded in the real ATLAS pipeline: stage order, worker dispatch,
artifact paths, and validation gates as implemented in src/atlas_api_jobs.py /
workflow_stage_*.py. Run:  python3 scripts/make_ssot_rtl_ppt.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── palette ──────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x16, 0x21, 0x3D)   # title bar / primary
BLUE   = RGBColor(0x2D, 0x6C, 0xDF)   # accent
TEAL   = RGBColor(0x17, 0x9E, 0x8A)   # SSOT
AMBER  = RGBColor(0xE0, 0x8A, 0x1E)   # RTL
GREY   = RGBColor(0x5B, 0x63, 0x72)
LIGHT  = RGBColor(0xF1, 0xF4, 0xF9)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1E, 0x24, 0x30)
RED    = RGBColor(0xC0, 0x3A, 0x2B)

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _font(run, size, *, bold=False, color=DARK, name="Calibri"):
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color; run.font.name = name


def box(slide, x, y, w, h, *, fill=None, line=None, line_w=1.0, rounded=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    return shp


def text(slide, x, y, w, h, lines, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """lines: list of (text, size, bold, color) or (text, size, bold, color, name)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, spec in enumerate(lines):
        t, size, bold, color = spec[0], spec[1], spec[2], spec[3]
        name = spec[4] if len(spec) > 4 else "Calibri"
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run(); run.text = t; _font(run, size, bold=bold, color=color, name=name)
        p.space_after = Pt(4)
    return tb


def header(slide, title, kicker, accent):
    box(slide, 0, 0, 13.333, 1.15, fill=NAVY)
    box(slide, 0, 1.15, 13.333, 0.07, fill=accent)
    text(slide, 0.55, 0.16, 12, 0.5, [(kicker, 12, True, accent)])
    text(slide, 0.55, 0.42, 12.2, 0.7, [(title, 27, True, WHITE)])


def chip(slide, x, y, w, label, color, *, h=0.62, sub=""):
    box(slide, x, y, w, h, fill=color, rounded=True)
    lines = [(label, 12.5, True, WHITE)]
    if sub:
        lines.append((sub, 8.5, False, RGBColor(0xEC, 0xF2, 0xFA)))
    text(slide, x, y + (0.05 if sub else 0.0), w, h, lines,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def arrow(slide, x, y, w=0.32):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(0.22))
    a.fill.solid(); a.fill.fore_color.rgb = GREY; a.line.fill.background(); a.shadow.inherit = False


# ── 1. Title ────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, 13.333, 7.5, fill=NAVY)
box(s, 0, 5.0, 13.333, 0.09, fill=TEAL)
box(s, 4.2, 5.0, 4.9, 0.09, fill=AMBER)
text(s, 0.9, 1.9, 11.5, 1.2, [("SSOT → RTL 생성 과정", 46, True, WHITE)])
text(s, 0.92, 3.05, 11.5, 0.7,
     [("ATLAS AI 설계 파이프라인 — 명세(SSOT)에서 합성가능 RTL까지", 18, False, RGBColor(0xB9, 0xC6, 0xDE))])
text(s, 0.92, 5.3, 11.5, 0.5,
     [("LLM 워커 디스패치 · 아티펙트 검증 게이트 · 상태 머신", 13, False, RGBColor(0x8A, 0x9A, 0xBC))])

# ── 2. Pipeline overview ─────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
header(s, "전체 파이프라인 안에서의 위치", "OVERVIEW", BLUE)
text(s, 0.55, 1.45, 12.2, 0.5,
     [("하나의 IP는 명세에서 검증까지 단계형 파이프라인을 거칩니다. 본 문서는 첫 두 스테이지를 다룹니다.", 13, False, GREY)])
stages = ["SSOT", "FL-model", "CL-model", "Equiv", "RTL", "Lint", "TB", "Sim"]
colors = [TEAL, GREY, GREY, GREY, AMBER, GREY, GREY, GREY]
x = 0.55; y = 2.5; w = 1.35
for i, (st, c) in enumerate(zip(stages, colors)):
    chip(s, x, y, w, st, c)
    if i < len(stages) - 1:
        arrow(s, x + w + 0.02, y + 0.20)
    x += w + 0.36
text(s, 0.55, 3.5, 12.2, 0.5,
     [("▲ 이 문서의 초점: SSOT(명세 생성) · RTL(합성가능 코드 생성)", 12, True, NAVY)])
# two focus cards
box(s, 0.55, 4.2, 6.0, 2.55, fill=LIGHT, rounded=True)
box(s, 0.55, 4.2, 0.12, 2.55, fill=TEAL)
text(s, 0.85, 4.45, 5.5, 2.2, [
    ("SSOT", 16, True, TEAL),
    ("Single Source of Truth — 설계의 단일 진실 명세.", 12, False, DARK),
    ("구조화된 YAML (20개 섹션: io_list, function_model,", 11, False, GREY),
    ("registers, fsm, timing, cdc ...)", 11, False, GREY),
    ("→ 이후 모든 생성·검증의 기준점.", 12, True, DARK),
])
box(s, 6.95, 4.2, 5.85, 2.55, fill=LIGHT, rounded=True)
box(s, 6.95, 4.2, 0.12, 2.55, fill=AMBER)
text(s, 7.25, 4.45, 5.4, 2.2, [
    ("RTL", 16, True, AMBER),
    ("SSOT를 입력으로 SystemVerilog 구현 생성.", 12, False, DARK),
    ("rtl/<ip>.sv + list/<ip>.f 산출, 컴파일·게이트 검증.", 11, False, GREY),
    ("→ lint/tb/sim 등 하탄 스테이지의 입력.", 12, True, DARK),
])

# ── 3. SSOT generation flow ──────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
header(s, "SSOT 생성 과정", "STAGE 1 — SSOT", TEAL)
steps = [
    ("1. create_ip", "scaffold 초기화", "yaml/<ip>.ssot.yaml\n(type: draft, TBD ...)", GREY),
    ("2. ssot-gen worker", "LLM이 명세 작성", "draft → 실제 명세\n(top_module, io, model...)", TEAL),
    ("3. validator", "check_ssot_disk.sh", "스키마·계약 검증", BLUE),
    ("4. gate", "pipeline state", "ssot = passed", TEAL),
]
x = 0.55; y = 1.7; w = 2.95; h = 2.0
for i, (t, sub, body, c) in enumerate(steps):
    box(s, x, y, w, h, fill=WHITE, line=c, line_w=1.75, rounded=True)
    box(s, x, y, w, 0.55, fill=c, rounded=True)
    text(s, x + 0.12, y + 0.04, w - 0.24, 0.5, [(t, 13, True, WHITE)], anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + 0.18, y + 0.7, w - 0.36, 1.2, [
        (sub, 12, True, DARK),
        (body, 10.5, False, GREY),
    ])
    if i < len(steps) - 1:
        arrow(s, x + w + 0.02, y + h/2 - 0.1, w=0.30)
    x += w + 0.34
text(s, 0.55, 4.1, 12.2, 0.5,
     [("핵심: create_ip이 먼저 TBD scaffold를 심고, ssot-gen이 그걸 실제 명세로 덮어씁니다.", 13, True, NAVY)])
box(s, 0.55, 4.75, 12.27, 1.95, fill=LIGHT, rounded=True)
text(s, 0.85, 4.95, 11.8, 1.7, [
    ("⚠  scaffold vs 실제 SSOT 구분 (검증 포인트)", 13, True, AMBER),
    ("•  scaffold: type: draft + 다수의 TBD  →  아직 생성 안 됨 (파일 존재≠완료)", 11.5, False, DARK),
    ("•  실제 SSOT: draft 마커 제거 + function_model 존재 + TBD 극소  →  passed", 11.5, False, DARK),
    ("•  실측: 실제 ssot-gen 1회 ≈ 수 분 (워커 LLM 추론 시간이 지배적)", 11.5, False, GREY),
])

# ── 4. SSOT structure ────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
header(s, "SSOT 구조 — 20개 섹션", "STAGE 1 — SSOT", TEAL)
text(s, 0.55, 1.4, 12.2, 0.5,
     [("yaml/<ip>.ssot.yaml 은 설계의 모든 면을 구조화해 담습니다. 주요 섹션:", 13, False, GREY)])
groups = [
    ("구조", ["top_module", "sub_modules", "decomposition", "rtl_contract", "parameters"], TEAL),
    ("인터페이스·동작", ["io_list", "features", "dataflow", "function_model", "cycle_model"], BLUE),
    ("클럭·상태·자원", ["clock_reset_domains", "cdc / rdc", "registers", "memory", "interrupts"], AMBER),
    ("구현 제약", ["fsm", "timing", "power", "security"], GREY),
]
x = 0.55; w = 2.98
for title_g, items, c in groups:
    box(s, x, 2.1, w, 4.4, fill=WHITE, line=c, line_w=1.5, rounded=True)
    box(s, x, 2.1, w, 0.5, fill=c, rounded=True)
    text(s, x + 0.1, 2.13, w - 0.2, 0.45, [(title_g, 12.5, True, WHITE)], anchor=MSO_ANCHOR.MIDDLE)
    lines = [(f"• {it}", 12, False, DARK) for it in items]
    text(s, x + 0.2, 2.8, w - 0.35, 3.5, lines)
    x += w + 0.13

# ── 5. RTL generation flow ───────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
header(s, "RTL 생성 과정", "STAGE 2 — RTL", AMBER)
steps = [
    ("1. 입력: SSOT", "passed 명세", "function_model /\nio_list / rtl_contract", TEAL),
    ("2. rtl-gen worker", "LLM이 SV 작성", "rtl/<ip>.sv\nlist/<ip>.f", AMBER),
    ("3. compile", "rtl_compile.json", "compile_rc == 0", BLUE),
    ("4. gate", "ssot-rtl.json", "placeholder 없음\n+ gate pass", AMBER),
]
x = 0.55; y = 1.7; w = 2.95; h = 2.0
for i, (t, sub, body, c) in enumerate(steps):
    box(s, x, y, w, h, fill=WHITE, line=c, line_w=1.75, rounded=True)
    box(s, x, y, w, 0.55, fill=c, rounded=True)
    text(s, x + 0.12, y + 0.04, w - 0.24, 0.5, [(t, 13, True, WHITE)], anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + 0.18, y + 0.7, w - 0.36, 1.2, [
        (sub, 12, True, DARK),
        (body, 10.5, False, GREY),
    ])
    if i < len(steps) - 1:
        arrow(s, x + w + 0.02, y + h/2 - 0.1, w=0.30)
    x += w + 0.34
text(s, 0.55, 4.1, 12.2, 0.5,
     [("산출물: rtl/<ip>.sv (구현) + list/<ip>.f (파일리스트) + rtl_compile.json (컴파일 결과)", 12.5, True, NAVY)])
box(s, 0.55, 4.75, 12.27, 1.95, fill=LIGHT, rounded=True)
text(s, 0.85, 4.95, 11.8, 1.7, [
    ("⚠  RTL 게이트 판정 (실제 검증 로직)", 13, True, AMBER),
    ("•  placeholder 체크: rtl/*.sv 에 TBD/TODO/PLACEHOLDER/stub 있으면 fail", 11.5, False, DARK),
    ("•  verdict 증거 있을 때만 판정: rtl_compile.json 또는 logs/stage_engine/ssot-rtl.json", 11.5, False, DARK),
    ("•  생성 전 scaffold(placeholder + list.f)는 'failed'가 아니라 'locked/idle' ← 수정 반영됨", 11.5, False, GREY),
])

# ── 6. How it is driven (orchestrator) ───────────────────────────────────
s = prs.slides.add_slide(BLANK)
header(s, "어떻게 구동되는가 — 조율 루프", "ENGINE", BLUE)
text(s, 0.55, 1.4, 12.2, 0.5,
     [("Orchestrator는 직접 만들지 않고 조율합니다. 실제 생성은 worker가 수행.", 13, False, GREY)])
loop = [
    ("read_pipeline_state", "현황 읽기", BLUE),
    ("dispatch_workflow", "worker에 ssot/rtl 시키기", TEAL),
    ("yield_run", "job 완료·유저까지 park", AMBER),
    ("(wake)", "worker가 벨 → 재개", GREY),
]
x = 0.8; y = 2.45; w = 2.7; h = 1.5
for i, (t, sub, c) in enumerate(loop):
    box(s, x, y, w, h, fill=WHITE, line=c, line_w=1.75, rounded=True)
    text(s, x + 0.15, y + 0.2, w - 0.3, 1.1, [
        (t, 13.5, True, c),
        (sub, 11, False, DARK),
    ], anchor=MSO_ANCHOR.MIDDLE)
    if i < len(loop) - 1:
        arrow(s, x + w + 0.0, y + h/2 - 0.1, w=0.28)
    x += w + 0.28
box(s, 0.8, 4.5, 11.7, 2.0, fill=LIGHT, rounded=True)
text(s, 1.1, 4.72, 11.1, 1.7, [
    ("worker(ssot-gen / rtl-gen) = 실제 일꾼", 13, True, NAVY),
    ("•  worker가 LLM으로 명세·코드 생성 → 디스크에 산출물 기록", 11.5, False, DARK),
    ("•  완료 시 jobs API가 notify_job_complete → park된 orchestrator 깨움", 11.5, False, DARK),
    ("•  상태는 read_pipeline_state가 DB(도는 job) + 디스크 산출물을 합쳐서 판정", 11.5, False, GREY),
])

# ── 7. State model ───────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
header(s, "상태 판정 — pipeline state", "ENGINE", BLUE)
text(s, 0.55, 1.4, 12.2, 0.5,
     [("GET /api/pipeline/state 가 각 스테이지의 상태를 돌려줍니다 (DB 우선, 파일 폴백).", 13, False, GREY)])
states = [
    ("idle", "안 돌음 · 의존 없음", GREY),
    ("locked", "상위 stage 미통과", GREY),
    ("ready", "의존 충족, 실행 가능", BLUE),
    ("running", "worker 동작 중", AMBER),
    ("passed", "게이트 통과", TEAL),
    ("failed", "검증 실패", RED),
]
x = 0.55; y = 2.2; w = 3.9; h = 1.1
for i, (st, desc, c) in enumerate(states):
    col = i % 3; row = i // 3
    bx = 0.55 + col * (w + 0.25); by = y + row * (h + 0.3)
    box(s, bx, by, w, h, fill=WHITE, line=c, line_w=2.0, rounded=True)
    box(s, bx, by, 0.16, h, fill=c)
    text(s, bx + 0.35, by + 0.1, w - 0.5, h - 0.2, [
        (st, 15, True, c),
        (desc, 11, False, DARK),
    ], anchor=MSO_ANCHOR.MIDDLE)
box(s, 0.55, 5.55, 12.27, 1.15, fill=LIGHT, rounded=True)
text(s, 0.85, 5.7, 11.8, 0.95, [
    ("→ SSOT/RTL 생성이 끝나면 해당 stage가 running → passed 로 전환.", 12.5, True, NAVY),
    ("   단, 산출물은 일찍 완성되어도 검증 꼬리(느린 모델) 때문에 passed 전환이 지연될 수 있음.", 11, False, GREY),
])

# ── 8. Summary ───────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, 13.333, 7.5, fill=NAVY)
text(s, 0.9, 0.7, 11.5, 0.8, [("요약", 30, True, WHITE)])
box(s, 0.9, 1.55, 5.7, 0.06, fill=TEAL)
rows = [
    ("SSOT", "create_ip(scaffold) → ssot-gen(LLM) → validator → passed.", TEAL),
    ("", "yaml/<ip>.ssot.yaml, 20개 섹션, draft≠실제 구분이 핵심.", RGBColor(0x9F,0xB2,0xD0)),
    ("RTL", "SSOT → rtl-gen(LLM) → rtl/<ip>.sv + compile + gate → passed.", AMBER),
    ("", "placeholder·verdict 기반 게이트로 진짜 생성만 통과.", RGBColor(0x9F,0xB2,0xD0)),
    ("엔진", "orchestrator: read→dispatch→yield. worker가 실제 생성.", BLUE),
    ("", "상태는 pipeline state(idle/running/passed/failed)로 가시화.", RGBColor(0x9F,0xB2,0xD0)),
]
y = 2.0
for label, body, c in rows:
    if label:
        text(s, 0.9, y, 2.2, 0.5, [(label, 17, True, c)])
        text(s, 3.0, y + 0.02, 9.3, 0.5, [(body, 14, False, WHITE)])
        y += 0.5
    else:
        text(s, 3.0, y, 9.3, 0.5, [(body, 11.5, False, c)])
        y += 0.55
text(s, 0.9, 6.6, 11.5, 0.5,
     [("ATLAS · SSOT → RTL generation · grounded in src/atlas_api_jobs.py pipeline", 11, False, RGBColor(0x6E,0x80,0xA4))])

import os
out = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "doc", "ssot_rtl_generation.pptx")
prs.save(out)
print("saved:", out, "| slides:", len(prs.slides._sldIdLst))
