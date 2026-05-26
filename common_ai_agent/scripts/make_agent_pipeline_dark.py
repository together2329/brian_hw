#!/usr/bin/env python3
"""ATLAS SSOT · RTL Generator — single dark/formal one-pager.
Pipeline + components + agent loops + TODO workflow + expected impact, all on
one slide. Run: python3 scripts/make_agent_pipeline_dark.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG     = RGBColor(0x0D, 0x14, 0x22)
PANEL  = RGBColor(0x18, 0x21, 0x36)
PANEL2 = RGBColor(0x12, 0x1A, 0x2C)
EDGE   = RGBColor(0x2B, 0x38, 0x52)
TEAL   = RGBColor(0x35, 0xC4, 0xAE)
AMBER  = RGBColor(0xEA, 0xA8, 0x3A)
BLUE   = RGBColor(0x5A, 0x97, 0xF5)
GREEN  = RGBColor(0x4F, 0xC2, 0x84)
VIOLET = RGBColor(0x9B, 0x8C, 0xF0)
RED    = RGBColor(0xD9, 0x5A, 0x4F)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
MUTE   = RGBColor(0x9A, 0xAA, 0xC6)
FAINT  = RGBColor(0x6B, 0x79, 0x95)
DARKTX = RGBColor(0x0D, 0x14, 0x22)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
s = prs.slides.add_slide(prs.slide_layouts[6])


def shape(kind, x, y, w, h, *, fill=None, line=None, line_w=1.0):
    sh = s.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.shadow.inherit = False
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb = line; sh.line.width = Pt(line_w)
    return sh


def text(x, y, w, h, lines, *, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=1):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, (t, sz, b, c) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; r = p.add_run(); r.text = t
        r.font.size = Pt(sz); r.font.bold = b; r.font.color.rgb = c; r.font.name = "Calibri"
        p.space_after = Pt(spacing)
    return tb


def arrow(x, y, w, c, h=0.18):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = c; a.line.fill.background(); a.shadow.inherit = False


# ── background ──
shape(MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 7.5, fill=BG)

# ── header ──
shape(MSO_SHAPE.RECTANGLE, 0.55, 0.32, 0.07, 0.85, fill=TEAL)
text(0.8, 0.26, 12, 0.3, [("ATLAS", 10.5, True, FAINT)], align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
text(0.78, 0.49, 12, 0.5, [("SSOT · RTL Generator", 23, True, WHITE)], align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
text(0.80, 0.97, 12.4, 0.3,
     [("An AI agent that turns existing DB knowledge and Q&A into a standardized spec (SSOT), then generates synthesizable RTL.",
       10.5, False, MUTE)], align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

# ── BAND 1: pipeline (hero) ──
text(0.78, 1.4, 12, 0.26, [("PIPELINE", 11, True, TEAL)], align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
AG = MSO_SHAPE.ROUNDED_RECTANGLE; DOC = MSO_SHAPE.FLOWCHART_DOCUMENT
y = 1.7; h = 1.4
nodes = [
    ("input", "Requirements", "INPUT", VIOLET, ["▸ AtlasDB (existing DB)", "▸ Deep Interview (Q&A)"]),
    ("agent", "SSOT-AGENT", "spec gen", TEAL, ["▸ ssot-gen (LLM)", "▸ verify_ssot.py"]),
    ("out", "SSOT", "OUTPUT", TEAL, ["▸ <ip>.ssot.yaml", "▸ 20-section spec"]),
    ("agent", "RTL-AGENT", "RTL gen", AMBER, ["▸ ssot_to_rtl.py (LLM)", "▸ build_gate.sh"]),
    ("out", "RTL", "OUTPUT", AMBER, ["▸ <ip>.sv", "▸ compile + gate"]),
]
widths = [2.6, 2.3, 2.4, 2.3, 2.4]; gap = 0.2
total = sum(widths) + gap*4; x = (13.333-total)/2
for i, (kind, label, sub, accent, comps) in enumerate(nodes):
    w = widths[i]
    if kind == "agent":
        shape(AG, x, y, w, h, fill=accent)
        text(x, y+0.1, w, 0.6, [("◇ AI AGENT", 8, True, DARKTX), (label, 13, True, DARKTX), (sub, 10, False, DARKTX)])
        text(x+0.2, y+0.86, w-0.36, 0.5, [(c, 8.5, False, DARKTX) for c in comps], align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    elif kind == "out":
        shape(DOC, x, y, w, h+0.08, fill=PANEL, line=accent, line_w=2.0)
        text(x, y+0.12, w, 0.55, [("OUTPUT", 8.5, True, accent), (label, 15, True, accent)])
        text(x+0.2, y+0.85, w-0.36, 0.5, [(c, 8.5, False, MUTE) for c in comps], align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    else:
        shape(MSO_SHAPE.RECTANGLE, x, y, w, h, fill=PANEL2, line=EDGE, line_w=1.5)
        shape(MSO_SHAPE.RECTANGLE, x, y, w, 0.07, fill=VIOLET)
        text(x, y+0.12, w, 0.55, [("INPUT", 8.5, True, VIOLET), (label, 13.5, True, WHITE)])
        text(x+0.2, y+0.85, w-0.36, 0.5, [(c, 8.5, False, MUTE) for c in comps], align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    if i < 4:
        arrow(x+w-0.01, y+h/2-0.09, gap+0.06, TEAL if i == 1 else (AMBER if i == 3 else FAINT))
    x += w + gap

# ── BAND 2: agent loops (compact, two lines) ──
shape(MSO_SHAPE.RECTANGLE, 0.0, 3.32, 13.333, 0.01, fill=EDGE)
text(0.78, 3.4, 12, 0.26, [("AGENT LOOPS", 11, True, AMBER)], align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
text(0.80, 3.68, 12.5, 0.3, [
    ("Orchestrator:  read_pipeline_state  →  dispatch_workflow  →  yield_run  ↻ wake on job/user", 11, False, MUTE)],
    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
text(0.80, 3.96, 12.5, 0.3, [
    ("Worker (ReAct):  Observe  →  Reason (LLM)  →  Act  ↻ until done  →  emit + gate", 11, False, MUTE)],
    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
text(0.80, 4.24, 12.5, 0.28, [
    ("tools — orch: read_pipeline_state · dispatch_workflow · ask_user      worker: read_file · replace_in_file · run_command · todo_write",
     8.5, False, FAINT)], align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

# ── BAND 3: TODO workflow (lifecycle strip) ──
shape(MSO_SHAPE.RECTANGLE, 0.0, 4.62, 13.333, 0.01, fill=EDGE)
text(0.78, 4.7, 12, 0.26, [("TODO WORKFLOW", 11, True, BLUE)], align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
states = [("PENDING", FAINT), ("IN_PROGRESS", AMBER), ("COMPLETED", BLUE), ("APPROVED ✓", GREEN)]
sw = 2.25; sg = 0.42; sx = 0.95; sy = 5.0; sh = 0.55
for i, (t, c) in enumerate(states):
    shape(MSO_SHAPE.ROUNDED_RECTANGLE, sx, sy, sw, sh, fill=PANEL, line=c, line_w=1.75)
    text(sx, sy, sw, sh, [(t, 11.5, True, c)])
    if i < 3:
        arrow(sx+sw+0.05, sy+sh/2-0.08, sg-0.1, c, h=0.16)
    sx += sw + sg
text(sx+0.05, sy, 2.9, sh, [("↩ rejected · ⏸ blocked", 9.5, True, RED), ("→ todo_events ledger", 8.5, False, FAINT)],
     align=PP_ALIGN.LEFT)
text(0.80, 5.62, 12.5, 0.26, [
    ("derived from SSOT (golden todos) & RTL gaps (derive_rtl_todos.py) · stored in workflow_todos + todo.json", 8.5, False, FAINT)],
    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

# ── BAND 4: expected impact ──
shape(MSO_SHAPE.RECTANGLE, 0.0, 5.98, 13.333, 0.01, fill=EDGE)
text(0.78, 6.04, 12, 0.26, [("EXPECTED IMPACT", 11, True, MUTE)], align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
impacts = [
    ("Design productivity", "spec & RTL authoring automated", TEAL),
    ("Standardization", "one SSOT → consistent IP, reuse", BLUE),
    ("Traceability", "spec→RTL + gates catch defects", GREEN),
    ("Knowledge asset", "DB & Q&A captured as formal spec", AMBER),
]
iw = 2.98; ig = 0.25; ix = (13.333-(4*iw+3*ig))/2; iy = 6.35; ih = 0.85
for i, (t_i, b_i, c) in enumerate(impacts):
    x = ix + i*(iw+ig)
    shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, iy, iw, ih, fill=PANEL, line=EDGE, line_w=1.0)
    shape(MSO_SHAPE.RECTANGLE, x, iy, 0.09, ih, fill=c)
    text(x+0.26, iy+0.1, iw-0.4, 0.32, [(t_i, 11, True, c)], align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    text(x+0.26, iy+0.42, iw-0.42, 0.4, [(b_i, 9, False, MUTE)], align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

import os
out = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "doc", "agent_pipeline_dark.pptx")
prs.save(out); print("saved:", out, "| slides:", len(prs.slides._sldIdLst))
