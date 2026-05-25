#!/usr/bin/env python3
"""Executive one-pager — agent/IO chain:
INPUT → SSOT-AGENT → OUTPUT(SSOT) → RTL-AGENT → OUTPUT(RTL).
Run: python3 scripts/make_agent_pipeline.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY  = RGBColor(0x14, 0x1E, 0x3A)
TEAL  = RGBColor(0x17, 0x9E, 0x8A)
AMBER = RGBColor(0xE0, 0x8A, 0x1E)
BLUE  = RGBColor(0x2D, 0x6C, 0xDF)
GREEN = RGBColor(0x2E, 0x9E, 0x5B)
GREY  = RGBColor(0x5B, 0x63, 0x72)
LGREY = RGBColor(0x8A, 0x93, 0xA3)
LIGHT = RGBColor(0xF1, 0xF4, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK  = RGBColor(0x1E, 0x24, 0x30)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
s = prs.slides.add_slide(prs.slide_layouts[6])


def shape(kind, x, y, w, h, *, fill=None, line=None, line_w=1.0):
    shp = s.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.shadow.inherit = False
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    return shp


def text(x, y, w, h, lines, *, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, (t, sz, b, c) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; r = p.add_run(); r.text = t
        r.font.size = Pt(sz); r.font.bold = b; r.font.color.rgb = c; r.font.name = "Calibri"
        p.space_after = Pt(2)
    return tb


def arrow(x, y, w=0.5, c=GREY):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(0.28))
    a.fill.solid(); a.fill.fore_color.rgb = c; a.line.fill.background(); a.shadow.inherit = False


# ── header ──
shape(MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 1.35, fill=NAVY)
shape(MSO_SHAPE.RECTANGLE, 0, 1.35, 13.333, 0.07, fill=TEAL)
text(0.55, 0.18, 12.3, 0.4, [("AI 반도체 IP 설계 자동화", 12, True, TEAL)], align=PP_ALIGN.LEFT)
text(0.53, 0.44, 12.3, 0.7, [("요구사항 입력 → AI 에이전트가 명세·RTL 자동 생성", 26, True, WHITE)], align=PP_ALIGN.LEFT)
# legend
text(8.7, 1.5, 4.4, 0.35, [("■ AGENT = AI가 실행   ▢ OUTPUT = 산출물", 10.5, True, GREY)], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.TOP)

# ── agent/IO chain ──
AG = MSO_SHAPE.ROUNDED_RECTANGLE
DOC = MSO_SHAPE.FLOWCHART_DOCUMENT
y = 2.35; h = 2.35
nodes = [
    # (kind, label, sub, fill, line, textcolor, kicker)
    ("input", "INPUT", "요구사항", None, GREY, DARK,
     "Import(문서·PDF 추출)\n+ Deep Interview(확정)"),
    ("agent", "SSOT-AGENT", "명세 생성", TEAL, None, WHITE, "AI"),
    ("out",   "OUTPUT", "SSOT", None, TEAL, TEAL,
     "yaml/<ip>.ssot.yaml\n(20개 섹션 표준 명세)"),
    ("agent", "RTL-AGENT", "RTL 생성", AMBER, None, WHITE, "AI"),
    ("out",   "OUTPUT", "RTL", None, AMBER, AMBER,
     "rtl/<ip>.sv\n(검증가능 구현)"),
]
widths = [2.55, 2.15, 2.45, 2.15, 2.45]
gap = 0.30
total = sum(widths) + gap*(len(nodes)-1)
x = (13.333 - total) / 2
for i, (kind, label, sub, fill, line, tc, body) in enumerate(nodes):
    w = widths[i]
    if kind == "agent":
        shape(AG, x, y, w, h, fill=fill)
        text(x, y+0.18, w, 0.45, [("◇ AI AGENT", 9.5, True, RGBColor(0xE6,0xEF,0xEC))])
        text(x, y+0.55, w, 1.0, [(label, 15, True, WHITE), (sub, 12, False, RGBColor(0xEC,0xF4,0xF1))])
    elif kind == "out":
        shape(DOC, x, y, w, h+0.15, fill=WHITE, line=line, line_w=2.5)
        text(x, y+0.2, w, 0.4, [("OUTPUT", 10, True, line)])
        text(x, y+0.6, w, 0.55, [(sub, 17, True, line)])
        text(x, y+1.35, w, 0.85, [(body, 10, False, GREY)])
    else:  # input
        shape(MSO_SHAPE.RECTANGLE, x, y, w, h, fill=LIGHT, line=GREY, line_w=1.75)
        shape(MSO_SHAPE.RECTANGLE, x, y, w, 0.1, fill=GREY)
        text(x, y+0.25, w, 0.4, [("INPUT", 11, True, GREY)])
        text(x, y+0.65, w, 0.5, [(sub, 16, True, DARK)])
        text(x, y+1.3, w, 0.9, [(body, 10.5, False, GREY)])
    if i < len(nodes)-1:
        ac = TEAL if i in (1,) else (AMBER if i == 3 else GREY)
        arrow(x + w - 0.03, y + h/2 - 0.14, w=gap+0.12, c=ac)
    x += w + gap

# small "produces" hints under agents
text((13.333-total)/2 + widths[0]+gap, y+h+0.05, widths[1]+gap+widths[2], 0.3,
     [("AI가 명세를 생성", 9.5, False, TEAL)])
xr = (13.333-total)/2 + sum(widths[:3]) + gap*3
text(xr, y+h+0.05, widths[3]+gap+widths[4], 0.3, [("AI가 RTL을 생성", 9.5, False, AMBER)])

# ── expected impact ──
text(0.55, 5.4, 12.3, 0.4, [("기대효과", 16, True, NAVY)], align=PP_ALIGN.LEFT)
shape(MSO_SHAPE.RECTANGLE, 0.55, 5.82, 1.7, 0.05, fill=AMBER)
impacts = [
    ("설계 기간 단축", "명세·코드 작성을 AI 자동화", TEAL),
    ("일관성·추적성", "SSOT 단일 명세에서 RTL 파생", BLUE),
    ("품질 자동 검증", "단계별 게이트로 결함 조기 차단", GREEN),
    ("사람은 핵심 결정만", "반복은 AI, 판단은 사람", AMBER),
]
iw = 2.98; ig = 0.25; ix = (13.333-(4*iw+3*ig))/2
iy = 6.0; ih = 1.05
for i, (t_i, b_i, c) in enumerate(impacts):
    x = ix + i*(iw+ig)
    shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, iy, iw, ih, fill=LIGHT)
    shape(MSO_SHAPE.RECTANGLE, x, iy, 0.12, ih, fill=c)
    text(x+0.3, iy+0.12, iw-0.45, 0.4, [(t_i, 12.5, True, c)], align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    text(x+0.3, iy+0.52, iw-0.45, 0.45, [(b_i, 10.5, False, DARK)], align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

import os
out = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "doc", "agent_pipeline_onepage.pptx")
prs.save(out); print("saved:", out)
